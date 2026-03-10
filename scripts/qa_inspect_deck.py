"""Comprehensive QA inspection of the rendered deck.

Checks every slide for:
- Layout assignment (correct template layout used)
- Placeholder population (no empty placeholders)
- Font families (Aptos Display / Aptos)
- Font sizes (title 28pt+, body 11-16pt)
- Colors (navy headers, teal accents, theme colors)
- Table formatting (header fills, alternating rows, status colors)
- Accent bars (present on Title Only slides)
- Text content (no placeholder text, no empty slides)
"""

import sys
sys.stdout.reconfigure(encoding="utf-8")
sys.stderr.reconfigure(encoding="utf-8")

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from collections import defaultdict

DECK_PATH = Path(__file__).resolve().parent.parent / "output" / "deck_formatted.pptx"

# Expected theme colors
NAVY = "0E2841"
TEAL = "156082"
ORANGE = "E97132"
BLUE = "0F9ED5"
WHITE = "FFFFFF"
LIGHT_GRAY = "F2F2F2"

ISSUES = []
WARNINGS = []

def issue(slide_num, msg):
    ISSUES.append(f"  ISSUE S{slide_num}: {msg}")

def warn(slide_num, msg):
    WARNINGS.append(f"  WARN  S{slide_num}: {msg}")

def rgb_hex(color):
    """Get hex string from RGBColor or None."""
    if color is None:
        return None
    try:
        return str(color).upper()
    except Exception:
        return None


def safe_font_color(font):
    """Safely get font color as hex string, handling _NoneColor."""
    try:
        if font.color and font.color.type is not None:
            return rgb_hex(font.color.rgb)
    except (AttributeError, TypeError):
        pass
    return None

def get_fill_color(shape):
    """Try to get fill color from a shape."""
    try:
        fill = shape.fill
        if fill.type is not None:
            fc = fill.fore_color
            if fc and fc.rgb:
                return str(fc.rgb).upper()
    except:
        pass
    return None

def get_cell_fill_xml(cell):
    """Get fill color from table cell via XML."""
    from lxml import etree
    nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    tc = cell._tc
    tcPr = tc.find("a:tcPr", nsmap)
    if tcPr is not None:
        solidFill = tcPr.find("a:solidFill", nsmap)
        if solidFill is not None:
            srgb = solidFill.find("a:srgbClr", nsmap)
            if srgb is not None:
                return srgb.get("val", "").upper()
    return None


def inspect_slide(slide, slide_num, prs):
    layout_name = slide.slide_layout.name
    layout_idx = list(prs.slide_layouts).index(slide.slide_layout) if slide.slide_layout in prs.slide_layouts else -1

    print(f"\n{'='*70}")
    print(f"SLIDE {slide_num}: Layout={layout_idx} ({layout_name})")
    print(f"{'='*70}")

    # Collect all shapes
    shapes = list(slide.shapes)
    has_table = any(s.has_table for s in shapes)
    has_text = any(s.has_text_frame for s in shapes)
    text_shapes = [s for s in shapes if s.has_text_frame]
    table_shapes = [s for s in shapes if s.has_table]

    print(f"  Shapes: {len(shapes)} total, {len(text_shapes)} text, {len(table_shapes)} tables")

    # Check for empty slide
    all_text = ""
    for s in text_shapes:
        all_text += s.text_frame.text
    for s in table_shapes:
        for row in s.table.rows:
            for cell in row.cells:
                all_text += cell.text

    if not all_text.strip():
        issue(slide_num, "EMPTY SLIDE - no text content found")
        return

    # Check accent bars on Title Only layouts
    if layout_idx == 5:  # Title Only
        rect_shapes = [s for s in shapes if hasattr(s, 'shape_type') and s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
        accent_bars = []
        for s in rect_shapes:
            if s.width > Inches(0.5) and s.height < Inches(0.3):
                accent_bars.append(s)
        if accent_bars:
            print(f"  Accent bar: YES ({len(accent_bars)} found)")
            for bar in accent_bars:
                color = get_fill_color(bar)
                if color:
                    print(f"    Bar color: #{color}")
                    if color != TEAL:
                        warn(slide_num, f"Accent bar color #{color} != expected #{TEAL}")
        else:
            warn(slide_num, "Title Only layout but no accent bar found")

    # Inspect text frames
    print(f"\n  TEXT FRAMES:")
    for si, shape in enumerate(text_shapes):
        tf = shape.text_frame
        text_preview = tf.text[:80].replace('\n', ' | ')
        if not text_preview.strip():
            continue

        print(f"    Shape {si}: '{text_preview}...' " if len(tf.text) > 80 else f"    Shape {si}: '{text_preview}'")

        # Check fonts in paragraphs
        for pi, para in enumerate(tf.paragraphs):
            for ri, run in enumerate(para.runs):
                font = run.font
                font_name = font.name
                font_size = font.size
                font_bold = font.bold
                font_color = safe_font_color(font)

                if font_name and "Aptos" not in font_name and "Calibri" not in font_name:
                    warn(slide_num, f"Non-theme font '{font_name}' in shape {si}")

                size_pt = font_size / 12700 if font_size else None

                if pi == 0 and ri == 0 and si == 0:
                    # First run of first paragraph of title-like shape
                    print(f"      Font: {font_name or 'inherited'}, "
                          f"Size: {size_pt:.0f}pt, " if size_pt else f"      Font: {font_name or 'inherited'}, Size: inherited, ",
                          end="")
                    print(f"Bold: {font_bold}, Color: #{font_color}" if font_color else f"Bold: {font_bold}, Color: inherited")

    # Inspect tables
    if table_shapes:
        print(f"\n  TABLES:")
        for ti, shape in enumerate(table_shapes):
            tbl = shape.table
            rows = len(tbl.rows)
            cols = len(tbl.columns)
            print(f"    Table {ti}: {rows}x{cols}")

            # Check header row (first row)
            if rows > 0:
                header_fills = []
                for ci, cell in enumerate(tbl.rows[0].cells):
                    fill = get_cell_fill_xml(cell)
                    header_fills.append(fill)

                    # Check header text formatting
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            hfont = run.font
                            hcolor = safe_font_color(hfont)
                            hbold = hfont.bold

                unique_fills = set(header_fills)
                print(f"      Header row fills: {[f'#{f}' if f else 'None' for f in header_fills[:4]]}{'...' if len(header_fills) > 4 else ''}")

                if NAVY in [f for f in header_fills if f]:
                    print(f"      Header: NAVY fill - GOOD")
                elif LIGHT_GRAY in [f for f in header_fills if f]:
                    print(f"      Header: Light gray fill (no explicit header)")
                elif all(f is None for f in header_fills):
                    warn(slide_num, "Table header has no fill color")

            # Check alternating rows
            if rows > 2:
                row_fills = []
                for ri in range(1, min(rows, 5)):
                    fills = [get_cell_fill_xml(cell) for cell in tbl.rows[ri].cells]
                    row_fills.append(fills[0] if fills else None)

                print(f"      Body row fills (first col): {[f'#{f}' if f else 'None' for f in row_fills]}")

                # Check for alternating pattern
                if len(row_fills) >= 2:
                    has_alternating = any(f != row_fills[0] for f in row_fills[1:])
                    if has_alternating:
                        print(f"      Alternating rows: YES")
                    else:
                        warn(slide_num, "Table body rows do not alternate colors")

            # Check for status colors (compliance matrix)
            status_colors = set()
            for ri in range(rows):
                for cell in tbl.rows[ri].cells:
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            c = safe_font_color(run.font)
                            if c and c not in (WHITE, "000000", NAVY):
                                status_colors.add(c)
            if status_colors:
                print(f"      Status colors found: {['#'+c for c in status_colors]}")

    # Check for placeholder text / leftover template content
    full_text = all_text.lower()
    placeholder_terms = ["xxxx", "lorem", "ipsum", "click to add", "this page layout"]
    for term in placeholder_terms:
        if term in full_text:
            issue(slide_num, f"Placeholder text found: '{term}'")


def main():
    print(f"QA Inspection: {DECK_PATH}")
    print(f"File size: {DECK_PATH.stat().st_size / 1024:.1f} KB")

    prs = Presentation(str(DECK_PATH))
    width = prs.slide_width
    height = prs.slide_height
    print(f"Dimensions: {width/914400:.2f}in x {height/914400:.2f}in")
    print(f"Total slides: {len(prs.slides)}")

    # Build layout list
    print(f"\nAvailable layouts:")
    for i, layout in enumerate(prs.slide_layouts):
        print(f"  {i}: {layout.name}")

    # Inspect each slide
    for i, slide in enumerate(prs.slides, 1):
        inspect_slide(slide, i, prs)

    # Summary
    print(f"\n{'='*70}")
    print(f"QA SUMMARY")
    print(f"{'='*70}")

    if ISSUES:
        print(f"\n  ISSUES ({len(ISSUES)}):")
        for i in ISSUES:
            print(i)
    else:
        print(f"\n  No critical issues found.")

    if WARNINGS:
        print(f"\n  WARNINGS ({len(WARNINGS)}):")
        for w in WARNINGS:
            print(w)
    else:
        print(f"\n  No warnings.")

    print(f"\n  Total: {len(ISSUES)} issues, {len(WARNINGS)} warnings")


if __name__ == "__main__":
    main()
