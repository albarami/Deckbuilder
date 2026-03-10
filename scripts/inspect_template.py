"""Deep inspection of Presentation6.pptx template.

Outputs for EVERY slide layout:
  - Layout name and index
  - Every placeholder: index, type, position, default font
  - Every shape: type, position, fill color, line color
  - Background fill
  - Theme colors
"""

import sys
from pathlib import Path
from lxml import etree

# Force UTF-8 output on Windows
sys.stdout.reconfigure(encoding="utf-8")
sys.stderr.reconfigure(encoding="utf-8")

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from pptx import Presentation
from pptx.util import Emu, Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn


def emu_to_inches(emu):
    if emu is None:
        return "N/A"
    return f"{emu / 914400:.2f}in"


def get_fill_info(fill):
    """Extract fill information from a shape or background."""
    try:
        if fill is None:
            return "None"
        ftype = fill.type
        if ftype is None:
            return "No fill"
        if str(ftype) == "SOLID (1)" or "SOLID" in str(ftype):
            try:
                rgb = fill.fore_color.rgb
                return f"Solid #{rgb}"
            except Exception:
                try:
                    theme_color = fill.fore_color.theme_color
                    brightness = getattr(fill.fore_color, 'brightness', None)
                    return f"Solid theme={theme_color} bright={brightness}"
                except Exception:
                    return "Solid (color unreadable)"
        return str(ftype)
    except Exception as e:
        return f"Error: {e}"


def inspect_placeholder(ph, indent="    "):
    """Detailed placeholder inspection."""
    pf = ph.placeholder_format
    print(f"{indent}Placeholder idx={pf.idx} type={pf.type}")
    print(f"{indent}  Position: left={emu_to_inches(ph.left)}, top={emu_to_inches(ph.top)}")
    print(f"{indent}  Size: width={emu_to_inches(ph.width)}, height={emu_to_inches(ph.height)}")

    if ph.has_text_frame:
        tf = ph.text_frame
        for para in tf.paragraphs:
            for run in para.runs:
                font = run.font
                fname = font.name or "(inherited)"
                fsize = f"{font.size}" if font.size else "(inherited)"
                fbold = font.bold if font.bold is not None else "(inherited)"
                fcolor = ""
                try:
                    if font.color and font.color.rgb:
                        fcolor = f" color=#{font.color.rgb}"
                except Exception:
                    pass
                print(f"{indent}  Font: {fname} size={fsize} bold={fbold}{fcolor}")
                if run.text.strip():
                    print(f"{indent}  Text: '{run.text[:60]}'")
                break
            break
    print()


def inspect_shape(shape, indent="    "):
    """Detailed shape inspection."""
    stype = shape.shape_type
    name = shape.name
    print(f"{indent}Shape: '{name}' type={stype}")
    print(f"{indent}  Position: left={emu_to_inches(shape.left)}, top={emu_to_inches(shape.top)}")
    print(f"{indent}  Size: width={emu_to_inches(shape.width)}, height={emu_to_inches(shape.height)}")

    try:
        fill = get_fill_info(shape.fill)
        print(f"{indent}  Fill: {fill}")
    except Exception:
        pass

    try:
        if shape.line and shape.line.color:
            lcolor = shape.line.color.rgb
            print(f"{indent}  Line: #{lcolor}")
    except Exception:
        pass
    print()


def main():
    template_path = Path("templates/Presentation6.pptx")
    prs = Presentation(str(template_path))

    print("=" * 80)
    print("TEMPLATE INSPECTION: Presentation6.pptx")
    print("=" * 80)
    print(f"\nSlide dimensions: width={emu_to_inches(prs.slide_width)}, height={emu_to_inches(prs.slide_height)}")

    # Theme colors - read from XML directly
    print("\n" + "=" * 80)
    print("THEME COLORS")
    print("=" * 80)
    try:
        master = prs.slide_masters[0]
        for rel in master.part.rels.values():
            if "theme" in str(rel.reltype):
                theme_part = rel.target_part
                theme_blob = theme_part.blob
                theme_xml = etree.fromstring(theme_blob)
                ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                clr_scheme = theme_xml.find('.//a:clrScheme', ns)
                if clr_scheme is not None:
                    print(f"  Color Scheme: '{clr_scheme.get('name')}'")
                    for child in clr_scheme:
                        tag = child.tag.split('}')[-1]
                        for sub in child:
                            val = sub.get('val') or sub.get('lastClr')
                            print(f"    {tag}: #{val}")
                font_scheme = theme_xml.find('.//a:fontScheme', ns)
                if font_scheme is not None:
                    print(f"\n  Font Scheme: '{font_scheme.get('name')}'")
                    major = font_scheme.find('.//a:majorFont', ns)
                    minor = font_scheme.find('.//a:minorFont', ns)
                    if major is not None:
                        latin = major.find('a:latin', ns)
                        if latin is not None:
                            print(f"    Major (Headings): {latin.get('typeface')}")
                        # Show complex script font too
                        cs = major.find('a:cs', ns)
                        if cs is not None:
                            print(f"    Major (CS/Arabic): {cs.get('typeface')}")
                    if minor is not None:
                        latin = minor.find('a:latin', ns)
                        if latin is not None:
                            print(f"    Minor (Body): {latin.get('typeface')}")
                        cs = minor.find('a:cs', ns)
                        if cs is not None:
                            print(f"    Minor (CS/Arabic): {cs.get('typeface')}")
                break
    except Exception as e:
        print(f"  Error reading theme: {e}")

    # Slide Masters
    print("\n" + "=" * 80)
    print("SLIDE MASTERS")
    print("=" * 80)
    for i, master in enumerate(prs.slide_masters):
        print(f"\n  Master {i}: {len(master.slide_layouts)} layouts")
        try:
            bg = master.background
            fill = get_fill_info(bg.fill)
            print(f"  Background: {fill}")
        except Exception as e:
            print(f"  Background: Error - {e}")

        for shape in master.shapes:
            if shape.has_text_frame and shape.text.strip():
                print(f"  Master shape: '{shape.name}' text='{shape.text[:50]}'")

    # Slide Layouts
    print("\n" + "=" * 80)
    print("SLIDE LAYOUTS (ALL)")
    print("=" * 80)

    for i, layout in enumerate(prs.slide_layouts):
        print(f"\n{'-' * 70}")
        print(f"LAYOUT {i}: '{layout.name}'")
        print(f"{'-' * 70}")

        try:
            bg = layout.background
            fill = get_fill_info(bg.fill)
            print(f"  Background: {fill}")
        except Exception:
            print(f"  Background: (inherits from master)")

        phs = list(layout.placeholders)
        if phs:
            print(f"\n  PLACEHOLDERS ({len(phs)}):")
            for ph in phs:
                inspect_placeholder(ph)
        else:
            print("  No placeholders")

        non_ph_shapes = [s for s in layout.shapes if not s.is_placeholder]
        if non_ph_shapes:
            print(f"  OTHER SHAPES ({len(non_ph_shapes)}):")
            for shape in non_ph_shapes:
                inspect_shape(shape)

    # Existing slides in the template
    print("\n" + "=" * 80)
    print("EXISTING SLIDES IN TEMPLATE")
    print("=" * 80)
    for i, slide in enumerate(prs.slides):
        layout = slide.slide_layout
        layout_idx = list(prs.slide_layouts).index(layout)
        print(f"\n  Slide {i+1}: layout='{layout.name}' (idx={layout_idx})")

        for shape in slide.shapes:
            if shape.is_placeholder:
                pf = shape.placeholder_format
                text = ""
                if shape.has_text_frame:
                    text = shape.text_frame.text[:80].replace('\n', ' | ')
                print(f"    PH idx={pf.idx} type={pf.type}: '{text}'")
            else:
                fill_info = ""
                try:
                    fill_info = get_fill_info(shape.fill)
                except Exception:
                    pass
                print(f"    Shape: '{shape.name}' type={shape.shape_type} fill={fill_info}")

    print("\n" + "=" * 80)
    print("INSPECTION COMPLETE")
    print("=" * 80)


if __name__ == "__main__":
    main()
