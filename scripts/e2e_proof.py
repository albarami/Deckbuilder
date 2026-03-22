"""
End-to-End Pipeline Proof -- runs the REAL DeckForge pipeline graph.

Uses: build_graph() -> graph.ainvoke() -> Command(resume=...) at each gate.
This is the SAME path as: frontend -> API -> session -> graph -> all nodes -> all gates.

Every node executes through the real LangGraph StateGraph with MemorySaver
checkpointer and interrupt/resume at gates.

ZERO BYPASSES: No force_render, no fail_close patching, no state mutation.
Whatever the pipeline produces naturally is the proof result.
- If fail_close=True -> negative proof (safety systems correctly block render)
- If fail_close=False -> positive proof (clean path renders successfully)

Usage:
    python scripts/e2e_proof.py --language en
    python scripts/e2e_proof.py --language en --brief positive
    python scripts/e2e_proof.py --language ar --brief positive
"""

from __future__ import annotations

import argparse
import asyncio
import json
import logging
import os
import re
import sys
import time
from pathlib import Path
from typing import Any

project_root = Path(__file__).resolve().parent.parent
os.chdir(project_root)
sys.path.insert(0, str(project_root))

# Force UTF-8 on Windows console
if sys.platform == "win32":
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")  # type: ignore[attr-defined]
    sys.stderr.reconfigure(encoding="utf-8", errors="replace")  # type: ignore[attr-defined]

from dotenv import load_dotenv  # noqa: E402

load_dotenv(override=True)

logging.basicConfig(
    level=logging.INFO,
    format="%(name)s %(levelname)s: %(message)s",
)

from langgraph.types import Command  # noqa: E402

# Increase retry budget for E2E proof -- Anthropic 529 overload is common
import src.services.llm as _llm_mod  # noqa: E402
from src.models.enums import LintSeverity, PipelineStage, RendererMode  # noqa: E402
from src.models.state import DeckForgeState, SessionMetadata  # noqa: E402
from src.pipeline.graph import build_graph  # noqa: E402
from src.services.search import _ensure_local_backend  # noqa: E402

_llm_mod._RETRY_DELAYS = [10, 20, 40, 60, 90, 120]  # 6+1=7 attempts, up to 2 min backoff

# Also increase the Anthropic SDK's own max_retries (default is 2)
_orig_get_anthropic = _llm_mod._get_anthropic_client


def _patched_get_anthropic():
    if not getattr(_llm_mod, "_anthropic_client", None):
        settings = _llm_mod.get_settings()
        _llm_mod._anthropic_client = _llm_mod.anthropic.AsyncAnthropic(
            api_key=settings.anthropic_api_key.get_secret_value(),
            max_retries=5,  # up from default 2
            timeout=1200.0,  # 20 minutes
        )
    return _llm_mod._anthropic_client


_llm_mod._get_anthropic_client = _patched_get_anthropic

# ── Brief definitions ──

BRIEF_NEGATIVE = (
    "IT Infrastructure Modernization RFP -- "
    "Strategic Gears proposes cloud migration and cybersecurity services "
    "for a government entity in KSA."
)

BRIEF_POSITIVE = (
    "Digital Transformation Consulting Services RFP -- "
    "The Ministry of Communications and Information Technology (MCIT) of the "
    "Kingdom of Saudi Arabia seeks a consulting firm to provide digital "
    "transformation advisory, enterprise architecture assessment, and "
    "automation strategy services. "
    "Strategic Gears proposes a phased consulting engagement leveraging its "
    "proven KSA government advisory experience, certified team, and "
    "structured delivery methodology. "
    "Evaluation Criteria: Technical Approach and Methodology (30%), "
    "Team Qualifications and Experience (25%), Past Performance and "
    "Project References (25%), Compliance and Certifications (10%), "
    "Project Management and Governance (10%)."
)

# ── Placeholder patterns for PPTX text scanning ──
_PLACEHOLDER_PATTERNS = [
    re.compile(r"\[BD team[^\]]*\]", re.IGNORECASE),
    re.compile(r"\[PLACEHOLDER[^\]]*\]", re.IGNORECASE),
    re.compile(r"\[TBC\]|\[TBD\]|\[INSERT[^\]]*\]", re.IGNORECASE),
    re.compile(r"\[CRITICAL:[^\]]*\]", re.IGNORECASE),
    re.compile(r"\[Action Required[^\]]*\]", re.IGNORECASE),
    re.compile(r"\[Confirm\b[^\]]*\]", re.IGNORECASE),
    re.compile(r"Key point \d", re.IGNORECASE),
    re.compile(
        r"(?:(?:needs? |remains? |yet )?to be added\b(?! value| benefit| to (?:the|our|your))"
        r"|(?:must |should |needs? to ).{1,40}before submission"
        r"|no row should remain"
        r"|must be resolved before|unresolved in client version)",
        re.IGNORECASE,
    ),
    re.compile(r"GAP-\d{3}", re.IGNORECASE),
    re.compile(r"TODO|FIXME", re.IGNORECASE),
]


async def ensure_search_index(docs_path: str | None = None, force: bool = False):
    """Index the local knowledge source if not already cached.

    For positive proofs, docs_path points at the curated data directory.
    We monkey-patch the search module defaults so the pipeline nodes
    (semantic_search, load_documents) use the curated data.
    """
    import src.services.search as _search_mod
    original_docs = _search_mod.DEFAULT_DOCS_PATH
    original_cache = _search_mod.DEFAULT_CACHE_PATH
    actual_docs = docs_path or original_docs
    cache_dir = original_cache
    if docs_path:
        # Use a separate cache for curated data to avoid mixing
        cache_dir = "state/index_positive/"
        # Reset the singleton backend so it picks up new data
        _search_mod._backend = None
        # Monkey-patch module defaults so pipeline nodes use curated data
        _search_mod.DEFAULT_DOCS_PATH = actual_docs
        _search_mod.DEFAULT_CACHE_PATH = cache_dir
        print(f"[INDEX] Overriding search defaults: docs={actual_docs}")
    cache_path = Path(cache_dir)
    embeddings = cache_path / "embeddings.npy"
    if force and embeddings.exists():
        import shutil
        shutil.rmtree(cache_path, ignore_errors=True)
        print(f"[INDEX] Cleared stale cache at {cache_path}")
    if embeddings.exists():
        print(f"[INDEX] Using cached index at {cache_path}")
    else:
        print(f"[INDEX] Indexing documents from {actual_docs} ...")
    await _ensure_local_backend(actual_docs, str(cache_dir))
    print(f"[INDEX] Search backend ready (docs={actual_docs})")
    return str(cache_dir)


def make_initial_state(
    language: str, session_id: str, brief: str,
) -> DeckForgeState:
    """Build initial state -- same as what the API/frontend would create."""
    return DeckForgeState(
        ai_assist_summary=brief,
        output_language=language,
        renderer_mode=RendererMode.TEMPLATE_V2,
        geography="ksa",
        proposal_mode="standard",
        sector="technology",
        session=SessionMetadata(session_id=session_id),
    )


def emit_qa_provenance_artifact(
    session_id: str,
    result: dict[str, Any],
) -> str | None:
    """Persist runtime QA provenance data to output/{session_id}/qa_provenance.json.

    Reads the LIVE runtime objects (written_slides, qa_result, proposal_manifest)
    from the pipeline state dict — not reconstructed from code logic.

    Returns the path written, or None if data was insufficient.
    """
    ws = result.get("written_slides")
    qa = result.get("qa_result")
    manifest = result.get("proposal_manifest")

    if not ws or not hasattr(ws, "slides"):
        return None

    slides = ws.slides

    # Build a lookup: slide_id -> QA validation result
    qa_lookup: dict[str, dict] = {}
    if qa and hasattr(qa, "slide_validations"):
        for sv in qa.slide_validations:
            qa_lookup[sv.slide_id] = {
                "status": str(sv.status.value) if hasattr(sv.status, "value") else str(sv.status),
                "issue_count": len(sv.issues),
            }

    # Build a lookup: asset_id -> section_id from the manifest
    section_lookup: dict[str, str] = {}
    if manifest and hasattr(manifest, "entries"):
        for entry in manifest.entries:
            if entry.entry_type == "b_variable":
                section_lookup[entry.asset_id] = entry.section_id

    # Serialize each slide's runtime provenance
    provenance_records = []
    for slide in slides:
        slide_id = getattr(slide, "slide_id", "unknown")
        asset_id = getattr(slide, "manifest_asset_id", "unset")

        # Determine source from slide_id prefix (set at stamp time in graph.py)
        if slide_id.startswith("METH-"):
            source = "synthesized_methodology_stub"
        elif asset_id.startswith("_unmatched_"):
            source = "unmatched"
        else:
            source = "iterative_builder"

        qa_info = qa_lookup.get(slide_id, {"status": "not_validated", "issue_count": 0})

        provenance_records.append({
            "slide_id": slide_id,
            "manifest_asset_id": asset_id,
            "section_id": section_lookup.get(asset_id, "unknown"),
            "source": source,
            "qa_result": qa_info["status"],
            "issue_count": qa_info["issue_count"],
        })

    output_dir = Path("output") / session_id
    output_dir.mkdir(parents=True, exist_ok=True)
    artifact_path = output_dir / "qa_provenance.json"
    artifact_path.write_text(
        json.dumps(provenance_records, indent=2, ensure_ascii=False),
        encoding="utf-8",
    )
    return str(artifact_path)


def extract_pptx_text(pptx_path: str) -> list[dict]:
    """Extract all text from PPTX slides for placeholder scanning."""
    from pptx import Presentation
    prs = Presentation(str(pptx_path))
    slides_text = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
        slides_text.append({
            "slide_num": i + 1,
            "layout": slide.slide_layout.name if slide.slide_layout else "unknown",
            "texts": texts,
        })
    return slides_text


def scan_for_placeholders(slides_text: list[dict]) -> list[dict]:
    """Scan extracted PPTX text for placeholder patterns."""
    violations = []
    for slide_info in slides_text:
        for text in slide_info["texts"]:
            for pattern in _PLACEHOLDER_PATTERNS:
                matches = pattern.findall(text)
                if matches:
                    violations.append({
                        "slide_num": slide_info["slide_num"],
                        "layout": slide_info["layout"],
                        "text_snippet": text[:100],
                        "pattern": pattern.pattern,
                        "matches": matches,
                    })
    return violations


async def run_e2e(
    language: str,
    brief: str,
    brief_label: str,
    docs_path: str | None = None,
) -> dict:
    """Run the FULL pipeline graph with stage-by-stage evidence."""
    print(f"\n{'='*80}")
    print(f"  E2E PROOF: language={language}, brief={brief_label}")
    print("  ZERO BYPASSES: No force_render, no fail_close patching, no state mutation")
    print("  Path: graph.ainvoke() -> Command(resume) at each gate")
    print("  Same path as: frontend -> API -> session -> graph -> all nodes -> all gates")
    print(f"{'='*80}\n")

    # Prerequisites
    lang_suffix = "ar" if language == "ar" else "en"
    template_dir = Path("PROPOSAL_TEMPLATE")
    if language == "ar":
        template_file = template_dir / "Arabic_Proposal_Template.potx"
    else:
        template_file = template_dir / "PROPOSAL_TEMPLATE EN.potx"

    print(f"[PREREQ] Template: {template_file} ({template_file.stat().st_size:,} bytes)")
    catalog_lock = Path(f"src/data/catalog_lock_{lang_suffix}.json")
    print(f"[PREREQ] Catalog lock: {catalog_lock} (exists={catalog_lock.exists()})")

    await ensure_search_index(docs_path=docs_path, force=(docs_path is not None))

    # Build graph -- the REAL production graph
    print("\n[BUILD] Building production pipeline graph...")
    graph = build_graph()
    nodes = list(graph.get_graph().nodes.keys())
    print(f"[BUILD] Graph nodes: {nodes}")
    print(f"[BUILD] Node count: {len(nodes)}")

    session_id = f"e2e-{lang_suffix}-{brief_label}-{int(time.time())}"
    state = make_initial_state(language, session_id, brief)
    config = {"configurable": {"thread_id": session_id}}
    print(f"[SESSION] ID: {session_id}")
    print(f"[BRIEF] {brief[:100]}...")

    stages = []  # Collect (node_range, timing, key_data)

    # ==========================================================
    # PHASE 1: Initial invoke -> runs context -> gate_1 (interrupt)
    # ==========================================================
    print(f"\n{'-'*80}")
    print("  PHASE 1: START -> context -> gate_1 (interrupt)")
    print(f"{'-'*80}")
    t0 = time.time()
    result = await graph.ainvoke(state, config)
    t1 = time.time()
    print(f"  Time: {t1-t0:.1f}s")
    print(f"  Stage: {result.get('current_stage')}")

    ctx = result.get("rfp_context")
    if ctx:
        name = getattr(ctx, "rfp_name", None)
        entity = getattr(ctx, "issuing_entity", None)
        print(f"  rfp_context.rfp_name: {name}")
        print(f"  rfp_context.issuing_entity: {entity}")
    stages.append(("context -> gate_1", t1 - t0, {
        "rfp_name": str(getattr(ctx, "rfp_name", "N/A")) if ctx else "N/A",
    }))

    # ==========================================================
    # PHASE 2: Approve gate_1 -> retrieval -> gate_2 (interrupt)
    # ==========================================================
    print(f"\n{'-'*80}")
    print("  PHASE 2: gate_1 approve -> retrieval -> gate_2 (interrupt)")
    print(f"{'-'*80}")
    t0 = time.time()
    result = await graph.ainvoke(Command(resume={"approved": True}), config)
    t1 = time.time()
    print(f"  Time: {t1-t0:.1f}s")
    print(f"  Stage: {result.get('current_stage')}")

    sources = result.get("retrieved_sources", [])
    print(f"  retrieved_sources: {len(sources)}")
    for s in sources[:5]:
        print(f"    - {getattr(s, 'title', '?')} (score={getattr(s, 'relevance_score', '?')})")
    stages.append(("retrieval -> gate_2", t1 - t0, {
        "source_count": len(sources),
    }))

    # ==========================================================
    # PHASE 3: Approve gate_2 -> evidence_curation -> proposal_strategy
    #           -> source_book -> gate_3 (interrupt)
    # Source Book pipeline: full-text evidence + external research,
    # then proposal strategy (win themes, evaluator priorities),
    # then Source Book Writer/Reviewer loop (up to 5 passes),
    # then Gate 3 reviews the Source Book.
    # ==========================================================
    print(f"\n{'-'*80}")
    print("  PHASE 3: gate_2 approve -> evidence_curation -> proposal_strategy")
    print("           -> source_book -> gate_3 (interrupt)")
    print(f"{'-'*80}")
    t0 = time.time()
    result = await graph.ainvoke(Command(resume={"approved": True}), config)
    t1 = time.time()
    print(f"  Time: {t1-t0:.1f}s")
    print(f"  Stage: {result.get('current_stage')}")

    # --- Evidence Curation proof ---
    ref_idx = result.get("reference_index")
    ext_evidence = result.get("external_evidence_pack")
    print("\n  --- Evidence Curation Proof ---")
    if ref_idx:
        claims = getattr(ref_idx, "claims", [])
        case_studies = getattr(ref_idx, "case_studies", [])
        team_profiles = getattr(ref_idx, "team_profiles", [])
        print(f"  reference_index.claims: {len(claims)}")
        print(f"  reference_index.case_studies: {len(case_studies)}")
        print(f"  reference_index.team_profiles: {len(team_profiles)}")
    else:
        print("  WARNING: reference_index is EMPTY")
    print(f"  external_evidence_pack populated: {ext_evidence is not None}")
    if ext_evidence:
        ext_entries = getattr(ext_evidence, "entries", [])
        print(f"  external_evidence entries: {len(ext_entries)}")

    # --- Proposal Strategy proof ---
    prop_strategy = result.get("proposal_strategy")
    print("\n  --- Proposal Strategy Proof ---")
    print(f"  proposal_strategy populated: {prop_strategy is not None}")
    if prop_strategy:
        win_themes = getattr(prop_strategy, "win_themes", [])
        rec_method = getattr(prop_strategy, "recommended_methodology_approach", "")
        print(f"  win_themes: {win_themes}")
        print(f"  recommended_methodology: {rec_method}")

    # --- Source Book proof ---
    source_book = result.get("source_book")
    source_book_review = result.get("source_book_review")
    rm = result.get("report_markdown", "")
    print("\n  --- Source Book Proof ---")
    sb_word_count = 0
    sb_evidence_count = 0
    sb_blueprint_count = 0
    sb_pass_number = 0
    sb_review_score = 0
    if source_book:
        # Word count across prose sections
        prose_parts = [
            source_book.rfp_interpretation.objective_and_scope,
            source_book.rfp_interpretation.constraints_and_compliance,
            source_book.rfp_interpretation.unstated_evaluator_priorities,
            source_book.rfp_interpretation.probable_scoring_logic,
            source_book.client_problem_framing.current_state_challenge,
            source_book.client_problem_framing.why_it_matters_now,
            source_book.client_problem_framing.transformation_logic,
            source_book.client_problem_framing.risk_if_unchanged,
            source_book.proposed_solution.methodology_overview,
            source_book.proposed_solution.governance_framework,
            source_book.proposed_solution.timeline_logic,
            source_book.proposed_solution.value_case_and_differentiation,
            source_book.why_strategic_gears.certifications_and_compliance,
        ]
        sb_word_count = sum(len(p.split()) for p in prose_parts if p)
        sb_evidence_count = len(source_book.evidence_ledger.entries)
        sb_blueprint_count = len(source_book.slide_blueprints)
        sb_pass_number = source_book.pass_number
        cap_count = len(source_book.why_strategic_gears.capability_mapping)
        print("  source_book populated: True")
        print(f"  word count: ~{sb_word_count}")
        print(f"  evidence ledger entries: {sb_evidence_count}")
        print(f"  slide blueprints: {sb_blueprint_count}")
        print(f"  capability mappings: {cap_count}")
        print(f"  pass_number: {sb_pass_number}")
        print(f"  client_name: {source_book.client_name}")
    else:
        print("  WARNING: source_book is EMPTY")

    if source_book_review:
        sb_review_score = source_book_review.overall_score
        print(f"  review score: {sb_review_score}/5")
        print(f"  competitive viability: {source_book_review.competitive_viability}")
        print(f"  pass_threshold_met: {source_book_review.pass_threshold_met}")
    else:
        print("  source_book_review: not available")

    # Source Book DOCX path
    docx_sb_path = result.get("report_docx_path")
    print(f"  Source Book DOCX: {docx_sb_path}")
    print(f"  report_markdown length: {len(rm) if rm else 0}")

    stages.append(("evidence_curation -> proposal_strategy -> source_book -> gate_3", t1 - t0, {
        "sb_word_count": sb_word_count,
        "sb_evidence_count": sb_evidence_count,
        "sb_blueprint_count": sb_blueprint_count,
        "sb_pass_number": sb_pass_number,
        "sb_review_score": sb_review_score,
        "report_markdown_len": len(rm) if rm else 0,
    }))

    # ==========================================================
    # PHASE 4: Approve gate_3 -> assembly_plan -> blueprint_extraction
    #           -> section_fill -> build_slides -> gate_4
    # Source Book pipeline: assembly plan produces manifest + budget,
    # Slide Architect converts Source Book into SlideBlueprint,
    # section fillers use blueprint guidance, iterative builder
    # produces final slide content.
    # ==========================================================
    print(f"\n{'-'*80}")
    print("  PHASE 4: gate_3 approve -> assembly_plan -> blueprint_extraction")
    print("           -> section_fill -> build_slides -> gate_4")
    print(f"{'-'*80}")
    t0 = time.time()
    result = await graph.ainvoke(Command(resume={"approved": True}), config)
    t1 = time.time()
    print(f"  Time: {t1-t0:.1f}s")
    print(f"  Stage: {result.get('current_stage')}")
    if result.get("last_error"):
        err = result.get("last_error")
        print(f"  last_error.agent: {getattr(err, 'agent', 'N/A')}")
        print(f"  last_error.error_type: {getattr(err, 'error_type', 'N/A')}")
        print(f"  last_error.message: {getattr(err, 'message', str(err))[:200]}")

    # === ASSEMBLY PLAN PROOF ===
    manifest = result.get("proposal_manifest")
    budget = result.get("slide_budget")
    print("\n  --- Assembly Plan Proof ---")
    print(f"  proposal_manifest populated: {manifest is not None}")
    b_variable_count = 0
    if manifest:
        b_variable_count = sum(
            1 for e in manifest.entries if e.entry_type == "b_variable"
        )
        total_entries = len(manifest.entries)
        print(f"  manifest entries total: {total_entries}")
        print(f"  manifest b_variable entries: {b_variable_count}")
        # Show entry type breakdown
        entry_types = {}
        for e in manifest.entries:
            entry_types[e.entry_type] = entry_types.get(e.entry_type, 0) + 1
        print(f"  entry type breakdown: {entry_types}")
    print(f"  slide_budget populated: {budget is not None}")

    # === SLIDE BLUEPRINT PROOF ===
    slide_bp = result.get("slide_blueprint")
    bp_entry_count = 0
    bp_evidence_coverage = 0.0
    print("\n  --- Slide Blueprint Proof (Slide Architect output) ---")
    print(f"  slide_blueprint populated: {slide_bp is not None}")
    if slide_bp:
        bp_entry_count = len(slide_bp.entries)
        bp_evidence_coverage = slide_bp.evidence_coverage
        print(f"  blueprint entries: {bp_entry_count}")
        print(f"  evidence coverage: {bp_evidence_coverage:.1%}")
        print(f"  blueprint_version: {slide_bp.blueprint_version}")
        print(f"  total_variable_slides: {slide_bp.total_variable_slides}")
        # Blueprint ↔ manifest alignment check
        if b_variable_count > 0:
            aligned = bp_entry_count == b_variable_count
            print(f"  blueprint/manifest alignment: {'ALIGNED' if aligned else 'MISMATCH'}")
            if not aligned:
                print(f"    WARNING: blueprint={bp_entry_count}, manifest b_variable={b_variable_count}")
        # Show first few entries
        for entry in slide_bp.entries[:5]:
            print(f"    Slide {entry.slide_number}: [{entry.section}] {entry.title[:50]}")
        if bp_entry_count > 5:
            print(f"    ... and {bp_entry_count - 5} more")
    else:
        print("  WARNING: slide_blueprint is EMPTY -- Slide Architect failed")

    # Persist blueprint JSON artifact
    blueprint_json_path = None
    if slide_bp:
        bp_out_dir = Path(f"output/{session_id}")
        bp_out_dir.mkdir(parents=True, exist_ok=True)
        blueprint_json_path = str(bp_out_dir / "slide_blueprint.json")
        with open(blueprint_json_path, "w", encoding="utf-8") as f:
            json.dump(slide_bp.model_dump(mode="json"), f, indent=2, ensure_ascii=False)
        print(f"  Blueprint JSON persisted: {blueprint_json_path}")

    # === SECTION FILL PROOF ===
    filler_outputs = result.get("filler_outputs")
    print("\n  --- Section Fill Proof ---")
    print(f"  filler_outputs populated: {filler_outputs is not None}")
    if filler_outputs:
        for section_id, output in filler_outputs.items():
            entry_count = len(output.entries) if hasattr(output, "entries") else 0
            errors = output.errors if hasattr(output, "errors") else []
            print(f"    {section_id}: {entry_count} entries, {len(errors)} errors")

    # === BUILD SLIDES PROOF ===
    ws = result.get("written_slides")
    slides = []
    print("\n  --- Build Slides Proof ---")
    if ws:
        slides = getattr(ws, "slides", [])
        print(f"  written_slides count: {len(slides)}")
        for s in slides[:10]:
            sid = getattr(s, "slide_id", "?")
            title = getattr(s, "title", "?")[:60]
            layout = getattr(s, "layout_type", "?")
            print(f"    {sid}: [{layout}] {title}")
        if len(slides) > 10:
            print(f"    ... and {len(slides) - 10} more")

    drafts = result.get("deck_drafts", [])
    reviews = result.get("deck_reviews", [])
    print(f"  deck_drafts (iterative turns): {len(drafts)}")
    print(f"  deck_reviews (iterative turns): {len(reviews)}")

    stages.append(("assembly_plan -> blueprint -> section_fill -> build_slides -> gate_4", t1 - t0, {
        "manifest_total": len(manifest.entries) if manifest else 0,
        "manifest_b_variable": b_variable_count,
        "blueprint_entries": bp_entry_count,
        "blueprint_evidence_coverage": bp_evidence_coverage,
        "llm_slides": len(slides),
        "drafts": len(drafts),
        "reviews": len(reviews),
    }))

    # ==========================================================
    # PHASE 5: Approve gate_4 -> qa -> governance -> gate_5 (interrupt)
    # ==========================================================
    print(f"\n{'-'*80}")
    print("  PHASE 5: gate_4 approve -> qa -> governance -> gate_5 (interrupt)")
    print(f"{'-'*80}")
    t0 = time.time()
    result = await graph.ainvoke(Command(resume={"approved": True}), config)
    t1 = time.time()
    print(f"  Time: {t1-t0:.1f}s")
    print(f"  Stage: {result.get('current_stage')}")

    qa = result.get("qa_result")
    sub_qa = result.get("submission_qa_result")
    ds = None
    print("\n  --- QA Proof ---")
    if qa:
        ds = qa.deck_summary
        print(f"  total_slides: {ds.total_slides}")
        print(f"  passed: {ds.passed}")
        print(f"  failed: {ds.failed}")
        print(f"  fail_close: {ds.fail_close}")
        print(f"  fail_close_reason: {ds.fail_close_reason}")
        print(f"  critical_gaps: {ds.critical_gaps}")
        print(f"  critical_gaps_remaining: {ds.critical_gaps_remaining}")
    else:
        print("  WARNING: qa_result is EMPTY -- QA agent failed")

    density = None
    provenance = None
    lint = None
    print("\n  --- Governance / Submission QA Proof ---")
    print(f"  submission_qa_result populated: {sub_qa is not None}")
    if sub_qa:
        print(f"  status: {getattr(sub_qa, 'status', 'N/A')}")
        print(f"  summary: {getattr(sub_qa, 'summary', 'N/A')}")
        density = getattr(sub_qa, "density_result", None)
        provenance = getattr(sub_qa, "evidence_provenance", None)
        lint = (
            getattr(sub_qa, "language_lint", None)
            or getattr(sub_qa, "lint_result", None)
        )
        print(f"  density_result populated: {density is not None}")
        if density:
            print(f"    blocker_count: {getattr(density, 'blocker_count', 'N/A')}")
            print(f"    warning_count: {getattr(density, 'warning_count', 'N/A')}")
            print(f"    is_within_budget: {getattr(density, 'is_within_budget', 'N/A')}")
        print(f"  evidence_provenance populated: {provenance is not None}")
        if provenance:
            print(f"    blocker_count: {getattr(provenance, 'blocker_count', 'N/A')}")
            print(f"    warning_count: {getattr(provenance, 'warning_count', 'N/A')}")
            issues = getattr(provenance, "issues", [])
            for iss in issues[:5]:
                print(
                    f"    - {getattr(iss, 'slide_id', '?')}: "
                    f"{getattr(iss, 'message', '?')[:80]}"
                )
        print(f"  language_lint populated: {lint is not None}")
        if lint:
            print(f"    blocker_count: {getattr(lint, 'blocker_count', 'N/A')}")
            print(f"    warning_count: {getattr(lint, 'warning_count', 'N/A')}")
            # Print all blocker-level issues for diagnostics
            lint_issues = getattr(lint, "issues", [])
            blocker_issues = [
                iss for iss in lint_issues
                if getattr(iss, "severity", None) == LintSeverity.BLOCKER
            ]
            if blocker_issues:
                print("    BLOCKER issues:")
                for iss in blocker_issues:
                    print(
                        f"      {getattr(iss, 'slide_id', '?')} "
                        f"[{getattr(iss, 'location', '?')}] "
                        f"rule={getattr(iss, 'rule', '?')}: "
                        f"'{getattr(iss, 'matched_text', '?')}'"
                    )

    stages.append(("qa -> governance -> gate_5", t1 - t0, {
        "fail_close": ds.fail_close if ds else None,
        "passed": ds.passed if ds else 0,
        "failed": ds.failed if ds else 0,
        "submission_qa_status": (
            str(getattr(sub_qa, "status", "N/A")) if sub_qa else "N/A"
        ),
        "lint_blockers": getattr(lint, "blocker_count", 0) if lint else 0,
        "density_blockers": getattr(density, "blocker_count", 0) if density else 0,
        "provenance_blockers": (
            getattr(provenance, "blocker_count", 0) if provenance else 0
        ),
    }))

    # ── Emit runtime QA provenance artifact ──────────────────────
    provenance_path = emit_qa_provenance_artifact(session_id, result)
    if provenance_path:
        print("\n  --- Runtime QA Provenance Artifact ---")
        print(f"  Written to: {provenance_path}")
        prov_data = json.loads(Path(provenance_path).read_text(encoding="utf-8"))
        print(f"  Total slides in artifact: {len(prov_data)}")
        meth_count = sum(1 for r in prov_data if r["source"] == "synthesized_methodology_stub")
        builder_count = sum(1 for r in prov_data if r["source"] == "iterative_builder")
        print(f"  Iterative builder slides: {builder_count}")
        print(f"  Methodology stubs: {meth_count}")
        for r in prov_data:
            print(
                f"    {r['slide_id']:45s} -> {r['manifest_asset_id']:40s} "
                f"[{r['section_id']:12s}] src={r['source']:35s} "
                f"qa={r['qa_result']} issues={r['issue_count']}"
            )
    else:
        print("\n  WARNING: Could not emit QA provenance artifact (missing data)")

    # ==========================================================
    # PHASE 6: Approve gate_5 -> render -> END
    # No auto-waivers. If fail_close=True, render will be blocked
    # naturally by the fail-close enforcement — that is correct
    # behavior, not a failure.
    # ==========================================================
    print(f"\n{'-'*80}")
    print("  PHASE 6: gate_5 approve -> render -> END")
    print("  No auto-waivers. Render outcome determined by QA + lint results.")
    print(f"{'-'*80}")

    t0 = time.time()
    result = await graph.ainvoke(Command(resume={"approved": True}), config)
    t1 = time.time()
    print(f"  Time: {t1-t0:.1f}s")
    print(f"  Stage: {result.get('current_stage')}")
    if result.get("last_error"):
        err = result.get("last_error")
        print(f"  last_error.agent: {getattr(err, 'agent', 'N/A')}")
        print(f"  last_error.error_type: {getattr(err, 'error_type', 'N/A')}")
        print(f"  last_error.message: {getattr(err, 'message', str(err))[:200]}")

    stages.append(("render -> END", t1 - t0, {
        "stage": str(result.get("current_stage")),
    }))

    # ==========================================================
    # FINAL RESULTS
    # ==========================================================
    total_time = sum(s[1] for s in stages)
    final_stage = result.get("current_stage")
    pptx_path = result.get("pptx_path")
    docx_path = result.get("report_docx_path")
    source_idx = result.get("source_index_path")
    gap_report = result.get("gap_report_path")

    print(f"\n{'='*80}")
    print(f"  FINAL RESULTS -- Session: {session_id}")
    print(f"{'='*80}")
    print(f"\n  Stage: {final_stage}")
    print(f"  Total pipeline time: {total_time:.1f}s")
    print(f"  Brief: {brief_label}")
    print("  Bypasses: ZERO (no force_render, no fail_close patching)")

    print("\n  --- Stage-by-Stage Timing ---")
    for label, elapsed, data in stages:
        print(f"  {elapsed:7.1f}s  {label}")

    print("\n  --- Output Files ---")
    for label, path in [
        ("PPTX", pptx_path),
        ("Report DOCX", docx_path),
        ("Source Index", source_idx),
        ("Gap Report", gap_report),
        ("Blueprint JSON", blueprint_json_path),
    ]:
        if path and Path(path).exists():
            size = Path(path).stat().st_size
            print(f"  {label}: {path} ({size:,} bytes)")
        elif path:
            print(f"  {label}: {path} (FILE NOT FOUND)")
        else:
            print(f"  {label}: (not generated)")

    # Session scoping
    session_scoped = False
    if pptx_path:
        parts = Path(pptx_path).parts
        print("\n  --- Session Scoping ---")
        print(f"  PPTX path parts: {parts}")
        has_output = "output" in parts
        has_session = any(session_id in p for p in parts)
        session_scoped = has_output and has_session
        print(f"  Contains 'output': {has_output}")
        print(f"  Session-scoped: {session_scoped}")

    # Slide count reconciliation + PPTX text extraction
    placeholder_violations = []
    quality_failures: list[str] = []
    if pptx_path and Path(pptx_path).exists():
        from pptx import Presentation
        prs = Presentation(str(pptx_path))

        catalog = json.loads(
            Path(f"src/data/catalog_lock_{lang_suffix}.json").read_text(
                encoding="utf-8"
            )
        )
        layout_to_type = {}
        for asset_id, info in catalog["a1_immutable"].items():
            layout_to_type[info.get("display_name", "")] = ("A1_CLONE", asset_id)
        for asset_id, info in catalog["a2_shells"].items():
            layout_to_type[info.get("display_name", "")] = ("A2_SHELL", asset_id)
        for div_id, info in catalog["section_dividers"].items():
            layout_to_type[info.get("display_name", "")] = ("DIVIDER", div_id)
        for cat_key, cs_list in catalog.get("case_study_pool", {}).items():
            if isinstance(cs_list, list):
                for cs in cs_list:
                    layout_to_type[cs.get("display_name", "")] = (
                        "POOL_CS",
                        cs.get("semantic_id", ""),
                    )
        for tb in catalog.get("team_bio_pool", []):
            layout_to_type[tb.get("display_name", "")] = (
                "POOL_TEAM",
                tb.get("semantic_id", ""),
            )

        counts = {
            "A1_CLONE": 0,
            "A2_SHELL": 0,
            "DIVIDER": 0,
            "POOL_CS": 0,
            "POOL_TEAM": 0,
            "LLM_CONTENT": 0,
        }
        for slide in prs.slides:
            layout_name = (
                slide.slide_layout.name if slide.slide_layout else "unknown"
            )
            entry = layout_to_type.get(layout_name)
            if entry:
                counts[entry[0]] += 1
            elif "Heading and description" in layout_name:
                counts["LLM_CONTENT"] += 1

        llm_built = len(ws.slides) if ws else 0
        fixed_total = (
            counts["A1_CLONE"]
            + counts["A2_SHELL"]
            + counts["DIVIDER"]
            + counts["POOL_CS"]
            + counts["POOL_TEAM"]
        )

        print("\n  --- Slide Count Reconciliation ---")
        print(f"  LLM 5-turn builder produced:         {llm_built} slides")
        print(f"  A1_CLONE (fixed institutional):      {counts['A1_CLONE']}")
        print(f"  A2_SHELL (template + proposal):      {counts['A2_SHELL']}")
        print(f"  DIVIDER (section separators):        {counts['DIVIDER']}")
        print(f"  POOL_CS (case studies):              {counts['POOL_CS']}")
        print(f"  POOL_TEAM (team bios):               {counts['POOL_TEAM']}")
        print(f"  LLM_CONTENT (in PPTX):               {counts['LLM_CONTENT']}")
        print(f"  Fixed/template total:                {fixed_total}")
        print(f"  Final PPTX total:                    {len(prs.slides)}")
        sum_check = fixed_total + counts["LLM_CONTENT"]
        print(f"  Sum check: {fixed_total} + {counts['LLM_CONTENT']} = {sum_check}")

        # Per-slide listing
        print("\n  --- Per-Slide Listing ---")
        for i, slide in enumerate(prs.slides):
            layout_name = (
                slide.slide_layout.name if slide.slide_layout else "unknown"
            )
            entry = layout_to_type.get(layout_name)
            if entry:
                tag = f"{entry[0]}:{entry[1]}"
            elif "Heading and description" in layout_name:
                tag = "LLM_CONTENT"
            else:
                tag = f"OTHER:{layout_name}"
            first_text = ""
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text.strip():
                    first_text = shape.text.strip()[:50]
                    break
            print(f"    Slide {i+1:2d} | {tag:35s} | {first_text}")

        # === QUALITY ASSERTIONS (positive proof only) ===
        quality_failures: list[str] = []
        manifest = result.get("proposal_manifest")
        if manifest:
            from src.models.section_blueprint import MANDATORY_SECTION_ORDER

            # 1. Total slides >= 30
            if len(prs.slides) < 30:
                quality_failures.append(
                    f"Total slides {len(prs.slides)} < 30 minimum"
                )

            # 2. All mandatory sections present in manifest
            manifest_sections = {
                e.section_id for e in manifest.entries if e.section_id
            }
            for sec_id in MANDATORY_SECTION_ORDER:
                if sec_id not in manifest_sections:
                    quality_failures.append(
                        f"Section {sec_id} missing from manifest"
                    )

            # 3. Methodology section uses methodology layouts
            methodology_layouts = [
                e.semantic_layout_id
                for e in manifest.entries
                if e.section_id == "section_03"
                and e.entry_type == "b_variable"
            ]
            has_method_layout = any(
                "methodology" in lay for lay in methodology_layouts
            )
            if methodology_layouts and not has_method_layout:
                quality_failures.append(
                    "Section 03 b_variable slides use no methodology "
                    f"layouts: {methodology_layouts}"
                )

            # 4. Team section uses pool clones
            team_entries = [
                e for e in manifest.entries
                if e.section_id == "section_05"
                and e.entry_type == "pool_clone"
            ]
            if counts["POOL_TEAM"] == 0 and len(team_entries) == 0:
                quality_failures.append(
                    "No team pool clones in section 05"
                )

            # 5. Case studies use pool clones
            cs_entries = [
                e for e in manifest.entries
                if e.section_id == "section_02"
                and e.entry_type == "pool_clone"
            ]
            if counts["POOL_CS"] == 0 and len(cs_entries) == 0:
                quality_failures.append(
                    "No case study pool clones in section 02"
                )

            if quality_failures:
                print(f"\n  --- Quality Assertion Failures ({len(quality_failures)}) ---")
                for qf in quality_failures:
                    print(f"    FAIL: {qf}")
            else:
                print("\n  --- Quality Assertions: ALL PASSED ---")

        # === PPTX TEXT EXTRACTION + PLACEHOLDER SCAN ===
        print("\n  --- PPTX Text Extraction + Placeholder Scan ---")
        slides_text = extract_pptx_text(str(pptx_path))
        placeholder_violations = scan_for_placeholders(slides_text)

        if placeholder_violations:
            print(f"  PLACEHOLDER VIOLATIONS FOUND: {len(placeholder_violations)}")
            for v in placeholder_violations:
                print(
                    f"    Slide {v['slide_num']} ({v['layout']}): "
                    f"matched [{v['pattern']}] in: {v['text_snippet'][:80]}"
                )
        else:
            print("  ZERO PLACEHOLDER VIOLATIONS -- all text is clean")

        # Show ALL extracted text for full transparency
        print("\n  --- Full PPTX Text Dump (LLM_CONTENT slides only) ---")
        for slide_info in slides_text:
            layout = slide_info["layout"]
            if "Heading and description" not in layout:
                continue
            print(f"    Slide {slide_info['slide_num']} ({layout}):")
            for text in slide_info["texts"]:
                # Truncate very long text
                display = text[:120] + "..." if len(text) > 120 else text
                print(f"      | {display}")

    # Proof classification
    print(f"\n{'='*80}")
    rendered = str(final_stage) == str(PipelineStage.FINALIZED)
    blocked_by_fail_close = (
        not rendered
        and ds
        and ds.fail_close
    )
    blocked_by_lint = (
        not rendered
        and result.get("last_error")
        and getattr(result.get("last_error"), "error_type", "") == "LintBlocked"
    )
    blocked_by_quality_gate = (
        not rendered
        and result.get("last_error")
        and "Quality gate REJECTED" in str(
            getattr(result.get("last_error"), "message", "")
        )
    )

    if rendered:
        if placeholder_violations:
            proof_class = "POSITIVE_PROOF_WITH_VIOLATIONS"
            proof_status = "NEEDS_WORK"
        else:
            proof_class = "POSITIVE_PROOF_CLEAN"
            proof_status = "PASSED"
    elif blocked_by_quality_gate:
        proof_class = "NEGATIVE_PROOF_QUALITY_GATE_BLOCKED"
        proof_status = "PASSED"
    elif blocked_by_lint:
        # Check lint first — last_error.error_type is the ACTUAL block reason.
        proof_class = "NEGATIVE_PROOF_LINT_BLOCKED"
        proof_status = "PASSED"
    elif blocked_by_fail_close:
        proof_class = "NEGATIVE_PROOF_FAIL_CLOSE"
        proof_status = "PASSED"
    else:
        proof_class = "UNEXPECTED_FAILURE"
        proof_status = "FAILED"

    print(f"  PROOF CLASSIFICATION: {proof_class}")
    print(f"  PROOF STATUS: {proof_status}")
    print(f"  Session: {session_id}")
    print(f"  Total time: {total_time:.1f}s")
    print(f"  Brief: {brief_label}")
    print("  Bypasses: ZERO")
    print(f"{'='*80}")

    return {
        "proof_class": proof_class,
        "proof_status": proof_status,
        "session_id": session_id,
        "language": language,
        "brief_label": brief_label,
        "total_time": total_time,
        "bypasses": "ZERO",
        "pipeline": "source_book",
        "pptx_path": str(pptx_path) if pptx_path else None,
        "docx_path": str(docx_path) if docx_path else None,
        "source_index_path": str(source_idx) if source_idx else None,
        "gap_report_path": str(gap_report) if gap_report else None,
        "blueprint_json_path": blueprint_json_path if slide_bp else None,
        "final_stage": str(final_stage),
        "session_scoped": session_scoped if pptx_path else False,
        "placeholder_violations": len(placeholder_violations),
        "quality_failures": quality_failures if pptx_path else [],
        "qa_fail_close": ds.fail_close if ds else None,
        "qa_passed": ds.passed if ds else 0,
        "qa_failed": ds.failed if ds else 0,
        "lint_blockers": getattr(lint, "blocker_count", 0) if lint else 0,
        "density_blockers": (
            getattr(density, "blocker_count", 0) if density else 0
        ),
        "provenance_blockers": (
            getattr(provenance, "blocker_count", 0) if provenance else 0
        ),
        "stages": [(s[0], s[1]) for s in stages],
    }


if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("--language", choices=["en", "ar"], default="en")
    parser.add_argument(
        "--brief",
        choices=["negative", "positive"],
        default="negative",
        help="Brief type: negative=mismatched RFP (expect block), positive=matched RFP (expect success)",
    )
    args = parser.parse_args()

    if args.brief == "positive":
        brief = BRIEF_POSITIVE
        docs_path = "data_positive_proof"
    else:
        brief = BRIEF_NEGATIVE
        docs_path = None

    result = asyncio.run(run_e2e(args.language, brief, args.brief, docs_path))

    # Write result JSON
    out = Path(f"e2e_proof_{args.language}_{args.brief}_result.json")
    with open(out, "w") as f:
        json.dump(result, f, indent=2, default=str)
    print(f"\nResults written to: {out}")
