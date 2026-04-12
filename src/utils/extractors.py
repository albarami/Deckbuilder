"""Document extractors for the Knowledge Layer.

Extracts structured content from PPTX, PDF, DOCX, and XLSX files.
Each extractor returns an ExtractedDocument with format-specific fields
populated (slides, pages, sections, sheets) plus a unified full_text.

Production patterns:
- Error isolation: corrupted files logged, never crash the batch
- Memory efficient: page-by-page / slide-by-slide extraction
- Content hash: SHA-256 of full_text for incremental indexing
"""

import hashlib
import logging
import os
from pathlib import Path

from src.models.enums import ExtractionQuality
from src.models.extraction import (
    ExtractedDocument,
    ExtractedPage,
    ExtractedSection,
    ExtractedSheet,
    ExtractedSlide,
    ExtractedTable,
)

logger = logging.getLogger(__name__)

_SUPPORTED_EXTENSIONS = {".pptx", ".pdf", ".docx", ".xlsx"}


def _normalize_arabic_text(text: str) -> str:
    """Normalize Arabic presentation forms and collapse extra whitespace.

    Many Arabic PDFs (especially from Etimad platform HTML-to-PDF renders)
    use Unicode Arabic Presentation Forms (FB50-FDFF, FE70-FEFF) instead of
    standard Arabic characters (0600-06FF). This makes keyword matching,
    LLM extraction, and NLP processing fail because the characters look
    identical visually but are different codepoints.

    NFKC normalization maps presentation forms to their standard equivalents.
    We also collapse multiple spaces (common in these PDFs) to single spaces
    while preserving line breaks.
    """
    import re
    import unicodedata

    if not text:
        return text

    # NFKC: maps presentation forms → standard Arabic, decomposes ligatures
    normalized = unicodedata.normalize("NFKC", text)

    # Collapse runs of 2+ spaces to single space (preserve newlines)
    normalized = re.sub(r"[ \t]{2,}", " ", normalized)

    # Strip leading/trailing spaces per line
    lines = normalized.splitlines()
    normalized = "\n".join(line.strip() for line in lines)

    return normalized


def _compute_hash(text: str) -> str:
    """Compute SHA-256 hex digest of text content."""
    return hashlib.sha256(text.encode("utf-8")).hexdigest()


def _file_size(filepath: str) -> int:
    """Get file size in bytes, 0 if not found."""
    try:
        return os.path.getsize(filepath)
    except OSError:
        return 0


def _make_error_doc(
    filepath: str,
    error: str,
    file_type: str = "",
) -> ExtractedDocument:
    """Create an ExtractedDocument representing an extraction failure."""
    path = Path(filepath)
    if not file_type:
        file_type = path.suffix.lstrip(".").lower()
    return ExtractedDocument(
        filepath=filepath,
        filename=path.name,
        file_type=file_type,
        file_size_bytes=_file_size(filepath),
        extraction_quality=ExtractionQuality.MANUAL_REVIEW_NEEDED,
        error=error,
    )


# ──────────────────────────────────────────────────────────────
# PPTX Extractor
# ──────────────────────────────────────────────────────────────


def extract_pptx(filepath: str) -> ExtractedDocument:
    """Extract structured content from a PPTX file.

    Preserves slide-level structure: title, body text, speaker notes,
    layout type, and tables per slide.
    """
    from pptx import Presentation

    path = Path(filepath)
    try:
        prs = Presentation(filepath)
    except Exception as e:
        return _make_error_doc(filepath, f"Failed to open PPTX: {e}", "pptx")

    slides: list[ExtractedSlide] = []
    full_parts: list[str] = []

    for idx, slide in enumerate(prs.slides, start=1):
        title = ""
        body_parts: list[str] = []
        tables: list[ExtractedTable] = []

        for shape in slide.shapes:
            # Extract title
            if shape.has_text_frame:
                if shape.is_placeholder:
                    ph_idx = shape.placeholder_format.idx
                    if ph_idx == 0:
                        title = shape.text_frame.text.strip()
                        continue
                # Body text from all text frames
                text = shape.text_frame.text.strip()
                if text:
                    body_parts.append(text)

            # Extract tables
            if shape.has_table:
                table_rows: list[list[str]] = []
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    table_rows.append(cells)
                tables.append(ExtractedTable(rows=table_rows))

        # Speaker notes
        speaker_notes = ""
        try:
            if slide.has_notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                speaker_notes = notes_frame.text.strip()
        except Exception:
            pass

        # Layout type
        layout_type = ""
        try:
            layout_type = slide.slide_layout.name
        except Exception:
            pass

        body_text = "\n".join(body_parts)
        slides.append(ExtractedSlide(
            slide_number=idx,
            title=title,
            body_text=body_text,
            speaker_notes=speaker_notes,
            layout_type=layout_type,
            tables=tables,
        ))

        # Build full text for this slide
        slide_text_parts = [p for p in [title, body_text, speaker_notes] if p]
        full_parts.append(f"[Slide {idx}] " + " | ".join(slide_text_parts))

    full_text = _normalize_arabic_text("\n\n".join(full_parts))

    return ExtractedDocument(
        filepath=filepath,
        filename=path.name,
        file_type="pptx",
        file_size_bytes=_file_size(filepath),
        content_hash=_compute_hash(full_text),
        slides=slides,
        full_text=full_text,
        extraction_quality=ExtractionQuality.CLEAN,
    )


# ──────────────────────────────────────────────────────────────
# PDF Extractor
# ──────────────────────────────────────────────────────────────


def extract_pdf(filepath: str) -> ExtractedDocument:
    """Extract text content from a PDF file, page by page.

    Uses the OCR service with smart fallback:
    1. Try PyPDF2 first (fast, free, works for digitally-born PDFs)
    2. If text is sparse (< 50 chars/page avg), fall back to Tesseract OCR
    3. If Tesseract unavailable, return degraded quality gracefully

    Handles both sync and async calling contexts:
    - From sync code: creates a new event loop with asyncio.run()
    - From async code: uses the existing loop to run the coroutine
    """
    import asyncio

    from src.services.ocr import extract_pdf_with_ocr

    path = Path(filepath)
    try:
        try:
            loop = asyncio.get_running_loop()
        except RuntimeError:
            loop = None

        if loop and loop.is_running():
            # Already inside an async context — run coroutine on the existing loop.
            # Since OCR backends are CPU-bound (pytesseract, pypdf2), they don't
            # actually need async I/O. We use a new thread to avoid blocking.
            import concurrent.futures
            with concurrent.futures.ThreadPoolExecutor(max_workers=1) as pool:
                ocr_result = pool.submit(
                    asyncio.run, extract_pdf_with_ocr(filepath, backend="auto")
                ).result()
        else:
            # No running loop — safe to use asyncio.run()
            ocr_result = asyncio.run(extract_pdf_with_ocr(filepath, backend="auto"))
    except Exception as e:
        return _make_error_doc(filepath, f"Failed to extract PDF: {e}", "pdf")

    # Normalize Arabic presentation forms → standard Arabic + collapse extra spaces
    full_text = _normalize_arabic_text(ocr_result.full_text)

    # Map OCR pages to ExtractedPage objects with normalized text
    pages = [
        ExtractedPage(page_number=p.page_number, text=_normalize_arabic_text(p.text))
        for p in ocr_result.pages
    ]

    # Map quality based on engine and confidence
    if ocr_result.engine_used == "pypdf2" and ocr_result.average_confidence >= 0.7:
        quality = ExtractionQuality.CLEAN
    elif ocr_result.engine_used == "tesseract" and ocr_result.average_confidence >= 0.7:
        quality = ExtractionQuality.PARTIAL_OCR
    elif not full_text.strip():
        quality = ExtractionQuality.DEGRADED
    else:
        quality = ExtractionQuality.DEGRADED

    return ExtractedDocument(
        filepath=filepath,
        filename=path.name,
        file_type="pdf",
        file_size_bytes=_file_size(filepath),
        content_hash=_compute_hash(full_text),
        pages=pages,
        full_text=full_text,
        extraction_quality=quality,
    )


# ──────────────────────────────────────────────────────────────
# DOCX Extractor
# ──────────────────────────────────────────────────────────────


def extract_docx(filepath: str) -> ExtractedDocument:
    """Extract structured content from a DOCX file.

    Groups content into sections based on heading paragraphs.
    """
    from docx import Document

    path = Path(filepath)
    try:
        doc = Document(filepath)
    except Exception as e:
        return _make_error_doc(filepath, f"Failed to open DOCX: {e}", "docx")

    sections: list[ExtractedSection] = []
    current_heading = ""
    current_level = 1
    current_paragraphs: list[str] = []

    def _flush_section() -> None:
        if current_heading or current_paragraphs:
            sections.append(ExtractedSection(
                heading=current_heading,
                text="\n".join(current_paragraphs).strip(),
                level=current_level,
            ))

    for paragraph in doc.paragraphs:
        style_name = paragraph.style.name if paragraph.style else ""
        text = paragraph.text.strip()

        if style_name.startswith("Heading"):
            _flush_section()
            current_heading = text
            try:
                current_level = int(style_name.split()[-1])
            except (ValueError, IndexError):
                current_level = 1
            current_paragraphs = []
        elif text:
            current_paragraphs.append(text)

    _flush_section()

    # Build full text
    full_parts: list[str] = []
    for section in sections:
        if section.heading:
            full_parts.append(section.heading)
        if section.text:
            full_parts.append(section.text)

    full_text = _normalize_arabic_text("\n\n".join(full_parts))

    return ExtractedDocument(
        filepath=filepath,
        filename=path.name,
        file_type="docx",
        file_size_bytes=_file_size(filepath),
        content_hash=_compute_hash(full_text),
        sections=sections,
        full_text=full_text,
        extraction_quality=ExtractionQuality.CLEAN,
    )


# ──────────────────────────────────────────────────────────────
# XLSX Extractor
# ──────────────────────────────────────────────────────────────


def extract_xlsx(filepath: str) -> ExtractedDocument:
    """Extract structured content from an XLSX file.

    Reads each sheet's headers and data rows as text.
    """
    from openpyxl import load_workbook  # type: ignore[import-untyped]

    path = Path(filepath)
    try:
        wb = load_workbook(filepath, read_only=True, data_only=True)
    except Exception as e:
        return _make_error_doc(filepath, f"Failed to open XLSX: {e}", "xlsx")

    sheets: list[ExtractedSheet] = []
    full_parts: list[str] = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = list(ws.iter_rows(values_only=True))

        if not rows:
            sheets.append(ExtractedSheet(sheet_name=sheet_name))
            continue

        # First row as headers
        headers = [str(cell) if cell is not None else "" for cell in rows[0]]
        data_rows = rows[1:]

        # Build text representation
        text_parts: list[str] = [" | ".join(headers)]
        for row in data_rows:
            cells = [str(cell) if cell is not None else "" for cell in row]
            text_parts.append(" | ".join(cells))

        text = "\n".join(text_parts)
        sheets.append(ExtractedSheet(
            sheet_name=sheet_name,
            headers=headers,
            row_count=len(data_rows),
            text=text,
        ))
        full_parts.append(f"[Sheet: {sheet_name}]\n{text}")

    wb.close()
    full_text = _normalize_arabic_text("\n\n".join(full_parts))

    return ExtractedDocument(
        filepath=filepath,
        filename=path.name,
        file_type="xlsx",
        file_size_bytes=_file_size(filepath),
        content_hash=_compute_hash(full_text),
        sheets=sheets,
        full_text=full_text,
        extraction_quality=ExtractionQuality.CLEAN,
    )


# ──────────────────────────────────────────────────────────────
# Dispatcher
# ──────────────────────────────────────────────────────────────


_EXTRACTOR_MAP = {
    ".pptx": extract_pptx,
    ".pdf": extract_pdf,
    ".docx": extract_docx,
    ".xlsx": extract_xlsx,
}


def extract_document(filepath: str) -> ExtractedDocument:
    """Extract content from any supported document format.

    Dispatches to the correct extractor based on file extension.
    Returns an ExtractedDocument with error set for unsupported
    or corrupted files.
    """
    path = Path(filepath)

    if not path.exists():
        return _make_error_doc(filepath, f"File not found: {filepath}")

    ext = path.suffix.lower()
    extractor = _EXTRACTOR_MAP.get(ext)

    if extractor is None:
        return _make_error_doc(
            filepath,
            f"Unsupported file type: {ext}",
            ext.lstrip("."),
        )

    try:
        return extractor(filepath)
    except Exception as e:
        logger.warning("Extraction failed for %s: %s", filepath, e)
        return _make_error_doc(filepath, str(e), ext.lstrip("."))


# ──────────────────────────────────────────────────────────────
# Batch Extraction
# ──────────────────────────────────────────────────────────────


def extract_directory(dir_path: str) -> list[ExtractedDocument]:
    """Extract all supported documents from a directory.

    Skips unsupported file types. Returns one ExtractedDocument
    per supported file. Errors are captured per-file, never crash
    the batch.
    """
    results: list[ExtractedDocument] = []
    directory = Path(dir_path)

    if not directory.is_dir():
        logger.warning("Not a directory: %s", dir_path)
        return results

    for filepath in sorted(directory.iterdir()):
        if not filepath.is_file():
            continue
        if filepath.suffix.lower() not in _SUPPORTED_EXTENSIONS:
            continue

        doc = extract_document(str(filepath))
        results.append(doc)

    return results
