"""Phase 1 — Template Manager.

POTX patcher, template loader, catalog lock consumer, slide cloner.
Template-hash enforcement on every load. All layout resolution by
semantic layout ID only — never by raw display names.

Consumes catalog lock artifacts from PHASE 0.
"""

from __future__ import annotations

import copy
import json
import logging
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from lxml import etree
from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.slide import Slide, SlideLayout

from src.services.template_auditor import (
    EN_POTX,
    AR_POTX,
    file_hash,
    patch_potx,
)

logger = logging.getLogger(__name__)

# ── Constants ─────────────────────────────────────────────────────────────

DATA_DIR = Path(__file__).resolve().parent.parent / "data"
CATALOG_LOCK_EN = DATA_DIR / "catalog_lock_en.json"
CATALOG_LOCK_AR = DATA_DIR / "catalog_lock_ar.json"


# ── Exceptions ────────────────────────────────────────────────────────────


class TemplateHashMismatch(RuntimeError):
    """Template file changed since catalog lock was generated."""


class SemanticIDNotFound(KeyError):
    """Semantic ID not present in catalog lock — fail-closed."""


class CatalogLockError(RuntimeError):
    """Catalog lock is invalid or missing required fields."""


# ── Data classes ──────────────────────────────────────────────────────────


@dataclass(frozen=True)
class AssetLocation:
    """Resolved location of a template asset."""

    slide_idx: int
    semantic_layout_id: str
    display_name: str
    shape_count: int = 0
    media_count: int = 0


@dataclass(frozen=True)
class LayoutLocation:
    """Resolved location of a template layout."""

    semantic_layout_id: str
    display_name: str
    master_idx: int
    placeholder_count: int
    placeholders: dict[str, str] = field(default_factory=dict)


@dataclass(frozen=True)
class ShellDef:
    """A2 shell definition with allowlist."""

    slide_idx: int
    semantic_layout_id: str
    display_name: str
    allowlist: dict[str, Any] = field(default_factory=dict)


# ── Catalog Lock Loader ──────────────────────────────────────────────────


def load_catalog_lock(path: Path) -> dict:
    """Load and validate a catalog lock file.

    Raises CatalogLockError if required fields are missing.
    """
    if not path.exists():
        raise CatalogLockError(f"Catalog lock not found: {path}")

    with open(path, encoding="utf-8") as f:
        lock = json.load(f)

    required = [
        "template_hash", "language", "a1_immutable", "a2_shells",
        "section_dividers", "case_study_pool", "team_bio_pool",
        "service_divider_pool", "layouts",
    ]
    missing = [k for k in required if k not in lock]
    if missing:
        raise CatalogLockError(f"Catalog lock missing fields: {missing}")

    return lock


# ── Placeholder index repair ────────────────────────────────────────────


_CORRUPT_IDX = 4294967295  # 0xFFFFFFFF — sentinel / corrupt value in some templates

_NSMAP = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}


def _repair_corrupt_placeholder_indices(slide: Slide) -> int:
    """Repair corrupt placeholder indices on a cloned slide.

    Some template slides (e.g., AR divider 02) have placeholder idx
    values of 4294967295 (0xFFFFFFFF) — a corrupt sentinel.  This
    function detects these and repairs them by matching the placeholder
    type (e.g., ``body``) against the slide layout's placeholders.

    Returns the number of repaired placeholders.
    """
    repaired = 0
    layout = slide.slide_layout

    # Build a map of layout placeholder types -> indices for repair lookup
    layout_type_to_idx: dict[str, int] = {}
    for ly_ph in layout.placeholders:
        ph_type = ly_ph.placeholder_format.type
        # Use the XML type string (e.g., 'body', 'title')
        ly_sp = ly_ph._element
        ly_nvPr = ly_sp.find(".//p:nvSpPr/p:nvPr/p:ph", _NSMAP)
        if ly_nvPr is not None:
            ly_type = ly_nvPr.get("type", "")
            ly_idx_str = ly_nvPr.get("idx", "")
            if ly_type and ly_idx_str:
                try:
                    layout_type_to_idx[ly_type] = int(ly_idx_str)
                except ValueError:
                    pass

    # Scan slide-level shapes for corrupt indices
    sp_tree = slide.shapes._spTree
    for sp in sp_tree:
        if etree.QName(sp.tag).localname != "sp":
            continue
        nvPr_ph = sp.find(".//p:nvSpPr/p:nvPr/p:ph", _NSMAP)
        if nvPr_ph is None:
            continue
        idx_str = nvPr_ph.get("idx", "")
        if not idx_str:
            continue
        try:
            idx_val = int(idx_str)
        except ValueError:
            continue
        if idx_val == _CORRUPT_IDX:
            ph_type = nvPr_ph.get("type", "")
            correct_idx = layout_type_to_idx.get(ph_type)
            if correct_idx is not None:
                nvPr_ph.set("idx", str(correct_idx))
                repaired += 1
                logger.info(
                    "Repaired corrupt placeholder idx %s -> %d (type=%s)",
                    idx_str, correct_idx, ph_type,
                )

    return repaired


# ── Template Manager ─────────────────────────────────────────────────────


class TemplateManager:
    """Manages template loading, hash enforcement, and slide operations.

    All layout resolution is by semantic layout ID from the catalog lock.
    Raw display names are never used for runtime resolution.
    """

    def __init__(
        self,
        potx_path: Path,
        catalog_lock_path: Path,
        patch_dir: Path | None = None,
    ):
        self._potx_path = Path(potx_path)
        self._lock = load_catalog_lock(catalog_lock_path)
        self._language = self._lock["language"]

        # Validate template hash — fail-closed on mismatch
        actual_hash = file_hash(self._potx_path)
        expected_hash = self._lock["template_hash"]
        if actual_hash != expected_hash:
            raise TemplateHashMismatch(
                f"Template file changed since catalog lock was generated. "
                f"Expected: {expected_hash}, Got: {actual_hash}. "
                f"Regenerate catalog lock with template_auditor."
            )

        # Patch POTX and load presentation
        self._pptx_path = patch_potx(self._potx_path, patch_dir)
        self._prs = Presentation(str(self._pptx_path))

        # Build internal indexes from catalog lock
        self._a1_assets: dict[str, AssetLocation] = {}
        self._a2_shells: dict[str, ShellDef] = {}
        self._dividers: dict[str, ShellDef] = {}
        self._layouts: dict[str, LayoutLocation] = {}
        self._layout_display_to_semantic: dict[str, str] = {}
        self._case_pool: dict[str, list[AssetLocation]] = {}
        self._team_pool: list[dict] = []
        self._service_divider_pool: list[AssetLocation] = []

        self._build_indexes()

        # Track original slide count for later cleanup
        self._original_slide_count = len(self._prs.slides)

        logger.info(
            "TemplateManager loaded: %s (%s), %d slides, %d layouts, "
            "%d A1, %d A2, %d dividers",
            self._language,
            self._potx_path.name,
            self._original_slide_count,
            len(self._layouts),
            len(self._a1_assets),
            len(self._a2_shells),
            len(self._dividers),
        )

    def _build_indexes(self) -> None:
        """Build all internal indexes from the catalog lock."""
        lock = self._lock

        # A1 immutable assets
        for sid, info in lock["a1_immutable"].items():
            self._a1_assets[sid] = AssetLocation(
                slide_idx=info["slide_idx"],
                semantic_layout_id=info.get("semantic_layout_id", ""),
                display_name=info.get("display_name", ""),
                shape_count=info.get("shape_count", 0),
                media_count=info.get("media_count", 0),
            )

        # A2 shells
        for sid, info in lock["a2_shells"].items():
            self._a2_shells[sid] = ShellDef(
                slide_idx=info["slide_idx"],
                semantic_layout_id=info.get("semantic_layout_id", ""),
                display_name=info.get("display_name", ""),
                allowlist=info.get("allowlist", {}),
            )

        # Section dividers
        for num, info in lock["section_dividers"].items():
            self._dividers[num] = ShellDef(
                slide_idx=info["slide_idx"],
                semantic_layout_id=info.get("semantic_layout_id", ""),
                display_name=info.get("display_name", ""),
                allowlist=info.get("allowlist", {}),
            )

        # Layouts
        for sem_id, info in lock["layouts"].items():
            self._layouts[sem_id] = LayoutLocation(
                semantic_layout_id=sem_id,
                display_name=info["display_name"],
                master_idx=info["master_idx"],
                placeholder_count=info["placeholder_count"],
                placeholders=info.get("placeholders", {}),
            )
            # Reverse map for internal use only
            self._layout_display_to_semantic[
                info["display_name"].lower().strip()
            ] = sem_id

        # Case study pool
        for cat, items in lock["case_study_pool"].items():
            self._case_pool[cat] = [
                AssetLocation(
                    slide_idx=item["slide_idx"],
                    semantic_layout_id=item.get("semantic_layout_id", ""),
                    display_name=item.get("display_name", ""),
                )
                for item in items
            ]

        # Team pool
        self._team_pool = lock["team_bio_pool"]

        # Service divider pool
        for item in lock["service_divider_pool"]:
            self._service_divider_pool.append(
                AssetLocation(
                    slide_idx=item["slide_idx"],
                    semantic_layout_id=item.get("semantic_layout_id", ""),
                    display_name=item.get("display_name", ""),
                )
            )

    # ── Properties ────────────────────────────────────────────────────

    @property
    def language(self) -> str:
        return self._language

    @property
    def presentation(self) -> Presentation:
        return self._prs

    @property
    def template_hash(self) -> str:
        return self._lock["template_hash"]

    @property
    def a1_asset_ids(self) -> list[str]:
        return sorted(self._a1_assets.keys())

    @property
    def a2_shell_ids(self) -> list[str]:
        return sorted(self._a2_shells.keys())

    @property
    def divider_numbers(self) -> list[str]:
        return sorted(self._dividers.keys())

    @property
    def layout_ids(self) -> list[str]:
        return sorted(self._layouts.keys())

    # ── Resolve by semantic ID (fail-closed) ─────────────────────────

    def resolve_a1(self, semantic_id: str) -> AssetLocation:
        """Resolve an A1 asset by semantic ID. Fail-closed if missing."""
        if semantic_id not in self._a1_assets:
            raise SemanticIDNotFound(
                f"A1 asset '{semantic_id}' not found in catalog lock. "
                f"Available: {sorted(self._a1_assets.keys())}"
            )
        return self._a1_assets[semantic_id]

    def resolve_a2(self, semantic_id: str) -> ShellDef:
        """Resolve an A2 shell by semantic ID. Fail-closed if missing."""
        if semantic_id not in self._a2_shells:
            raise SemanticIDNotFound(
                f"A2 shell '{semantic_id}' not found in catalog lock. "
                f"Available: {sorted(self._a2_shells.keys())}"
            )
        return self._a2_shells[semantic_id]

    def resolve_divider(self, number: str) -> ShellDef:
        """Resolve a section divider by number. Fail-closed if missing."""
        if number not in self._dividers:
            raise SemanticIDNotFound(
                f"Divider '{number}' not found in catalog lock. "
                f"Available: {sorted(self._dividers.keys())}"
            )
        return self._dividers[number]

    def resolve_layout(self, semantic_layout_id: str) -> LayoutLocation:
        """Resolve a layout by semantic layout ID. Fail-closed if missing."""
        if semantic_layout_id not in self._layouts:
            raise SemanticIDNotFound(
                f"Layout '{semantic_layout_id}' not found in catalog lock. "
                f"Available: {sorted(self._layouts.keys())}"
            )
        return self._layouts[semantic_layout_id]

    # ── Layout object resolution ─────────────────────────────────────

    def get_layout_object(self, semantic_layout_id: str) -> SlideLayout:
        """Get the python-pptx SlideLayout object for a semantic layout ID.

        Resolves through the catalog lock, never by raw display name.
        """
        loc = self.resolve_layout(semantic_layout_id)
        display_name = loc.display_name

        # Find the layout object by display name (internal resolution only)
        for master in self._prs.slide_masters:
            for layout in master.slide_layouts:
                if layout.name == display_name:
                    return layout

        raise SemanticIDNotFound(
            f"Layout object for '{semantic_layout_id}' "
            f"(display: '{display_name}') not found in presentation."
        )

    # ── Slide operations ─────────────────────────────────────────────

    def add_slide_from_layout(self, semantic_layout_id: str) -> Slide:
        """Add a new blank slide using an official template layout.

        Resolves the layout by semantic layout ID from the catalog lock.
        Returns the newly added Slide object.
        """
        layout = self.get_layout_object(semantic_layout_id)
        slide = self._prs.slides.add_slide(layout)
        return slide

    def clone_slide(self, source_slide_idx: int) -> Slide:
        """Clone a slide from the template preserving master/layout linkage.

        Uses XML-level copy to preserve all shapes, relationships,
        and visual structure. The cloned slide is appended to the
        presentation.

        If the source slide has no explicit shapes (all content is
        inherited from the slide layout — common for A1 institutional
        slides), the new slide is created from the same layout and
        its layout-inherited shapes are preserved as-is.
        """
        if source_slide_idx >= len(self._prs.slides):
            raise IndexError(
                f"Slide index {source_slide_idx} out of range "
                f"({len(self._prs.slides)} slides)"
            )

        source = self._prs.slides[source_slide_idx]
        layout = source.slide_layout

        # Check if source has explicit shapes (beyond layout inheritance)
        source_tree = source.shapes._spTree
        source_shape_tags = [
            c for c in source_tree
            if etree.QName(c.tag).localname in ("sp", "pic", "graphicFrame", "grpSp", "cxnSp")
        ]

        # Add a new slide with the same layout
        new_slide = self._prs.slides.add_slide(layout)

        if source_shape_tags:
            # Source has explicit shapes — deep copy them into the new slide
            sp_tree = new_slide.shapes._spTree
            for child in list(sp_tree):
                tag = etree.QName(child.tag).localname
                if tag in ("sp", "pic", "graphicFrame", "grpSp", "cxnSp"):
                    sp_tree.remove(child)

            for child in source_shape_tags:
                new_elem = copy.deepcopy(child)
                sp_tree.append(new_elem)
        # else: source has no explicit shapes — layout-inherited content
        # is automatically provided by add_slide(layout), so keep it as-is

        # Repair corrupt placeholder indices (0xFFFFFFFF = 4294967295)
        # Some template slides (e.g., AR divider 02) have corrupted
        # placeholder idx values in the XML.  Repair by matching against
        # the layout's placeholder index for the same type.
        _repair_corrupt_placeholder_indices(new_slide)

        # Copy slide-level relationships (images, media)
        for rel in source.part.rels.values():
            if rel.reltype not in (RT.SLIDE_LAYOUT,):
                # Copy the relationship target
                try:
                    if hasattr(rel, "target_part") and rel.target_part is not None:
                        new_slide.part.rels.get_or_add(
                            rel.reltype, rel.target_part
                        )
                except Exception:
                    # Some relationships can't be copied; skip gracefully
                    pass

        return new_slide

    def clone_a1(self, semantic_id: str) -> Slide:
        """Clone an A1 immutable house slide by semantic ID.

        Structurally clones the slide preserving all visual structure.
        """
        loc = self.resolve_a1(semantic_id)
        return self.clone_slide(loc.slide_idx)

    def clone_a2(self, semantic_id: str) -> Slide:
        """Clone an A2 shell slide by semantic ID.

        Returns the cloned slide (not yet sanitized — call shell_sanitizer
        before injection).
        """
        loc = self.resolve_a2(semantic_id)
        return self.clone_slide(loc.slide_idx)

    def clone_divider(self, number: str) -> Slide:
        """Clone a section divider slide by number (e.g., '01')."""
        loc = self.resolve_divider(number)
        return self.clone_slide(loc.slide_idx)

    # ── Pool access ──────────────────────────────────────────────────

    def get_case_study_pool(self) -> dict[str, list[AssetLocation]]:
        """Get case study pool (disjoint from service dividers)."""
        return dict(self._case_pool)

    def get_team_pool(self) -> list[dict]:
        """Get team bio pool with family normalization."""
        return list(self._team_pool)

    def get_service_divider_pool(self) -> list[AssetLocation]:
        """Get service divider pool (disjoint from case studies)."""
        return list(self._service_divider_pool)

    # ── Slide cleanup ────────────────────────────────────────────────

    def remove_original_slides(self) -> int:
        """Remove all original template slides, keeping only rendered slides.

        After rendering, the presentation contains original template slides
        (indices 0..N-1) followed by cloned/added slides (indices N..).
        This method removes the original slides so only rendered output remains.

        Returns the number of slides removed.
        """
        count = self._original_slide_count
        if count <= 0:
            return 0

        xml_slides = self._prs.slides._sldIdLst
        ns = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"

        removed = 0
        for _ in range(count):
            if len(xml_slides) == 0:
                break
            sld_id = xml_slides[0]
            r_id = sld_id.get(f"{ns}id")
            if r_id:
                try:
                    self._prs.part.drop_rel(r_id)
                except KeyError:
                    pass  # Relationship may already be gone
            xml_slides.remove(sld_id)
            removed += 1

        logger.info("Removed %d original template slides", removed)
        return removed

    # ── Output ───────────────────────────────────────────────────────

    def save(self, output_path: Path) -> Path:
        """Save the presentation to a file."""
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        self._prs.save(str(output_path))
        return output_path
