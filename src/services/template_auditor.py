"""PHASE 0 — Template Intelligence Audit.

Produces machine-readable JSON artifacts from official POTX templates.
No renderer_v2 coding is allowed until this audit exists and is validated.

Usage:
    python -m src.services.template_auditor [output_base_dir]
"""

from __future__ import annotations

import hashlib
import json
import logging
import re
import shutil
import tempfile
import zipfile
from datetime import UTC, datetime
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

logger = logging.getLogger(__name__)

# ── Constants ──────────────────────────────────────────────────────────────

POTX_CT = (
    "application/vnd.openxmlformats-officedocument"
    ".presentationml.template.main+xml"
)
PPTX_CT = (
    "application/vnd.openxmlformats-officedocument"
    ".presentationml.presentation.main+xml"
)

TEMPLATE_DIR = Path(r"C:\Projects\Deckbuilder\PROPOSAL_TEMPLATE")
EXAMPLE_DIR = Path(r"C:\Projects\Deckbuilder\real_example")
EN_POTX = TEMPLATE_DIR / "PROPOSAL_TEMPLATE EN.potx"
AR_POTX = TEMPLATE_DIR / "Arabic_Proposal_Template.potx"

# ── Classification regex patterns ──────────────────────────────────────────

# Actual template layout name patterns (from audit discovery)
_RE_DIVIDER = re.compile(r"^0[1-9]$")  # "01", "02", ..., "09"
_RE_METH = re.compile(r"^Meth", re.I)  # Handles template typos: Methodology/Methdology/Methodolgy
_RE_TEAM = re.compile(r"team\s*member", re.I)
_RE_CASE = re.compile(r"services\s*-\s*(cases|detailed\s*case)", re.I)
_RE_COVER = re.compile(r"^proposal\s*cover$", re.I)
_RE_TOC = re.compile(r"(toc|agenda)", re.I)
_RE_INTRO = re.compile(r"^introduction\s*message$", re.I)
_RE_LEAD = re.compile(r"^our\s*leadership$", re.I)
_RE_KNOW = re.compile(r"know\s*more", re.I)
_RE_CONTACT = re.compile(r"^contact$", re.I)
_RE_MAIN_COVER = re.compile(r"^main\s*cover$", re.I)
_RE_OUR_SERVICES = re.compile(r"^our\s*services$", re.I)

# KSA context layout names (Master 1, immutable)
_KSA_LAYOUT_NAMES = {
    "ksa",
    "pillars of the vision",
    "vision realization programs numbers",
    "vision realization programs",
}

# Company profile layout names (Master 3, immutable)
_COMPANY_PROFILE_NAMES: dict[str, str] = {
    "overview": "overview",
    "what drives us": "what_drives_us",
    "at a glance": "at_a_glance",
    "why strategic gears": "why_sg",
    "deep experience 1/4": "deep_experience_1",
    "deep experience 2/4": "deep_experience_2",
    "deep experience 3/4": "deep_experience_3",
    "deep experience 4/4": "deep_experience_4",
    "a house of expertise": "expertise",
    "vast network": "vast_network",
    "purpose beyond business 1/2": "purpose_1",
    "purpose beyond business 2/2": "purpose_2",
}

# Service category divider layout names (Master 6, immutable pool)
_SERVICE_CATEGORY_NAMES = {
    "strategy",
    "organizational excellence",
    "marketing",
    "digital, cloud, and ai",
    "people advisory",
    "deals advisory",
    "research",
}

# ── Semantic Layout ID mapping ────────────────────────────────────────────
# Maps layout display names (lowercased) to stable, language-neutral IDs.
# Runtime code resolves layouts ONLY by semantic layout ID from catalog lock.
# Display names exist as audit metadata only.

_LAYOUT_SEMANTIC_MAP: dict[str, str] = {
    # Master 0 — fixed-purpose layouts
    "main cover": "main_cover",
    "proposal cover": "proposal_cover",
    "toc / agenda": "toc_table",
    "know more page": "know_more",
    "contact": "contact",
    # Master 1 — general content + KSA context
    "introduction message": "intro_message",
    "ksa": "ksa_context",
    "pillars of the vision": "vision_pillars",
    "vision realization programs numbers": "vision_programs_numbers",
    "vision realization programs": "vision_programs",
    "heading only": "content_heading_only",
    "heading and description": "content_heading_desc",
    "heading description and content box": "content_heading_desc_box",
    "heading and content": "content_heading_content",
    "heading and 4-boxes": "content_heading_4boxes",
    # Master 3 — company profile
    "overview": "overview",
    "what drives us": "what_drives_us",
    "at a glance": "at_a_glance",
    "why strategic gears": "why_sg",
    "deep experience 1/4": "deep_experience_1",
    "deep experience 2/4": "deep_experience_2",
    "deep experience 3/4": "deep_experience_3",
    "deep experience 4/4": "deep_experience_4",
    "a house of expertise": "expertise",
    "vast network": "vast_network",
    "purpose beyond business 1/2": "purpose_1",
    "purpose beyond business 2/2": "purpose_2",
    # Master 5
    "our leadership": "our_leadership",
    "two team members": "team_two_members",
    # Master 6 — services
    "our services": "services_overview",
    "strategy": "svc_strategy",
    "organizational excellence": "svc_organizational_excellence",
    "marketing": "svc_marketing",
    "digital, cloud, and ai": "svc_digital_cloud_ai",
    "digital cloud and ai": "svc_digital_cloud_ai",
    "people advisory": "svc_people_advisory",
    "deals advisory": "svc_deals_advisory",
    "research": "svc_research",
    "services - cases": "case_study_cases",
    "services - detailed case": "case_study_detailed",
}


def _layout_name_to_semantic_id(name: str) -> str:
    """Map a layout display name to a stable semantic layout ID.

    Uses exact match first, then methodology pattern, divider pattern,
    and finally falls back to a deterministic slug.
    """
    nl = name.lower().strip()

    # Exact match from known layouts
    if nl in _LAYOUT_SEMANTIC_MAP:
        return _LAYOUT_SEMANTIC_MAP[nl]

    # Section dividers: "01" through "09"
    if _RE_DIVIDER.match(nl):
        return f"section_divider_{nl}"

    # Methodology variants (handles typos: Methodology/Methdology/Methodolgy)
    if _RE_METH.search(name):
        variant = _meth_variant(name)
        return f"methodology_{variant}"

    # Fallback: deterministic slug from display name
    slug = re.sub(r"[^a-z0-9]+", "_", nl).strip("_")
    return f"layout_{slug}"


# ── POTX Patching ──────────────────────────────────────────────────────────


def patch_potx(potx_path: Path, out_dir: Path | None = None) -> Path:
    """Patch .potx content-type to make it openable by python-pptx."""
    potx_path = Path(potx_path)
    if out_dir is None:
        out_dir = Path(tempfile.mkdtemp(prefix="potx_"))
    out = out_dir / (potx_path.stem.replace(" ", "_") + ".pptx")
    with zipfile.ZipFile(potx_path, "r") as zin:
        with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename == "[Content_Types].xml":
                    data = data.replace(
                        POTX_CT.encode(), PPTX_CT.encode()
                    )
                zout.writestr(item, data)
    return out


def file_hash(path: Path) -> str:
    """SHA256 of file."""
    h = hashlib.sha256()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(1 << 16), b""):
            h.update(chunk)
    return f"sha256:{h.hexdigest()}"


# ── Extraction helpers ─────────────────────────────────────────────────────


def _ph_type(ph) -> str:
    """Placeholder type as string."""
    try:
        raw = str(ph.placeholder_format.type)
        # Format: "TITLE (15)" or "MSO_PLACEHOLDER_TYPE.TITLE (15)" etc.
        return raw.split("(")[0].split(".")[-1].strip()
    except Exception:
        return "UNKNOWN"


def _shape_text(shape) -> str:
    """Extract text from a shape (text frame or table)."""
    parts: list[str] = []
    if shape.has_text_frame:
        parts.extend(p.text for p in shape.text_frame.paragraphs)
    if shape.has_table:
        for row in shape.table.rows:
            parts.extend(cell.text for cell in row.cells)
    return "\n".join(parts).strip()


def _media_count(slide) -> int:
    """Count picture shapes on a slide."""
    n = 0
    for s in slide.shapes:
        if s.shape_type == MSO_SHAPE_TYPE.PICTURE:
            n += 1
    return n


# ── Effective slide structure (master-aware) ──────────────────────────────


def build_effective_slide_structure(prs: Presentation) -> list[dict]:
    """Build master-aware effective slide structure for each slide.

    For every slide, captures three levels of content:
    - slide_local: shapes/text defined directly on the slide
    - layout_inherited: shapes/text from the slide layout
    - master_inherited: shapes/text from the slide master

    This is the ground-truth for understanding what a rendered slide
    actually shows, since many template slides have shape_count=0 but
    inherit rich visual structure from masters/layouts.
    """
    # Build layout -> master map
    lm_map: dict[str, int] = {}
    for mi, master in enumerate(prs.slide_masters):
        for sl in master.slide_layouts:
            lm_map[sl.name] = mi

    result: list[dict] = []
    for si, slide in enumerate(prs.slides):
        layout = slide.slide_layout
        master = layout.slide_master
        ln = layout.name
        sem_layout_id = _layout_name_to_semantic_id(ln)

        # Slide-local shapes
        slide_shapes = _extract_shape_info(slide.shapes)

        # Layout-level shapes (inherited if slide doesn't override)
        layout_shapes = _extract_shape_info(layout.shapes)

        # Master-level shapes
        master_shapes = _extract_shape_info(master.shapes)

        # Effective placeholder inventory (slide > layout > master)
        effective_phs: dict[int, dict] = {}
        # Start with master placeholders
        for ph in master.placeholders:
            effective_phs[ph.placeholder_format.idx] = {
                "idx": ph.placeholder_format.idx,
                "type": _ph_type(ph),
                "source": "master",
                "has_text": bool(ph.has_text_frame and ph.text_frame.text.strip()),
                "text_preview": (ph.text_frame.text[:200] if ph.has_text_frame else ""),
            }
        # Override with layout placeholders
        for ph in layout.placeholders:
            effective_phs[ph.placeholder_format.idx] = {
                "idx": ph.placeholder_format.idx,
                "type": _ph_type(ph),
                "source": "layout",
                "has_text": bool(ph.has_text_frame and ph.text_frame.text.strip()),
                "text_preview": (ph.text_frame.text[:200] if ph.has_text_frame else ""),
            }
        # Override with slide placeholders
        for ph in slide.placeholders:
            effective_phs[ph.placeholder_format.idx] = {
                "idx": ph.placeholder_format.idx,
                "type": _ph_type(ph),
                "source": "slide",
                "has_text": bool(ph.has_text_frame and ph.text_frame.text.strip()),
                "text_preview": (ph.text_frame.text[:200] if ph.has_text_frame else ""),
            }

        # All effective text regions across all levels
        effective_text_regions: list[dict] = []
        for src_label, shapes in [("slide", slide.shapes), ("layout", layout.shapes), ("master", master.shapes)]:
            for shape in shapes:
                text = ""
                if shape.has_text_frame:
                    text = shape.text_frame.text[:300]
                if shape.has_table:
                    cells = []
                    for row in shape.table.rows:
                        for cell in row.cells:
                            ct = cell.text.strip()
                            if ct:
                                cells.append(ct[:100])
                    text = " | ".join(cells[:10])
                if text.strip():
                    effective_text_regions.append({
                        "source": src_label,
                        "shape_name": shape.name,
                        "is_placeholder": shape.is_placeholder,
                        "ph_idx": (shape.placeholder_format.idx if shape.is_placeholder else None),
                        "text_preview": text.strip()[:200],
                        "has_table": shape.has_table,
                    })

        result.append({
            "slide_idx": si,
            "layout_name": ln,
            "semantic_layout_id": sem_layout_id,
            "master_idx": lm_map.get(ln, -1),
            "slide_local_shape_count": len(list(slide.shapes)),
            "layout_shape_count": len(list(layout.shapes)),
            "master_shape_count": len(list(master.shapes)),
            "effective_placeholder_count": len(effective_phs),
            "effective_placeholders": sorted(effective_phs.values(), key=lambda p: p["idx"]),
            "effective_text_region_count": len(effective_text_regions),
            "effective_text_regions": effective_text_regions,
            "slide_local_shapes": slide_shapes,
            "layout_shapes": layout_shapes,
            "master_shapes": master_shapes,
        })

    return result


def _extract_shape_info(shapes) -> list[dict]:
    """Extract shape metadata for effective structure reporting."""
    result: list[dict] = []
    for shape in shapes:
        info: dict[str, Any] = {
            "name": shape.name,
            "is_placeholder": shape.is_placeholder,
            "has_text_frame": shape.has_text_frame,
            "has_table": shape.has_table,
        }
        if shape.is_placeholder:
            info["ph_idx"] = shape.placeholder_format.idx
            info["ph_type"] = _ph_type(shape)
        if shape.has_text_frame:
            info["text_preview"] = shape.text_frame.text[:200]
        if shape.has_table:
            cells = []
            for row in shape.table.rows:
                for cell in row.cells:
                    ct = cell.text.strip()
                    if ct:
                        cells.append(ct[:100])
            info["table_cell_preview"] = cells[:10]
        # Bounding box
        try:
            info["left"] = shape.left
            info["top"] = shape.top
            info["width"] = shape.width
            info["height"] = shape.height
        except Exception:
            pass
        result.append(info)
    return result


# ── Layout catalog ─────────────────────────────────────────────────────────


def catalog_layouts(prs: Presentation) -> list[dict]:
    """Full layout catalog with placeholder details and semantic layout IDs."""
    result: list[dict] = []
    for mi, master in enumerate(prs.slide_masters):
        for layout in master.slide_layouts:
            phs = []
            for ph in layout.placeholders:
                phs.append({
                    "idx": ph.placeholder_format.idx,
                    "type": _ph_type(ph),
                    "left": ph.left,
                    "top": ph.top,
                    "width": ph.width,
                    "height": ph.height,
                })
            sem_id = _layout_name_to_semantic_id(layout.name)
            result.append({
                "name": layout.name,
                "semantic_layout_id": sem_id,
                "master_idx": mi,
                "shape_count": len(layout.shapes),
                "placeholder_count": len(phs),
                "placeholders": phs,
            })
    return result


# ── Slide audit ────────────────────────────────────────────────────────────


def audit_all_slides(prs: Presentation) -> list[dict]:
    """Walk every slide: metadata, placeholders, text content, shape details."""
    # Build layout -> master index map
    lm_map: dict[str, int] = {}
    for mi, master in enumerate(prs.slide_masters):
        for sl in master.slide_layouts:
            lm_map[sl.name] = mi

    slides: list[dict] = []
    for si, slide in enumerate(prs.slides):
        ln = slide.slide_layout.name
        sem_layout_id = _layout_name_to_semantic_id(ln)
        # Placeholders
        phs = []
        for ph in slide.placeholders:
            phs.append({
                "idx": ph.placeholder_format.idx,
                "type": _ph_type(ph),
                "has_text": bool(
                    ph.has_text_frame and ph.text_frame.text.strip()
                ),
                "text": (
                    ph.text_frame.text[:300] if ph.has_text_frame else ""
                ),
            })
        # Full text from every shape
        text: dict[str, str] = {}
        for shape in slide.shapes:
            t = _shape_text(shape)
            if t:
                key = (
                    f"ph_{shape.placeholder_format.idx}"
                    if shape.is_placeholder
                    else shape.name
                )
                text[key] = t[:500]

        # Shape details for A2 allowlist construction
        shapes_detail: list[dict] = []
        for shape in slide.shapes:
            detail: dict[str, Any] = {
                "name": shape.name,
                "is_placeholder": shape.is_placeholder,
                "has_text_frame": shape.has_text_frame,
                "has_table": shape.has_table,
            }
            if shape.is_placeholder:
                detail["ph_idx"] = shape.placeholder_format.idx
            if shape.has_text_frame:
                detail["text_preview"] = shape.text_frame.text[:200]
            if shape.has_table:
                cell_texts = []
                for row in shape.table.rows:
                    for cell in row.cells:
                        ct = cell.text.strip()
                        if ct:
                            cell_texts.append(ct[:100])
                detail["table_cell_preview"] = cell_texts[:10]
            shapes_detail.append(detail)

        slides.append({
            "slide_idx": si,
            "layout_name": ln,
            "semantic_layout_id": sem_layout_id,
            "master_idx": lm_map.get(ln, -1),
            "shape_count": len(slide.shapes),
            "placeholder_count": len(phs),
            "placeholders": phs,
            "media_count": _media_count(slide),
            "text": text,
            "shapes_detail": shapes_detail,
            "classification": "unclassified",
            "semantic_id": None,
        })
    return slides


# ── Classification ─────────────────────────────────────────────────────────


def classify_slides(slides: list[dict]) -> list[dict]:
    """Classify each slide and assign semantic IDs.

    Uses layout name + master_idx as primary classifiers.
    """
    for s in slides:
        ln = s["layout_name"]
        ln_low = ln.lower().strip()
        mi = s["master_idx"]

        # ── Section dividers: exact "01"-"09" (Master 4) ──────────────
        if _RE_DIVIDER.match(ln.strip()):
            s["classification"] = "a2_shell"
            s["semantic_id"] = f"section_{ln.strip()}"
            s["approved_placeholders"] = [0, 10]
            continue

        # ── Proposal cover ─────────────────────────────────────────────
        if _RE_COVER.match(ln.strip()):
            s["classification"] = "a2_shell"
            s["semantic_id"] = "proposal_cover"
            s["approved_placeholders"] = [1, 10, 11, 12]
            continue

        # ── Main cover (Master 0) ──────────────────────────────────────
        if _RE_MAIN_COVER.match(ln.strip()):
            s["classification"] = "a1_immutable"
            s["semantic_id"] = "main_cover"
            continue

        # ── ToC / Agenda ───────────────────────────────────────────────
        if _RE_TOC.search(ln):
            s["classification"] = "a2_shell"
            s["semantic_id"] = "toc_agenda"
            s["approved_placeholders"] = [0, 10]
            continue

        # ── Introduction message ───────────────────────────────────────
        if _RE_INTRO.match(ln.strip()):
            s["classification"] = "a2_shell"
            s["semantic_id"] = "intro_message"
            s["approved_placeholders"] = [0, 1, 13, 14, 15, 16, 17]
            continue

        # ── KSA context layouts (Master 1, specific names) ────────────
        if ln_low in _KSA_LAYOUT_NAMES:
            s["classification"] = "a1_immutable"
            if ln_low == "ksa":
                s["semantic_id"] = "ksa_context"
            elif ln_low == "pillars of the vision":
                s["semantic_id"] = "vision_pillars"
            elif "numbers" in ln_low:
                s["semantic_id"] = "vision_programs_numbers"
            else:
                s["semantic_id"] = "vision_programs"
            continue

        # ── Company profile (Master 3, by layout name) ────────────────
        if ln_low in _COMPANY_PROFILE_NAMES:
            s["classification"] = "a1_immutable"
            s["semantic_id"] = _COMPANY_PROFILE_NAMES[ln_low]
            continue

        # ── Leadership ─────────────────────────────────────────────────
        if _RE_LEAD.match(ln.strip()):
            s["classification"] = "a1_immutable"
            s["semantic_id"] = "our_leadership"
            continue

        # ── Know More ──────────────────────────────────────────────────
        if _RE_KNOW.search(ln):
            s["classification"] = "a1_immutable"
            s["semantic_id"] = "know_more"
            continue

        # ── Contact ────────────────────────────────────────────────────
        if _RE_CONTACT.match(ln.strip()):
            s["classification"] = "a1_immutable"
            s["semantic_id"] = "contact"
            continue

        # ── Our Services (Master 6, immutable overview) ────────────────
        if _RE_OUR_SERVICES.match(ln.strip()):
            s["classification"] = "a1_immutable"
            s["semantic_id"] = "services_overview"
            continue

        # ── Service category dividers (Master 6, disjoint from case studies) ─
        if ln_low in _SERVICE_CATEGORY_NAMES:
            s["classification"] = "pool_service_divider"
            s["semantic_id"] = f"svc_divider_{ln_low.replace(' ', '_').replace(',', '')}"
            s["service_category"] = ln_low
            continue

        # ── Team bios ──────────────────────────────────────────────────
        if _RE_TEAM.search(ln):
            s["classification"] = "pool_team_bio"
            names: list[str] = []
            for txt in s["text"].values():
                for line in txt.split("\n"):
                    line = line.strip()
                    if 3 < len(line) < 60 and not line.startswith(
                        "\u2022"
                    ):
                        names.append(line)
            s["semantic_id"] = f"team_bio_{s['slide_idx']}"
            s["member_names"] = names[:6]
            # Normalize team family from layout name
            s["team_family"] = _detect_team_family(ln, s)
            continue

        # ── Case studies (Services - Cases / Detailed Case) ───────────
        if _RE_CASE.search(ln):
            s["classification"] = "pool_case_study"
            s["semantic_id"] = f"case_study_{s['slide_idx']}"
            s["category"] = "general"
            continue

        # ── Methodology example content ────────────────────────────────
        if _RE_METH.search(ln):
            s["classification"] = "example_content"
            s["semantic_id"] = f"methodology_example_{s['slide_idx']}"
            continue

        # ── General content layouts (Master 1, variable) ──────────────
        if mi == 1:
            s["classification"] = "example_content"
            s["semantic_id"] = f"content_example_{s['slide_idx']}"
            continue

        # ── Fallback ───────────────────────────────────────────────────
        s["classification"] = "unclassified"
        s["semantic_id"] = f"slide_{s['slide_idx']}"

    return slides


def build_a2_allowlists(
    slides: list[dict], prs: Presentation | None = None
) -> None:
    """Build effective allowlists for A2 shells.

    Effective = considers ALL text-bearing regions across three levels:
    slide-local, layout-inherited, master-inherited.

    Sets ``allowlist`` on each A2 slide dict with:
    - approved_injection_placeholders (from classification)
    - candidate_preserved_text_regions (all text shapes by source level)
    - candidate_preserved_tables (all table shapes by source level)
    - must_clear_regions (text regions NOT on any approved list)
    - unresolved_regions (regions that need human review)

    Human reviewer validates during PHASE 0 review.
    """
    for s in slides:
        if s.get("classification") != "a2_shell":
            continue

        approved_ph = s.get("approved_placeholders", [])
        slide_idx = s["slide_idx"]

        # Collect text regions from all three levels
        all_text_regions: list[dict] = []
        all_table_regions: list[dict] = []

        # 1. Slide-local shapes
        for shape in s.get("shapes_detail", []):
            if shape.get("is_placeholder"):
                continue
            if shape.get("has_table"):
                all_table_regions.append({
                    "shape_name": shape["name"],
                    "source": "slide",
                    "cell_preview": shape.get("table_cell_preview", []),
                })
            elif shape.get("has_text_frame"):
                txt = shape.get("text_preview", "").strip()
                if txt:
                    all_text_regions.append({
                        "shape_name": shape["name"],
                        "source": "slide",
                        "text_preview": txt[:200],
                    })

        # 2. Layout-inherited and master-inherited shapes (if prs available)
        if prs is not None and slide_idx < len(prs.slides):
            slide_obj = prs.slides[slide_idx]
            layout = slide_obj.slide_layout
            master = layout.slide_master

            for src_label, shapes in [("layout", layout.shapes), ("master", master.shapes)]:
                for shape in shapes:
                    if shape.is_placeholder:
                        continue
                    if shape.has_table:
                        cells = []
                        for row in shape.table.rows:
                            for cell in row.cells:
                                ct = cell.text.strip()
                                if ct:
                                    cells.append(ct[:100])
                        if cells:
                            all_table_regions.append({
                                "shape_name": shape.name,
                                "source": src_label,
                                "cell_preview": cells[:10],
                            })
                    elif shape.has_text_frame:
                        txt = shape.text_frame.text.strip()
                        if txt:
                            all_text_regions.append({
                                "shape_name": shape.name,
                                "source": src_label,
                                "text_preview": txt[:200],
                            })

        # Classify: regions to preserve vs clear vs unresolved
        # Slide-local non-placeholder text = must_clear (unless approved)
        # Layout/master text = candidate_preserved (institutional branding)
        must_clear: list[dict] = []
        candidate_preserved: list[dict] = []
        for region in all_text_regions:
            if region["source"] == "slide":
                must_clear.append(region)
            else:
                candidate_preserved.append(region)

        preserved_tables: list[dict] = []
        must_clear_tables: list[dict] = []
        for table in all_table_regions:
            if table["source"] == "slide":
                must_clear_tables.append(table)
            else:
                preserved_tables.append(table)

        s["allowlist"] = {
            "shell_id": s.get("semantic_id", ""),
            "approved_injection_placeholders": {
                idx: f"ph_{idx}" for idx in approved_ph
            },
            "candidate_preserved_text_regions": candidate_preserved,
            "candidate_preserved_tables": preserved_tables,
            "must_clear_text_regions": must_clear,
            "must_clear_tables": must_clear_tables,
            "effective_text_region_count": len(all_text_regions),
            "effective_table_count": len(all_table_regions),
            "source_breakdown": {
                "slide": len([r for r in all_text_regions if r["source"] == "slide"]),
                "layout": len([r for r in all_text_regions if r["source"] == "layout"]),
                "master": len([r for r in all_text_regions if r["source"] == "master"]),
            },
        }


# ── Specialized detections ─────────────────────────────────────────────────


def detect_section_dividers(slides: list[dict]) -> list[dict]:
    """Extract section divider details."""
    return [
        {
            "slide_idx": s["slide_idx"],
            "divider_number": s["semantic_id"].replace("section_", ""),
            "layout_name": s["layout_name"],
            "title": s["text"].get("ph_0", ""),
            "subtitle": s["text"].get("ph_10", ""),
        }
        for s in slides
        if s.get("classification") == "a2_shell"
        and s.get("semantic_id", "").startswith("section_")
    ]


def detect_methodology_family(layouts: list[dict]) -> list[dict]:
    """Detect methodology-related layouts with semantic layout IDs."""
    result: list[dict] = []
    for lay in layouts:
        if _RE_METH.search(lay["name"]):
            variant = _meth_variant(lay["name"])
            result.append({
                "layout_name": lay["name"],
                "semantic_layout_id": f"methodology_{variant}",
                "placeholder_count": lay["placeholder_count"],
                "placeholders": lay["placeholders"],
                "variant": variant,
            })
    return result


def _detect_team_family(layout_name: str, slide: dict) -> str:
    """Detect team bio family variant from layout name and slide content.

    Returns one of: team_two_members, team_single_member, team_variant.
    """
    nl = layout_name.lower().strip()
    if "two" in nl or "2" in nl:
        return "team_two_members"
    if "single" in nl or "one" in nl or "1" in nl:
        return "team_single_member"
    # Heuristic: count member names detected to infer variant
    member_count = len(slide.get("member_names", []))
    if member_count <= 1:
        return "team_single_member"
    if member_count == 2:
        return "team_two_members"
    return "team_variant"


def _meth_variant(name: str) -> str:
    nl = name.lower()
    if "overview" in nl and "-4-" in nl:
        return "overview_4"
    if "overview" in nl and "-3-" in nl:
        return "overview_3"
    if "focused" in nl and "-4-" in nl:
        return "focused_4"
    if "focused" in nl and "-3-" in nl:
        return "focused_3"
    if "detail" in nl:
        return "detail"
    if "overview" in nl:
        return "overview"
    if "focused" in nl:
        return "focused"
    return "unknown"


def detect_team_pool(slides: list[dict]) -> list[dict]:
    """Team bio slide pool with family normalization."""
    return [
        {
            "slide_idx": s["slide_idx"],
            "layout_name": s["layout_name"],
            "semantic_id": s.get("semantic_id", ""),
            "team_family": s.get("team_family", "team_variant"),
            "member_names": s.get("member_names", []),
        }
        for s in slides
        if s.get("classification") == "pool_team_bio"
    ]


def detect_case_study_pool(slides: list[dict]) -> list[dict]:
    """Case study slide pool (disjoint from service dividers)."""
    return [
        {
            "slide_idx": s["slide_idx"],
            "layout_name": s["layout_name"],
            "semantic_id": s.get("semantic_id", ""),
            "category": s.get("category", "general"),
            "text_preview": list(s["text"].values())[:3],
        }
        for s in slides
        if s.get("classification") == "pool_case_study"
    ]


def detect_service_divider_pool(slides: list[dict]) -> list[dict]:
    """Service category divider slides (disjoint from case studies)."""
    return [
        {
            "slide_idx": s["slide_idx"],
            "layout_name": s["layout_name"],
            "semantic_id": s.get("semantic_id", ""),
            "service_category": s.get("service_category", ""),
        }
        for s in slides
        if s.get("classification") == "pool_service_divider"
    ]


def extract_forbidden_fingerprints(slides: list[dict]) -> dict:
    """Extract text fingerprints from example/forbidden template content.

    These are distinctive phrases from template example slides (methodology
    examples, Film Sector, Tadawul, etc.) that must NEVER appear in generated
    output. Used by anti-leak tests.

    Returns dict with:
    - forbidden_phrases: list of distinctive 4+ word phrases
    - source_slides: which slides contributed fingerprints
    - keyword_set: unique keywords for fast matching
    """
    forbidden_phrases: list[str] = []
    source_slides: list[dict] = []
    keyword_set: set[str] = set()

    # Known forbidden project names
    _FORBIDDEN_PROJECTS = {
        "film sector", "tadawul", "saudi exchange", "film authority",
        "general authority for media regulation",
    }

    for s in slides:
        cls = s.get("classification", "")
        if cls not in ("example_content",):
            continue

        all_text = " ".join(v for v in s["text"].values() if v.strip())
        if not all_text.strip():
            continue

        # Extract distinctive phrases (sentences/fragments > 4 words)
        phrases: list[str] = []
        for line in all_text.split("\n"):
            line = line.strip()
            words = line.split()
            if len(words) >= 4 and len(line) > 20:
                # Take up to 15 words as fingerprint
                fp = " ".join(words[:15])
                phrases.append(fp)

        # Check for known forbidden project names
        text_lower = all_text.lower()
        for proj in _FORBIDDEN_PROJECTS:
            if proj in text_lower:
                keyword_set.add(proj)

        if phrases:
            forbidden_phrases.extend(phrases[:5])  # Top 5 per slide
            source_slides.append({
                "slide_idx": s["slide_idx"],
                "semantic_id": s.get("semantic_id", ""),
                "classification": cls,
                "phrase_count": len(phrases),
            })

    # Also check pool_case_study slides for Tadawul/Film references
    for s in slides:
        if s.get("classification") != "pool_case_study":
            continue
        all_text = " ".join(v for v in s["text"].values() if v.strip())
        text_lower = all_text.lower()
        for proj in _FORBIDDEN_PROJECTS:
            if proj in text_lower:
                keyword_set.add(proj)
                # Extract phrases from contaminated case studies too
                for line in all_text.split("\n"):
                    line = line.strip()
                    words = line.split()
                    if len(words) >= 4 and len(line) > 20:
                        forbidden_phrases.append(" ".join(words[:15]))
                source_slides.append({
                    "slide_idx": s["slide_idx"],
                    "semantic_id": s.get("semantic_id", ""),
                    "classification": "pool_case_study",
                    "phrase_count": 0,
                    "contamination": "forbidden_project_name",
                })

    return {
        "forbidden_phrases": sorted(set(forbidden_phrases)),
        "forbidden_phrase_count": len(set(forbidden_phrases)),
        "source_slide_count": len(source_slides),
        "source_slides": source_slides,
        "keyword_set": sorted(keyword_set),
    }


def detect_contact_closing(slides: list[dict]) -> dict:
    """Contact and closing slide indices."""
    return {
        "know_more": [
            s["slide_idx"]
            for s in slides
            if s.get("semantic_id") == "know_more"
        ],
        "contact": [
            s["slide_idx"]
            for s in slides
            if s.get("semantic_id") == "contact"
        ],
    }


# ── Parity audit ───────────────────────────────────────────────────────────


def run_parity_audit(
    en_layouts: list[dict],
    ar_layouts: list[dict],
    en_slides: list[dict],
    ar_slides: list[dict],
) -> dict:
    """Compare EN and AR templates for structural parity."""
    checks: list[dict] = []
    critical = 0
    warnings = 0

    def _check(name: str, en_val: Any, ar_val: Any, sev: str = "critical"):
        nonlocal critical, warnings
        match = en_val == ar_val
        if not match:
            if sev == "critical":
                critical += 1
            else:
                warnings += 1
        checks.append({
            "check": name,
            "en": en_val,
            "ar": ar_val,
            "match": match,
            "severity": sev,
        })

    _check("slide_count", len(en_slides), len(ar_slides))
    _check("layout_count", len(en_layouts), len(ar_layouts))

    # Semantic layout ID set parity (critical — must resolve same IDs)
    en_sem_ids = sorted({
        lay.get(
            "semantic_layout_id",
            _layout_name_to_semantic_id(lay["name"]),
        )
        for lay in en_layouts
    })
    ar_sem_ids = sorted({
        lay.get(
            "semantic_layout_id",
            _layout_name_to_semantic_id(lay["name"]),
        )
        for lay in ar_layouts
    })
    _check("semantic_layout_id_set", en_sem_ids, ar_sem_ids)

    # Layout placeholder index parity (matched by semantic layout ID)
    en_ph = {
        lay.get("semantic_layout_id", lay["name"]): sorted(
            p["idx"] for p in lay["placeholders"]
        )
        for lay in en_layouts
    }
    ar_ph = {
        lay.get("semantic_layout_id", lay["name"]): sorted(
            p["idx"] for p in lay["placeholders"]
        )
        for lay in ar_layouts
    }
    common = set(en_ph) & set(ar_ph)
    for sem_id in sorted(common):
        if en_ph[sem_id] != ar_ph[sem_id]:
            checks.append({
                "check": f"placeholder_parity:{sem_id}",
                "en": en_ph[sem_id],
                "ar": ar_ph[sem_id],
                "match": False,
                "severity": "critical",
            })
            critical += 1

    _check(
        "layout_name_count",
        len(en_ph),
        len(ar_ph),
        "warning",
    )

    # Layout sequence (may differ in language-specific naming)
    en_seq = [s.get("semantic_layout_id", s["layout_name"]) for s in en_slides]
    ar_seq = [s.get("semantic_layout_id", s["layout_name"]) for s in ar_slides]
    _check("semantic_layout_sequence", en_seq, ar_seq, "warning")

    # Section divider parity
    en_div = sorted(
        s["semantic_id"]
        for s in en_slides
        if s.get("semantic_id", "").startswith("section_")
    )
    ar_div = sorted(
        s["semantic_id"]
        for s in ar_slides
        if s.get("semantic_id", "").startswith("section_")
    )
    _check("section_dividers", en_div, ar_div)

    # Media count parity
    en_media = [s["media_count"] for s in en_slides]
    ar_media = [s["media_count"] for s in ar_slides]
    _check("media_counts", en_media, ar_media, "warning")

    return {
        "overall_pass": critical == 0,
        "critical_failures": critical,
        "warnings": warnings,
        "checks": checks,
    }


# ── Grammar extraction ─────────────────────────────────────────────────────


def extract_grammar(
    slides: list[dict], dividers: list[dict]
) -> dict[str, dict]:
    """Extract template grammar artifacts from slide text content."""
    grammar: dict[str, Any] = {}

    # Section naming grammar
    section_naming: dict[str, dict] = {}
    for d in dividers:
        section_naming[d["divider_number"]] = {
            "title": d["title"],
            "subtitle": d["subtitle"],
        }
    grammar["section_naming_grammar"] = section_naming

    # Intro letter grammar
    intro = [
        s for s in slides if s.get("semantic_id") == "intro_message"
    ]
    if intro:
        s = intro[0]
        grammar["intro_letter_grammar"] = {
            "greeting_pattern": s["text"].get("ph_0", ""),
            "body_text": s["text"].get("ph_1", ""),
            "subtitle": s["text"].get("ph_13", ""),
            "footer_fields": {
                k: v
                for k, v in s["text"].items()
                if k.startswith("ph_1")
                and k not in ("ph_1", "ph_10", "ph_13")
            },
        }

    # Methodology framing
    meth = [
        s
        for s in slides
        if s.get("classification") == "example_content"
        and "methodology" in s.get("semantic_id", "")
    ]
    if meth:
        all_meth = "\n".join(
            v for s in meth for v in s["text"].values()
        )
        grammar["methodology_framing"] = {
            "example_slide_count": len(meth),
            "activity_prefix_patterns": _find_patterns(all_meth, [
                r"The team will\b",
                r"سيقوم الفريق",
                r"We will\b",
                r"The project\b",
            ]),
            "phase_keywords": _find_patterns(all_meth, [
                r"Phase \d",
                r"المرحلة",
                r"Deliverable",
                r"مخرج",
            ]),
            "full_text_preview": all_meth[:2000],
        }

    # Timeline / deliverables grammar
    timeline = [
        s
        for s in slides
        if s.get("classification") == "example_content"
        and any(
            kw in v.lower()
            for v in s["text"].values()
            for kw in ("deliverable", "timeline", "مخرج")
        )
    ]
    if timeline:
        grammar["timeline_grammar"] = {
            "example_text": {
                k: v for s in timeline for k, v in s["text"].items()
            },
            "column_patterns": [
                "Item",
                "Description",
                "Unit",
                "Quantity",
                "Timeline",
            ],
        }

    # Team bio grammar
    team = [
        s for s in slides if s.get("classification") == "pool_team_bio"
    ]
    if team:
        sample = team[0]
        grammar["team_bio_grammar"] = {
            "pool_size": len(team),
            "layout_name": sample["layout_name"],
            "placeholder_count": sample["placeholder_count"],
            "sample_text": {
                k: v[:200] for k, v in sample["text"].items()
            },
            "bio_structure": [
                "experience_level",
                "name",
                "title",
                "career_summary",
                "sectors",
                "education",
            ],
        }

    # Case study grammar
    cases = [
        s
        for s in slides
        if s.get("classification") == "pool_case_study"
        and _RE_CASE.search(s["layout_name"])
    ]
    if cases:
        sample = cases[0]
        grammar["case_study_grammar"] = {
            "pool_size": len(cases),
            "layout_name": sample["layout_name"],
            "framing": "Challenge -> Approach -> Impact",
            "fields": ["client", "project_title", "outcome", "services"],
            "sample_text": {
                k: v[:200] for k, v in sample["text"].items()
            },
        }

    # Governance grammar (structural — may not be explicit in template)
    grammar["governance_grammar"] = {
        "tier_names": [
            "Steering Committee",
            "Operational Committee",
            "Delivery Team",
            "Escalation Path",
        ],
        "tier_structure": {
            "name": "",
            "composition": "",
            "frequency": "",
            "responsibility": "",
        },
    }

    # Closing grammar
    closing = [
        s
        for s in slides
        if s.get("semantic_id") in ("know_more", "contact")
    ]
    if closing:
        grammar["closing_grammar"] = {
            "slides": [
                {
                    "semantic_id": s["semantic_id"],
                    "text": s["text"],
                }
                for s in closing
            ],
        }

    return grammar


def _find_patterns(text: str, patterns: list[str]) -> list[str]:
    """Return patterns that match in text."""
    return [p for p in patterns if re.search(p, text, re.I)]


# ── Catalog lock ───────────────────────────────────────────────────────────


def generate_catalog_lock(
    slides: list[dict],
    layouts: list[dict],
    template_hash: str,
    language: str,
) -> dict:
    """Catalog lock: semantic IDs -> validated template positions.

    Maps both semantic asset IDs (slides) and semantic layout IDs (layouts).
    Runtime code resolves everything through this lock.
    """
    a1: dict[str, dict] = {}
    a2: dict[str, dict] = {}
    dividers: dict[str, dict] = {}
    case_pool: dict[str, list[dict]] = {}
    team_pool: list[dict] = []
    service_divider_pool: list[dict] = []

    for s in slides:
        cls = s.get("classification", "")
        sid = s.get("semantic_id", "")
        sem_layout = s.get("semantic_layout_id", "")

        if cls == "a1_immutable":
            a1[sid] = {
                "slide_idx": s["slide_idx"],
                "semantic_layout_id": sem_layout,
                "display_name": s["layout_name"],
                "shape_count": s["shape_count"],
                "media_count": s["media_count"],
            }
        elif cls == "a2_shell":
            allowlist = s.get("allowlist", {})
            entry = {
                "slide_idx": s["slide_idx"],
                "semantic_layout_id": sem_layout,
                "display_name": s["layout_name"],
                "allowlist": allowlist,
            }
            if sid.startswith("section_"):
                dividers[sid.replace("section_", "")] = entry
            else:
                a2[sid] = entry
        elif cls == "pool_case_study":
            cat = s.get("category", "general")
            case_pool.setdefault(cat, []).append({
                "slide_idx": s["slide_idx"],
                "semantic_layout_id": sem_layout,
                "display_name": s["layout_name"],
                "semantic_id": sid,
            })
        elif cls == "pool_team_bio":
            team_pool.append({
                "slide_idx": s["slide_idx"],
                "semantic_layout_id": sem_layout,
                "display_name": s["layout_name"],
                "semantic_id": sid,
                "team_family": s.get("team_family", "team_variant"),
                "member_names": s.get("member_names", []),
            })
        elif cls == "pool_service_divider":
            service_divider_pool.append({
                "slide_idx": s["slide_idx"],
                "semantic_layout_id": sem_layout,
                "display_name": s["layout_name"],
                "semantic_id": sid,
                "service_category": s.get("service_category", ""),
            })

    # Layout catalog keyed by semantic layout ID (not display name)
    layout_cat: dict[str, dict] = {}
    for lay in layouts:
        sem_id = lay.get(
            "semantic_layout_id",
            _layout_name_to_semantic_id(lay["name"]),
        )
        layout_cat[sem_id] = {
            "display_name": lay["name"],
            "master_idx": lay["master_idx"],
            "placeholder_count": lay["placeholder_count"],
            "placeholders": {
                str(p["idx"]): p["type"] for p in lay["placeholders"]
            },
        }

    return {
        "template_hash": template_hash,
        "language": language,
        "generated_at": datetime.now(UTC).isoformat(),
        "a1_immutable": a1,
        "a2_shells": a2,
        "section_dividers": dividers,
        "case_study_pool": case_pool,
        "team_bio_pool": team_pool,
        "service_divider_pool": service_divider_pool,
        "layouts": layout_cat,
    }


# ── Example alignment ──────────────────────────────────────────────────────


def run_example_alignment(
    catalog_lock: dict, example_dir: Path
) -> dict:
    """Check alignment between catalog lock and real example decks."""
    report: dict[str, Any] = {"examples": [], "summary": {}}

    for ex_dir in sorted(example_dir.iterdir()):
        if not ex_dir.is_dir() or not ex_dir.name.startswith("example_"):
            continue
        pptx_files = list(ex_dir.glob("*.pptx"))
        if not pptx_files:
            continue

        pptx_path = pptx_files[0]
        try:
            prs = Presentation(str(pptx_path))
        except Exception as e:
            report["examples"].append({
                "name": ex_dir.name,
                "error": str(e),
            })
            continue

        layout_names = [s.slide_layout.name for s in prs.slides]
        catalog_layouts_set = set(catalog_lock.get("layouts", {}).keys())
        example_layouts_set = set(layout_names)

        # Divider flow
        divider_flow: list[dict] = []
        for i, ln in enumerate(layout_names):
            if _RE_DIVIDER.search(ln):
                m = re.search(r"(0[1-9])", ln)
                if m:
                    divider_flow.append({
                        "slide_idx": i,
                        "divider": m.group(1),
                    })

        report["examples"].append({
            "name": ex_dir.name,
            "file": pptx_path.name,
            "slide_count": len(prs.slides),
            "layout_usage": {
                ln: layout_names.count(ln)
                for ln in sorted(set(layout_names))
            },
            "layouts_in_common": len(
                catalog_layouts_set & example_layouts_set
            ),
            "layouts_only_in_example": sorted(
                example_layouts_set - catalog_layouts_set
            ),
            "section_divider_flow": divider_flow,
        })

    report["summary"] = {
        "examples_analyzed": len(report["examples"]),
        "catalog_layout_count": len(
            catalog_lock.get("layouts", {})
        ),
    }
    return report


# ── JSON writer ────────────────────────────────────────────────────────────


def _write_json(path: Path, data: Any) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    with open(path, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2, default=str)


# ── House slide registry builder ───────────────────────────────────────────


def _build_house_registry(slides: list[dict]) -> dict:
    """Group classified slides into registry structure."""
    reg: dict[str, Any] = {
        "a1": {},
        "a2": {},
        "example": {},
        "pool_team": [],
        "pool_case": [],
        "pool_service_divider": [],
        "unclassified": [],
    }
    for s in slides:
        cls = s.get("classification", "")
        sid = s.get("semantic_id", "")
        entry = {
            "slide_idx": s["slide_idx"],
            "layout_name": s["layout_name"],
            "semantic_layout_id": s.get("semantic_layout_id", ""),
            "semantic_id": sid,
            "shape_count": s["shape_count"],
            "media_count": s["media_count"],
        }
        if cls == "a1_immutable":
            reg["a1"][sid] = entry
        elif cls == "a2_shell":
            reg["a2"][sid] = {
                **entry,
                "approved_placeholders": s.get(
                    "approved_placeholders", []
                ),
                "allowlist": s.get("allowlist", {}),
            }
        elif cls == "example_content":
            reg["example"][sid] = entry
        elif cls == "pool_team_bio":
            reg["pool_team"].append({
                **entry,
                "team_family": s.get("team_family", "team_variant"),
                "member_names": s.get("member_names", []),
            })
        elif cls == "pool_case_study":
            reg["pool_case"].append(entry)
        elif cls == "pool_service_divider":
            reg["pool_service_divider"].append({
                **entry,
                "service_category": s.get("service_category", ""),
            })
        else:
            reg["unclassified"].append(entry)
    return reg


# ── Main runner ────────────────────────────────────────────────────────────


def run_full_audit(
    en_potx: Path = EN_POTX,
    ar_potx: Path = AR_POTX,
    example_dir: Path = EXAMPLE_DIR,
    output_base: Path | None = None,
) -> dict[str, Path]:
    """Run complete PHASE 0 audit. Returns artifact name → file path."""
    if output_base is None:
        output_base = Path.cwd()

    audit_dir = output_base / "audit"
    grammar_dir = output_base / "src" / "data" / "template_grammar"
    data_dir = output_base / "src" / "data"
    audit_dir.mkdir(parents=True, exist_ok=True)
    grammar_dir.mkdir(parents=True, exist_ok=True)

    artifacts: dict[str, Path] = {}
    patch_dir = Path(tempfile.mkdtemp(prefix="phase0_"))

    try:
        # ── Patch & load templates ─────────────────────────────────────
        print("=" * 60)
        print("PHASE 0 — Template Intelligence Audit")
        print("=" * 60)

        print("\n[1/16] Patching POTX templates...")
        en_pptx = patch_potx(en_potx, patch_dir)
        ar_pptx = patch_potx(ar_potx, patch_dir)
        en_hash = file_hash(en_potx)
        ar_hash = file_hash(ar_potx)

        print("[2/16] Loading EN template...")
        en_prs = Presentation(str(en_pptx))
        print(f"  ->{len(en_prs.slides)} slides loaded")

        print("[3/16] Loading AR template...")
        ar_prs = Presentation(str(ar_pptx))
        print(f"  ->{len(ar_prs.slides)} slides loaded")

        # ── Layout catalogs ────────────────────────────────────────────
        print("[4/16] Cataloging layouts...")
        en_layouts = catalog_layouts(en_prs)
        ar_layouts = catalog_layouts(ar_prs)
        _write_json(audit_dir / "layout_catalog_en.json", en_layouts)
        _write_json(audit_dir / "layout_catalog_ar.json", ar_layouts)
        artifacts["layout_catalog_en"] = audit_dir / "layout_catalog_en.json"
        artifacts["layout_catalog_ar"] = audit_dir / "layout_catalog_ar.json"
        print(f"  ->EN: {len(en_layouts)} layouts, AR: {len(ar_layouts)} layouts")

        # ── Placeholder maps ───────────────────────────────────────────
        print("[5/16] Building placeholder maps...")
        en_ph_map = {
            lay["name"]: lay["placeholders"] for lay in en_layouts
        }
        ar_ph_map = {
            lay["name"]: lay["placeholders"] for lay in ar_layouts
        }
        _write_json(audit_dir / "placeholder_map_en.json", en_ph_map)
        _write_json(audit_dir / "placeholder_map_ar.json", ar_ph_map)
        artifacts["placeholder_map_en"] = audit_dir / "placeholder_map_en.json"
        artifacts["placeholder_map_ar"] = audit_dir / "placeholder_map_ar.json"

        # ── Slide audit + classification ───────────────────────────────
        print("[6/16] Auditing & classifying slides...")
        en_slides = classify_slides(audit_all_slides(en_prs))
        ar_slides = classify_slides(audit_all_slides(ar_prs))

        # ── A2 allowlists ─────────────────────────────────────────────
        print("[7/16] Building effective A2 shell allowlists...")
        build_a2_allowlists(en_slides, en_prs)
        build_a2_allowlists(ar_slides, ar_prs)
        a2_count_en = sum(
            1 for s in en_slides if s.get("classification") == "a2_shell"
        )
        a2_count_ar = sum(
            1 for s in ar_slides if s.get("classification") == "a2_shell"
        )
        print(f"  ->EN: {a2_count_en} shells with allowlists")
        print(f"  ->AR: {a2_count_ar} shells with allowlists")

        # ── House slide registries ─────────────────────────────────────
        print("[8/16] Building house slide registries...")
        en_house = _build_house_registry(en_slides)
        ar_house = _build_house_registry(ar_slides)
        _write_json(
            audit_dir / "house_slide_registry_en.json", en_house
        )
        _write_json(
            audit_dir / "house_slide_registry_ar.json", ar_house
        )
        artifacts["house_slide_registry_en"] = (
            audit_dir / "house_slide_registry_en.json"
        )
        artifacts["house_slide_registry_ar"] = (
            audit_dir / "house_slide_registry_ar.json"
        )
        print(
            f"  ->EN: {len(en_house['a1'])} A1, "
            f"{len(en_house['a2'])} A2, "
            f"{len(en_house['pool_team'])} team, "
            f"{len(en_house['pool_case'])} case, "
            f"{len(en_house['pool_service_divider'])} svc-div, "
            f"{len(en_house['example'])} example"
        )
        print(
            f"  ->AR: {len(ar_house['a1'])} A1, "
            f"{len(ar_house['a2'])} A2, "
            f"{len(ar_house['pool_team'])} team, "
            f"{len(ar_house['pool_case'])} case, "
            f"{len(ar_house['pool_service_divider'])} svc-div, "
            f"{len(ar_house['example'])} example"
        )

        # ── Specialized detections ─────────────────────────────────────
        print("[9/16] Detecting structures...")
        en_dividers = detect_section_dividers(en_slides)
        ar_dividers = detect_section_dividers(ar_slides)
        _write_json(
            audit_dir / "section_divider_map.json",
            {"en": en_dividers, "ar": ar_dividers},
        )
        artifacts["section_divider_map"] = (
            audit_dir / "section_divider_map.json"
        )

        meth_family = detect_methodology_family(en_layouts)
        _write_json(audit_dir / "methodology_family.json", meth_family)
        artifacts["methodology_family"] = (
            audit_dir / "methodology_family.json"
        )

        en_team = detect_team_pool(en_slides)
        _write_json(audit_dir / "team_pool.json", en_team)
        artifacts["team_pool"] = audit_dir / "team_pool.json"

        en_cases = detect_case_study_pool(en_slides)
        _write_json(audit_dir / "case_study_pool.json", en_cases)
        artifacts["case_study_pool"] = audit_dir / "case_study_pool.json"

        svc_div = detect_service_divider_pool(en_slides)
        _write_json(audit_dir / "service_divider_pool.json", svc_div)
        artifacts["service_divider_pool"] = (
            audit_dir / "service_divider_pool.json"
        )

        contact = detect_contact_closing(en_slides)
        _write_json(audit_dir / "contact_closing_map.json", contact)
        artifacts["contact_closing_map"] = (
            audit_dir / "contact_closing_map.json"
        )

        print(
            f"  ->Dividers: EN={len(en_dividers)} AR={len(ar_dividers)}"
        )
        print(f"  ->Methodology layouts: {len(meth_family)}")
        print(f"  ->Team pool: {len(en_team)} slides")
        print(f"  ->Case study pool: {len(en_cases)} slides")
        print(f"  ->Service divider pool: {len(svc_div)} slides")

        # ── Text extraction ────────────────────────────────────────────
        print("[10/16] Extracting template text...")
        _write_json(
            audit_dir / "template_text_extraction_en.json",
            {str(s["slide_idx"]): s["text"] for s in en_slides},
        )
        _write_json(
            audit_dir / "template_text_extraction_ar.json",
            {str(s["slide_idx"]): s["text"] for s in ar_slides},
        )
        artifacts["template_text_extraction_en"] = (
            audit_dir / "template_text_extraction_en.json"
        )
        artifacts["template_text_extraction_ar"] = (
            audit_dir / "template_text_extraction_ar.json"
        )

        # ── Parity report ──────────────────────────────────────────────
        print("[11/16] Running EN<->AR parity audit...")
        parity = run_parity_audit(
            en_layouts, ar_layouts, en_slides, ar_slides
        )
        _write_json(audit_dir / "parity_report.json", parity)
        artifacts["parity_report"] = audit_dir / "parity_report.json"
        print(
            f"  ->Pass: {parity['overall_pass']}  "
            f"Critical: {parity['critical_failures']}  "
            f"Warnings: {parity['warnings']}"
        )

        # ── Grammar artifacts ──────────────────────────────────────────
        print("[12/16] Extracting template grammar...")
        grammar = extract_grammar(en_slides, en_dividers)
        for name, data in grammar.items():
            path = grammar_dir / f"{name}.json"
            _write_json(path, data)
            artifacts[f"grammar_{name}"] = path
        print(f"  ->{len(grammar)} grammar artifacts")

        # ── Catalog locks ──────────────────────────────────────────────
        print("[13/16] Generating catalog locks...")
        en_lock = generate_catalog_lock(
            en_slides, en_layouts, en_hash, "en"
        )
        ar_lock = generate_catalog_lock(
            ar_slides, ar_layouts, ar_hash, "ar"
        )
        _write_json(data_dir / "catalog_lock_en.json", en_lock)
        _write_json(data_dir / "catalog_lock_ar.json", ar_lock)
        artifacts["catalog_lock_en"] = data_dir / "catalog_lock_en.json"
        artifacts["catalog_lock_ar"] = data_dir / "catalog_lock_ar.json"

        # ── Effective slide structure (master-aware) ─────────────────
        print("[14/16] Building effective slide structure (master-aware)...")
        en_eff = build_effective_slide_structure(en_prs)
        ar_eff = build_effective_slide_structure(ar_prs)
        _write_json(audit_dir / "effective_slide_structure_en.json", en_eff)
        _write_json(audit_dir / "effective_slide_structure_ar.json", ar_eff)
        artifacts["effective_slide_structure_en"] = (
            audit_dir / "effective_slide_structure_en.json"
        )
        artifacts["effective_slide_structure_ar"] = (
            audit_dir / "effective_slide_structure_ar.json"
        )
        # Count slides with inherited-only content
        inherited_only = sum(
            1 for s in en_eff
            if s["slide_local_shape_count"] == 0
            and (s["layout_shape_count"] > 0 or s["master_shape_count"] > 0)
        )
        print(f"  ->EN: {len(en_eff)} slides, {inherited_only} inherited-only")

        # ── Anti-leak fingerprints ────────────────────────────────────
        print("[15/16] Extracting forbidden example fingerprints...")
        en_fingerprints = extract_forbidden_fingerprints(en_slides)
        ar_fingerprints = extract_forbidden_fingerprints(ar_slides)
        _write_json(
            audit_dir / "forbidden_example_fingerprints_en.json",
            en_fingerprints,
        )
        _write_json(
            audit_dir / "forbidden_example_fingerprints_ar.json",
            ar_fingerprints,
        )
        artifacts["forbidden_fingerprints_en"] = (
            audit_dir / "forbidden_example_fingerprints_en.json"
        )
        artifacts["forbidden_fingerprints_ar"] = (
            audit_dir / "forbidden_example_fingerprints_ar.json"
        )
        print(
            f"  ->EN: {en_fingerprints['forbidden_phrase_count']} phrases "
            f"from {en_fingerprints['source_slide_count']} slides"
        )
        print(
            f"  ->AR: {ar_fingerprints['forbidden_phrase_count']} phrases "
            f"from {ar_fingerprints['source_slide_count']} slides"
        )

        # ── Pool disjointness proof ──────────────────────────────────
        print("[16/16] Verifying pool disjointness...")
        case_indices = {s["slide_idx"] for s in en_slides if s.get("classification") == "pool_case_study"}
        svc_indices = {s["slide_idx"] for s in en_slides if s.get("classification") == "pool_service_divider"}
        team_indices = {s["slide_idx"] for s in en_slides if s.get("classification") == "pool_team_bio"}
        overlap_case_svc = case_indices & svc_indices
        overlap_case_team = case_indices & team_indices
        overlap_svc_team = svc_indices & team_indices
        disjointness = {
            "case_study_indices": sorted(case_indices),
            "service_divider_indices": sorted(svc_indices),
            "team_bio_indices": sorted(team_indices),
            "overlap_case_svc": sorted(overlap_case_svc),
            "overlap_case_team": sorted(overlap_case_team),
            "overlap_svc_team": sorted(overlap_svc_team),
            "all_pools_disjoint": (
                len(overlap_case_svc) == 0
                and len(overlap_case_team) == 0
                and len(overlap_svc_team) == 0
            ),
        }
        _write_json(audit_dir / "pool_disjointness_proof.json", disjointness)
        artifacts["pool_disjointness_proof"] = audit_dir / "pool_disjointness_proof.json"
        disjoint_str = "PASS" if disjointness["all_pools_disjoint"] else "FAIL"
        print(f"  ->Pools disjoint: {disjoint_str}")
        if not disjointness["all_pools_disjoint"]:
            print(f"  ->OVERLAP case/svc: {overlap_case_svc}")
            print(f"  ->OVERLAP case/team: {overlap_case_team}")
            print(f"  ->OVERLAP svc/team: {overlap_svc_team}")

        # ── Example alignment ──────────────────────────────────────────
        print("[bonus] Running example alignment...")
        alignment = run_example_alignment(en_lock, example_dir)
        _write_json(
            audit_dir / "example_alignment_report.json", alignment
        )
        artifacts["example_alignment"] = (
            audit_dir / "example_alignment_report.json"
        )
        print(
            f"  ->{alignment['summary']['examples_analyzed']} "
            "examples analyzed"
        )

        # ── Summary ────────────────────────────────────────────────────
        summary = {
            "phase": "PHASE_0",
            "status": "COMPLETE",
            "timestamp": datetime.now(UTC).isoformat(),
            "en_template": str(en_potx),
            "ar_template": str(ar_potx),
            "en_hash": en_hash,
            "ar_hash": ar_hash,
            "en_slides": len(en_slides),
            "ar_slides": len(ar_slides),
            "en_layouts": len(en_layouts),
            "ar_layouts": len(ar_layouts),
            "parity_pass": parity["overall_pass"],
            "artifact_count": len(artifacts),
            "artifacts": {k: str(v) for k, v in artifacts.items()},
            "classification_summary": {
                "en": {
                    "a1_immutable": len(en_house["a1"]),
                    "a2_shell": len(en_house["a2"]),
                    "pool_team_bio": len(en_house["pool_team"]),
                    "pool_case_study": len(en_house["pool_case"]),
                    "pool_service_divider": len(en_house["pool_service_divider"]),
                    "example_content": len(en_house["example"]),
                    "unclassified": len(en_house["unclassified"]),
                },
                "ar": {
                    "a1_immutable": len(ar_house["a1"]),
                    "a2_shell": len(ar_house["a2"]),
                    "pool_team_bio": len(ar_house["pool_team"]),
                    "pool_case_study": len(ar_house["pool_case"]),
                    "pool_service_divider": len(ar_house["pool_service_divider"]),
                    "example_content": len(ar_house["example"]),
                    "unclassified": len(ar_house["unclassified"]),
                },
            },
        }
        _write_json(audit_dir / "audit_summary.json", summary)

        print("\n" + "=" * 60)
        print("PHASE 0 — Template Intelligence Audit COMPLETE")
        print("=" * 60)
        print(f"Artifacts: {len(artifacts)} files")
        print(
            f"EN: {len(en_slides)} slides, {len(en_layouts)} layouts"
        )
        print(
            f"AR: {len(ar_slides)} slides, {len(ar_layouts)} layouts"
        )
        stat = "PASS" if parity["overall_pass"] else "FAIL"
        print(f"Parity: {stat}")
        print(f"  Critical: {parity['critical_failures']}")
        print(f"  Warnings: {parity['warnings']}")
        print(f"Output: {audit_dir}")
        return artifacts

    finally:
        shutil.rmtree(patch_dir, ignore_errors=True)


# ── Entry point ────────────────────────────────────────────────────────────

if __name__ == "__main__":
    import sys

    base = Path(sys.argv[1]) if len(sys.argv) > 1 else Path.cwd()
    run_full_audit(output_base=base)
