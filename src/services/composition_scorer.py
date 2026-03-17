"""Composition Scorer — Profile-Driven Visual QA.

Post-render quality-assurance layer that extracts shape geometry from
a rendered PPTX and evaluates composition rules.  All scoring behavior
is driven by a ``ProfileConfig`` obtained from ``scorer_profiles.py``.

Dispatch flow:
  1. ``extract_shapes(pptx_path)`` → ``list[ShapeInfo]``  (Layer 1)
  2. ``score_composition(shapes, slides, *, profile=...)`` → ``CompositionResult``  (Layer 2)

The active ``ScorerProfile`` determines brand fonts, font-size thresholds,
margin expectations, and template-fidelity toggles.  Legacy profile is the
default — identical to pre-v2 behavior.  Template-v2 profile is isolated.
"""

from __future__ import annotations

import logging
from dataclasses import dataclass
from enum import StrEnum
from pathlib import Path
from typing import Any

from src.services.scorer_profiles import (
    ProfileConfig,
    ScorerProfile,
    get_profile,
)

logger = logging.getLogger(__name__)


# ── Severity ───────────────────────────────────────────────────────────


class ViolationSeverity(StrEnum):
    """How serious a composition violation is."""

    BLOCKER = "blocker"
    WARNING = "warning"
    INFO = "info"


# ── Data classes ───────────────────────────────────────────────────────


@dataclass(frozen=True)
class ShapeInfo:
    """Extracted geometric data for one shape on one slide."""

    slide_index: int
    slide_id: str
    shape_name: str
    shape_type: str           # "placeholder", "textbox", "table", "picture", etc.
    left_in: float            # inches from left edge
    top_in: float             # inches from top edge
    width_in: float
    height_in: float
    text: str = ""
    font_name: str = ""
    font_size_pt: float = 0.0
    is_placeholder: bool = False
    placeholder_idx: int = -1


@dataclass(frozen=True)
class Violation:
    """One composition rule violation."""

    rule: str                 # e.g. "overlap_severe", "font_min_soft"
    message: str
    severity: ViolationSeverity
    slide_id: str = ""
    shape_name: str = ""


@dataclass(frozen=True)
class SlideScore:
    """Aggregated score for one slide."""

    slide_id: str
    violations: tuple[Violation, ...] = ()

    @property
    def blocker_count(self) -> int:
        return sum(1 for v in self.violations if v.severity == ViolationSeverity.BLOCKER)


@dataclass(frozen=True)
class CompositionResult:
    """Aggregated result of scoring an entire deck."""

    slide_scores: tuple[SlideScore, ...] = ()
    profile_used: ScorerProfile = ScorerProfile.LEGACY

    @property
    def blocker_count(self) -> int:
        return sum(s.blocker_count for s in self.slide_scores)

    @property
    def total_violations(self) -> int:
        return sum(len(s.violations) for s in self.slide_scores)


# ── Layer 1: Shape extraction ─────────────────────────────────────────


def extract_shapes(pptx_path: str | Path) -> list[ShapeInfo]:
    """Extract shape geometry from a rendered PPTX file.

    Uses python-pptx to read every shape on every slide and produce
    a flat list of ``ShapeInfo`` records.
    """
    from pptx import Presentation

    prs = Presentation(str(pptx_path))
    shapes_out: list[ShapeInfo] = []

    for slide_idx, slide in enumerate(prs.slides):
        slide_id = f"slide_{slide_idx}"

        for shape in slide.shapes:
            left_in = shape.left / 914400 if shape.left else 0.0
            top_in = shape.top / 914400 if shape.top else 0.0
            width_in = shape.width / 914400 if shape.width else 0.0
            height_in = shape.height / 914400 if shape.height else 0.0

            text = ""
            font_name = ""
            font_size_pt = 0.0

            if shape.has_text_frame:
                text = shape.text_frame.text
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.name:
                            font_name = run.font.name
                        if run.font.size:
                            font_size_pt = run.font.size / 12700  # EMU to pt

            is_ph = shape.is_placeholder if hasattr(shape, "is_placeholder") else False
            ph_idx = -1
            if is_ph and hasattr(shape, "placeholder_format"):
                ph_idx = shape.placeholder_format.idx

            shape_type = "shape"
            if is_ph:
                shape_type = "placeholder"
            elif shape.has_table:
                shape_type = "table"
            elif hasattr(shape, "image"):
                shape_type = "picture"
            elif shape.has_text_frame:
                shape_type = "textbox"

            shapes_out.append(ShapeInfo(
                slide_index=slide_idx,
                slide_id=slide_id,
                shape_name=shape.name or f"Shape_{slide_idx}_{len(shapes_out)}",
                shape_type=shape_type,
                left_in=left_in,
                top_in=top_in,
                width_in=width_in,
                height_in=height_in,
                text=text,
                font_name=font_name,
                font_size_pt=font_size_pt,
                is_placeholder=is_ph,
                placeholder_idx=ph_idx,
            ))

    return shapes_out


# ── Layer 2: Rule engine (profile-driven) ─────────────────────────────


def score_composition(
    shapes: list[ShapeInfo],
    slides: list[Any],
    *,
    profile: ScorerProfile = ScorerProfile.LEGACY,
) -> CompositionResult:
    """Score extracted shapes against composition rules.

    Parameters
    ----------
    shapes : list[ShapeInfo]
        Flat list of shape data from ``extract_shapes()``.
    slides : list
        Original slide objects (SlideObject or similar) for context.
    profile : ScorerProfile
        Which scorer profile to use.  Default is LEGACY for backward
        compatibility.

    Returns
    -------
    CompositionResult
        Per-slide scores and blocker count.
    """
    config = get_profile(profile)

    # Group shapes by slide
    slide_shapes: dict[str, list[ShapeInfo]] = {}
    for s in shapes:
        slide_shapes.setdefault(s.slide_id, []).append(s)

    slide_scores: list[SlideScore] = []
    for slide_id, shape_list in slide_shapes.items():
        violations = _evaluate_slide(shape_list, slide_id, config)
        slide_scores.append(SlideScore(
            slide_id=slide_id,
            violations=tuple(violations),
        ))

    return CompositionResult(
        slide_scores=tuple(slide_scores),
        profile_used=profile,
    )


def _evaluate_slide(
    shapes: list[ShapeInfo],
    slide_id: str,
    config: ProfileConfig,
) -> list[Violation]:
    """Evaluate all composition rules for one slide."""
    violations: list[Violation] = []

    violations.extend(_check_font_compliance(shapes, slide_id, config))
    violations.extend(_check_overlap(shapes, slide_id, config))
    violations.extend(_check_bounds(shapes, slide_id, config))

    if config.enforce_template_fidelity:
        violations.extend(_check_template_fidelity(shapes, slide_id, config))

    return violations


# ── Rule: Font compliance ──────────────────────────────────────────────


def _check_font_compliance(
    shapes: list[ShapeInfo],
    slide_id: str,
    config: ProfileConfig,
) -> list[Violation]:
    """Check font name and size against profile thresholds."""
    violations: list[Violation] = []

    for s in shapes:
        if not s.font_name and not s.text:
            continue

        # Font family check
        if s.font_name and s.font_name not in config.brand_fonts and s.font_name not in config.heading_fonts:
            # Skip decorative elements in v2 profile
            if config.classify_template_native_decorative and _is_template_decorative(s):
                continue

            violations.append(Violation(
                rule="font_brand",
                message=(
                    f"Shape '{s.shape_name}' uses font '{s.font_name}', "
                    f"expected one of {config.brand_fonts + config.heading_fonts}"
                ),
                severity=ViolationSeverity.WARNING,
                slide_id=slide_id,
                shape_name=s.shape_name,
            ))

        # Font size check (only for shapes with text)
        if s.font_size_pt > 0 and s.text:
            min_pt = config.hard_floor_pt

            if s.font_size_pt < min_pt:
                violations.append(Violation(
                    rule="font_min_soft",
                    message=(
                        f"Shape '{s.shape_name}' font {s.font_size_pt}pt "
                        f"below minimum {min_pt}pt"
                    ),
                    severity=ViolationSeverity.BLOCKER,
                    slide_id=slide_id,
                    shape_name=s.shape_name,
                ))

    return violations


# ── Rule: Overlap detection ────────────────────────────────────────────


def _check_overlap(
    shapes: list[ShapeInfo],
    slide_id: str,
    config: ProfileConfig,
) -> list[Violation]:
    """Detect shapes overlapping beyond the profile threshold."""
    violations: list[Violation] = []
    threshold = config.overlap_severe_threshold_in

    text_shapes = [s for s in shapes if s.text and s.shape_type in ("placeholder", "textbox")]

    for i, a in enumerate(text_shapes):
        for b in text_shapes[i + 1:]:
            overlap = _compute_overlap(a, b)
            if overlap > threshold:
                violations.append(Violation(
                    rule="overlap_severe",
                    message=(
                        f"Shapes '{a.shape_name}' and '{b.shape_name}' "
                        f"overlap by {overlap:.2f}in (threshold {threshold}in)"
                    ),
                    severity=ViolationSeverity.BLOCKER,
                    slide_id=slide_id,
                    shape_name=a.shape_name,
                ))

    return violations


def _compute_overlap(a: ShapeInfo, b: ShapeInfo) -> float:
    """Compute vertical overlap between two shapes in inches."""
    a_top, a_bottom = a.top_in, a.top_in + a.height_in
    b_top, b_bottom = b.top_in, b.top_in + b.height_in

    overlap_top = max(a_top, b_top)
    overlap_bottom = min(a_bottom, b_bottom)
    vertical_overlap = max(0.0, overlap_bottom - overlap_top)

    # Also check horizontal overlap
    a_left, a_right = a.left_in, a.left_in + a.width_in
    b_left, b_right = b.left_in, b.left_in + b.width_in

    h_overlap_left = max(a_left, b_left)
    h_overlap_right = min(a_right, b_right)
    horizontal_overlap = max(0.0, h_overlap_right - h_overlap_left)

    if horizontal_overlap > 0 and vertical_overlap > 0:
        return vertical_overlap  # return vertical overlap amount

    return 0.0


# ── Rule: Bounds / margin check ────────────────────────────────────────


def _check_bounds(
    shapes: list[ShapeInfo],
    slide_id: str,
    config: ProfileConfig,
) -> list[Violation]:
    """Check shapes against margin boundaries."""
    violations: list[Violation] = []
    min_left = config.bounds_margin_left_min_in

    for s in shapes:
        if not s.text:
            continue
        if s.left_in < min_left and s.left_in > 0:
            # Skip very small offsets (template chrome)
            if s.left_in < min_left * 0.5:
                continue
            violations.append(Violation(
                rule="bounds_left_margin",
                message=(
                    f"Shape '{s.shape_name}' left edge at {s.left_in:.2f}in, "
                    f"minimum is {min_left}in"
                ),
                severity=ViolationSeverity.WARNING,
                slide_id=slide_id,
                shape_name=s.shape_name,
            ))

    return violations


# ── Rule: Template fidelity (v2 only) ─────────────────────────────────


def _check_template_fidelity(
    shapes: list[ShapeInfo],
    slide_id: str,
    config: ProfileConfig,
) -> list[Violation]:
    """V2-only: verify shapes come from template layouts, not programmatic creation.

    Heuristic: non-placeholder, non-table shapes with text that are not
    classified as template-native decorative elements are suspicious.
    """
    violations: list[Violation] = []

    for s in shapes:
        if s.shape_type == "textbox" and s.text:
            if not _is_template_decorative(s):
                violations.append(Violation(
                    rule="template_fidelity",
                    message=(
                        f"Non-template textbox '{s.shape_name}' detected "
                        f"with text: '{s.text[:50]}'"
                    ),
                    severity=ViolationSeverity.WARNING,
                    slide_id=slide_id,
                    shape_name=s.shape_name,
                ))

    return violations


def _is_template_decorative(shape: ShapeInfo) -> bool:
    """Classify a shape as template-native decorative content.

    Template-native elements include: headers, footers, slide numbers,
    decorative bars, copyright notices, etc.  These are preserved from
    the template and should not trigger fidelity violations.
    """
    # Shapes at very top or very bottom are likely template chrome
    if shape.top_in < 0.3 or shape.top_in > 6.5:
        return True
    # Very narrow shapes are likely decorative bars/lines
    if shape.height_in < 0.15 or shape.width_in < 0.3:
        return True
    # Shapes with no meaningful text
    if len(shape.text.strip()) < 3:
        return True
    return False
