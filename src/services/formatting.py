"""Formatting utilities for PPTX and DOCX rendering.

Provides:
  - Professional PPTX table formatting (navy headers, alternating rows, status colors)
  - PPTX text formatting (bold key phrases, bullet levels, font sizing)
  - Pipe-table detection and parsing for slide tables
  - Layout-specific formatters (STAT_CALLOUT, AGENDA, TEAM, COMPLIANCE_MATRIX, TIMELINE)
  - Markdown-to-DOCX conversion via markdown-it-py
"""

from __future__ import annotations

import re
from typing import Any

from docx.shared import Inches, Pt
from markdown_it import MarkdownIt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches as PptxInches
from pptx.util import Pt as PptxPt

# ── Color Constants ───────────────────────────────────────────

# Primary palette from template theme
NAVY = RGBColor(0x0E, 0x28, 0x41)         # dk2 — primary dark
TEAL = RGBColor(0x15, 0x60, 0x82)         # accent1 — primary accent
ORANGE = RGBColor(0xE9, 0x71, 0x32)       # accent2 — warm accent
GREEN = RGBColor(0x19, 0x6B, 0x24)        # accent3
BLUE = RGBColor(0x0F, 0x9E, 0xD5)         # accent4 — bright blue
DARK_TEAL = RGBColor(0x46, 0x78, 0x86)    # hlink color

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)   # Alternating row fill
MID_GRAY = RGBColor(0xD9, 0xD9, 0xD9)     # Borders
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)     # Body text color

# Status colors
STATUS_GREEN = RGBColor(0x19, 0x6B, 0x24)
STATUS_AMBER = RGBColor(0xE9, 0x71, 0x32)
STATUS_RED = RGBColor(0xC0, 0x39, 0x2B)

# Font Constants
FONT_BODY = "Aptos"
FONT_HEADING = "Aptos Display"

# Legacy tuple exports for backward compatibility
ACCENT1_RGB = (0x15, 0x60, 0x82)  # Deep Blue  #156082
ACCENT2_RGB = (0xE9, 0x71, 0x32)  # Orange     #E97132

# Match "key — detail" where — is U+2014 (em-dash)
KEY_DETAIL_RE = re.compile(r"^(.+?)\s*\u2014\s*(.+)$")

# Match "• text" (bullet prefix)
BULLET_PREFIX_RE = re.compile(r"^[\u2022\u2023\u25CF\u25CB]\s*(.+)$")


# ==============================================================
# PIPE-TABLE UTILITIES
# ==============================================================


def is_pipe_table(elements: list[str]) -> bool:
    """Return True if 2+ non-empty elements contain pipe separators."""
    pipe_count = sum(1 for e in elements if e.strip() and " | " in e)
    return pipe_count >= 2


def parse_pipe_table(
    elements: list[str],
) -> tuple[list[str], list[list[str]]]:
    """Parse pipe-separated text elements into (headers, rows).

    The first pipe-containing element becomes the header row.
    Separator-only rows (all dashes) are filtered out.
    Non-pipe elements are silently skipped.
    """
    headers: list[str] = []
    rows: list[list[str]] = []

    for elem in elements:
        if not elem.strip() or " | " not in elem:
            continue
        cells = [c.strip() for c in elem.split(" | ")]
        # Skip separator rows like "---|---"
        if all(set(c) <= {"-", ":"} for c in cells):
            continue
        if not headers:
            headers = cells
        else:
            rows.append(cells)

    return headers, rows


# ==============================================================
# TABLE CELL STYLING HELPERS
# ==============================================================


def _set_cell_fill(cell: Any, color: RGBColor) -> None:
    """Set solid fill on a table cell via XML."""
    tcPr = cell._tc.get_or_add_tcPr()
    # Remove existing shading
    for existing in tcPr.findall(qn("a:solidFill")):
        tcPr.remove(existing)
    solidFill = tcPr.makeelement(qn("a:solidFill"), {})
    srgb = solidFill.makeelement(
        qn("a:srgbClr"),
        {"val": f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"},
    )
    solidFill.append(srgb)
    tcPr.append(solidFill)


def _set_cell_margins(cell: Any, top: int = 36000, bottom: int = 36000,
                      left: int = 72000, right: int = 72000) -> None:
    """Set cell margins (in EMU). Default = ~0.05in top/bottom, ~0.08in left/right."""
    tcPr = cell._tc.get_or_add_tcPr()
    mar = tcPr.makeelement(qn("a:tcMar"), {})
    for side, val in [("top", top), ("bottom", bottom), ("left", left), ("right", right)]:
        elem = mar.makeelement(qn(f"a:{side}"), {"w": str(val), "type": "emu"})
        mar.append(elem)
    tcPr.append(mar)


def _remove_table_banding(table: Any) -> None:
    """Remove PowerPoint's default banding and first-row styling."""
    tbl = table._tbl
    tblPr = tbl.tblPr
    if tblPr is not None:
        tblPr.set("bandRow", "0")
        tblPr.set("bandCol", "0")
        tblPr.set("firstRow", "0")
        tblPr.set("lastRow", "0")
        tblPr.set("firstCol", "0")
        tblPr.set("lastCol", "0")


def _set_cell_border(cell: Any, sides: list[str] | None = None,
                     color: RGBColor | None = None, width: int = 6350) -> None:
    """Set borders on specific sides of a cell.

    Args:
        cell: python-pptx table cell
        sides: List of sides to set borders for ("top", "bottom", "left", "right")
                Defaults to all four sides.
        color: Border color (defaults to MID_GRAY)
        width: Border width in EMU (default 6350 = ~0.5pt)
    """
    if sides is None:
        sides = ["top", "bottom", "left", "right"]
    if color is None:
        color = MID_GRAY

    border_map = {
        "top": "a:top",
        "bottom": "a:bottom",
        "left": "a:left",
        "right": "a:right",
    }

    tcPr = cell._tc.get_or_add_tcPr()

    # Remove existing tcBorders if any
    for existing in tcPr.findall(qn("a:tcBorders")):
        tcPr.remove(existing)

    borders_elem = tcPr.makeelement(qn("a:tcBorders"), {})
    color_str = f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"

    for side in sides:
        ns_tag = border_map.get(side)
        if ns_tag:
            border = borders_elem.makeelement(
                qn(ns_tag), {"w": str(width), "cap": "flat", "cmpd": "sng", "algn": "ctr"}
            )
            solid_fill = border.makeelement(qn("a:solidFill"), {})
            srgb = solid_fill.makeelement(qn("a:srgbClr"), {"val": color_str})
            solid_fill.append(srgb)
            border.append(solid_fill)
            prstDash = border.makeelement(qn("a:prstDash"), {"val": "solid"})
            border.append(prstDash)
            borders_elem.append(border)

    tcPr.append(borders_elem)


def _format_cell_text(cell: Any, text: str, font_name: str = FONT_BODY,
                      font_size_pt: int = 11, bold: bool = False,
                      color: RGBColor | None = None,
                      alignment: Any = PP_ALIGN.LEFT) -> None:
    """Format a table cell with styled text."""
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text.strip()
    run.font.name = font_name
    run.font.size = PptxPt(font_size_pt)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def _detect_status(text: str) -> str | None:
    """Detect status indicators in text for color coding."""
    lower = text.lower().strip()
    # Green indicators
    if any(kw in lower for kw in ["compliant", "yes", "ready", "pass", "met", "available"]):
        return "green"
    if lower.startswith(("\u2713",)):  # checkmark
        return "green"
    # Amber indicators
    if any(kw in lower for kw in ["pending", "partial", "in progress", "action required", "\u26a0"]):
        return "amber"
    # Red indicators
    if any(kw in lower for kw in ["no", "fail", "missing", "not met", "non-compliant", "\u2717"]):
        return "red"
    return None


# ==============================================================
# PPTX TEXT FORMATTING
# ==============================================================


def format_text_frame(
    tf: Any,
    elements: list[str],
    font_size_pt: int = 14,
) -> None:
    """Fill a PPTX text frame with formatted bullet-point elements.

    Formatting rules:
      - Lines starting with bullet char -> paragraph.level=1, smaller font
      - Lines matching "key -- detail" (em-dash) -> bold run for key part
      - Font: Aptos, size per level, dark text color
    """
    tf.clear()
    first_para = True

    for text in elements:
        if not text.strip():
            continue

        if first_para:
            p = tf.paragraphs[0]
            first_para = False
        else:
            p = tf.add_paragraph()

        # Detect bullet prefix -> sub-bullet
        level = 0
        m_bullet = BULLET_PREFIX_RE.match(text)
        if m_bullet:
            text = m_bullet.group(1)
            level = 1

        p.level = level
        size = PptxPt(font_size_pt if level == 0 else font_size_pt - 2)
        p.space_after = PptxPt(4)

        # Detect key--detail pattern -> bold key
        m_key = KEY_DETAIL_RE.match(text)
        if m_key:
            key_run = p.add_run()
            key_run.text = m_key.group(1)
            key_run.font.bold = True
            key_run.font.size = size
            key_run.font.name = FONT_BODY
            key_run.font.color.rgb = DARK_TEXT

            detail_run = p.add_run()
            detail_run.text = " \u2014 " + m_key.group(2)
            detail_run.font.bold = False
            detail_run.font.size = size
            detail_run.font.name = FONT_BODY
            detail_run.font.color.rgb = DARK_TEXT
        else:
            run = p.add_run()
            run.text = text
            run.font.size = size
            run.font.name = FONT_BODY
            run.font.color.rgb = DARK_TEXT


def format_stat_callout(tf: Any, elements: list[str]) -> None:
    """Format STAT_CALLOUT layout: large bold stat + supporting text.

    First element: 44pt bold teal (the big number/stat)
    Remaining: 14pt dark text below
    """
    tf.clear()
    first_para = True

    for i, text in enumerate(elements):
        if not text.strip():
            continue

        if first_para:
            p = tf.paragraphs[0]
            first_para = False
        else:
            p = tf.add_paragraph()

        p.alignment = PP_ALIGN.CENTER

        if i == 0:
            # Big stat number
            p.space_after = PptxPt(12)
            run = p.add_run()
            run.text = text
            run.font.name = FONT_HEADING
            run.font.size = PptxPt(44)
            run.font.bold = True
            run.font.color.rgb = TEAL
        else:
            # Supporting text
            p.space_after = PptxPt(4)
            # Support key--detail pattern
            m_key = KEY_DETAIL_RE.match(text)
            if m_key:
                key_run = p.add_run()
                key_run.text = m_key.group(1)
                key_run.font.bold = True
                key_run.font.size = PptxPt(14)
                key_run.font.name = FONT_BODY
                key_run.font.color.rgb = DARK_TEXT

                detail_run = p.add_run()
                detail_run.text = " \u2014 " + m_key.group(2)
                detail_run.font.bold = False
                detail_run.font.size = PptxPt(14)
                detail_run.font.name = FONT_BODY
                detail_run.font.color.rgb = DARK_TEXT
            else:
                run = p.add_run()
                run.text = text
                run.font.size = PptxPt(14)
                run.font.name = FONT_BODY
                run.font.color.rgb = DARK_TEXT


def format_agenda(tf: Any, elements: list[str]) -> None:
    """Format AGENDA layout: auto-numbered list with spacing and teal numbers."""
    tf.clear()
    num = 0
    first_para = True

    for text in elements:
        if not text.strip():
            continue

        num += 1

        if first_para:
            p = tf.paragraphs[0]
            first_para = False
        else:
            p = tf.add_paragraph()

        p.space_after = PptxPt(8)

        # Add teal number
        num_run = p.add_run()
        # Check if already has number prefix
        has_num = re.match(r"^\d+[.)]\s*", text)
        if has_num:
            num_run.text = has_num.group(0)
            text = text[has_num.end():]
        else:
            num_run.text = f"{num}.  "
        num_run.font.size = PptxPt(16)
        num_run.font.name = FONT_HEADING
        num_run.font.bold = True
        num_run.font.color.rgb = TEAL

        # Add text
        text_run = p.add_run()
        text_run.text = text
        text_run.font.size = PptxPt(14)
        text_run.font.name = FONT_BODY
        text_run.font.color.rgb = DARK_TEXT


# ==============================================================
# PPTX TABLE RENDERING — PROFESSIONAL STYLED TABLES
# ==============================================================


def add_pptx_table(
    slide: Any,
    elements: list[str],
    left_inches: float = 0.92,
    top_inches: float = 2.0,
    width_inches: float = 11.5,
) -> Any | None:
    """Render pipe-separated text elements as a professionally styled PowerPoint table.

    Features:
      - Navy header row (#0E2841) with white bold text
      - Alternating light gray (#F2F2F2) / white body rows
      - Status column color coding (green/amber/red indicators)
      - Proper cell margins and borders
      - Aptos 11pt body, 12pt header
    """
    headers, rows = parse_pipe_table(elements)
    if not headers:
        return None

    num_rows = len(rows) + 1  # +1 for header
    num_cols = len(headers)

    left = PptxInches(left_inches)
    top = PptxInches(top_inches)
    width = PptxInches(width_inches)
    row_height = PptxInches(0.38)
    height = row_height * num_rows

    table_shape = slide.shapes.add_table(
        num_rows, num_cols, left, top, width, height,
    )
    table = table_shape.table

    # Remove default banding
    _remove_table_banding(table)

    # ── Header row ──
    for c, text in enumerate(headers):
        cell = table.cell(0, c)
        _set_cell_fill(cell, NAVY)
        _format_cell_text(
            cell, text,
            font_size_pt=12, bold=True, color=WHITE,
            alignment=PP_ALIGN.LEFT,
        )
        _set_cell_margins(cell)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # ── Data rows with alternating fills ──
    for r, row_data in enumerate(rows, start=1):
        fill_color = LIGHT_GRAY if r % 2 == 0 else WHITE
        for c in range(min(len(row_data), num_cols)):
            cell = table.cell(r, c)
            _set_cell_fill(cell, fill_color)
            _set_cell_margins(cell)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            cell_text = row_data[c].strip()

            # Check for status color coding
            status = _detect_status(cell_text)
            if status:
                text_color = {
                    "green": STATUS_GREEN,
                    "amber": STATUS_AMBER,
                    "red": STATUS_RED,
                }.get(status, DARK_TEXT)
                _format_cell_text(
                    cell, cell_text,
                    font_size_pt=11, bold=True, color=text_color,
                )
            else:
                _format_cell_text(
                    cell, cell_text,
                    font_size_pt=11, color=DARK_TEXT,
                )

    # ── Set column widths proportionally ──
    total_width = width
    col_widths = _calculate_column_widths(headers, rows, num_cols, total_width)
    for c in range(num_cols):
        table.columns[c].width = col_widths[c]

    return table_shape


def _calculate_column_widths(
    headers: list[str],
    rows: list[list[str]],
    num_cols: int,
    total_width: int,
) -> list[int]:
    """Calculate proportional column widths based on content length."""
    # Estimate character counts per column
    char_counts = []
    for c in range(num_cols):
        header_len = len(headers[c]) if c < len(headers) else 5
        data_lens = [len(row[c]) if c < len(row) else 0 for row in rows]
        avg_data_len = sum(data_lens) / max(len(data_lens), 1) if data_lens else 5
        char_counts.append(max(header_len, avg_data_len, 3))

    total_chars = sum(char_counts)
    widths = []
    for count in char_counts:
        proportion = count / total_chars
        # Minimum 8% width, maximum 50% width per column
        proportion = max(0.08, min(0.50, proportion))
        widths.append(int(total_width * proportion))

    # Adjust to match total
    diff = total_width - sum(widths)
    if widths:
        widths[-1] += diff

    return widths


def add_styled_table_from_elements(
    slide: Any,
    elements: list[str],
    has_header: bool = True,
    column_headers: list[str] | None = None,
    left_inches: float = 0.92,
    top_inches: float = 2.0,
    width_inches: float = 11.5,
) -> Any | None:
    """Render non-pipe elements as a styled table.

    Treats each element as a row with em-dash separated columns.
    Used for TEAM, TIMELINE layouts where content is in "key -- detail" format.

    Args:
        column_headers: If provided, prepend a navy header row with these
            labels before the data rows. Overrides has_header for the
            header styling.
    """
    # Parse elements into rows
    parsed_rows: list[list[str]] = []
    for elem in elements:
        if not elem.strip():
            continue
        # Try splitting by em-dash first
        m = KEY_DETAIL_RE.match(elem)
        if m:
            parsed_rows.append([m.group(1).strip(), m.group(2).strip()])
        else:
            # Try splitting by pipe
            if " | " in elem:
                parsed_rows.append([c.strip() for c in elem.split(" | ")])
            else:
                parsed_rows.append([elem.strip()])

    if not parsed_rows:
        return None

    num_cols = max(len(r) for r in parsed_rows)
    # Ensure num_cols matches column_headers if provided
    if column_headers:
        num_cols = max(num_cols, len(column_headers))
    num_rows = len(parsed_rows)
    # Extra row for explicit column headers
    if column_headers:
        num_rows += 1

    left = PptxInches(left_inches)
    top = PptxInches(top_inches)
    width = PptxInches(width_inches)
    row_height = PptxInches(0.40)
    height = row_height * num_rows

    table_shape = slide.shapes.add_table(
        num_rows, num_cols, left, top, width, height,
    )
    table = table_shape.table
    _remove_table_banding(table)

    # Render explicit column headers as navy row 0
    data_row_offset = 0
    if column_headers:
        data_row_offset = 1
        for c in range(num_cols):
            cell = table.cell(0, c)
            _set_cell_fill(cell, NAVY)
            _set_cell_margins(cell)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            header_text = column_headers[c] if c < len(column_headers) else ""
            _format_cell_text(
                cell, header_text,
                font_size_pt=12, bold=True, color=WHITE,
            )

    for r, row_data in enumerate(parsed_rows):
        table_row = r + data_row_offset
        is_header = (not column_headers) and has_header and r == 0
        fill_color = NAVY if is_header else (LIGHT_GRAY if r % 2 == 0 else WHITE)
        text_color = WHITE if is_header else DARK_TEXT

        for c in range(num_cols):
            cell = table.cell(table_row, c)
            _set_cell_fill(cell, fill_color)
            _set_cell_margins(cell)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            cell_text = row_data[c].strip() if c < len(row_data) else ""
            _format_cell_text(
                cell, cell_text,
                font_size_pt=12 if is_header else 11,
                bold=is_header or c == 0,
                color=text_color,
            )

    # Adjust first column width (label column = 30%)
    if num_cols >= 2:
        first_col_width = int(width * 0.30)
        remaining_width = int((width - first_col_width) / max(num_cols - 1, 1))
        table.columns[0].width = first_col_width
        for c in range(1, num_cols):
            table.columns[c].width = remaining_width

    return table_shape


# ==============================================================
# COMPARISON LAYOUT FORMATTING
# ==============================================================


def format_comparison_body(
    ph: Any,
    elements: list[str],
    font_size_pt: int = 13,
) -> None:
    """Format a body placeholder in the Comparison layout as a sub-header.

    Used for idx=1 and idx=3 (the column headers in Comparison layout).
    """
    ph.text = ""
    tf = ph.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT

    # Use first element as the sub-header text
    text = elements[0] if elements else ""
    run = p.add_run()
    run.text = text
    run.font.name = FONT_HEADING
    run.font.size = PptxPt(font_size_pt)
    run.font.bold = True
    run.font.color.rgb = TEAL


# ==============================================================
# PROPOSAL COVER (LAYOUT 11) FORMATTING
# ==============================================================


def format_cover_subtitle(ph: Any, text: str) -> None:
    """Format the project name on the Proposal Cover layout (idx=1, SUBTITLE)."""
    ph.text = ""
    tf = ph.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = FONT_HEADING
    run.font.size = PptxPt(28)
    run.font.bold = True
    run.font.color.rgb = WHITE


def format_cover_body(ph: Any, text: str, font_size: int = 18) -> None:
    """Format a body field on the Proposal Cover layout (idx=10 or idx=11)."""
    ph.text = ""
    tf = ph.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = FONT_BODY
    run.font.size = PptxPt(font_size)
    run.font.color.rgb = WHITE


# ==============================================================
# TITLE SLIDE (LAYOUT 0) CLOSING FORMAT
# ==============================================================


def format_closing_slide(slide: Any, title: str, elements: list[str]) -> None:
    """Format a closing slide (Layout 0 - Title Slide).

    CENTER_TITLE(0): Big thank you / closing headline
    SUBTITLE(1): Key closing points
    """
    # Title
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = ""
            tf = ph.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = title
            run.font.name = FONT_HEADING
            run.font.size = PptxPt(36)
            run.font.bold = True
            break

    # Subtitle with closing points
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            ph.text = ""
            tf = ph.text_frame
            tf.clear()
            first = True
            for text in elements:
                if not text.strip():
                    continue
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                p.alignment = PP_ALIGN.CENTER
                p.space_after = PptxPt(6)

                m_key = KEY_DETAIL_RE.match(text)
                if m_key:
                    key_run = p.add_run()
                    key_run.text = m_key.group(1)
                    key_run.font.bold = True
                    key_run.font.size = PptxPt(16)
                    key_run.font.name = FONT_BODY

                    detail_run = p.add_run()
                    detail_run.text = " \u2014 " + m_key.group(2)
                    detail_run.font.size = PptxPt(16)
                    detail_run.font.name = FONT_BODY
                else:
                    run = p.add_run()
                    run.text = text
                    run.font.size = PptxPt(16)
                    run.font.name = FONT_BODY
            break


# ==============================================================
# TEAL ACCENT BAR (decorative element for Title Only layouts)
# ==============================================================


def add_accent_bar(slide: Any, left_inches: float = 0.92,
                   top_inches: float = 1.85, width_inches: float = 1.0,
                   height_inches: float = 0.04) -> None:
    """Add a small teal accent bar below the title area."""
    from pptx.util import Inches as I
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        I(left_inches), I(top_inches), I(width_inches), I(height_inches),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = TEAL
    shape.line.fill.background()  # No border


# ==============================================================
# DOCX MARKDOWN RENDERING
# ==============================================================


def _add_inline_to_paragraph(para: Any, inline_token: Any) -> None:
    """Walk inline token children and add formatted runs to a paragraph."""
    if not inline_token.children:
        if inline_token.content:
            para.add_run(inline_token.content)
        return

    bold = False
    italic = False

    for child in inline_token.children:
        if child.type == "strong_open":
            bold = True
        elif child.type == "strong_close":
            bold = False
        elif child.type == "em_open":
            italic = True
        elif child.type == "em_close":
            italic = False
        elif child.type == "text":
            run = para.add_run(child.content)
            if bold:
                run.bold = True
            if italic:
                run.italic = True
        elif child.type == "code_inline":
            run = para.add_run(child.content)
            run.font.name = "Consolas"
            run.font.size = Pt(9)
        elif child.type == "softbreak":
            para.add_run("\n")
        elif child.type == "hardbreak":
            para.add_run("\n")


def render_markdown_to_docx(doc: Any, markdown_text: str) -> None:
    """Parse markdown and render into a python-docx Document.

    Handles: headings (h2, h3), bold, italic, bullet lists,
    numbered lists, blockquotes, tables.
    """
    md = MarkdownIt("commonmark").enable("table")
    tokens = md.parse(markdown_text)

    # State
    in_blockquote = False
    in_bullet_list = False
    in_ordered_list = False
    in_table = False
    is_thead = False
    table_rows: list[list[str]] = []
    table_header: list[str] = []
    current_row: list[str] = []
    skipped_first_heading = False

    i = 0
    while i < len(tokens):
        token = tokens[i]

        # -- Headings --
        if token.type == "heading_open":
            level = int(token.tag[1])  # h2 -> 2, h3 -> 3
            if i + 1 < len(tokens) and tokens[i + 1].type == "inline":
                heading_text = tokens[i + 1].content
                if not skipped_first_heading and level == 2:
                    skipped_first_heading = True
                    i += 3
                    continue
                doc.add_heading(heading_text, level=level)
                i += 3
                continue

        # -- Blockquotes --
        elif token.type == "blockquote_open":
            in_blockquote = True
        elif token.type == "blockquote_close":
            in_blockquote = False

        # -- Bullet lists --
        elif token.type == "bullet_list_open":
            in_bullet_list = True
        elif token.type == "bullet_list_close":
            in_bullet_list = False

        # -- Ordered lists --
        elif token.type == "ordered_list_open":
            in_ordered_list = True
        elif token.type == "ordered_list_close":
            in_ordered_list = False

        # -- Table --
        elif token.type == "table_open":
            in_table = True
            table_rows = []
            table_header = []
        elif token.type == "thead_open":
            is_thead = True
        elif token.type == "thead_close":
            is_thead = False
        elif token.type == "tr_open":
            current_row = []
        elif token.type in ("th_close", "td_close"):
            pass
        elif token.type == "tr_close":
            if is_thead:
                table_header = current_row
            else:
                table_rows.append(current_row)
            current_row = []
        elif token.type == "table_close":
            in_table = False
            if table_header:
                num_cols = len(table_header)
                num_data_rows = len(table_rows)
                table = doc.add_table(
                    rows=num_data_rows + 1, cols=num_cols,
                )
                table.style = "Table Grid"
                for c, text in enumerate(table_header):
                    cell = table.rows[0].cells[c]
                    p = cell.paragraphs[0]
                    run = p.add_run(text)
                    run.bold = True
                for r, row_data in enumerate(table_rows):
                    for c in range(min(len(row_data), num_cols)):
                        table.rows[r + 1].cells[c].text = row_data[c]

        # -- Inline content (inside th/td in tables) --
        elif token.type == "inline" and in_table:
            current_row.append(token.content)

        # -- Paragraphs with inline content --
        elif token.type == "paragraph_open":
            if i + 1 < len(tokens) and tokens[i + 1].type == "inline":
                inline_token = tokens[i + 1]

                if in_bullet_list:
                    p = doc.add_paragraph(style="List Bullet")
                elif in_ordered_list:
                    p = doc.add_paragraph(style="List Number")
                elif in_blockquote:
                    p = doc.add_paragraph()
                    p.paragraph_format.left_indent = Inches(0.5)
                    p.paragraph_format.right_indent = Inches(0.3)
                else:
                    p = doc.add_paragraph()

                _add_inline_to_paragraph(p, inline_token)
                i += 2
                i += 1
                continue

        i += 1
