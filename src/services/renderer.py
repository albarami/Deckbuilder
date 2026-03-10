"""Design Agent -- deterministic PPTX renderer and DOCX exporter.

No LLM involvement. Loads the Strategic Gears master template
(Presentation6.pptx) and populates slides from validated SlideObjects.

Template layouts (Presentation6.pptx):
  0: Title Slide       -- CENTER_TITLE(0) + SUBTITLE(1)
  1: Title and Content  -- TITLE(0) + OBJECT(1)
  2: Section Header     -- TITLE(0) + BODY(1)
  3: Two Content        -- TITLE(0) + OBJECT(1) + OBJECT(2)
  4: Comparison         -- TITLE(0) + BODY(1) + OBJECT(2) + BODY(3) + OBJECT(4)
  5: Title Only         -- TITLE(0)
  6: Blank              -- (footer only)
  7: Content with Caption -- TITLE(0) + OBJECT(1, right 6.75in) + BODY(2, left 4.30in)
  8: Picture with Caption -- TITLE(0) + PICTURE(1) + BODY(2)
  9: Title and Vertical Text
 10: Vertical Title and Text
 11: Proposal Cover     -- SUBTITLE(1) + BODY(10,11) + PICTURE(12)
 12: ToC / Agenda       -- TITLE(0) + TABLE(10)
"""

import datetime
import logging
from pathlib import Path
from typing import Any

from docx import Document
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches

from src.models.enums import Language, LayoutType, RenderStatus
from src.models.report import ResearchReport
from src.models.slides import SlideObject
from src.services.formatting import (
    add_accent_bar,
    add_pptx_table,
    add_styled_table_from_elements,
    format_agenda,
    format_closing_slide,
    format_comparison_body,
    format_cover_body,
    format_cover_subtitle,
    format_stat_callout,
    format_text_frame,
    is_pipe_table,
    render_markdown_to_docx,
)

logger = logging.getLogger(__name__)

# ----------------------------------------------------------------
# Layout mapping -- DeckForge LayoutType -> template layout index
# ----------------------------------------------------------------

LAYOUT_MAP: dict[LayoutType, int] = {
    LayoutType.TITLE: 11,              # Proposal Cover (with photo BG)
    LayoutType.AGENDA: 12,             # ToC / Agenda (with table placeholder)
    LayoutType.SECTION: 2,             # Section Header
    LayoutType.CONTENT_1COL: 1,        # Title and Content
    LayoutType.CONTENT_2COL: 3,        # Two Content
    LayoutType.DATA_CHART: 1,          # Title and Content (chart in OBJECT)
    LayoutType.FRAMEWORK: 1,           # Title and Content (formatted body)
    LayoutType.COMPARISON: 4,          # Comparison (4 body areas)
    LayoutType.STAT_CALLOUT: 5,        # Title Only (custom textbox)
    LayoutType.TEAM: 5,               # Title Only (styled table)
    LayoutType.TIMELINE: 5,           # Title Only (styled table)
    LayoutType.COMPLIANCE_MATRIX: 5,  # Title Only (styled table)
    LayoutType.CLOSING: 0,             # Title Slide (bookend)
}

_DEFAULT_LAYOUT_IDX = 1  # "Title and Content" as fallback


# ----------------------------------------------------------------
# RenderResult
# ----------------------------------------------------------------


class RenderResult:
    """Result of PPTX rendering."""

    def __init__(
        self,
        pptx_path: str,
        slide_count: int,
        render_log: list[dict[str, Any]],
    ) -> None:
        self.pptx_path = pptx_path
        self.slide_count = slide_count
        self.render_log = render_log


# ----------------------------------------------------------------
# Layout-specific populate functions
# ----------------------------------------------------------------


def _get_layout_index(layout_type: LayoutType) -> int:
    """Get template layout index for a LayoutType."""
    return LAYOUT_MAP.get(layout_type, _DEFAULT_LAYOUT_IDX)


def _populate_title_slide(slide: Any, slide_obj: SlideObject) -> None:
    """Populate Layout 11 (Proposal Cover).

    idx=1 (SUBTITLE): Project name / proposal title
    idx=10 (BODY): Client name / issuing entity
    idx=11 (BODY): Date
    idx=12 (PICTURE): Client logo placeholder (skipped)
    """
    elements = slide_obj.body_content.text_elements if slide_obj.body_content else []

    # Parse title components from elements
    project_name = slide_obj.title
    client_name = ""
    date_str = datetime.date.today().strftime("%d %B %Y")
    extra_lines = []

    for elem in elements:
        lower = elem.lower()
        if "issuing entity" in lower or "client" in lower.split("\u2014")[0] if "\u2014" in lower else "":
            # Extract value after em-dash
            parts = elem.split("\u2014", 1)
            if len(parts) > 1:
                client_name = parts[1].strip()
            else:
                client_name = elem.strip()
        elif "date" in lower or "format" in lower:
            parts = elem.split("\u2014", 1)
            if len(parts) > 1:
                date_str = parts[1].strip()
        elif "rfp" in lower and "\u2014" in elem:
            parts = elem.split("\u2014", 1)
            if len(parts) > 1:
                project_name = parts[1].strip()
        else:
            extra_lines.append(elem)

    # If no client extracted, try to extract from elements
    if not client_name and elements:
        for elem in elements:
            if "\u2014" in elem:
                parts = elem.split("\u2014", 1)
                key = parts[0].strip().lower()
                if any(kw in key for kw in ["issuing", "entity", "client"]):
                    client_name = parts[1].strip()
                    break

    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 1:  # SUBTITLE - Project name
            format_cover_subtitle(ph, project_name)
        elif idx == 10:  # BODY - Client name
            format_cover_body(ph, client_name or "Strategic Gears Consulting", font_size=18)
        elif idx == 11:  # BODY - Date
            format_cover_body(ph, date_str, font_size=14)


def _populate_agenda_slide(slide: Any, slide_obj: SlideObject) -> None:
    """Populate Layout 12 (ToC / Agenda).

    idx=0 (TITLE): "Table of Contents" or custom title
    idx=10 (TABLE): Table placeholder - we add a textbox instead
                    since table placeholders are tricky.
    """
    # Set title
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = slide_obj.title
            break

    elements = slide_obj.body_content.text_elements if slide_obj.body_content else []
    if not elements:
        return

    # Add formatted agenda as textbox (Layout 12 has complex BG shape already)
    from pptx.util import Inches as I
    left = I(0.92)
    top = I(2.0)
    width = I(8.0)
    height = I(4.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    format_agenda(tf, elements)


def _populate_content_1col(slide: Any, slide_obj: SlideObject) -> None:
    """Populate Layout 1 (Title and Content) or Layout 7 (Content with Caption).

    idx=0 (TITLE): Slide title
    idx=1 (OBJECT): Body content with formatted bullets
    """
    _set_title(slide, slide_obj.title)

    elements = slide_obj.body_content.text_elements if slide_obj.body_content else []
    if not elements:
        return

    # Check if pipe-table data -> render as table
    if is_pipe_table(elements):
        add_pptx_table(slide, elements)
        return

    # Find body placeholder
    body_ph = _find_body_placeholder(slide)
    if body_ph is not None:
        format_text_frame(body_ph.text_frame, elements)
    else:
        _add_body_textbox(slide, elements)


def _populate_content_2col(slide: Any, slide_obj: SlideObject) -> None:
    """Populate Layout 3 (Two Content).

    idx=0 (TITLE): Slide title
    idx=1 (OBJECT): Left column
    idx=2 (OBJECT): Right column
    """
    _set_title(slide, slide_obj.title)

    elements = slide_obj.body_content.text_elements if slide_obj.body_content else []
    if not elements:
        return

    # Pipe tables override columns
    if is_pipe_table(elements):
        add_pptx_table(slide, elements)
        return

    mid = len(elements) // 2

    col_phs: dict[int, Any] = {}
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx in (1, 2) and ph.placeholder_format.type in (2, 7):
            col_phs[idx] = ph

    for col_idx, items in [(1, elements[:mid]), (2, elements[mid:])]:
        ph = col_phs.get(col_idx)
        if ph is None or not items:
            continue
        format_text_frame(ph.text_frame, items)


def _populate_comparison(slide: Any, slide_obj: SlideObject) -> None:
    """Populate Layout 4 (Comparison).

    idx=0 (TITLE): Slide title
    idx=1 (BODY): Left sub-header
    idx=2 (OBJECT): Left content
    idx=3 (BODY): Right sub-header
    idx=4 (OBJECT): Right content

    If content has pipe tables, renders as a table instead.
    """
    _set_title(slide, slide_obj.title)

    elements = slide_obj.body_content.text_elements if slide_obj.body_content else []
    if not elements:
        return

    # Pipe table -> render as full-width table
    if is_pipe_table(elements):
        add_pptx_table(slide, elements, top_inches=2.0)
        return

    # Split elements into two halves
    mid = len(elements) // 2
    left_elems = elements[:mid]
    right_elems = elements[mid:]

    placeholders: dict[int, Any] = {}
    for ph in slide.placeholders:
        placeholders[ph.placeholder_format.idx] = ph

    # Left sub-header (idx=1, BODY) - use first left element as header
    if 1 in placeholders and left_elems:
        format_comparison_body(placeholders[1], [left_elems[0]])
        left_content = left_elems[1:] if len(left_elems) > 1 else left_elems
    else:
        left_content = left_elems

    # Right sub-header (idx=3, BODY) - use first right element as header
    if 3 in placeholders and right_elems:
        format_comparison_body(placeholders[3], [right_elems[0]])
        right_content = right_elems[1:] if len(right_elems) > 1 else right_elems
    else:
        right_content = right_elems

    # Left content (idx=2, OBJECT)
    if 2 in placeholders and left_content:
        format_text_frame(placeholders[2].text_frame, left_content, font_size_pt=13)

    # Right content (idx=4, OBJECT)
    if 4 in placeholders and right_content:
        format_text_frame(placeholders[4].text_frame, right_content, font_size_pt=13)


def _populate_stat_callout(slide: Any, slide_obj: SlideObject) -> None:
    """Populate Layout 5 (Title Only) for STAT_CALLOUT.

    idx=0 (TITLE): Slide title
    Custom textbox: Large stat number + supporting points
    """
    _set_title(slide, slide_obj.title)

    elements = slide_obj.body_content.text_elements if slide_obj.body_content else []
    if not elements:
        return

    # Add accent bar
    add_accent_bar(slide)

    # Custom stat callout textbox
    left = Inches(1.5)
    top = Inches(2.2)
    width = Inches(10.0)
    height = Inches(4.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    format_stat_callout(tf, elements)


def _populate_team(slide: Any, slide_obj: SlideObject) -> None:
    """Populate Layout 5 (Title Only) for TEAM layout.

    Renders as a styled table with navy header row.
    """
    _set_title(slide, slide_obj.title)

    elements = slide_obj.body_content.text_elements if slide_obj.body_content else []
    if not elements:
        return

    # Add accent bar
    add_accent_bar(slide)

    # If pipe-table data, use styled table
    if is_pipe_table(elements):
        add_pptx_table(slide, elements, top_inches=2.2)
    else:
        # Render as em-dash separated table
        add_styled_table_from_elements(
            slide, elements, has_header=False, top_inches=2.2,
        )


def _populate_timeline(slide: Any, slide_obj: SlideObject) -> None:
    """Populate Layout 5 (Title Only) for TIMELINE layout.

    Renders as a styled table with phase/period rows.
    """
    _set_title(slide, slide_obj.title)

    elements = slide_obj.body_content.text_elements if slide_obj.body_content else []
    if not elements:
        return

    # Add accent bar
    add_accent_bar(slide)

    # If pipe-table data, use styled table
    if is_pipe_table(elements):
        add_pptx_table(slide, elements, top_inches=2.2)
    else:
        # Render as em-dash separated table
        add_styled_table_from_elements(
            slide, elements, has_header=False, top_inches=2.2,
        )


def _populate_compliance_matrix(slide: Any, slide_obj: SlideObject) -> None:
    """Populate Layout 5 (Title Only) for COMPLIANCE_MATRIX.

    Always renders as a styled table with status color coding.
    """
    _set_title(slide, slide_obj.title)

    elements = slide_obj.body_content.text_elements if slide_obj.body_content else []
    if not elements:
        return

    # Add accent bar
    add_accent_bar(slide)

    # Compliance matrix always renders as table
    if is_pipe_table(elements):
        add_pptx_table(slide, elements, top_inches=2.2)
    else:
        add_styled_table_from_elements(
            slide, elements, has_header=False, top_inches=2.2,
        )


def _populate_framework(slide: Any, slide_obj: SlideObject) -> None:
    """Populate Layout 1 (Title and Content) for FRAMEWORK.

    Uses the standard content placeholder with formatted bullets.
    """
    _set_title(slide, slide_obj.title)

    elements = slide_obj.body_content.text_elements if slide_obj.body_content else []
    if not elements:
        return

    # Check for pipe table
    if is_pipe_table(elements):
        add_pptx_table(slide, elements)
        return

    body_ph = _find_body_placeholder(slide)
    if body_ph is not None:
        format_text_frame(body_ph.text_frame, elements)
    else:
        _add_body_textbox(slide, elements)


def _populate_closing(slide: Any, slide_obj: SlideObject) -> None:
    """Populate Layout 0 (Title Slide) for CLOSING.

    Bookend that matches the opening. Uses CENTER_TITLE + SUBTITLE.
    """
    elements = slide_obj.body_content.text_elements if slide_obj.body_content else []
    format_closing_slide(slide, slide_obj.title, elements)


def _populate_section(slide: Any, slide_obj: SlideObject) -> None:
    """Populate Layout 2 (Section Header).

    idx=0 (TITLE): Large section heading
    idx=1 (BODY): Optional descriptive text
    """
    _set_title(slide, slide_obj.title)

    elements = slide_obj.body_content.text_elements if slide_obj.body_content else []
    if not elements:
        return

    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1 and ph.has_text_frame:
            format_text_frame(ph.text_frame, elements, font_size_pt=16)
            return


# ----------------------------------------------------------------
# Generic helpers
# ----------------------------------------------------------------


def _set_title(slide: Any, title: str) -> None:
    """Set the title placeholder text (idx=0)."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = title
            return


def _find_body_placeholder(slide: Any) -> Any | None:
    """Find the best placeholder for body content.

    Priority:
      1. idx=1 with type BODY (2) or OBJECT (7) -- standard content layouts
      2. idx=1 with type SUBTITLE (4) -- Title Slide layout
      3. idx=1 with any text frame -- catch-all
    """
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1 and ph.placeholder_format.type in (2, 7):
            return ph
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1 and ph.has_text_frame:
            return ph
    return None


def _add_body_textbox(
    slide: Any,
    elements: list[str],
) -> None:
    """Add a text box with body content when no placeholder exists."""
    left = Inches(0.92)
    top = Inches(2.0)
    width = Inches(11.5)
    height = Inches(4.76)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    format_text_frame(tf, elements)


def _add_speaker_notes(slide: Any, notes: str) -> None:
    """Add speaker notes to the slide."""
    if not notes:
        return
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes


def _apply_rtl(slide: Any) -> None:
    """Apply RTL text direction to all paragraphs in the slide."""
    ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            pPr = paragraph._p.get_or_add_pPr()
            pPr.set(f"{{{ns}}}rtl", "1")


def _add_chart(
    slide: Any,
    slide_obj: SlideObject,
    log_entry: dict[str, Any],
) -> None:
    """Add a chart to the slide from chart_spec."""
    if not slide_obj.chart_spec:
        return

    spec = slide_obj.chart_spec
    supported = {"bar", "line", "pie"}

    if spec.type not in supported:
        log_entry["status"] = RenderStatus.WARNING
        log_entry["message"] += (
            f" Chart type '{spec.type}' not yet supported -- skipped."
        )
        return

    if not spec.y_axis or not spec.y_axis.get("values"):
        log_entry["status"] = RenderStatus.WARNING
        log_entry["message"] += " Chart has no data -- skipped."
        return

    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    chart_type_map = {
        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE,
        "pie": XL_CHART_TYPE.PIE,
    }

    chart_data = CategoryChartData()
    categories = spec.x_axis.get("values", []) if spec.x_axis else []
    chart_data.categories = categories
    chart_data.add_series(
        spec.y_axis.get("label", "Data"),
        spec.y_axis["values"],
    )

    x = Inches(1.5)
    y = Inches(2.0)
    cx = Inches(8.0)
    cy = Inches(4.5)

    slide.shapes.add_chart(
        chart_type_map[spec.type],
        x, y, cx, cy,
        chart_data,
    )


# ----------------------------------------------------------------
# Layout dispatch table
# ----------------------------------------------------------------

_LAYOUT_HANDLERS: dict[LayoutType, Any] = {
    LayoutType.TITLE: _populate_title_slide,
    LayoutType.AGENDA: _populate_agenda_slide,
    LayoutType.SECTION: _populate_section,
    LayoutType.CONTENT_1COL: _populate_content_1col,
    LayoutType.CONTENT_2COL: _populate_content_2col,
    LayoutType.DATA_CHART: _populate_content_1col,  # Chart handled separately
    LayoutType.FRAMEWORK: _populate_framework,
    LayoutType.COMPARISON: _populate_comparison,
    LayoutType.STAT_CALLOUT: _populate_stat_callout,
    LayoutType.TEAM: _populate_team,
    LayoutType.TIMELINE: _populate_timeline,
    LayoutType.COMPLIANCE_MATRIX: _populate_compliance_matrix,
    LayoutType.CLOSING: _populate_closing,
}


# ----------------------------------------------------------------
# PPTX Renderer
# ----------------------------------------------------------------


async def render_pptx(
    slides: list[SlideObject],
    template_path: str,
    output_path: str,
    language: Language = Language.EN,
) -> RenderResult:
    """Render validated slides into branded PPTX using the SG template.

    Args:
        slides: QA-validated SlideObjects from the pipeline.
        template_path: Path to Presentation6.pptx master template.
        output_path: Path to write the rendered .pptx file.
        language: Output language (EN or AR for RTL).

    Returns:
        RenderResult with pptx_path, slide_count, and per-slide render_log.
    """
    prs = Presentation(template_path)

    # Remove existing slides from template
    for sldId in list(prs.slides._sldIdLst):
        rId = sldId.get(qn("r:id"))
        if rId:
            prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(sldId)

    render_log: list[dict[str, Any]] = []

    for slide_obj in slides:
        log_entry: dict[str, Any] = {
            "slide_id": slide_obj.slide_id,
            "status": RenderStatus.SUCCESS,
            "message": f"Rendered {slide_obj.layout_type} slide.",
        }

        # Select layout
        layout_idx = _get_layout_index(slide_obj.layout_type)
        layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(layout)

        # Dispatch to layout-specific handler
        handler = _LAYOUT_HANDLERS.get(slide_obj.layout_type)
        if handler:
            try:
                handler(slide, slide_obj)
            except Exception as e:
                logger.warning(
                    "Layout handler for %s failed: %s. Falling back.",
                    slide_obj.layout_type, e,
                )
                log_entry["status"] = RenderStatus.WARNING
                log_entry["message"] += f" Handler error: {e}. Used fallback."
                # Fallback: basic title + body
                _set_title(slide, slide_obj.title)
                if slide_obj.body_content and slide_obj.body_content.text_elements:
                    body_ph = _find_body_placeholder(slide)
                    if body_ph:
                        format_text_frame(
                            body_ph.text_frame,
                            slide_obj.body_content.text_elements,
                        )
                    else:
                        _add_body_textbox(
                            slide, slide_obj.body_content.text_elements,
                        )
        else:
            # No specific handler -- generic populate
            _set_title(slide, slide_obj.title)
            if slide_obj.body_content and slide_obj.body_content.text_elements:
                body_ph = _find_body_placeholder(slide)
                if body_ph:
                    format_text_frame(
                        body_ph.text_frame,
                        slide_obj.body_content.text_elements,
                    )
                else:
                    _add_body_textbox(
                        slide, slide_obj.body_content.text_elements,
                    )

        # Add chart if specified
        _add_chart(slide, slide_obj, log_entry)

        # Add speaker notes
        _add_speaker_notes(slide, slide_obj.speaker_notes)

        # Apply RTL for Arabic
        if language == Language.AR:
            _apply_rtl(slide)

        render_log.append(log_entry)

    # Ensure output directory exists
    Path(output_path).parent.mkdir(parents=True, exist_ok=True)

    # Save
    prs.save(output_path)

    # Verification: reopen and confirm slide count matches
    verify_prs = Presentation(output_path)
    actual_count = len(verify_prs.slides)
    if actual_count != len(slides):
        logger.error(
            "PPTX verification failed: expected %d slides, found %d in %s",
            len(slides),
            actual_count,
            output_path,
        )

    return RenderResult(
        pptx_path=output_path,
        slide_count=actual_count,
        render_log=render_log,
    )


# ----------------------------------------------------------------
# DOCX Exporter -- Research Report
# ----------------------------------------------------------------


async def export_report_docx(
    report: ResearchReport,
    output_path: str,
    language: Language = Language.EN,
) -> str:
    """Export Research Report as .docx file."""
    doc = Document()
    doc.add_heading(report.title, level=0)

    for section in report.sections:
        doc.add_heading(section.heading, level=1)
        render_markdown_to_docx(doc, section.content_markdown)

    if report.all_gaps:
        doc.add_heading("Evidence Gaps", level=1)
        table = doc.add_table(rows=1, cols=4)
        table.style = "Table Grid"
        hdr = table.rows[0].cells
        hdr[0].text = "Gap ID"
        hdr[1].text = "Description"
        hdr[2].text = "Severity"
        hdr[3].text = "Action Required"

        for gap in report.all_gaps:
            row = table.add_row().cells
            row[0].text = gap.gap_id
            row[1].text = gap.description
            row[2].text = str(gap.severity)
            row[3].text = gap.action_required

    if report.source_index:
        doc.add_heading("Source Index", level=1)
        table = doc.add_table(rows=1, cols=3)
        table.style = "Table Grid"
        hdr = table.rows[0].cells
        hdr[0].text = "Claim ID"
        hdr[1].text = "Document"
        hdr[2].text = "Path"

        for src in report.source_index:
            row = table.add_row().cells
            row[0].text = src.claim_id
            row[1].text = src.document_title
            row[2].text = src.sharepoint_path

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    doc.save(output_path)
    return output_path
