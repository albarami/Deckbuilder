"""Integration tests for the LangGraph pipeline — all LLM calls mocked."""

import os
import tempfile
from pathlib import Path
from unittest.mock import AsyncMock, patch

import pytest

_TEMPLATE_DIR = Path("PROPOSAL_TEMPLATE")
_TEMPLATE_EN = _TEMPLATE_DIR / "PROPOSAL_TEMPLATE EN.potx"


def _can_open_template() -> bool:
    """Check if python-pptx can actually open the template file."""
    if not _TEMPLATE_EN.exists():
        return False
    try:
        from pptx import Presentation
        Presentation(str(_TEMPLATE_EN))
        return True
    except Exception:
        return False


_HAS_TEMPLATES = _can_open_template()
_skip_no_template = pytest.mark.skipif(
    not _HAS_TEMPLATES,
    reason="PROPOSAL_TEMPLATE .potx not available or not openable by python-pptx",
)
from langgraph.types import Command

from src.models.claims import ClaimObject, GapObject, ReferenceIndex, SourceManifestEntry
from src.models.common import BilingualText
from src.models.enums import PipelineStage, RendererMode
from src.models.iterative import DeckDraft, DeckReview, SlideCritique, SlideText
from src.models.qa import DeckValidationSummary, QAResult, SlideValidation
from src.models.report import ReportSection, ResearchReport
from src.models.retrieval import (
    RankedSourcesOutput,
    RetrievalQueries,
    RetrievalSummary,
    SearchQuery,
)
from src.models.rfp import RFPContext
from src.models.slides import (
    BodyContent,
    SlideObject,
    SlideOutline,
    WrittenSlides,
)
from src.models.state import DeckForgeState, RetrievedSource
from src.services.llm import LLMResponse

# ──────────────────────────────────────────────────────────────
# Fixtures — mock LLM responses for every agent in the pipeline
# ──────────────────────────────────────────────────────────────

def _rfp_context() -> RFPContext:
    return RFPContext(
        rfp_name=BilingualText(en="SAP Support Renewal"),
        issuing_entity=BilingualText(en="SIDF"),
        mandate=BilingualText(en="Renew SAP licenses for 24 months."),
    )


def _retrieval_queries() -> RetrievalQueries:
    return RetrievalQueries(
        search_queries=[
            SearchQuery(
                query="SAP support experience",
                strategy="rfp_aligned",
                language="en",
                priority="high",
            ),
        ],
        retrieval_summary=RetrievalSummary(
            total_queries=1,
            by_strategy={"rfp_aligned": 1},
        ),
    )


def _ranked_sources() -> RankedSourcesOutput:
    return RankedSourcesOutput(
        ranked_sources=[
            RetrievedSource(
                doc_id="DOC-001",
                title="SAP Case Study",
                relevance_score=95,
            ),
        ],
    )


def _ranked_sources_two_docs() -> RankedSourcesOutput:
    return RankedSourcesOutput(
        ranked_sources=[
            RetrievedSource(
                doc_id="DOC-001",
                title="SAP Case Study",
                relevance_score=95,
            ),
            RetrievedSource(
                doc_id="DOC-002",
                title="SAP Delivery Framework",
                relevance_score=89,
            ),
        ],
    )


def _reference_index() -> ReferenceIndex:
    return ReferenceIndex(
        claims=[
            ClaimObject(
                claim_id="CLM-001",
                claim_text="Delivered SAP migration in 6 months",
                source_doc_id="DOC-001",
                source_location="Slide 3",
                evidence_span="Completed migration in 6 months",
                sensitivity_tag="capability",
                category="project_reference",
                confidence=0.95,
            ),
        ],
        gaps=[
            GapObject(
                gap_id="GAP-001",
                description="No ISO 27001 evidence",
                rfp_criterion="Compliance",
                severity="medium",
                action_required="Request certification docs",
            ),
        ],
        source_manifest=[
            SourceManifestEntry(
                doc_id="DOC-001",
                title="SAP Case Study",
                sharepoint_path="/test_docs/sap_case.pdf",
            ),
        ],
    )


def _research_report() -> ResearchReport:
    return ResearchReport(
        title="SAP Support Renewal — Research Report",
        language="en",
        sections=[
            ReportSection(
                section_id="SEC-001",
                heading="Executive Summary",
                content_markdown=(
                    "Strategic Gears has deep SAP expertise."
                    " [Ref: CLM-001]"
                ),
            ),
        ],
        full_markdown="# SAP Support Renewal\n\nResearch report content.",
    )


def _research_report_without_full_markdown() -> ResearchReport:
    report = _research_report()
    report.full_markdown = ""
    return report


def _slide_outline() -> SlideOutline:
    return SlideOutline(
        slides=[
            SlideObject(
                slide_id="S-001",
                title="Executive Summary",
                layout_type="TITLE",
            ),
            SlideObject(
                slide_id="S-002",
                title="Our SAP Experience",
                layout_type="CONTENT_1COL",
            ),
        ],
    )


def _written_slides() -> WrittenSlides:
    return WrittenSlides(
        slides=[
            SlideObject(
                slide_id="S-001",
                title="Executive Summary",
                layout_type="TITLE",
                body_content=BodyContent(
                    text_elements=["Strategic Gears — SAP expertise"],
                ),
            ),
            SlideObject(
                slide_id="S-002",
                title="Our SAP Experience",
                layout_type="CONTENT_1COL",
                body_content=BodyContent(
                    text_elements=["6-month SAP migration [Ref: CLM-001]"],
                ),
            ),
        ],
    )


def _qa_result() -> QAResult:
    return QAResult(
        slide_validations=[
            SlideValidation(slide_id="S-001", status="PASS"),
            SlideValidation(slide_id="S-002", status="PASS"),
        ],
        deck_summary=DeckValidationSummary(
            total_slides=2,
            passed=2,
            failed=0,
        ),
    )


def _llm_response(parsed: object) -> LLMResponse:
    """Wrap any parsed model in an LLMResponse."""
    return LLMResponse(
        parsed=parsed,
        input_tokens=1000,
        output_tokens=200,
        model="mock-model",
        latency_ms=500.0,
    )


_SEARCH_RESULTS = [
    {
        "doc_id": "DOC-001",
        "title": "SAP Case Study",
        "excerpt": "Delivered SAP migration in 6 months.",
        "metadata": {},
        "search_score": 0.95,
    },
]


def _deck_draft() -> DeckDraft:
    return DeckDraft(
        slides=[
            SlideText(slide_number=1, title="Executive Summary",
                      bullets=["Strategic Gears — SAP expertise"],
                      evidence_level="sourced"),
            SlideText(slide_number=2, title="Our SAP Experience",
                      bullets=["6-month SAP migration [Ref: CLM-001]"],
                      evidence_level="sourced"),
        ],
        turn_number=1,
        mode="strict",
    )


def _deck_review() -> DeckReview:
    return DeckReview(
        critiques=[
            SlideCritique(slide_number=1, score=4, issues=[]),
            SlideCritique(slide_number=2, score=3, issues=["Add more detail"]),
        ],
        overall_score=4,
        turn_number=2,
    )


def _refined_deck_draft() -> DeckDraft:
    return DeckDraft(
        slides=[
            SlideText(slide_number=1, title="Executive Summary",
                      bullets=["Strategic Gears — SAP expertise"],
                      evidence_level="sourced"),
            SlideText(slide_number=2, title="Our SAP Experience",
                      bullets=["Completed SAP migration in 6 months [Ref: CLM-001]"],
                      evidence_level="sourced"),
        ],
        turn_number=3,
        mode="strict",
    )


def _final_deck_review() -> DeckReview:
    return DeckReview(
        critiques=[
            SlideCritique(slide_number=1, score=5, issues=[]),
            SlideCritique(slide_number=2, score=4, issues=[]),
        ],
        overall_score=5,
        turn_number=4,
    )


def _make_input_state() -> DeckForgeState:
    """Minimal DeckForgeState for pipeline intake.

    Uses LEGACY renderer because these integration tests don't have
    the official .potx template or catalog lock files available.
    """
    return DeckForgeState(
        ai_assist_summary="SAP Support Renewal RFP from SIDF",
        output_language="en",
        renderer_mode=RendererMode.LEGACY,
    )


# ──────────────────────────────────────────────────────────────
# Tests
# ──────────────────────────────────────────────────────────────


@_skip_no_template
@pytest.mark.asyncio
async def test_full_pipeline_happy_path() -> None:
    """Mock all LLM calls, auto-approve all gates, verify state at end."""
    from src.pipeline.graph import build_graph

    graph = build_graph()
    state = _make_input_state()

    # Patch call_llm at each agent's module level
    agent_patches = {
        "src.agents.context.agent.call_llm": _llm_response(_rfp_context()),
        "src.agents.retrieval.planner.call_llm": _llm_response(
            _retrieval_queries()
        ),
        "src.agents.retrieval.ranker.call_llm": _llm_response(
            _ranked_sources()
        ),
        "src.agents.analysis.agent.call_llm": _llm_response(
            _reference_index()
        ),
        "src.agents.research.agent.call_llm": _llm_response(
            _research_report()
        ),
        # 5-turn iterative builder agents (M10)
        "src.agents.draft.agent.call_llm": _llm_response(
            _deck_draft()
        ),
        "src.agents.review.agent.call_llm": _llm_response(
            _deck_review()
        ),
        "src.agents.refine.agent.call_llm": _llm_response(
            _refined_deck_draft()
        ),
        "src.agents.final_review.agent.call_llm": _llm_response(
            _final_deck_review()
        ),
        "src.agents.presentation.agent.call_llm": _llm_response(
            _written_slides()
        ),
        "src.agents.qa.agent.call_llm": _llm_response(_qa_result()),
    }

    patches = []
    for path, return_val in agent_patches.items():
        m = AsyncMock(return_value=return_val)
        p = patch(path, m)
        patches.append(p)

    # Patch search and document loading at graph module level
    search_patch = patch(
        "src.pipeline.graph.semantic_search",
        new_callable=AsyncMock,
        return_value=_SEARCH_RESULTS,
    )
    load_docs_patch = patch(
        "src.services.search.load_full_documents",
        new_callable=AsyncMock,
        return_value=[
            {
                "doc_id": "DOC-001",
                "title": "SAP Case Study",
                "content_text": "SAP migration content",
                "metadata": {},
            },
        ],
    )
    patches.append(search_patch)
    patches.append(load_docs_patch)

    for p in patches:
        p.start()
    try:
        config = {"configurable": {"thread_id": "test-happy"}}
        result = await graph.ainvoke(state, config)

        # Resume through all 5 gates
        for _ in range(5):
            result = await graph.ainvoke(
                Command(resume={"approved": True}), config
            )

        # Pipeline should complete — render is the last node
        assert result["qa_result"] is not None
        assert result["rfp_context"] is not None
        assert result["written_slides"] is not None
        assert result["session"].total_llm_calls >= 10
        # Render node should have produced PPTX and DOCX
        assert result.get("pptx_path") is not None
        assert result.get("report_docx_path") is not None
        assert result["current_stage"] == PipelineStage.FINALIZED
    finally:
        for p in patches:
            p.stop()


@pytest.mark.asyncio
async def test_pipeline_stops_at_gate() -> None:
    """Pipeline pauses at gate_1 with interrupt — state preserved."""
    from src.pipeline.graph import build_graph

    graph = build_graph()
    state = _make_input_state()

    with patch(
        "src.agents.context.agent.call_llm", new_callable=AsyncMock
    ) as mock_llm:
        mock_llm.return_value = _llm_response(_rfp_context())

        config = {"configurable": {"thread_id": "test-gate-stop"}}
        result = await graph.ainvoke(state, config)

        # Should have interrupted at gate_1
        assert "__interrupt__" in result
        assert result["rfp_context"] is not None
        assert result["current_stage"] == PipelineStage.CONTEXT_REVIEW


@pytest.mark.asyncio
async def test_pipeline_error_handling() -> None:
    """Context agent fails — verify error state."""
    from src.pipeline.graph import build_graph
    from src.services.llm import LLMError

    graph = build_graph()
    state = _make_input_state()

    with patch(
        "src.agents.context.agent.call_llm", new_callable=AsyncMock
    ) as mock_llm:
        mock_llm.side_effect = LLMError(
            model="gpt-5.4",
            attempts=4,
            last_error=TimeoutError("timed out"),
        )

        config = {"configurable": {"thread_id": "test-error"}}
        result = await graph.ainvoke(state, config)

        assert result["current_stage"] == PipelineStage.ERROR
        assert len(result["errors"]) == 1
        assert result["errors"][0].agent == "context_agent"


def test_gate_three_summary_uses_report_sections_when_markdown_missing() -> None:
    """Gate 3 summary should still acknowledge a generated report without full markdown."""
    from src.pipeline.graph import _gate_3_summary

    state = _make_input_state()
    state.research_report = _research_report_without_full_markdown()
    state.report_markdown = ""

    summary = _gate_3_summary(state)

    assert "Research report ready" in summary
    assert "1 sections" in summary


@pytest.mark.asyncio
async def test_retrieval_chain() -> None:
    """Verify planner -> search -> ranker chaining in retrieval node."""
    from src.pipeline.graph import build_graph

    graph = build_graph()
    state = _make_input_state()

    with (
        patch(
            "src.agents.context.agent.call_llm",
            new_callable=AsyncMock,
            return_value=_llm_response(_rfp_context()),
        ),
        patch(
            "src.agents.retrieval.planner.call_llm",
            new_callable=AsyncMock,
            return_value=_llm_response(_retrieval_queries()),
        ) as mock_planner,
        patch(
            "src.agents.retrieval.ranker.call_llm",
            new_callable=AsyncMock,
            return_value=_llm_response(_ranked_sources()),
        ) as mock_ranker,
        patch(
            "src.pipeline.graph.semantic_search",
            new_callable=AsyncMock,
            return_value=_SEARCH_RESULTS,
        ) as mock_search,
    ):
        config = {"configurable": {"thread_id": "test-retrieval"}}
        # Run — context executes, gate_1 interrupts
        result = await graph.ainvoke(state, config)
        # Resume past gate_1
        result = await graph.ainvoke(
            Command(resume={"approved": True}), config
        )

        # Should have interrupted at gate_2 now
        assert "__interrupt__" in result
        assert result["current_stage"] == PipelineStage.SOURCE_REVIEW
        assert len(result["retrieved_sources"]) >= 1

        # Verify all three steps in the chain ran
        mock_planner.assert_called_once()
        mock_search.assert_called_once()
        mock_ranker.assert_called_once()


@pytest.mark.asyncio
async def test_gate_two_resume_applies_selected_source_ids() -> None:
    """Gate 2 resume should persist reviewer-selected source IDs into graph state."""
    from src.pipeline.graph import build_graph

    graph = build_graph()
    state = _make_input_state()

    async def _assembly_plan_run(state: DeckForgeState) -> dict:
        """Mock assembly plan agent that returns a minimal result."""
        return {
            "current_stage": PipelineStage.ANALYSIS,
            "assembly_plan": None,
            "methodology_blueprint": None,
            "slide_budget": None,
            "sector": "technology",
            "geography": "ksa",
            "proposal_mode": "standard",
            "session": state.session,
        }

    with (
        patch(
            "src.agents.context.agent.call_llm",
            new_callable=AsyncMock,
            return_value=_llm_response(_rfp_context()),
        ),
        patch(
            "src.agents.retrieval.planner.call_llm",
            new_callable=AsyncMock,
            return_value=_llm_response(_retrieval_queries()),
        ),
        patch(
            "src.agents.retrieval.ranker.call_llm",
            new_callable=AsyncMock,
            return_value=_llm_response(_ranked_sources_two_docs()),
        ),
        patch(
            "src.pipeline.graph.semantic_search",
            new_callable=AsyncMock,
            return_value=_SEARCH_RESULTS,
        ),
        patch(
            "src.services.search.load_full_documents",
            new_callable=AsyncMock,
            return_value=[
                {"doc_id": "DOC-001", "title": "SAP Case Study", "content_text": "Case study"},
                {"doc_id": "DOC-002", "title": "SAP Delivery Framework", "content_text": "Framework"},
            ],
        ),
        patch(
            "src.pipeline.graph.assembly_plan_agent.run",
            new_callable=AsyncMock,
            side_effect=_assembly_plan_run,
        ),
    ):
        config = {"configurable": {"thread_id": "test-gate-2-selection"}}

        result = await graph.ainvoke(state, config)
        assert "__interrupt__" in result

        result = await graph.ainvoke(
            Command(resume={"approved": True}),
            config,
        )
        assert "__interrupt__" in result

        result = await graph.ainvoke(
            Command(
                resume={
                    "approved": True,
                    "modifications": {
                        "included_sources": ["DOC-002"],
                    },
                }
            ),
            config,
        )

        assert "__interrupt__" in result
        assert result["approved_source_ids"] == ["DOC-002"]
        assert result["current_stage"] == PipelineStage.ANALYSIS


@pytest.mark.asyncio
async def test_state_persistence() -> None:
    """Verify state saves to JSON and reloads correctly."""
    from src.pipeline.graph import load_state, save_state

    state = _make_input_state()
    state.rfp_context = _rfp_context()
    state.current_stage = PipelineStage.CONTEXT_REVIEW

    with tempfile.TemporaryDirectory() as tmpdir:
        path = os.path.join(tmpdir, "session.json")

        # Save
        save_state(state, path)
        assert os.path.exists(path)

        # Load
        loaded = load_state(path)
        assert loaded.current_stage == PipelineStage.CONTEXT_REVIEW
        assert loaded.rfp_context is not None
        assert loaded.rfp_context.rfp_name.en == "SAP Support Renewal"
        assert loaded.ai_assist_summary == (
            "SAP Support Renewal RFP from SIDF"
        )


@_skip_no_template
@pytest.mark.asyncio
async def test_full_pipeline_with_render() -> None:
    """Full pipeline through all gates → render produces .pptx and .docx."""
    from pptx import Presentation as PptxPresentation

    from src.pipeline.dry_run import get_dry_run_patches
    from src.pipeline.graph import build_graph

    graph = build_graph()
    state = _make_input_state()

    patches = get_dry_run_patches()
    for p in patches:
        p.start()
    try:
        config = {"configurable": {"thread_id": "test-render"}}
        result = await graph.ainvoke(state, config)  # type: ignore[arg-type]

        # Resume through all 5 gates
        for _ in range(5):
            result = await graph.ainvoke(
                Command(resume={"approved": True}), config  # type: ignore[arg-type]
            )

        # Verify render output
        assert result["current_stage"] == PipelineStage.FINALIZED
        pptx_path = result.get("pptx_path")
        docx_path = result.get("report_docx_path")

        assert pptx_path is not None
        assert os.path.exists(pptx_path)
        assert docx_path is not None
        assert os.path.exists(docx_path)

        # Verify PPTX has correct slide count (2 from dry-run mocks)
        prs = PptxPresentation(pptx_path)
        assert len(prs.slides) == 2
    finally:
        for p in patches:
            p.stop()
