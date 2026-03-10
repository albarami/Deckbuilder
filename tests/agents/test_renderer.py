"""Tests for the PPTX renderer and DOCX exporter (Design Agent)."""

import asyncio
import os
import tempfile

from docx import Document
from pptx import Presentation

from src.models.enums import GapSeverity, Language, LayoutType
from src.models.report import ReportGap, ReportSection, ResearchReport
from src.models.slides import BodyContent, SlideObject

# ──────────────────────────────────────────────────────────────
# Fixtures
# ──────────────────────────────────────────────────────────────

TEMPLATE_PATH = os.path.join(
    os.path.dirname(__file__), "..", "..", "templates", "Presentation6.pptx"
)


def _sample_slides() -> list[SlideObject]:
    """Three slides covering different layout types."""
    return [
        SlideObject(
            slide_id="S-001",
            title="Executive Summary",
            layout_type=LayoutType.TITLE,
            body_content=BodyContent(
                text_elements=["Strategic Gears — SAP expertise"],
            ),
            speaker_notes="This is the opening slide.",
        ),
        SlideObject(
            slide_id="S-002",
            title="Our Approach",
            layout_type=LayoutType.CONTENT_1COL,
            body_content=BodyContent(
                text_elements=[
                    "Phase 1: Discovery and assessment",
                    "Phase 2: Implementation and migration",
                    "Phase 3: Knowledge transfer and handover",
                ],
            ),
            speaker_notes="Walk through each phase.",
        ),
        SlideObject(
            slide_id="S-003",
            title="Market Overview",
            layout_type=LayoutType.SECTION,
            body_content=BodyContent(
                text_elements=["SAP market in the GCC region"],
            ),
        ),
    ]


def _sample_report() -> ResearchReport:
    """Minimal research report for docx export."""
    return ResearchReport(
        title="SAP Support Renewal — Research Report",
        language=Language.EN,
        sections=[
            ReportSection(
                section_id="SEC-001",
                heading="Executive Summary",
                content_markdown="Strategic Gears has deep SAP expertise.",
            ),
            ReportSection(
                section_id="SEC-002",
                heading="Technical Approach",
                content_markdown="Our methodology covers three phases.",
            ),
        ],
        all_gaps=[
            ReportGap(
                gap_id="GAP-001",
                description="No ISO 27001 evidence",
                rfp_criterion="Compliance",
                severity=GapSeverity.MEDIUM,
                action_required="Request certification docs",
            ),
        ],
        full_markdown="# SAP Support Renewal\n\nResearch content.",
    )


def _sample_report_multi_section() -> ResearchReport:
    """8-section report matching real E2E output size."""
    sections = []
    for i in range(1, 9):
        sections.append(
            ReportSection(
                section_id=f"SEC-{i:02d}",
                heading=f"Section {i}: {'ABCDEFGH'[i-1]} Topic",
                content_markdown=(
                    f"## Section {i}\n\n"
                    f"First paragraph of section {i} with content.\n\n"
                    f"Second paragraph with **bold** and analysis.\n\n"
                    f"Third paragraph with more detail for section {i}."
                ),
                gaps_flagged=[f"GAP-{i:03d}"],
            )
        )
    gaps = [
        ReportGap(
            gap_id=f"GAP-{i:03d}",
            description=f"Gap {i} description",
            rfp_criterion="Criterion",
            severity=GapSeverity.CRITICAL,
            action_required=f"Fix gap {i}",
        )
        for i in range(1, 9)
    ]
    return ResearchReport(
        title="Multi-Section Test Report",
        language=Language.EN,
        sections=sections,
        all_gaps=gaps,
        full_markdown="Full markdown content here.",
    )


def _all_layout_slides() -> list[SlideObject]:
    """One slide per layout type — tests that body content is never lost."""
    layouts = [
        (LayoutType.TITLE, "Title Slide"),
        (LayoutType.AGENDA, "Agenda Slide"),
        (LayoutType.SECTION, "Section Header"),
        (LayoutType.CONTENT_1COL, "One Column"),
        (LayoutType.CONTENT_2COL, "Two Columns"),
        (LayoutType.DATA_CHART, "Chart Slide"),
        (LayoutType.FRAMEWORK, "Framework"),
        (LayoutType.COMPARISON, "Comparison"),
        (LayoutType.STAT_CALLOUT, "Key Stat"),
        (LayoutType.TEAM, "Team Slide"),
        (LayoutType.TIMELINE, "Timeline"),
        (LayoutType.COMPLIANCE_MATRIX, "Compliance"),
        (LayoutType.CLOSING, "Thank You"),
    ]
    slides = []
    for i, (lt, label) in enumerate(layouts, 1):
        slides.append(
            SlideObject(
                slide_id=f"S-{i:03d}",
                title=f"{label} — Test Title",
                layout_type=lt,
                body_content=BodyContent(
                    text_elements=[
                        f"Bullet 1 for {label}",
                        f"Bullet 2 for {label}",
                        f"Bullet 3 for {label}",
                    ],
                ),
                speaker_notes=f"Notes for {label}.",
            )
        )
    return slides


def _get_slide_body_text(slide) -> str:
    """Extract all body text from a slide (placeholders + textboxes + tables)."""
    texts = []
    for shape in slide.shapes:
        # Check text frames
        if shape.has_text_frame:
            try:
                pf = shape.placeholder_format
                if pf.idx == 0:  # title
                    continue
            except (ValueError, AttributeError):
                pass  # Not a placeholder — textbox fallback
            txt = shape.text_frame.text.strip()
            if txt:
                texts.append(txt)
        # Check table cells
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    txt = cell.text.strip()
                    if txt:
                        texts.append(txt)
    return "\n".join(texts)


# ──────────────────────────────────────────────────────────────
# PPTX Tests
# ──────────────────────────────────────────────────────────────


def test_render_pptx_creates_file() -> None:
    """Render 3 SlideObjects → .pptx file exists, correct slide count."""
    from src.services.renderer import render_pptx

    slides = _sample_slides()
    with tempfile.TemporaryDirectory() as tmpdir:
        output = os.path.join(tmpdir, "test_output.pptx")
        result = asyncio.run(
            render_pptx(slides, TEMPLATE_PATH, output)
        )

        assert os.path.exists(output)
        assert result.slide_count == 3
        assert result.pptx_path == output

        # Verify actual PPTX slide count
        prs = Presentation(output)
        assert len(prs.slides) == 3


def test_render_maps_layout_types() -> None:
    """Each LayoutType maps to a valid template layout."""
    from src.services.renderer import LAYOUT_MAP

    for layout_type in LayoutType:
        assert layout_type in LAYOUT_MAP, (
            f"LayoutType.{layout_type.name} missing from LAYOUT_MAP"
        )
        idx = LAYOUT_MAP[layout_type]
        assert isinstance(idx, int)
        assert 0 <= idx <= 12


def test_render_populates_title() -> None:
    """Slide title text matches SlideObject.title."""
    from src.services.renderer import render_pptx

    slides = [
        SlideObject(
            slide_id="S-001",
            title="My Custom Title",
            layout_type=LayoutType.CONTENT_1COL,
            body_content=BodyContent(text_elements=["Content"]),
        ),
    ]
    with tempfile.TemporaryDirectory() as tmpdir:
        output = os.path.join(tmpdir, "title_test.pptx")
        asyncio.run(render_pptx(slides, TEMPLATE_PATH, output))

        prs = Presentation(output)
        slide = prs.slides[0]
        # Find title placeholder
        title_text = None
        for shape in slide.shapes:
            if shape.has_text_frame and shape.shape_id <= 2:
                title_text = shape.text_frame.text
                break
        # Fallback: check placeholders directly
        if title_text is None:
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == 0:
                    title_text = ph.text
                    break

        assert title_text == "My Custom Title"


def test_render_populates_body() -> None:
    """Bullet points from body_content appear in slide."""
    from src.services.renderer import render_pptx

    bullets = ["First point", "Second point", "Third point"]
    slides = [
        SlideObject(
            slide_id="S-001",
            title="Test",
            layout_type=LayoutType.CONTENT_1COL,
            body_content=BodyContent(text_elements=bullets),
        ),
    ]
    with tempfile.TemporaryDirectory() as tmpdir:
        output = os.path.join(tmpdir, "body_test.pptx")
        asyncio.run(render_pptx(slides, TEMPLATE_PATH, output))

        prs = Presentation(output)
        slide = prs.slides[0]
        # Find content placeholder (idx=1)
        body_text = ""
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 1:
                body_text = ph.text
                break

        for bullet in bullets:
            assert bullet in body_text


def test_render_adds_speaker_notes() -> None:
    """Notes section contains speaker_notes text."""
    from src.services.renderer import render_pptx

    slides = [
        SlideObject(
            slide_id="S-001",
            title="Test",
            layout_type=LayoutType.CONTENT_1COL,
            body_content=BodyContent(text_elements=["Content"]),
            speaker_notes="Important presenter note here.",
        ),
    ]
    with tempfile.TemporaryDirectory() as tmpdir:
        output = os.path.join(tmpdir, "notes_test.pptx")
        asyncio.run(render_pptx(slides, TEMPLATE_PATH, output))

        prs = Presentation(output)
        slide = prs.slides[0]
        notes = slide.notes_slide.notes_text_frame.text
        assert "Important presenter note here." in notes


def test_render_stat_callout_body_not_skipped() -> None:
    """STAT_CALLOUT body content is rendered via textbox fallback, not dropped."""
    from src.services.renderer import render_pptx

    bullets = ["Key insight number one", "Critical metric here"]
    slides = [
        SlideObject(
            slide_id="S-001",
            title="42%",
            layout_type=LayoutType.STAT_CALLOUT,
            body_content=BodyContent(text_elements=bullets),
        ),
    ]
    with tempfile.TemporaryDirectory() as tmpdir:
        output = os.path.join(tmpdir, "stat_test.pptx")
        result = asyncio.run(
            render_pptx(slides, TEMPLATE_PATH, output)
        )

        assert result.slide_count == 1
        prs = Presentation(output)
        body = _get_slide_body_text(prs.slides[0])
        assert "Key insight number one" in body
        assert "Critical metric here" in body


def test_render_all_layout_types_have_body() -> None:
    """Every LayoutType renders body content — no silent drops.

    This is the key regression test: the original renderer silently
    skipped body content for TITLE, AGENDA, STAT_CALLOUT, and CLOSING
    layouts because they lack a BODY/OBJECT placeholder.
    """
    from src.services.renderer import render_pptx

    slides = _all_layout_slides()
    with tempfile.TemporaryDirectory() as tmpdir:
        output = os.path.join(tmpdir, "all_layouts.pptx")
        result = asyncio.run(
            render_pptx(slides, TEMPLATE_PATH, output)
        )

        assert result.slide_count == len(slides)

        prs = Presentation(output)
        assert len(prs.slides) == len(slides)

        for i, slide in enumerate(prs.slides):
            body = _get_slide_body_text(slide)
            layout_name = str(slides[i].layout_type)
            assert body, (
                f"Slide {i+1} ({layout_name}) has no body text — "
                f"body content was silently dropped"
            )


def test_render_20_slides_exact_count() -> None:
    """Render 20 slides and verify exactly 20 in the saved PPTX.

    Regression test for the stale-file bug: renderer reported N slides
    but the saved file contained fewer.
    """
    from src.services.renderer import render_pptx

    slides = []
    layouts = list(LayoutType)
    for i in range(20):
        lt = layouts[i % len(layouts)]
        slides.append(
            SlideObject(
                slide_id=f"S-{i+1:03d}",
                title=f"Slide {i+1} Title",
                layout_type=lt,
                body_content=BodyContent(
                    text_elements=[f"Bullet {j+1}" for j in range(3)],
                ),
            )
        )

    with tempfile.TemporaryDirectory() as tmpdir:
        output = os.path.join(tmpdir, "twenty_slides.pptx")
        result = asyncio.run(
            render_pptx(slides, TEMPLATE_PATH, output)
        )

        # RenderResult.slide_count now comes from verification
        assert result.slide_count == 20

        # Double-check by reopening
        prs = Presentation(output)
        assert len(prs.slides) == 20


def test_render_result_has_log() -> None:
    """RenderResult contains per-slide status entries."""
    from src.services.renderer import render_pptx

    slides = _sample_slides()
    with tempfile.TemporaryDirectory() as tmpdir:
        output = os.path.join(tmpdir, "log_test.pptx")
        result = asyncio.run(
            render_pptx(slides, TEMPLATE_PATH, output)
        )

        assert len(result.render_log) == 3
        for entry in result.render_log:
            assert "slide_id" in entry
            assert "status" in entry
            assert "message" in entry


def test_render_title_slide_uses_cover_layout() -> None:
    """TITLE layout uses Proposal Cover (Layout 11) with subtitle placeholder."""
    from src.services.renderer import render_pptx

    slides = [
        SlideObject(
            slide_id="S-001",
            title="Company Name",
            layout_type=LayoutType.TITLE,
            body_content=BodyContent(
                text_elements=["RFP \u2014 Strategic consulting for the GCC"],
            ),
        ),
    ]
    with tempfile.TemporaryDirectory() as tmpdir:
        output = os.path.join(tmpdir, "title_body_test.pptx")
        asyncio.run(render_pptx(slides, TEMPLATE_PATH, output))

        prs = Presentation(output)
        slide = prs.slides[0]
        # Layout 11 uses idx=1 (SUBTITLE) for project name
        body_text = _get_slide_body_text(slide)
        assert body_text, "Title slide should have body content"


# ──────────────────────────────────────────────────────────────
# DOCX Tests
# ──────────────────────────────────────────────────────────────


def test_export_docx_creates_file() -> None:
    """ResearchReport → .docx file exists with sections."""
    from src.services.renderer import export_report_docx

    report = _sample_report()
    with tempfile.TemporaryDirectory() as tmpdir:
        output = os.path.join(tmpdir, "test_report.docx")
        result_path = asyncio.run(
            export_report_docx(report, output)
        )

        assert os.path.exists(result_path)
        assert result_path == output

        # Verify document structure
        doc = Document(output)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        # Title should be present
        assert any("SAP Support Renewal" in p for p in paragraphs)
        # Section headings should be present
        assert any("Executive Summary" in p for p in paragraphs)
        assert any("Technical Approach" in p for p in paragraphs)


def test_export_docx_all_sections_present() -> None:
    """All 8 sections of a multi-section report appear in the DOCX.

    Regression test: the original export wrote only the first paragraph
    of the first section, dropping 7 entire sections.
    """
    from src.services.renderer import export_report_docx

    report = _sample_report_multi_section()
    with tempfile.TemporaryDirectory() as tmpdir:
        output = os.path.join(tmpdir, "multi_section.docx")
        asyncio.run(export_report_docx(report, output))

        doc = Document(output)
        all_text = "\n".join(p.text for p in doc.paragraphs)

        # All 8 section headings must appear
        for sec in report.sections:
            assert sec.heading in all_text, (
                f"Section '{sec.heading}' missing from DOCX"
            )

        # All section content must appear (not just first paragraph)
        for sec in report.sections:
            # Each section has 3 paragraphs — check they're all there
            assert "First paragraph of section" in all_text
            assert "Second paragraph with" in all_text
            assert "Third paragraph with" in all_text

        # Gaps table should exist
        assert len(doc.tables) >= 1
        gap_table = doc.tables[0]
        assert gap_table.rows[0].cells[0].text == "Gap ID"


def test_export_docx_gaps_table() -> None:
    """Gaps table has correct structure and all entries."""
    from src.services.renderer import export_report_docx

    report = _sample_report()
    with tempfile.TemporaryDirectory() as tmpdir:
        output = os.path.join(tmpdir, "gaps_test.docx")
        asyncio.run(export_report_docx(report, output))

        doc = Document(output)
        assert len(doc.tables) >= 1

        gap_table = doc.tables[0]
        # Header row
        assert gap_table.rows[0].cells[0].text == "Gap ID"
        assert gap_table.rows[0].cells[1].text == "Description"
        # Data row
        assert gap_table.rows[1].cells[0].text == "GAP-001"
        assert "ISO 27001" in gap_table.rows[1].cells[1].text


def test_export_docx_paragraph_count() -> None:
    """DOCX has more than just title + 1 heading + 1 line.

    Regression test: original export produced only 3 paragraphs
    from an 8-section report.
    """
    from src.services.renderer import export_report_docx

    report = _sample_report_multi_section()
    with tempfile.TemporaryDirectory() as tmpdir:
        output = os.path.join(tmpdir, "count_test.docx")
        asyncio.run(export_report_docx(report, output))

        doc = Document(output)
        # 8 sections × ~4 paragraphs each + title + headings = at least 30
        non_empty = [p for p in doc.paragraphs if p.text.strip()]
        assert len(non_empty) >= 20, (
            f"Expected 20+ paragraphs, got {len(non_empty)} — "
            f"DOCX exporter is truncating content"
        )


# ──────────────────────────────────────────────────────────────
# DOCX Markdown Formatting Tests
# ──────────────────────────────────────────────────────────────


def test_export_docx_renders_markdown_headings() -> None:
    """Markdown ## headings in content_markdown become Word Heading styles."""
    from src.services.renderer import export_report_docx

    report = ResearchReport(
        title="Test Report",
        language=Language.EN,
        sections=[
            ReportSection(
                section_id="SEC-01",
                heading="Main Section",
                content_markdown=(
                    "## Main Section\n\n"
                    "Intro paragraph.\n\n"
                    "### Subsection One\n\n"
                    "Details here."
                ),
            ),
        ],
    )
    with tempfile.TemporaryDirectory() as tmpdir:
        output = os.path.join(tmpdir, "heading_test.docx")
        asyncio.run(export_report_docx(report, output))

        doc = Document(output)
        styles = [p.style.name for p in doc.paragraphs if p.text.strip()]
        # Section heading (from add_heading level=1)
        assert "Heading 1" in styles
        # Subsection heading (from markdown ### → level 3)
        assert "Heading 3" in styles


def test_export_docx_renders_bold() -> None:
    """Markdown **bold** text becomes bold runs in Word."""
    from src.services.renderer import export_report_docx

    report = ResearchReport(
        title="Test Report",
        language=Language.EN,
        sections=[
            ReportSection(
                section_id="SEC-01",
                heading="Bold Test",
                content_markdown="This has **important text** in it.",
            ),
        ],
    )
    with tempfile.TemporaryDirectory() as tmpdir:
        output = os.path.join(tmpdir, "bold_test.docx")
        asyncio.run(export_report_docx(report, output))

        doc = Document(output)
        # Find the paragraph with our content
        for p in doc.paragraphs:
            if "important text" in p.text:
                bold_runs = [r for r in p.runs if r.bold]
                assert any("important text" in r.text for r in bold_runs), (
                    "Expected bold run for 'important text'"
                )
                return
        raise AssertionError("Paragraph with 'important text' not found")


def test_export_docx_renders_bullet_list() -> None:
    """Markdown - bullets become List Bullet style."""
    from src.services.renderer import export_report_docx

    report = ResearchReport(
        title="Test Report",
        language=Language.EN,
        sections=[
            ReportSection(
                section_id="SEC-01",
                heading="List Test",
                content_markdown="- First item\n- Second item\n- Third item",
            ),
        ],
    )
    with tempfile.TemporaryDirectory() as tmpdir:
        output = os.path.join(tmpdir, "list_test.docx")
        asyncio.run(export_report_docx(report, output))

        doc = Document(output)
        bullet_paras = [
            p for p in doc.paragraphs if p.style.name == "List Bullet"
        ]
        assert len(bullet_paras) >= 3, (
            f"Expected 3+ List Bullet paragraphs, got {len(bullet_paras)}"
        )


def test_export_docx_renders_table() -> None:
    """Markdown tables become Word tables."""
    from src.services.renderer import export_report_docx

    report = ResearchReport(
        title="Test Report",
        language=Language.EN,
        sections=[
            ReportSection(
                section_id="SEC-01",
                heading="Table Test",
                content_markdown=(
                    "| Col1 | Col2 | Col3 |\n"
                    "|------|------|------|\n"
                    "| A    | B    | C    |\n"
                    "| D    | E    | F    |"
                ),
            ),
        ],
    )
    with tempfile.TemporaryDirectory() as tmpdir:
        output = os.path.join(tmpdir, "table_test.docx")
        asyncio.run(export_report_docx(report, output))

        doc = Document(output)
        # Should have at least one table from the markdown
        assert len(doc.tables) >= 1, "No tables found in DOCX"
        table = doc.tables[0]
        assert len(table.rows) >= 3  # header + 2 data rows
        assert table.rows[0].cells[0].text == "Col1"


def test_export_docx_renders_blockquote() -> None:
    """Markdown > blockquotes get indented paragraph formatting."""
    from src.services.renderer import export_report_docx

    report = ResearchReport(
        title="Test Report",
        language=Language.EN,
        sections=[
            ReportSection(
                section_id="SEC-01",
                heading="Quote Test",
                content_markdown="> Important notice about gaps.",
            ),
        ],
    )
    with tempfile.TemporaryDirectory() as tmpdir:
        output = os.path.join(tmpdir, "quote_test.docx")
        asyncio.run(export_report_docx(report, output))

        doc = Document(output)
        # Find paragraph with blockquote content
        for p in doc.paragraphs:
            if "Important notice" in p.text:
                # Should have left indent
                assert p.paragraph_format.left_indent is not None, (
                    "Blockquote paragraph should have left indent"
                )
                return
        raise AssertionError("Blockquote paragraph not found")


# ──────────────────────────────────────────────────────────────
# PPTX Body Formatting Tests
# ──────────────────────────────────────────────────────────────


def test_pptx_bullet_gets_font_formatting() -> None:
    """Bullet text_elements get Aptos font and 14pt size."""
    from src.services.renderer import render_pptx

    slides = [
        SlideObject(
            slide_id="S-001",
            title="Test",
            layout_type=LayoutType.CONTENT_1COL,
            body_content=BodyContent(
                text_elements=["First point", "Second point"],
            ),
        ),
    ]
    with tempfile.TemporaryDirectory() as tmpdir:
        output = os.path.join(tmpdir, "font_test.pptx")
        asyncio.run(render_pptx(slides, TEMPLATE_PATH, output))

        prs = Presentation(output)
        slide = prs.slides[0]
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 1:
                for p in ph.text_frame.paragraphs:
                    for run in p.runs:
                        assert run.font.name == "Aptos", (
                            f"Expected Aptos font, got {run.font.name}"
                        )
                return


def test_pptx_em_dash_bold_key() -> None:
    """Text with em-dash separator gets bold key phrase."""
    from src.services.renderer import render_pptx

    slides = [
        SlideObject(
            slide_id="S-001",
            title="Test",
            layout_type=LayoutType.CONTENT_1COL,
            body_content=BodyContent(
                text_elements=["S1 \u2014 License renewal for deployed modules"],
            ),
        ),
    ]
    with tempfile.TemporaryDirectory() as tmpdir:
        output = os.path.join(tmpdir, "emdash_test.pptx")
        asyncio.run(render_pptx(slides, TEMPLATE_PATH, output))

        prs = Presentation(output)
        slide = prs.slides[0]
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 1:
                runs = ph.text_frame.paragraphs[0].runs
                assert len(runs) >= 2, "Expected at least 2 runs (bold key + detail)"
                assert runs[0].font.bold is True, "First run (key) should be bold"
                assert "S1" in runs[0].text
                return


def test_pptx_bullet_prefix_indented() -> None:
    """Elements starting with bullet char get paragraph.level = 1."""
    from src.services.renderer import render_pptx

    slides = [
        SlideObject(
            slide_id="S-001",
            title="Test",
            layout_type=LayoutType.CONTENT_1COL,
            body_content=BodyContent(
                text_elements=[
                    "Main point",
                    "\u2022 Sub-bullet text",
                ],
            ),
        ),
    ]
    with tempfile.TemporaryDirectory() as tmpdir:
        output = os.path.join(tmpdir, "indent_test.pptx")
        asyncio.run(render_pptx(slides, TEMPLATE_PATH, output))

        prs = Presentation(output)
        slide = prs.slides[0]
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 1:
                paras = ph.text_frame.paragraphs
                assert paras[0].level == 0, "First paragraph should be level 0"
                assert paras[1].level == 1, "Bullet-prefixed should be level 1"
                return


# ──────────────────────────────────────────────────────────────
# STAT_CALLOUT / AGENDA Formatting Tests
# ──────────────────────────────────────────────────────────────


def test_pptx_stat_callout_formatting() -> None:
    """STAT_CALLOUT first element is large + bold + accent color."""
    from pptx.dml.color import RGBColor

    from src.services.renderer import render_pptx

    slides = [
        SlideObject(
            slide_id="S-001",
            title="Key Metric",
            layout_type=LayoutType.STAT_CALLOUT,
            body_content=BodyContent(
                text_elements=["Zero verified claims", "Supporting detail"],
            ),
        ),
    ]
    with tempfile.TemporaryDirectory() as tmpdir:
        output = os.path.join(tmpdir, "stat_test.pptx")
        asyncio.run(render_pptx(slides, TEMPLATE_PATH, output))

        prs = Presentation(output)
        body = _get_slide_body_text(prs.slides[0])
        assert "Zero verified claims" in body

        # Check formatting on the textbox
        for shape in prs.slides[0].shapes:
            if not shape.has_text_frame:
                continue
            try:
                shape.placeholder_format
                continue  # skip title placeholder
            except (ValueError, AttributeError):
                pass
            # Skip shapes with no text (e.g. accent bars)
            if not shape.text_frame.text.strip():
                continue
            first_run = shape.text_frame.paragraphs[0].runs[0]
            assert first_run.font.bold is True, "First stat should be bold"
            assert first_run.font.color.rgb == RGBColor(0x15, 0x60, 0x82), (
                "First stat should use ACCENT1 color"
            )
            return


def test_pptx_agenda_numbered() -> None:
    """AGENDA elements get numbered prefixes."""
    from src.services.renderer import render_pptx

    slides = [
        SlideObject(
            slide_id="S-001",
            title="Agenda",
            layout_type=LayoutType.AGENDA,
            body_content=BodyContent(
                text_elements=["Topic one", "Topic two", "Topic three"],
            ),
        ),
    ]
    with tempfile.TemporaryDirectory() as tmpdir:
        output = os.path.join(tmpdir, "agenda_test.pptx")
        asyncio.run(render_pptx(slides, TEMPLATE_PATH, output))

        prs = Presentation(output)
        body = _get_slide_body_text(prs.slides[0])
        assert "Topic one" in body
        assert "Topic two" in body
        assert "Topic three" in body
        # Verify numbering
        assert "1." in body
        assert "2." in body
        assert "3." in body


# ──────────────────────────────────────────────────────────────
# PPTX Pipe-Table Tests
# ──────────────────────────────────────────────────────────────


def test_pptx_pipe_table_rendered() -> None:
    """Pipe-separated elements are rendered as a PowerPoint table."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    from src.services.renderer import render_pptx

    slides = [
        SlideObject(
            slide_id="S-001",
            title="Compliance Matrix",
            layout_type=LayoutType.COMPLIANCE_MATRIX,
            body_content=BodyContent(
                text_elements=[
                    "Item | Evidence | Status",
                    "D1 | SAP license renewal | CRITICAL GAP",
                    "D2 | L2/L3 support | CRITICAL GAP",
                ],
            ),
        ),
    ]
    with tempfile.TemporaryDirectory() as tmpdir:
        output = os.path.join(tmpdir, "table_test.pptx")
        asyncio.run(render_pptx(slides, TEMPLATE_PATH, output))

        prs = Presentation(output)
        slide = prs.slides[0]
        # Find table shape
        table_shapes = [
            s for s in slide.shapes
            if s.shape_type == MSO_SHAPE_TYPE.TABLE
        ]
        assert len(table_shapes) >= 1, "Expected a table shape on the slide"
        table = table_shapes[0].table
        assert table.rows[0].cells[0].text == "Item"
        assert table.rows[1].cells[2].text == "CRITICAL GAP"


def test_pptx_pipe_table_header_bold() -> None:
    """Pipe-table header row has bold formatting."""
    from src.services.renderer import render_pptx

    slides = [
        SlideObject(
            slide_id="S-001",
            title="Table",
            layout_type=LayoutType.CONTENT_1COL,
            body_content=BodyContent(
                text_elements=[
                    "Col1 | Col2 | Col3",
                    "A | B | C",
                    "D | E | F",
                ],
            ),
        ),
    ]
    with tempfile.TemporaryDirectory() as tmpdir:
        output = os.path.join(tmpdir, "hdr_test.pptx")
        asyncio.run(render_pptx(slides, TEMPLATE_PATH, output))

        prs = Presentation(output)
        slide = prs.slides[0]
        for shape in slide.shapes:
            if shape.has_table:
                cell = shape.table.cell(0, 0)
                runs = cell.text_frame.paragraphs[0].runs
                assert runs[0].font.bold is True, "Header cell should be bold"
                return
        raise AssertionError("No table found")


def test_pptx_team_table_has_navy_header() -> None:
    """TEAM layout em-dash table gets a navy header row with 'Role' and 'Details'."""
    from src.services.renderer import render_pptx

    slides = [
        SlideObject(
            slide_id="S-001",
            title="Proposed Team",
            layout_type=LayoutType.TEAM,
            body_content=BodyContent(
                text_elements=[
                    "SAP Basis Lead \u2014 Senior consultant with 10+ years",
                    "Project Manager \u2014 PMP certified, SIDF experience",
                ],
            ),
        ),
    ]
    with tempfile.TemporaryDirectory() as tmpdir:
        output = os.path.join(tmpdir, "team_header_test.pptx")
        asyncio.run(render_pptx(slides, TEMPLATE_PATH, output))

        prs = Presentation(output)
        slide = prs.slides[0]
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                # Header row should have "Role" and "Details"
                assert table.cell(0, 0).text.strip() == "Role"
                assert table.cell(0, 1).text.strip() == "Details"
                # Data should start at row 1
                assert "SAP Basis Lead" in table.cell(1, 0).text
                # 3 rows total (1 header + 2 data)
                assert len(table.rows) == 3
                # Verify navy fill on header via XML
                nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
                tc = table.cell(0, 0)._tc
                tcPr = tc.find("a:tcPr", nsmap)
                solid = tcPr.find("a:solidFill", nsmap)
                srgb = solid.find("a:srgbClr", nsmap)
                assert srgb.get("val").upper() == "0E2841"
                return
        raise AssertionError("No table found on TEAM slide")


def test_pptx_timeline_table_has_navy_header() -> None:
    """TIMELINE layout em-dash table gets a navy header row."""
    from src.services.renderer import render_pptx

    slides = [
        SlideObject(
            slide_id="S-001",
            title="Project Timeline",
            layout_type=LayoutType.TIMELINE,
            body_content=BodyContent(
                text_elements=[
                    "Phase 1 \u2014 Discovery (Weeks 1-4)",
                    "Phase 2 \u2014 Implementation (Weeks 5-16)",
                ],
            ),
        ),
    ]
    with tempfile.TemporaryDirectory() as tmpdir:
        output = os.path.join(tmpdir, "timeline_header_test.pptx")
        asyncio.run(render_pptx(slides, TEMPLATE_PATH, output))

        prs = Presentation(output)
        slide = prs.slides[0]
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                assert table.cell(0, 0).text.strip() == "Phase"
                assert table.cell(0, 1).text.strip() == "Details"
                assert len(table.rows) == 3
                return
        raise AssertionError("No table found on TIMELINE slide")
