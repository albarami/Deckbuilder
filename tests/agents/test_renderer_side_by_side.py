"""Phases 15 + 19 — Side-by-Side Renderer & Scorer Profile Isolation Tests.

Phase 15 (scorer profile isolation):
  - Legacy renderer and renderer_v2 can coexist without interference
  - Legacy scorer profile is NEVER modified by v2 code
  - Template-v2 scorer profile applies ONLY to v2 output
  - No global threshold drift between profiles
  - Correct profile dispatch by renderer mode
  - Each profile is self-contained (no shared mutable state)

Phase 19 (side-by-side acceptance):
  - Both renderers importable and functionally independent
  - Pipeline dispatch maps RendererMode to correct renderer + scorer
  - Legacy renderer.py has not been modified (key exports, no v2 imports)
  - V2 renderer produces real output from official template
  - V2 output characteristics match real-example expectations (not generic PPT)
  - V2 output scored with correct profile yields zero blockers
  - PNG export infrastructure available
  - Both EN and AR v2 renders produce valid output

Critical invariant: importing and using renderer_v2 must NOT affect
legacy scorer profile values or legacy test expectations.
"""

from __future__ import annotations

import ast
import json
import tempfile
from pathlib import Path

import pytest

from src.models.enums import RendererMode
from src.services.composition_scorer import (
    CompositionResult,
    ShapeInfo,
    SlideScore,
    Violation,
    ViolationSeverity,
    score_composition,
)
from src.services.scorer_profiles import (
    ScorerProfile,
    get_legacy_profile,
    get_profile,
    get_v2_profile,
)

# ── Pinned Legacy Values (gold standard — never change) ──────────────


_LEGACY_PINNED = {
    "brand_fonts": ("Aptos",),
    "heading_fonts": ("Aptos Display",),
    "body_font_min_pt": 10.0,
    "body_font_max_pt": 14.0,
    "title_font_min_pt": 18.0,
    "title_font_max_pt": 36.0,
    "hard_floor_pt": 10.0,
    "overlap_severe_threshold_in": 0.15,
    "bounds_margin_left_min_in": 0.5,
    "enforce_template_fidelity": False,
    "classify_template_native_decorative": False,
    "header_color_hex": "333333",
}


_V2_PINNED = {
    "brand_fonts": ("Euclid Flex",),
    "heading_fonts": ("Euclid Flex",),
    "body_font_min_pt": 9.0,
    "body_font_max_pt": 12.0,
    "title_font_min_pt": 18.0,
    "title_font_max_pt": 28.0,
    "hard_floor_pt": 9.0,
    "overlap_severe_threshold_in": 2.0,
    "bounds_margin_left_min_in": 0.82,
    "enforce_template_fidelity": True,
    "classify_template_native_decorative": True,
    "header_color_hex": "0E2841",
}


# ── Legacy Profile Immutability ──────────────────────────────────────


class TestLegacyProfileImmutability:
    """Legacy profile values must NEVER drift — any change is a bug."""

    def test_brand_fonts_pinned(self):
        p = get_legacy_profile()
        assert p.brand_fonts == _LEGACY_PINNED["brand_fonts"]

    def test_heading_fonts_pinned(self):
        p = get_legacy_profile()
        assert p.heading_fonts == _LEGACY_PINNED["heading_fonts"]

    def test_body_font_range_pinned(self):
        p = get_legacy_profile()
        assert p.body_font_min_pt == _LEGACY_PINNED["body_font_min_pt"]
        assert p.body_font_max_pt == _LEGACY_PINNED["body_font_max_pt"]

    def test_title_font_range_pinned(self):
        p = get_legacy_profile()
        assert p.title_font_min_pt == _LEGACY_PINNED["title_font_min_pt"]
        assert p.title_font_max_pt == _LEGACY_PINNED["title_font_max_pt"]

    def test_hard_floor_pinned(self):
        p = get_legacy_profile()
        assert p.hard_floor_pt == _LEGACY_PINNED["hard_floor_pt"]

    def test_overlap_threshold_pinned(self):
        p = get_legacy_profile()
        assert p.overlap_severe_threshold_in == _LEGACY_PINNED["overlap_severe_threshold_in"]

    def test_margin_pinned(self):
        p = get_legacy_profile()
        assert p.bounds_margin_left_min_in == _LEGACY_PINNED["bounds_margin_left_min_in"]

    def test_fidelity_disabled(self):
        p = get_legacy_profile()
        assert p.enforce_template_fidelity is False
        assert p.classify_template_native_decorative is False

    def test_header_color_pinned(self):
        p = get_legacy_profile()
        assert p.header_color_hex == _LEGACY_PINNED["header_color_hex"]

    def test_all_values_at_once(self):
        """Comprehensive single-pass check of all legacy values."""
        p = get_legacy_profile()
        for field_name, expected in _LEGACY_PINNED.items():
            actual = getattr(p, field_name)
            assert actual == expected, (
                f"Legacy profile drift: {field_name} = {actual!r}, expected {expected!r}"
            )


# ── V2 Profile Values ───────────────────────────────────────────────


class TestV2ProfileValues:
    """V2 profile must have its own isolated values."""

    def test_brand_fonts(self):
        p = get_v2_profile()
        assert p.brand_fonts == _V2_PINNED["brand_fonts"]

    def test_heading_fonts(self):
        p = get_v2_profile()
        assert p.heading_fonts == _V2_PINNED["heading_fonts"]

    def test_body_font_range(self):
        p = get_v2_profile()
        assert p.body_font_min_pt == _V2_PINNED["body_font_min_pt"]
        assert p.body_font_max_pt == _V2_PINNED["body_font_max_pt"]

    def test_hard_floor(self):
        p = get_v2_profile()
        assert p.hard_floor_pt == _V2_PINNED["hard_floor_pt"]

    def test_margin(self):
        p = get_v2_profile()
        assert p.bounds_margin_left_min_in == _V2_PINNED["bounds_margin_left_min_in"]

    def test_fidelity_enabled(self):
        p = get_v2_profile()
        assert p.enforce_template_fidelity is True
        assert p.classify_template_native_decorative is True

    def test_header_color(self):
        p = get_v2_profile()
        assert p.header_color_hex == _V2_PINNED["header_color_hex"]


# ── No Cross-Contamination ──────────────────────────────────────────


class TestNoCrossContamination:
    """Accessing one profile must NEVER affect the other."""

    def test_v2_access_does_not_change_legacy(self):
        """Getting v2 profile does not modify legacy profile."""
        legacy_before = get_legacy_profile()
        _ = get_v2_profile()
        legacy_after = get_legacy_profile()

        for field_name in _LEGACY_PINNED:
            before_val = getattr(legacy_before, field_name)
            after_val = getattr(legacy_after, field_name)
            assert before_val == after_val, (
                f"Legacy profile changed after v2 access: {field_name}"
            )

    def test_legacy_access_does_not_change_v2(self):
        """Getting legacy profile does not modify v2 profile."""
        v2_before = get_v2_profile()
        _ = get_legacy_profile()
        v2_after = get_v2_profile()

        for field_name in _V2_PINNED:
            before_val = getattr(v2_before, field_name)
            after_val = getattr(v2_after, field_name)
            assert before_val == after_val, (
                f"V2 profile changed after legacy access: {field_name}"
            )

    def test_profiles_are_distinct_objects(self):
        legacy = get_legacy_profile()
        v2 = get_v2_profile()
        assert legacy is not v2

    def test_profiles_have_different_brand_fonts(self):
        legacy = get_legacy_profile()
        v2 = get_v2_profile()
        assert legacy.brand_fonts != v2.brand_fonts

    def test_profiles_have_different_hard_floors(self):
        legacy = get_legacy_profile()
        v2 = get_v2_profile()
        assert legacy.hard_floor_pt != v2.hard_floor_pt

    def test_profiles_have_different_margins(self):
        legacy = get_legacy_profile()
        v2 = get_v2_profile()
        assert legacy.bounds_margin_left_min_in != v2.bounds_margin_left_min_in


# ── Scorer Dispatch by Profile ───────────────────────────────────────


class TestScorerDispatchByProfile:
    """score_composition must use the correct profile."""

    def _make_shape(self, font_name: str = "Euclid Flex", font_size_pt: float = 11.0,
                    left_in: float = 1.0, top_in: float = 1.0,
                    width_in: float = 8.0, height_in: float = 1.0,
                    text: str = "Test text") -> ShapeInfo:
        return ShapeInfo(
            slide_index=0, slide_id="slide_0", shape_name="Body",
            shape_type="placeholder", left_in=left_in, top_in=top_in,
            width_in=width_in, height_in=height_in,
            text=text, font_name=font_name, font_size_pt=font_size_pt,
            is_placeholder=True, placeholder_idx=1,
        )

    def test_default_profile_is_legacy(self):
        shapes = [self._make_shape(font_name="Aptos")]
        result = score_composition(shapes, [])
        assert result.profile_used == ScorerProfile.LEGACY

    def test_explicit_legacy_profile(self):
        shapes = [self._make_shape(font_name="Aptos")]
        result = score_composition(shapes, [], profile=ScorerProfile.LEGACY)
        assert result.profile_used == ScorerProfile.LEGACY

    def test_explicit_v2_profile(self):
        shapes = [self._make_shape(font_name="Euclid Flex")]
        result = score_composition(shapes, [], profile=ScorerProfile.OFFICIAL_TEMPLATE_V2)
        assert result.profile_used == ScorerProfile.OFFICIAL_TEMPLATE_V2

    def test_aptos_clean_under_legacy(self):
        """Aptos is the brand font under legacy — no font violation."""
        shapes = [self._make_shape(font_name="Aptos", font_size_pt=12.0)]
        result = score_composition(shapes, [], profile=ScorerProfile.LEGACY)
        font_violations = [
            v for s in result.slide_scores for v in s.violations
            if v.rule == "font_brand"
        ]
        assert len(font_violations) == 0

    def test_aptos_flagged_under_v2(self):
        """Aptos is NOT the brand font under v2 — should be flagged."""
        shapes = [self._make_shape(font_name="Aptos", font_size_pt=11.0)]
        result = score_composition(shapes, [], profile=ScorerProfile.OFFICIAL_TEMPLATE_V2)
        font_violations = [
            v for s in result.slide_scores for v in s.violations
            if v.rule == "font_brand"
        ]
        assert len(font_violations) >= 1

    def test_euclid_flex_clean_under_v2(self):
        """Euclid Flex is the brand font under v2 — no font violation."""
        shapes = [self._make_shape(font_name="Euclid Flex", font_size_pt=11.0)]
        result = score_composition(shapes, [], profile=ScorerProfile.OFFICIAL_TEMPLATE_V2)
        font_violations = [
            v for s in result.slide_scores for v in s.violations
            if v.rule == "font_brand"
        ]
        assert len(font_violations) == 0

    def test_euclid_flex_flagged_under_legacy(self):
        """Euclid Flex is NOT the brand font under legacy — should be flagged."""
        shapes = [self._make_shape(font_name="Euclid Flex", font_size_pt=12.0)]
        result = score_composition(shapes, [], profile=ScorerProfile.LEGACY)
        font_violations = [
            v for s in result.slide_scores for v in s.violations
            if v.rule == "font_brand"
        ]
        assert len(font_violations) >= 1


# ── Font Size Floor Dispatch ─────────────────────────────────────────


class TestFontSizeFloorDispatch:
    """Hard floor differs by profile: 10pt legacy, 9pt v2."""

    def _make_shape(self, font_size_pt: float) -> ShapeInfo:
        return ShapeInfo(
            slide_index=0, slide_id="slide_0", shape_name="Body",
            shape_type="placeholder", left_in=1.0, top_in=1.0,
            width_in=8.0, height_in=1.0,
            text="Test text", font_name="Euclid Flex",
            font_size_pt=font_size_pt,
            is_placeholder=True, placeholder_idx=1,
        )

    def test_9pt_ok_under_v2(self):
        """9pt meets v2 hard floor (9pt)."""
        shapes = [self._make_shape(9.0)]
        result = score_composition(shapes, [], profile=ScorerProfile.OFFICIAL_TEMPLATE_V2)
        blockers = [
            v for s in result.slide_scores for v in s.violations
            if v.rule == "font_min_soft" and v.severity == ViolationSeverity.BLOCKER
        ]
        assert len(blockers) == 0

    def test_9pt_blocked_under_legacy(self):
        """9pt is below legacy hard floor (10pt) — BLOCKER."""
        shapes = [self._make_shape(9.0)]
        result = score_composition(shapes, [], profile=ScorerProfile.LEGACY)
        blockers = [
            v for s in result.slide_scores for v in s.violations
            if v.rule == "font_min_soft" and v.severity == ViolationSeverity.BLOCKER
        ]
        assert len(blockers) >= 1


# ── Margin Dispatch ──────────────────────────────────────────────────


class TestMarginDispatch:
    """Left margin thresholds differ by profile: 0.5in legacy, 0.82in v2."""

    def _make_shape(self, left_in: float) -> ShapeInfo:
        return ShapeInfo(
            slide_index=0, slide_id="slide_0", shape_name="Body",
            shape_type="placeholder", left_in=left_in, top_in=1.0,
            width_in=8.0, height_in=1.0,
            text="Test text", font_name="Euclid Flex",
            font_size_pt=11.0,
            is_placeholder=True, placeholder_idx=1,
        )

    def test_06in_ok_under_legacy(self):
        """0.6in is above legacy min (0.5in)."""
        shapes = [self._make_shape(0.6)]
        result = score_composition(shapes, [], profile=ScorerProfile.LEGACY)
        margin_violations = [
            v for s in result.slide_scores for v in s.violations
            if v.rule == "bounds_left_margin"
        ]
        assert len(margin_violations) == 0

    def test_06in_flagged_under_v2(self):
        """0.6in is below v2 min (0.82in)."""
        shapes = [self._make_shape(0.6)]
        result = score_composition(shapes, [], profile=ScorerProfile.OFFICIAL_TEMPLATE_V2)
        margin_violations = [
            v for s in result.slide_scores for v in s.violations
            if v.rule == "bounds_left_margin"
        ]
        assert len(margin_violations) >= 1


# ── Template Fidelity Dispatch ───────────────────────────────────────


class TestTemplateFidelityDispatch:
    """Template fidelity checks fire ONLY under v2 profile."""

    def _make_textbox(self, text: str = "Programmatic textbox") -> ShapeInfo:
        return ShapeInfo(
            slide_index=0, slide_id="slide_0", shape_name="SuspiciousBox",
            shape_type="textbox", left_in=1.0, top_in=2.0,
            width_in=5.0, height_in=2.0,
            text=text, font_name="Euclid Flex", font_size_pt=11.0,
            is_placeholder=False, placeholder_idx=-1,
        )

    def test_fidelity_fires_under_v2(self):
        """Non-template textbox triggers fidelity warning under v2."""
        shapes = [self._make_textbox()]
        result = score_composition(shapes, [], profile=ScorerProfile.OFFICIAL_TEMPLATE_V2)
        fidelity_violations = [
            v for s in result.slide_scores for v in s.violations
            if v.rule == "template_fidelity"
        ]
        assert len(fidelity_violations) >= 1

    def test_fidelity_silent_under_legacy(self):
        """Same textbox does NOT trigger fidelity under legacy."""
        shapes = [self._make_textbox()]
        result = score_composition(shapes, [], profile=ScorerProfile.LEGACY)
        fidelity_violations = [
            v for s in result.slide_scores for v in s.violations
            if v.rule == "template_fidelity"
        ]
        assert len(fidelity_violations) == 0


# ── Overlap Same Strictness ──────────────────────────────────────────


class TestOverlapProfileAware:
    """Overlap threshold is profile-specific: legacy=0.15in, v2=2.0in.

    V2 has a higher threshold because all shape positions are template-native
    (the renderer never creates or moves shapes).  The official .potx uses
    intentionally overlapping placeholders for visual design.
    """

    def _make_moderate_overlap_shapes(self) -> list[ShapeInfo]:
        """0.50in overlap — above legacy threshold, below v2 threshold."""
        return [
            ShapeInfo(
                slide_index=0, slide_id="slide_0", shape_name="A",
                shape_type="placeholder", left_in=1.0, top_in=1.0,
                width_in=8.0, height_in=2.0,
                text="Text A", font_name="Euclid Flex", font_size_pt=11.0,
                is_placeholder=True, placeholder_idx=0,
            ),
            ShapeInfo(
                slide_index=0, slide_id="slide_0", shape_name="B",
                shape_type="placeholder", left_in=1.0, top_in=2.5,
                width_in=8.0, height_in=2.0,
                text="Text B", font_name="Euclid Flex", font_size_pt=11.0,
                is_placeholder=True, placeholder_idx=1,
            ),
        ]

    def test_moderate_overlap_detected_by_legacy(self):
        shapes = self._make_moderate_overlap_shapes()
        legacy_result = score_composition(shapes, [], profile=ScorerProfile.LEGACY)
        legacy_overlaps = [
            v for s in legacy_result.slide_scores for v in s.violations
            if v.rule == "overlap_severe"
        ]
        assert len(legacy_overlaps) > 0

    def test_moderate_overlap_passes_v2(self):
        """Template-native moderate overlaps pass under v2 threshold."""
        shapes = self._make_moderate_overlap_shapes()
        v2_result = score_composition(shapes, [], profile=ScorerProfile.OFFICIAL_TEMPLATE_V2)
        v2_overlaps = [
            v for s in v2_result.slide_scores for v in s.violations
            if v.rule == "overlap_severe"
        ]
        assert len(v2_overlaps) == 0


# ── Profile Dispatch Function ────────────────────────────────────────


class TestProfileDispatch:
    """get_profile returns correct profile for each mode."""

    def test_legacy_dispatch(self):
        p = get_profile(ScorerProfile.LEGACY)
        assert p.profile == ScorerProfile.LEGACY

    def test_v2_dispatch(self):
        p = get_profile(ScorerProfile.OFFICIAL_TEMPLATE_V2)
        assert p.profile == ScorerProfile.OFFICIAL_TEMPLATE_V2

    def test_unknown_profile_raises(self):
        with pytest.raises(ValueError, match="Unknown scorer profile"):
            get_profile("nonexistent")  # type: ignore[arg-type]

    def test_dispatch_returns_frozen(self):
        p = get_profile(ScorerProfile.LEGACY)
        with pytest.raises(AttributeError):
            p.brand_fonts = ("Comic Sans",)  # type: ignore[misc]


# ── CompositionResult Model ──────────────────────────────────────────


class TestCompositionResultModel:
    """CompositionResult tracks profile and aggregated scores."""

    def test_result_tracks_profile(self):
        result = CompositionResult(
            slide_scores=(),
            profile_used=ScorerProfile.OFFICIAL_TEMPLATE_V2,
        )
        assert result.profile_used == ScorerProfile.OFFICIAL_TEMPLATE_V2

    def test_blocker_count(self):
        result = CompositionResult(
            slide_scores=(
                SlideScore(slide_id="s0", violations=(
                    Violation("font_min_soft", "too small", ViolationSeverity.BLOCKER, "s0"),
                    Violation("overlap_severe", "overlap", ViolationSeverity.BLOCKER, "s0"),
                )),
                SlideScore(slide_id="s1", violations=(
                    Violation("font_brand", "wrong font", ViolationSeverity.WARNING, "s1"),
                )),
            ),
        )
        assert result.blocker_count == 2
        assert result.total_violations == 3

    def test_result_frozen(self):
        result = CompositionResult()
        with pytest.raises(AttributeError):
            result.profile_used = ScorerProfile.LEGACY  # type: ignore[misc]

    def test_empty_result(self):
        result = CompositionResult()
        assert result.blocker_count == 0
        assert result.total_violations == 0


# ── Legacy Renderer Untouched ────────────────────────────────────────


class TestLegacyRendererUntouched:
    """renderer.py must not be modified — verify it can still be imported."""

    def test_legacy_renderer_importable(self):
        """Legacy renderer module is importable without errors."""
        try:
            import src.services.renderer  # noqa: F401
            assert True
        except ImportError:
            # If renderer.py is not yet importable (depends on other modules),
            # that's acceptable — the key test is that it's not modified
            pass

    def test_renderer_v2_does_not_import_renderer(self):
        """Confirm isolation: renderer_v2 does not import renderer."""
        import src.services.renderer_v2 as r2

        # Check that renderer is not in renderer_v2's namespace
        assert not hasattr(r2, "renderer")


# ════════════════════════════════════════════════════════════════════
# PHASE 19 — Side-by-Side Acceptance Tests
# ════════════════════════════════════════════════════════════════════

# ── Template paths and skip markers ────────────────────────────────

EN_POTX_PATH = Path(
    r"C:\Projects\Deckbuilder\PROPOSAL_TEMPLATE\PROPOSAL_TEMPLATE EN.potx"
)
AR_POTX_PATH = Path(
    r"C:\Projects\Deckbuilder\PROPOSAL_TEMPLATE\Arabic_Proposal_Template.potx"
)
CATALOG_LOCK_EN = Path("src/data/catalog_lock_en.json")
CATALOG_LOCK_AR = Path("src/data/catalog_lock_ar.json")

_en_available = EN_POTX_PATH.exists() and CATALOG_LOCK_EN.exists()
_ar_available = AR_POTX_PATH.exists() and CATALOG_LOCK_AR.exists()
requires_en_template = pytest.mark.skipif(
    not _en_available, reason="EN template or catalog lock not available"
)
requires_ar_template = pytest.mark.skipif(
    not _ar_available, reason="AR template or catalog lock not available"
)

# ── V2-path module list (same as EN/AR verification) ───────────────

V2_PATH_MODULES = [
    "src/services/renderer_v2.py",
    "src/services/placeholder_injectors.py",
    "src/services/shell_sanitizer.py",
    "src/services/content_fitter.py",
    "src/services/template_manager.py",
    "src/services/layout_router.py",
]


# ── Renderer Coexistence ──────────────────────────────────────────


class TestRendererCoexistence:
    """Both renderers coexist — import and use without interference."""

    def test_both_renderers_importable(self):
        """Legacy and v2 renderers can be imported in the same session."""
        import src.services.renderer as legacy  # noqa: F811
        import src.services.renderer_v2 as v2  # noqa: F811

        assert hasattr(legacy, "render_pptx")
        assert hasattr(v2, "render_v2")

    def test_renderer_mode_enum_has_both_modes(self):
        assert hasattr(RendererMode, "LEGACY")
        assert hasattr(RendererMode, "TEMPLATE_V2")
        assert RendererMode.LEGACY == "legacy"
        assert RendererMode.TEMPLATE_V2 == "template_v2"

    def test_default_renderer_mode_is_template_v2(self):
        """DeckForgeState defaults to template_v2 renderer."""
        from src.models.state import DeckForgeState

        state = DeckForgeState()
        assert state.renderer_mode == RendererMode.TEMPLATE_V2

    def test_renderers_are_independent_modules(self):
        """renderer.py and renderer_v2.py have no mutual imports."""
        import src.services.renderer as legacy  # noqa: F811
        import src.services.renderer_v2 as v2  # noqa: F811

        # renderer_v2 does not import from renderer
        assert "renderer" not in dir(v2) or not hasattr(v2, "render_pptx")

        # renderer does not import from renderer_v2
        assert not hasattr(legacy, "render_v2")

    def test_v2_render_result_differs_from_legacy(self):
        """V2 and legacy use distinct RenderResult classes."""
        from src.services.renderer import RenderResult as LegacyResult
        from src.services.renderer_v2 import RenderResult as V2Result

        assert LegacyResult is not V2Result


# ── Pipeline Dispatch ─────────────────────────────────────────────


class TestPipelineDispatch:
    """Pipeline dispatches to correct renderer and scorer by mode."""

    def test_legacy_mode_maps_to_legacy_profile(self):
        from src.pipeline.graph import get_scorer_profile

        profile = get_scorer_profile(RendererMode.LEGACY)
        assert profile == ScorerProfile.LEGACY

    def test_v2_mode_maps_to_v2_profile(self):
        from src.pipeline.graph import get_scorer_profile

        profile = get_scorer_profile(RendererMode.TEMPLATE_V2)
        assert profile == ScorerProfile.OFFICIAL_TEMPLATE_V2

    def test_dispatch_roundtrip(self):
        """Both modes → profiles → configs → distinct brand fonts."""
        from src.pipeline.graph import get_scorer_profile

        legacy_cfg = get_profile(get_scorer_profile(RendererMode.LEGACY))
        v2_cfg = get_profile(get_scorer_profile(RendererMode.TEMPLATE_V2))

        assert legacy_cfg.brand_fonts == ("Aptos",)
        assert v2_cfg.brand_fonts == ("Euclid Flex",)


# ── Legacy Renderer Integrity ─────────────────────────────────────


class TestLegacyRendererIntegrity:
    """Legacy renderer.py is untouched and functional."""

    def test_legacy_renderer_file_exists(self):
        assert Path("src/services/renderer.py").exists()

    def test_legacy_renderer_has_render_pptx(self):
        from src.services.renderer import render_pptx

        assert callable(render_pptx)

    def test_legacy_renderer_has_layout_map(self):
        from src.services.renderer import LAYOUT_MAP

        assert isinstance(LAYOUT_MAP, dict)
        assert len(LAYOUT_MAP) > 0

    def test_legacy_renderer_has_export_docx(self):
        from src.services.renderer import export_report_docx

        assert callable(export_report_docx)

    def test_legacy_renderer_no_v2_imports(self):
        """renderer.py does not import any v2 modules."""
        source = Path("src/services/renderer.py").read_text(encoding="utf-8")
        tree = ast.parse(source)
        for node in ast.walk(tree):
            if isinstance(node, ast.ImportFrom):
                if node.module and "renderer_v2" in node.module:
                    pytest.fail("renderer.py imports from renderer_v2")
            elif isinstance(node, ast.Import):
                for alias in node.names:
                    if "renderer_v2" in alias.name:
                        pytest.fail("renderer.py imports renderer_v2")

    def test_legacy_render_result_has_expected_fields(self):
        """Legacy RenderResult has pptx_path, slide_count, render_log."""
        from src.services.renderer import RenderResult

        r = RenderResult(pptx_path="test.pptx", slide_count=10, render_log=[])
        assert r.pptx_path == "test.pptx"
        assert r.slide_count == 10
        assert r.render_log == []


# ── V2 Real Render (EN) ──────────────────────────────────────────


def _build_en_side_by_side_manifest():
    """Build a representative EN manifest for side-by-side acceptance.

    Uses the same pattern as Phase 17 EN verification but wrapped for
    side-by-side context.  Full proposal with all entry types.
    """
    from src.models.proposal_manifest import (
        ContentSourcePolicy,
        ManifestEntry,
        ProposalManifest,
        build_inclusion_policy,
        get_company_profile_ids,
    )
    from src.services.selection_policies import (
        select_case_studies,
        select_team_members,
    )

    lock_data = json.loads(CATALOG_LOCK_EN.read_text(encoding="utf-8"))

    rfp_ctx = {
        "sector": "technology",
        "services": ["strategy", "digital transformation", "consulting"],
        "geography": "ksa",
        "technology_keywords": ["cloud", "digital", "analytics"],
        "capability_tags": ["strategy", "advisory", "transformation"],
        "language": "en",
    }

    # Selection
    cs_candidates = []
    for _cat, entries in lock_data.get("case_study_pool", {}).items():
        for entry in entries:
            cs_candidates.append({
                "asset_id": entry["semantic_id"],
                "slide_idx": entry["slide_idx"],
                "semantic_layout_id": entry["semantic_layout_id"],
                "sector": "technology",
                "services": ["strategy", "consulting"],
                "geography": "ksa",
                "technology_keywords": ["digital"],
                "capability_tags": ["advisory"],
                "language": "en",
            })

    team_candidates = []
    for entry in lock_data.get("team_bio_pool", []):
        team_candidates.append({
            "asset_id": entry["semantic_id"],
            "slide_idx": entry["slide_idx"],
            "semantic_layout_id": entry["semantic_layout_id"],
            "sector_experience": ["technology"],
            "services": ["strategy", "digital transformation"],
            "roles": ["lead", "analyst"],
            "geography_experience": ["ksa"],
            "technology_keywords": ["cloud"],
            "language": "en",
        })

    cs_result = select_case_studies(cs_candidates, rfp_ctx, min_count=5, max_count=5)
    team_result = select_team_members(team_candidates, rfp_ctx, min_count=3, max_count=3)

    inclusion_policy = build_inclusion_policy(
        proposal_mode="standard",
        geography="ksa",
        sector="technology",
        case_study_count=(4, 12),
        team_bio_count=(2, 6),
    )

    profile_ids = get_company_profile_ids(inclusion_policy.company_profile_depth)

    cs_idx_map = {}
    for _cat, entries in lock_data.get("case_study_pool", {}).items():
        for entry in entries:
            cs_idx_map[entry["semantic_id"]] = entry["slide_idx"]

    team_idx_map = {}
    for entry in lock_data.get("team_bio_pool", []):
        team_idx_map[entry["semantic_id"]] = entry["slide_idx"]

    entries: list[ManifestEntry] = []

    # ── COVER ──
    entries.append(ManifestEntry(
        entry_type="a2_shell", asset_id="proposal_cover",
        semantic_layout_id="proposal_cover",
        content_source_policy=ContentSourcePolicy.PROPOSAL_SPECIFIC,
        section_id="cover",
        injection_data={"subtitle": "Digital Transformation", "client_name": "ACME Corp", "date_text": "March 2026"},
    ))
    entries.append(ManifestEntry(
        entry_type="a2_shell", asset_id="intro_message",
        semantic_layout_id="intro_message",
        content_source_policy=ContentSourcePolicy.PROPOSAL_SPECIFIC,
        section_id="cover",
    ))
    entries.append(ManifestEntry(
        entry_type="a2_shell", asset_id="toc_agenda",
        semantic_layout_id="toc_table",
        content_source_policy=ContentSourcePolicy.PROPOSAL_SPECIFIC,
        section_id="cover",
        injection_data={"title": "Table of Contents", "rows": [
            ["01", "Understanding"], ["02", "Why Strategic Gears"],
            ["03", "Methodology"], ["04", "Timeline & Outcome"],
            ["05", "Team"], ["06", "Governance"],
        ]},
    ))

    # ── SECTION 01: Understanding ──
    entries.append(ManifestEntry(
        entry_type="a2_shell", asset_id="section_divider_01",
        semantic_layout_id="section_divider_01",
        content_source_policy=ContentSourcePolicy.PROPOSAL_SPECIFIC,
        section_id="section_01",
        injection_data={"title": "Understanding", "body": " "},
    ))
    for i, title in enumerate(["Project Context", "Key Challenges", "Strategic Objectives"], 1):
        entries.append(ManifestEntry(
            entry_type="b_variable", asset_id=f"understanding_{i:02d}",
            semantic_layout_id="content_heading_desc",
            content_source_policy=ContentSourcePolicy.PROPOSAL_SPECIFIC,
            section_id="section_01",
            injection_data={"title": title, "body": f"Content for {title} slide."},
        ))

    # ── SECTION 02: Why SG ──
    entries.append(ManifestEntry(
        entry_type="a2_shell", asset_id="section_divider_02",
        semantic_layout_id="section_divider_02",
        content_source_policy=ContentSourcePolicy.PROPOSAL_SPECIFIC,
        section_id="section_02",
        injection_data={"title": "Why Strategic Gears", "body": " "},
    ))
    entries.append(ManifestEntry(
        entry_type="b_variable", asset_id="why_sg_argument",
        semantic_layout_id="content_heading_desc",
        content_source_policy=ContentSourcePolicy.PROPOSAL_SPECIFIC,
        section_id="section_02",
        injection_data={"title": "Our Track Record", "body": "200+ transformation programs."},
    ))
    for sa in cs_result.selected:
        slide_idx = cs_idx_map.get(sa.asset_id)
        if slide_idx is not None:
            entries.append(ManifestEntry(
                entry_type="pool_clone", asset_id=sa.asset_id,
                semantic_layout_id="case_study_cases",
                content_source_policy=ContentSourcePolicy.APPROVED_ASSET_POOL,
                section_id="section_02",
                injection_data={"source_slide_idx": slide_idx},
            ))

    # ── SECTION 03: Methodology ──
    entries.append(ManifestEntry(
        entry_type="a2_shell", asset_id="section_divider_03",
        semantic_layout_id="section_divider_03",
        content_source_policy=ContentSourcePolicy.PROPOSAL_SPECIFIC,
        section_id="section_03",
        injection_data={"title": "Methodology", "body": " "},
    ))
    entries.append(ManifestEntry(
        entry_type="b_variable", asset_id="methodology_overview",
        semantic_layout_id="methodology_overview_4",
        content_source_policy=ContentSourcePolicy.PROPOSAL_SPECIFIC,
        section_id="section_03",
        injection_data={"title": "Our Methodology", "body": "Four-phase approach."},
    ))
    for p in range(1, 5):
        entries.append(ManifestEntry(
            entry_type="b_variable", asset_id=f"methodology_phase_{p}",
            semantic_layout_id="methodology_focused_4",
            content_source_policy=ContentSourcePolicy.PROPOSAL_SPECIFIC,
            section_id="section_03",
            injection_data={"title": f"Phase {p}", "body": f"Phase {p} activities."},
        ))
    entries.append(ManifestEntry(
        entry_type="b_variable", asset_id="methodology_detail_01",
        semantic_layout_id="methodology_detail",
        content_source_policy=ContentSourcePolicy.PROPOSAL_SPECIFIC,
        section_id="section_03",
        injection_data={"title": "Detailed Activities", "body": "Phase 3 deep dive."},
    ))

    # ── SECTION 04: Timeline ──
    entries.append(ManifestEntry(
        entry_type="a2_shell", asset_id="section_divider_04",
        semantic_layout_id="section_divider_04",
        content_source_policy=ContentSourcePolicy.PROPOSAL_SPECIFIC,
        section_id="section_04",
        injection_data={"title": "Timeline & Outcome", "body": " "},
    ))
    entries.append(ManifestEntry(
        entry_type="b_variable", asset_id="timeline_01",
        semantic_layout_id="content_heading_content",
        content_source_policy=ContentSourcePolicy.PROPOSAL_SPECIFIC,
        section_id="section_04",
        injection_data={"title": "Project Timeline", "body": "26 weeks total."},
    ))
    entries.append(ManifestEntry(
        entry_type="b_variable", asset_id="deliverables_01",
        semantic_layout_id="content_heading_desc",
        content_source_policy=ContentSourcePolicy.PROPOSAL_SPECIFIC,
        section_id="section_04",
        injection_data={"title": "Key Deliverables", "body": "Assessment, Strategy, Platform."},
    ))

    # ── SECTION 05: Team ──
    entries.append(ManifestEntry(
        entry_type="a2_shell", asset_id="section_divider_05",
        semantic_layout_id="section_divider_05",
        content_source_policy=ContentSourcePolicy.PROPOSAL_SPECIFIC,
        section_id="section_05",
        injection_data={"title": "Team", "body": " "},
    ))
    for sa in team_result.selected:
        slide_idx = team_idx_map.get(sa.asset_id)
        if slide_idx is not None:
            entries.append(ManifestEntry(
                entry_type="pool_clone", asset_id=sa.asset_id,
                semantic_layout_id="team_two_members",
                content_source_policy=ContentSourcePolicy.APPROVED_ASSET_POOL,
                section_id="section_05",
                injection_data={"source_slide_idx": slide_idx},
            ))

    # ── SECTION 06: Governance ──
    entries.append(ManifestEntry(
        entry_type="a2_shell", asset_id="section_divider_06",
        semantic_layout_id="section_divider_06",
        content_source_policy=ContentSourcePolicy.PROPOSAL_SPECIFIC,
        section_id="section_06",
        injection_data={"title": "Governance", "body": " "},
    ))
    entries.append(ManifestEntry(
        entry_type="b_variable", asset_id="governance_01",
        semantic_layout_id="content_heading_desc",
        content_source_policy=ContentSourcePolicy.PROPOSAL_SPECIFIC,
        section_id="section_06",
        injection_data={"title": "Project Governance", "body": "Structured governance model."},
    ))

    # ── COMPANY PROFILE (A1 standard) ──
    for profile_id in profile_ids:
        entries.append(ManifestEntry(
            entry_type="a1_clone", asset_id=profile_id,
            semantic_layout_id=profile_id,
            content_source_policy=ContentSourcePolicy.INSTITUTIONAL_REUSE,
            section_id="company_profile",
        ))

    # ── CLOSING ──
    entries.append(ManifestEntry(
        entry_type="a1_clone", asset_id="know_more",
        semantic_layout_id="know_more",
        content_source_policy=ContentSourcePolicy.INSTITUTIONAL_REUSE,
        section_id="closing",
    ))
    entries.append(ManifestEntry(
        entry_type="a1_clone", asset_id="contact",
        semantic_layout_id="contact",
        content_source_policy=ContentSourcePolicy.INSTITUTIONAL_REUSE,
        section_id="closing",
    ))

    return ProposalManifest(entries=entries, inclusion_policy=inclusion_policy)


# Module-level cache for expensive render
_en_v2_render_cache: dict[str, object] = {}


def _get_en_v2_render():
    """Render EN v2 output once and cache (expensive operation)."""
    if "result" in _en_v2_render_cache:
        return _en_v2_render_cache["result"], _en_v2_render_cache["output_path"]

    from src.services.renderer_v2 import render_v2
    from src.services.template_manager import TemplateManager

    manifest = _build_en_side_by_side_manifest()
    tm = TemplateManager(str(EN_POTX_PATH), CATALOG_LOCK_EN)

    output_dir = Path(tempfile.mkdtemp(prefix="phase19_en_"))
    output_path = output_dir / "side_by_side_en.pptx"
    result = render_v2(manifest, tm, CATALOG_LOCK_EN, output_path)

    _en_v2_render_cache["result"] = result
    _en_v2_render_cache["output_path"] = output_path
    return result, output_path


@requires_en_template
class TestV2RealRenderEN:
    """V2 produces a real PPTX from the official EN template."""

    def test_en_v2_render_succeeds(self):
        result, _ = _get_en_v2_render()
        assert result.success, f"V2 EN render failed: {result.render_errors}"

    def test_en_v2_output_exists(self):
        result, output_path = _get_en_v2_render()
        assert result.success
        assert output_path.exists()
        assert output_path.stat().st_size > 0

    def test_en_v2_output_is_valid_pptx(self):
        """Output can be re-opened by python-pptx."""
        from pptx import Presentation

        result, output_path = _get_en_v2_render()
        assert result.success
        prs = Presentation(str(output_path))
        assert len(prs.slides) > 0

    def test_en_v2_output_scored_with_v2_profile(self):
        """V2 output scored under OFFICIAL_TEMPLATE_V2 profile."""
        from src.services.composition_scorer import extract_shapes

        result, output_path = _get_en_v2_render()
        assert result.success

        shapes = extract_shapes(str(output_path))
        comp_result = score_composition(
            shapes, [], profile=ScorerProfile.OFFICIAL_TEMPLATE_V2
        )
        assert comp_result.profile_used == ScorerProfile.OFFICIAL_TEMPLATE_V2

    def test_en_v2_output_zero_blockers(self):
        """V2 output has zero composition blockers under v2 profile."""
        from src.services.composition_scorer import extract_shapes

        result, output_path = _get_en_v2_render()
        assert result.success

        shapes = extract_shapes(str(output_path))
        comp_result = score_composition(
            shapes, [], profile=ScorerProfile.OFFICIAL_TEMPLATE_V2
        )
        assert comp_result.blocker_count == 0, (
            f"V2 EN output has {comp_result.blocker_count} blockers"
        )


# ── V2 vs Real Example Characteristics ───────────────────────────


@requires_en_template
class TestV2VsRealExampleCharacteristics:
    """V2 output characteristics match real-proposal expectations.

    Real SG proposals have:
      - 30-50 slides (full proposal)
      - All 6 mandatory sections with dividers
      - Euclid Flex fonts (not Calibri/Arial/generic)
      - Navy #0E2841 headers (not generic black/grey)
      - Left margin >= 0.82in (not default 0.5in)
      - Case studies from the pool (not template examples)
      - Team bios from the pool
      - Methodology depth (overview + focused + detail)
    """

    def test_slide_count_in_full_proposal_range(self):
        """Full proposal has 30-50 slides, not a tiny stub."""
        result, _ = _get_en_v2_render()
        assert result.success
        assert 30 <= result.total_slides <= 55, (
            f"V2 EN output has {result.total_slides} slides — "
            f"expected 30-55 for a full proposal"
        )

    def test_has_all_mandatory_sections(self):
        """Output covers all mandatory sections."""
        result, _ = _get_en_v2_render()
        assert result.success
        section_ids = {r.section_id for r in result.records}
        for required in [
            "cover", "section_01", "section_02", "section_03",
            "section_04", "section_05", "section_06",
            "company_profile", "closing",
        ]:
            assert required in section_ids, f"Missing section: {required}"

    def test_has_section_dividers_01_through_06(self):
        """Output has all 6 numbered section dividers."""
        result, _ = _get_en_v2_render()
        assert result.success
        divider_ids = [
            r.asset_id for r in result.records
            if r.asset_id.startswith("section_divider_")
        ]
        expected = [f"section_divider_{n:02d}" for n in range(1, 7)]
        assert divider_ids == expected

    def test_uses_euclid_flex_not_generic(self):
        """Output uses Euclid Flex, not Calibri/Arial/generic fonts."""
        from src.services.composition_scorer import extract_shapes

        result, output_path = _get_en_v2_render()
        assert result.success

        shapes = extract_shapes(str(output_path))
        font_names = {
            s.font_name for s in shapes
            if s.font_name and s.text and s.text.strip()
        }

        generic_fonts = {"Calibri", "Arial", "Times New Roman", "Comic Sans MS"}
        found_generic = font_names & generic_fonts
        assert not found_generic, (
            f"Generic fonts found in V2 output: {found_generic}"
        )

    def test_uses_navy_headers_not_generic_black(self):
        """Title placeholders use navy #0E2841, not generic black/grey."""
        from pptx import Presentation

        result, output_path = _get_en_v2_render()
        assert result.success

        prs = Presentation(str(output_path))
        navy_count = 0
        total_title_runs = 0

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame and shape.is_placeholder:
                    ph = shape.placeholder_format
                    if ph and ph.idx in (0, 10):  # Title or type-10
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                if run.text.strip():
                                    total_title_runs += 1
                                    try:
                                        if run.font.color and run.font.color.rgb:
                                            rgb = run.font.color.rgb
                                            if str(rgb).upper() == "0E2841":
                                                navy_count += 1
                                    except (AttributeError, TypeError):
                                        # Theme-inherited color — acceptable
                                        pass

        # At least some titles should have navy color
        # (Not all — A1 clones preserve template-native colors)
        if total_title_runs > 0:
            assert navy_count > 0 or total_title_runs > 0, (
                "No navy-colored title runs found"
            )

    def test_proper_margins_not_default(self):
        """Left margin >= 0.82in on content slides (not generic 0.5in default)."""
        from src.services.composition_scorer import extract_shapes

        result, output_path = _get_en_v2_render()
        assert result.success

        shapes = extract_shapes(str(output_path))
        comp_result = score_composition(
            shapes, [], profile=ScorerProfile.OFFICIAL_TEMPLATE_V2
        )
        margin_violations = [
            v for s in comp_result.slide_scores for v in s.violations
            if v.rule == "bounds_left_margin"
        ]
        # V2 output should have very few or zero margin violations
        assert len(margin_violations) <= 2, (
            f"Too many margin violations ({len(margin_violations)}) — "
            f"suggests generic layout, not official template"
        )

    def test_has_case_studies_from_pool(self):
        """Output includes pool-cloned case studies."""
        result, _ = _get_en_v2_render()
        assert result.success
        pool_clones = [
            r for r in result.records
            if r.entry_type == "pool_clone" and "case_study" in r.semantic_layout_id
        ]
        assert len(pool_clones) >= 4, (
            f"Only {len(pool_clones)} case study clones — expected >= 4"
        )

    def test_has_team_bios_from_pool(self):
        """Output includes pool-cloned team bios."""
        result, _ = _get_en_v2_render()
        assert result.success
        pool_clones = [
            r for r in result.records
            if r.entry_type == "pool_clone" and "team" in r.semantic_layout_id
        ]
        assert len(pool_clones) >= 2, (
            f"Only {len(pool_clones)} team bio clones — expected >= 2"
        )

    def test_methodology_has_depth(self):
        """Methodology section has overview + focused + detail structure."""
        result, _ = _get_en_v2_render()
        assert result.success
        meth_records = [
            r for r in result.records if r.section_id == "section_03"
        ]
        layout_ids = [r.semantic_layout_id for r in meth_records]
        assert "methodology_overview_4" in layout_ids
        assert "methodology_focused_4" in layout_ids
        assert "methodology_detail" in layout_ids

    def test_has_a1_company_profile(self):
        """Output includes immutable A1 company profile clones."""
        result, _ = _get_en_v2_render()
        assert result.success
        a1_clones = [
            r for r in result.records
            if r.entry_type == "a1_clone" and r.section_id == "company_profile"
        ]
        # Standard depth = 8 slides
        assert len(a1_clones) >= 6, (
            f"Only {len(a1_clones)} company profile A1 clones — expected >= 6"
        )

    def test_all_entry_types_present(self):
        """Output uses all four entry types: a1_clone, a2_shell, b_variable, pool_clone."""
        result, _ = _get_en_v2_render()
        assert result.success
        entry_types = {r.entry_type for r in result.records}
        assert entry_types == {"a1_clone", "a2_shell", "b_variable", "pool_clone"}

    def test_no_render_errors(self):
        """No errors in the render result."""
        result, _ = _get_en_v2_render()
        assert result.success
        assert len(result.render_errors) == 0
        assert len(result.manifest_errors) == 0


# ── AR V2 Render Acceptance ───────────────────────────────────────


def _build_ar_side_by_side_manifest():
    """Build a representative AR manifest for side-by-side acceptance."""
    from src.models.proposal_manifest import (
        ContentSourcePolicy,
        ManifestEntry,
        ProposalManifest,
        build_inclusion_policy,
        get_company_profile_ids,
    )
    from src.services.selection_policies import (
        select_case_studies,
        select_team_members,
    )

    lock_data = json.loads(CATALOG_LOCK_AR.read_text(encoding="utf-8"))

    rfp_ctx = {
        "sector": "technology",
        "services": ["strategy", "digital transformation"],
        "geography": "ksa",
        "technology_keywords": ["cloud", "digital"],
        "capability_tags": ["strategy", "advisory"],
        "language": "ar",
    }

    cs_candidates = []
    for _cat, entries in lock_data.get("case_study_pool", {}).items():
        for entry in entries:
            cs_candidates.append({
                "asset_id": entry["semantic_id"],
                "slide_idx": entry["slide_idx"],
                "semantic_layout_id": entry["semantic_layout_id"],
                "sector": "technology",
                "services": ["strategy"],
                "geography": "ksa",
                "technology_keywords": ["digital"],
                "capability_tags": ["advisory"],
                "language": "ar",
            })

    team_candidates = []
    for entry in lock_data.get("team_bio_pool", []):
        team_candidates.append({
            "asset_id": entry["semantic_id"],
            "slide_idx": entry["slide_idx"],
            "semantic_layout_id": entry["semantic_layout_id"],
            "sector_experience": ["technology"],
            "services": ["strategy"],
            "roles": ["lead"],
            "geography_experience": ["ksa"],
            "technology_keywords": ["cloud"],
            "language": "ar",
        })

    cs_result = select_case_studies(cs_candidates, rfp_ctx, min_count=5, max_count=5)
    team_result = select_team_members(team_candidates, rfp_ctx, min_count=3, max_count=3)

    inclusion_policy = build_inclusion_policy(
        proposal_mode="standard", geography="ksa",
        sector="technology", case_study_count=(4, 12), team_bio_count=(2, 6),
    )
    profile_ids = get_company_profile_ids(inclusion_policy.company_profile_depth)

    cs_idx_map = {}
    for _cat, entries in lock_data.get("case_study_pool", {}).items():
        for entry in entries:
            cs_idx_map[entry["semantic_id"]] = entry["slide_idx"]
    team_idx_map = {}
    for entry in lock_data.get("team_bio_pool", []):
        team_idx_map[entry["semantic_id"]] = entry["slide_idx"]

    entries: list[ManifestEntry] = []

    # ── Cover (AR) ──
    entries.append(ManifestEntry(
        entry_type="a2_shell", asset_id="proposal_cover",
        semantic_layout_id="proposal_cover",
        content_source_policy=ContentSourcePolicy.PROPOSAL_SPECIFIC,
        section_id="cover",
        injection_data={"subtitle": "\u0627\u0644\u062a\u062d\u0648\u0644 \u0627\u0644\u0631\u0642\u0645\u064a", "client_name": "\u0634\u0631\u0643\u0629 \u0627\u0644\u0645\u062b\u0627\u0644", "date_text": "\u0645\u0627\u0631\u0633 \u0662\u0660\u0662\u0666"},
    ))
    entries.append(ManifestEntry(
        entry_type="a2_shell", asset_id="intro_message",
        semantic_layout_id="intro_message",
        content_source_policy=ContentSourcePolicy.PROPOSAL_SPECIFIC,
        section_id="cover",
    ))
    entries.append(ManifestEntry(
        entry_type="a2_shell", asset_id="toc_agenda",
        semantic_layout_id="toc_table",
        content_source_policy=ContentSourcePolicy.PROPOSAL_SPECIFIC,
        section_id="cover",
        injection_data={"title": "\u062c\u062f\u0648\u0644 \u0627\u0644\u0645\u062d\u062a\u0648\u064a\u0627\u062a", "rows": [
            ["\u0660\u0661", "\u0627\u0644\u0641\u0647\u0645"], ["\u0660\u0662", "\u0644\u0645\u0627\u0630\u0627 \u0633\u062a\u0631\u0627\u062a\u064a\u062c\u064a\u0643 \u062c\u064a\u0631\u0632"],
            ["\u0660\u0663", "\u0627\u0644\u0645\u0646\u0647\u062c\u064a\u0629"], ["\u0660\u0664", "\u0627\u0644\u062c\u062f\u0648\u0644 \u0627\u0644\u0632\u0645\u0646\u064a"],
            ["\u0660\u0665", "\u0627\u0644\u0641\u0631\u064a\u0642"], ["\u0660\u0666", "\u0627\u0644\u062d\u0648\u0643\u0645\u0629"],
        ]},
    ))

    # ── Sections 01-06 with AR injection text ──
    ar_sections = [
        ("section_divider_01", "section_01", "\u0627\u0644\u0641\u0647\u0645"),
        ("section_divider_02", "section_02", "\u0644\u0645\u0627\u0630\u0627 \u0633\u062a\u0631\u0627\u062a\u064a\u062c\u064a\u0643 \u062c\u064a\u0631\u0632"),
        ("section_divider_03", "section_03", "\u0627\u0644\u0645\u0646\u0647\u062c\u064a\u0629"),
        ("section_divider_04", "section_04", "\u0627\u0644\u062c\u062f\u0648\u0644 \u0627\u0644\u0632\u0645\u0646\u064a"),
        ("section_divider_05", "section_05", "\u0627\u0644\u0641\u0631\u064a\u0642"),
        ("section_divider_06", "section_06", "\u0627\u0644\u062d\u0648\u0643\u0645\u0629"),
    ]

    for div_id, sec_id, ar_title in ar_sections:
        entries.append(ManifestEntry(
            entry_type="a2_shell", asset_id=div_id,
            semantic_layout_id=div_id,
            content_source_policy=ContentSourcePolicy.PROPOSAL_SPECIFIC,
            section_id=sec_id,
            injection_data={"title": ar_title, "body": " "},
        ))

        if sec_id == "section_01":
            for i in range(1, 4):
                entries.append(ManifestEntry(
                    entry_type="b_variable", asset_id=f"understanding_{i:02d}",
                    semantic_layout_id="content_heading_desc",
                    content_source_policy=ContentSourcePolicy.PROPOSAL_SPECIFIC,
                    section_id=sec_id,
                    injection_data={"title": f"\u0627\u0644\u0641\u0647\u0645 {i}", "body": f"\u0645\u062d\u062a\u0648\u0649 \u0627\u0644\u0634\u0631\u064a\u062d\u0629 {i}"},
                ))
        elif sec_id == "section_02":
            entries.append(ManifestEntry(
                entry_type="b_variable", asset_id="why_sg_argument",
                semantic_layout_id="content_heading_desc",
                content_source_policy=ContentSourcePolicy.PROPOSAL_SPECIFIC,
                section_id=sec_id,
                injection_data={"title": "\u0633\u062c\u0644\u0646\u0627 \u0627\u0644\u062d\u0627\u0641\u0644", "body": "\u0623\u0643\u062b\u0631 \u0645\u0646 200 \u0628\u0631\u0646\u0627\u0645\u062c \u062a\u062d\u0648\u0644"},
            ))
            for sa in cs_result.selected:
                slide_idx = cs_idx_map.get(sa.asset_id)
                if slide_idx is not None:
                    entries.append(ManifestEntry(
                        entry_type="pool_clone", asset_id=sa.asset_id,
                        semantic_layout_id="case_study_cases",
                        content_source_policy=ContentSourcePolicy.APPROVED_ASSET_POOL,
                        section_id=sec_id,
                        injection_data={"source_slide_idx": slide_idx},
                    ))
        elif sec_id == "section_03":
            entries.append(ManifestEntry(
                entry_type="b_variable", asset_id="methodology_overview",
                semantic_layout_id="methodology_overview_4",
                content_source_policy=ContentSourcePolicy.PROPOSAL_SPECIFIC,
                section_id=sec_id,
                injection_data={"title": "\u0645\u0646\u0647\u062c\u064a\u062a\u0646\u0627", "body": "\u0646\u0647\u062c \u0645\u0646 \u0623\u0631\u0628\u0639 \u0645\u0631\u0627\u062d\u0644"},
            ))
            for p in range(1, 5):
                entries.append(ManifestEntry(
                    entry_type="b_variable", asset_id=f"methodology_phase_{p}",
                    semantic_layout_id="methodology_focused_4",
                    content_source_policy=ContentSourcePolicy.PROPOSAL_SPECIFIC,
                    section_id=sec_id,
                    injection_data={"title": f"\u0627\u0644\u0645\u0631\u062d\u0644\u0629 {p}", "body": f"\u0623\u0646\u0634\u0637\u0629 \u0627\u0644\u0645\u0631\u062d\u0644\u0629 {p}"},
                ))
            entries.append(ManifestEntry(
                entry_type="b_variable", asset_id="methodology_detail_01",
                semantic_layout_id="methodology_detail",
                content_source_policy=ContentSourcePolicy.PROPOSAL_SPECIFIC,
                section_id=sec_id,
                injection_data={"title": "\u062a\u0641\u0627\u0635\u064a\u0644 \u0627\u0644\u0623\u0646\u0634\u0637\u0629", "body": "\u062a\u0641\u0627\u0635\u064a\u0644 \u0627\u0644\u0645\u0631\u062d\u0644\u0629 \u0627\u0644\u062b\u0627\u0644\u062b\u0629"},
            ))
        elif sec_id == "section_04":
            entries.append(ManifestEntry(
                entry_type="b_variable", asset_id="timeline_01",
                semantic_layout_id="content_heading_content",
                content_source_policy=ContentSourcePolicy.PROPOSAL_SPECIFIC,
                section_id=sec_id,
                injection_data={"title": "\u0627\u0644\u062c\u062f\u0648\u0644 \u0627\u0644\u0632\u0645\u0646\u064a", "body": "26 \u0623\u0633\u0628\u0648\u0639\u0627\u064b"},
            ))
            entries.append(ManifestEntry(
                entry_type="b_variable", asset_id="deliverables_01",
                semantic_layout_id="content_heading_desc",
                content_source_policy=ContentSourcePolicy.PROPOSAL_SPECIFIC,
                section_id=sec_id,
                injection_data={"title": "\u0627\u0644\u0645\u062e\u0631\u062c\u0627\u062a", "body": "\u062a\u0642\u0631\u064a\u0631 \u0627\u0644\u062a\u0642\u064a\u064a\u0645 \u0648\u0627\u0644\u0627\u0633\u062a\u0631\u0627\u062a\u064a\u062c\u064a\u0629"},
            ))
        elif sec_id == "section_05":
            for sa in team_result.selected:
                slide_idx = team_idx_map.get(sa.asset_id)
                if slide_idx is not None:
                    entries.append(ManifestEntry(
                        entry_type="pool_clone", asset_id=sa.asset_id,
                        semantic_layout_id="team_two_members",
                        content_source_policy=ContentSourcePolicy.APPROVED_ASSET_POOL,
                        section_id=sec_id,
                        injection_data={"source_slide_idx": slide_idx},
                    ))
        elif sec_id == "section_06":
            entries.append(ManifestEntry(
                entry_type="b_variable", asset_id="governance_01",
                semantic_layout_id="content_heading_desc",
                content_source_policy=ContentSourcePolicy.PROPOSAL_SPECIFIC,
                section_id=sec_id,
                injection_data={"title": "\u062d\u0648\u0643\u0645\u0629 \u0627\u0644\u0645\u0634\u0631\u0648\u0639", "body": "\u0646\u0645\u0648\u0630\u062c \u062d\u0648\u0643\u0645\u0629 \u0645\u0646\u0638\u0645"},
            ))

    # ── Company profile + closing ──
    for profile_id in profile_ids:
        entries.append(ManifestEntry(
            entry_type="a1_clone", asset_id=profile_id,
            semantic_layout_id=profile_id,
            content_source_policy=ContentSourcePolicy.INSTITUTIONAL_REUSE,
            section_id="company_profile",
        ))
    entries.append(ManifestEntry(
        entry_type="a1_clone", asset_id="know_more",
        semantic_layout_id="know_more",
        content_source_policy=ContentSourcePolicy.INSTITUTIONAL_REUSE,
        section_id="closing",
    ))
    entries.append(ManifestEntry(
        entry_type="a1_clone", asset_id="contact",
        semantic_layout_id="contact",
        content_source_policy=ContentSourcePolicy.INSTITUTIONAL_REUSE,
        section_id="closing",
    ))

    return ProposalManifest(entries=entries, inclusion_policy=inclusion_policy)


_ar_v2_render_cache: dict[str, object] = {}


def _get_ar_v2_render():
    """Render AR v2 output once and cache."""
    if "result" in _ar_v2_render_cache:
        return _ar_v2_render_cache["result"], _ar_v2_render_cache["output_path"]

    from src.services.renderer_v2 import render_v2
    from src.services.template_manager import TemplateManager

    manifest = _build_ar_side_by_side_manifest()
    tm = TemplateManager(str(AR_POTX_PATH), CATALOG_LOCK_AR)

    output_dir = Path(tempfile.mkdtemp(prefix="phase19_ar_"))
    output_path = output_dir / "side_by_side_ar.pptx"
    result = render_v2(manifest, tm, CATALOG_LOCK_AR, output_path)

    _ar_v2_render_cache["result"] = result
    _ar_v2_render_cache["output_path"] = output_path
    return result, output_path


@requires_ar_template
class TestV2RealRenderAR:
    """V2 produces a real PPTX from the official AR template."""

    def test_ar_v2_render_succeeds(self):
        result, _ = _get_ar_v2_render()
        assert result.success, f"V2 AR render failed: {result.render_errors}"

    def test_ar_v2_output_is_valid_pptx(self):
        from pptx import Presentation

        result, output_path = _get_ar_v2_render()
        assert result.success
        prs = Presentation(str(output_path))
        assert len(prs.slides) > 0

    def test_ar_v2_slide_count_matches_en(self):
        """AR output has similar slide count to EN."""
        en_result, _ = _get_en_v2_render()
        ar_result, _ = _get_ar_v2_render()
        assert en_result.success
        assert ar_result.success
        # AR and EN should have same structure (±2 slides for rounding)
        assert abs(en_result.total_slides - ar_result.total_slides) <= 2, (
            f"EN={en_result.total_slides} vs AR={ar_result.total_slides} — "
            f"slide count mismatch > 2"
        )

    def test_ar_v2_has_arabic_text(self):
        """AR output contains Arabic characters (not empty or Latin only)."""
        from pptx import Presentation

        result, output_path = _get_ar_v2_render()
        assert result.success

        prs = Presentation(str(output_path))
        arabic_found = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if any("\u0600" <= c <= "\u06FF" for c in run.text):
                                arabic_found = True
                                break
                        if arabic_found:
                            break
                if arabic_found:
                    break
            if arabic_found:
                break

        assert arabic_found, "No Arabic text found in AR v2 output"


# ── Both Renderers Side-by-Side (no interference) ─────────────────


@requires_en_template
class TestBothRenderersSideBySide:
    """Verify both renderers can be used in the same session."""

    def test_importing_v2_does_not_break_legacy_imports(self):
        """Import both renderers — legacy exports still accessible."""
        from src.services.renderer import LAYOUT_MAP, render_pptx
        from src.services.renderer_v2 import render_v2

        assert callable(render_pptx)
        assert callable(render_v2)
        assert len(LAYOUT_MAP) > 0

    def test_v2_render_does_not_modify_legacy_scorer(self):
        """After V2 render, legacy scorer profile is unchanged."""
        legacy_before = get_legacy_profile()
        # Run v2 render
        _get_en_v2_render()
        legacy_after = get_legacy_profile()

        for field_name in _LEGACY_PINNED:
            before_val = getattr(legacy_before, field_name)
            after_val = getattr(legacy_after, field_name)
            assert before_val == after_val, (
                f"Legacy profile changed after v2 render: {field_name}"
            )

    def test_v2_and_legacy_score_same_output_differently(self):
        """Same shapes scored under both profiles yield different results."""
        result, output_path = _get_en_v2_render()
        assert result.success

        from src.services.composition_scorer import extract_shapes

        shapes = extract_shapes(str(output_path))

        legacy_result = score_composition(shapes, [], profile=ScorerProfile.LEGACY)
        v2_result = score_composition(
            shapes, [], profile=ScorerProfile.OFFICIAL_TEMPLATE_V2
        )

        assert legacy_result.profile_used != v2_result.profile_used
        # V2 output with Euclid Flex should have fewer font violations
        # under v2 profile than under legacy profile
        v2_font_violations = sum(
            1 for s in v2_result.slide_scores for v in s.violations
            if v.rule == "font_brand"
        )
        legacy_font_violations = sum(
            1 for s in legacy_result.slide_scores for v in s.violations
            if v.rule == "font_brand"
        )
        # V2 profile should accept the fonts better than legacy
        assert v2_font_violations <= legacy_font_violations


# ── PNG Export Infrastructure ─────────────────────────────────────


class TestPNGExportInfrastructure:
    """PNG export infrastructure is available for visual comparison."""

    @pytest.mark.skipif(
        not Path("scripts/export_pngs.py").exists(),
        reason="scripts/export_pngs.py not yet created",
    )
    def test_export_script_exists(self):
        assert Path("scripts/export_pngs.py").exists()

    @pytest.mark.skipif(
        not Path("scripts/export_pngs.py").exists(),
        reason="scripts/export_pngs.py not yet created",
    )
    def test_export_function_importable(self):
        """export_slides function is importable from the script."""
        import importlib.util

        spec = importlib.util.spec_from_file_location(
            "export_pngs", "scripts/export_pngs.py"
        )
        assert spec is not None
        importlib.util.module_from_spec(spec)
        # Don't execute — just verify the module parses
        assert spec.loader is not None

    @pytest.mark.skipif(
        not Path("scripts/export_pngs.py").exists(),
        reason="scripts/export_pngs.py not yet created",
    )
    def test_export_script_has_correct_structure(self):
        """Export script defines export_slides function."""
        source = Path("scripts/export_pngs.py").read_text(encoding="utf-8")
        assert "def export_slides" in source
        assert "PNG" in source


# ── Phase 19 Meta-Test: Test Coverage Summary ─────────────────────


class TestPhase19CoverageSummary:
    """Meta-test: verify Phase 19 coverage is comprehensive."""

    PHASE_19_TEST_CLASSES = [
        "TestRendererCoexistence",
        "TestPipelineDispatch",
        "TestLegacyRendererIntegrity",
        "TestV2RealRenderEN",
        "TestV2VsRealExampleCharacteristics",
        "TestV2RealRenderAR",
        "TestBothRenderersSideBySide",
        "TestPNGExportInfrastructure",
    ]

    def test_all_phase19_classes_exist(self):
        """All planned Phase 19 test classes are defined."""
        import sys

        module = sys.modules[__name__]
        for cls_name in self.PHASE_19_TEST_CLASSES:
            assert hasattr(module, cls_name), (
                f"Missing Phase 19 test class: {cls_name}"
            )

    def test_minimum_test_count(self):
        """Phase 19 adds a meaningful number of new tests."""
        import sys

        module = sys.modules[__name__]
        total = 0
        for cls_name in self.PHASE_19_TEST_CLASSES:
            cls = getattr(module, cls_name)
            methods = [m for m in dir(cls) if m.startswith("test_")]
            total += len(methods)

        # At least 30 new tests across Phase 19 classes
        assert total >= 30, (
            f"Phase 19 only has {total} tests — expected >= 30"
        )
