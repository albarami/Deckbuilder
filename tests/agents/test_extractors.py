"""Tests for document extractors — PPTX, PDF, DOCX, XLSX."""

import os
import tempfile

from src.models.extraction import ExtractedDocument

# ──────────────────────────────────────────────────────────────
# Fixtures — create minimal test files programmatically
# ──────────────────────────────────────────────────────────────


def _create_test_pptx(path: str) -> None:
    """Create a minimal 1-slide PPTX for testing."""
    from pptx import Presentation

    prs = Presentation()
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    slide.placeholders[0].text = "Test Title"
    slide.placeholders[1].text = "Bullet point one\nBullet point two"
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "Speaker notes here"
    prs.save(path)


def _create_test_pdf(path: str) -> None:
    """Create a minimal 1-page PDF for testing."""
    from PyPDF2 import PdfWriter

    writer = PdfWriter()
    writer.add_blank_page(width=612, height=792)
    # PyPDF2 can't add text to blank pages easily,
    # so we create a PDF with metadata only
    writer.add_metadata({"/Title": "Test PDF"})
    with open(path, "wb") as f:
        writer.write(f)


def _create_test_docx(path: str) -> None:
    """Create a minimal 1-paragraph DOCX for testing."""
    from docx import Document

    doc = Document()
    doc.add_heading("Test Heading", level=1)
    doc.add_paragraph("This is test paragraph content.")
    doc.add_heading("Second Section", level=2)
    doc.add_paragraph("More content in section two.")
    doc.save(path)


def _create_test_xlsx(path: str) -> None:
    """Create a minimal 1-sheet XLSX for testing."""
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Data"
    ws.append(["Name", "Value", "Status"])
    ws.append(["Alpha", "100", "Active"])
    ws.append(["Beta", "200", "Inactive"])
    wb.save(path)


# ──────────────────────────────────────────────────────────────
# PPTX Extractor Tests
# ──────────────────────────────────────────────────────────────


def test_extract_pptx_basic() -> None:
    """Extract a 1-slide PPTX → ExtractedDocument with 1 slide."""
    from src.utils.extractors import extract_pptx

    with tempfile.TemporaryDirectory() as tmpdir:
        path = os.path.join(tmpdir, "test.pptx")
        _create_test_pptx(path)

        doc = extract_pptx(path)

        assert isinstance(doc, ExtractedDocument)
        assert doc.file_type == "pptx"
        assert doc.filename == "test.pptx"
        assert len(doc.slides) == 1
        assert doc.slides[0].slide_number == 1
        assert "Test Title" in doc.slides[0].title
        assert "Bullet point" in doc.slides[0].body_text
        assert "Speaker notes" in doc.slides[0].speaker_notes
        assert len(doc.full_text) > 0
        assert doc.content_hash != ""
        assert doc.error is None


def test_extract_pptx_preserves_structure() -> None:
    """PPTX extractor preserves slide-level structure."""
    from pptx import Presentation

    from src.utils.extractors import extract_pptx

    with tempfile.TemporaryDirectory() as tmpdir:
        path = os.path.join(tmpdir, "multi.pptx")
        prs = Presentation()
        for i in range(3):
            layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(layout)
            slide.placeholders[0].text = f"Slide {i + 1} Title"
            slide.placeholders[1].text = f"Content for slide {i + 1}"
        prs.save(path)

        doc = extract_pptx(path)

        assert len(doc.slides) == 3
        assert doc.slides[0].title == "Slide 1 Title"
        assert doc.slides[2].title == "Slide 3 Title"
        assert doc.slides[1].slide_number == 2


# ──────────────────────────────────────────────────────────────
# PDF Extractor Tests
# ──────────────────────────────────────────────────────────────


def test_extract_pdf_basic() -> None:
    """Extract a PDF → ExtractedDocument with pages."""
    from src.utils.extractors import extract_pdf

    with tempfile.TemporaryDirectory() as tmpdir:
        path = os.path.join(tmpdir, "test.pdf")
        _create_test_pdf(path)

        doc = extract_pdf(path)

        assert isinstance(doc, ExtractedDocument)
        assert doc.file_type == "pdf"
        assert doc.filename == "test.pdf"
        assert len(doc.pages) >= 1
        assert doc.pages[0].page_number == 1
        assert doc.content_hash != ""
        assert doc.error is None


# ──────────────────────────────────────────────────────────────
# DOCX Extractor Tests
# ──────────────────────────────────────────────────────────────


def test_extract_docx_basic() -> None:
    """Extract a DOCX → ExtractedDocument with sections."""
    from src.utils.extractors import extract_docx

    with tempfile.TemporaryDirectory() as tmpdir:
        path = os.path.join(tmpdir, "test.docx")
        _create_test_docx(path)

        doc = extract_docx(path)

        assert isinstance(doc, ExtractedDocument)
        assert doc.file_type == "docx"
        assert doc.filename == "test.docx"
        assert len(doc.sections) >= 2
        assert doc.sections[0].heading == "Test Heading"
        assert doc.sections[0].level == 1
        assert "test paragraph" in doc.sections[0].text
        assert doc.sections[1].heading == "Second Section"
        assert doc.sections[1].level == 2
        assert len(doc.full_text) > 0
        assert doc.content_hash != ""


# ──────────────────────────────────────────────────────────────
# XLSX Extractor Tests
# ──────────────────────────────────────────────────────────────


def test_extract_xlsx_basic() -> None:
    """Extract an XLSX → ExtractedDocument with sheet data."""
    from src.utils.extractors import extract_xlsx

    with tempfile.TemporaryDirectory() as tmpdir:
        path = os.path.join(tmpdir, "test.xlsx")
        _create_test_xlsx(path)

        doc = extract_xlsx(path)

        assert isinstance(doc, ExtractedDocument)
        assert doc.file_type == "xlsx"
        assert doc.filename == "test.xlsx"
        assert len(doc.sheets) == 1
        assert doc.sheets[0].sheet_name == "Data"
        assert doc.sheets[0].headers == ["Name", "Value", "Status"]
        assert doc.sheets[0].row_count == 2
        assert "Alpha" in doc.full_text
        assert doc.content_hash != ""


# ──────────────────────────────────────────────────────────────
# Dispatcher Tests
# ──────────────────────────────────────────────────────────────


def test_extract_document_dispatch_pptx() -> None:
    """extract_document dispatches to correct extractor by extension."""
    from src.utils.extractors import extract_document

    with tempfile.TemporaryDirectory() as tmpdir:
        path = os.path.join(tmpdir, "test.pptx")
        _create_test_pptx(path)

        doc = extract_document(path)

        assert doc.file_type == "pptx"
        assert len(doc.slides) == 1


def test_extract_document_dispatch_docx() -> None:
    """extract_document dispatches DOCX correctly."""
    from src.utils.extractors import extract_document

    with tempfile.TemporaryDirectory() as tmpdir:
        path = os.path.join(tmpdir, "test.docx")
        _create_test_docx(path)

        doc = extract_document(path)

        assert doc.file_type == "docx"
        assert len(doc.sections) >= 2


def test_extract_document_unsupported() -> None:
    """Unsupported file type returns document with error."""
    from src.utils.extractors import extract_document

    with tempfile.TemporaryDirectory() as tmpdir:
        path = os.path.join(tmpdir, "test.txt")
        with open(path, "w") as f:
            f.write("plain text")

        doc = extract_document(path)

        assert doc.error is not None
        assert "unsupported" in doc.error.lower()


# ──────────────────────────────────────────────────────────────
# Error Handling Tests
# ──────────────────────────────────────────────────────────────


def test_extract_corrupted_file() -> None:
    """Corrupted file returns document with error, no crash."""
    from src.utils.extractors import extract_document

    with tempfile.TemporaryDirectory() as tmpdir:
        path = os.path.join(tmpdir, "corrupt.pptx")
        with open(path, "wb") as f:
            f.write(b"this is not a valid pptx file")

        doc = extract_document(path)

        assert doc.error is not None
        assert doc.filename == "corrupt.pptx"
        assert doc.file_type == "pptx"


def test_extract_nonexistent_file() -> None:
    """Nonexistent file returns document with error."""
    from src.utils.extractors import extract_document

    doc = extract_document("/nonexistent/file.pptx")

    assert doc.error is not None


# ──────────────────────────────────────────────────────────────
# Content Hash Consistency
# ──────────────────────────────────────────────────────────────


def test_content_hash_consistent() -> None:
    """Same file produces same content_hash on re-extraction."""
    from src.utils.extractors import extract_document

    with tempfile.TemporaryDirectory() as tmpdir:
        path = os.path.join(tmpdir, "test.docx")
        _create_test_docx(path)

        doc1 = extract_document(path)
        doc2 = extract_document(path)

        assert doc1.content_hash == doc2.content_hash
        assert doc1.content_hash != ""


def test_content_hash_is_sha256() -> None:
    """Content hash is a valid SHA-256 hex string."""
    from src.utils.extractors import extract_document

    with tempfile.TemporaryDirectory() as tmpdir:
        path = os.path.join(tmpdir, "test.docx")
        _create_test_docx(path)

        doc = extract_document(path)

        # SHA-256 hex is 64 characters
        assert len(doc.content_hash) == 64
        # Valid hex
        int(doc.content_hash, 16)


# ──────────────────────────────────────────────────────────────
# Batch Extraction Tests
# ──────────────────────────────────────────────────────────────


def test_extract_directory_batch() -> None:
    """extract_directory processes multiple files in a directory."""
    from src.utils.extractors import extract_directory

    with tempfile.TemporaryDirectory() as tmpdir:
        _create_test_pptx(os.path.join(tmpdir, "a.pptx"))
        _create_test_docx(os.path.join(tmpdir, "b.docx"))
        _create_test_xlsx(os.path.join(tmpdir, "c.xlsx"))
        # Also put an unsupported file
        with open(os.path.join(tmpdir, "readme.txt"), "w") as f:
            f.write("not a document")

        results = extract_directory(tmpdir)

        # Should process 3 supported files (skip .txt)
        assert len(results) == 3
        types = {d.file_type for d in results}
        assert types == {"pptx", "docx", "xlsx"}
        assert all(d.error is None for d in results)
