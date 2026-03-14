"""Tests for Phase 1 — template_manager.py.

Tests POTX patching, catalog lock loading, hash enforcement,
semantic ID resolution (fail-closed), slide cloning, and layout access.
"""

from __future__ import annotations

import json
import shutil
import tempfile
from pathlib import Path
from unittest.mock import patch

import pytest
from pptx import Presentation as PptxPresentation

from src.services.template_auditor import EN_POTX, file_hash
from src.services.template_manager import (
    AssetLocation,
    CatalogLockError,
    LayoutLocation,
    SemanticIDNotFound,
    ShellDef,
    TemplateHashMismatch,
    TemplateManager,
    load_catalog_lock,
)

DATA_DIR = Path(__file__).resolve().parent.parent.parent / "src" / "data"
CATALOG_LOCK_EN = DATA_DIR / "catalog_lock_en.json"

# Skip all tests if template file is not available
pytestmark = pytest.mark.skipif(
    not EN_POTX.exists() or not CATALOG_LOCK_EN.exists(),
    reason="Template or catalog lock not available",
)


@pytest.fixture(scope="module")
def patch_dir():
    """Shared temp dir for POTX patching across all tests in this module."""
    d = Path(tempfile.mkdtemp(prefix="test_tm_"))
    yield d
    shutil.rmtree(d, ignore_errors=True)


@pytest.fixture(scope="module")
def tm(patch_dir):
    """Shared TemplateManager instance for the module."""
    return TemplateManager(EN_POTX, CATALOG_LOCK_EN, patch_dir)


# ── Catalog lock loading ─────────────────────────────────────────────────


class TestCatalogLockLoading:
    def test_load_valid_lock(self):
        lock = load_catalog_lock(CATALOG_LOCK_EN)
        assert "template_hash" in lock
        assert "a1_immutable" in lock
        assert "layouts" in lock

    def test_load_missing_file(self, tmp_path):
        with pytest.raises(CatalogLockError, match="not found"):
            load_catalog_lock(tmp_path / "nonexistent.json")

    def test_load_incomplete_lock(self, tmp_path):
        bad_lock = tmp_path / "bad.json"
        bad_lock.write_text('{"template_hash": "x"}', encoding="utf-8")
        with pytest.raises(CatalogLockError, match="missing fields"):
            load_catalog_lock(bad_lock)


# ── Template hash enforcement ────────────────────────────────────────────


class TestHashEnforcement:
    def test_valid_hash_accepted(self, tm):
        """TemplateManager loads successfully with matching hash."""
        assert tm.template_hash.startswith("sha256:")

    def test_mismatched_hash_rejected(self, tmp_path):
        """TemplateManager fails-closed on hash mismatch."""
        # Create a catalog lock with wrong hash
        lock = load_catalog_lock(CATALOG_LOCK_EN)
        lock["template_hash"] = "sha256:0000000000000000000000000000000000000000"
        bad_lock = tmp_path / "bad_hash.json"
        bad_lock.write_text(
            json.dumps(lock, default=str), encoding="utf-8"
        )
        with pytest.raises(TemplateHashMismatch, match="changed since"):
            TemplateManager(EN_POTX, bad_lock, tmp_path)


# ── Semantic ID resolution (fail-closed) ─────────────────────────────────


class TestSemanticResolution:
    def test_resolve_a1_known(self, tm):
        loc = tm.resolve_a1("ksa_context")
        assert isinstance(loc, AssetLocation)
        assert loc.slide_idx >= 0
        assert loc.semantic_layout_id == "ksa_context"

    def test_resolve_a1_missing_fails_closed(self, tm):
        with pytest.raises(SemanticIDNotFound, match="not found"):
            tm.resolve_a1("nonexistent_asset")

    def test_resolve_a2_known(self, tm):
        loc = tm.resolve_a2("proposal_cover")
        assert isinstance(loc, ShellDef)
        assert loc.slide_idx >= 0
        assert loc.allowlist  # Must have allowlist

    def test_resolve_a2_missing_fails_closed(self, tm):
        with pytest.raises(SemanticIDNotFound, match="not found"):
            tm.resolve_a2("nonexistent_shell")

    def test_resolve_divider_known(self, tm):
        loc = tm.resolve_divider("01")
        assert isinstance(loc, ShellDef)
        assert loc.slide_idx >= 0

    def test_resolve_divider_missing_fails_closed(self, tm):
        with pytest.raises(SemanticIDNotFound, match="not found"):
            tm.resolve_divider("99")

    def test_resolve_layout_known(self, tm):
        loc = tm.resolve_layout("methodology_overview_4")
        assert isinstance(loc, LayoutLocation)
        assert loc.placeholder_count > 0

    def test_resolve_layout_missing_fails_closed(self, tm):
        with pytest.raises(SemanticIDNotFound, match="not found"):
            tm.resolve_layout("nonexistent_layout")

    def test_all_a1_ids_resolvable(self, tm):
        """Every A1 ID in the catalog lock is resolvable."""
        for sid in tm.a1_asset_ids:
            loc = tm.resolve_a1(sid)
            assert loc.slide_idx >= 0

    def test_all_a2_ids_resolvable(self, tm):
        """Every A2 ID in the catalog lock is resolvable."""
        for sid in tm.a2_shell_ids:
            loc = tm.resolve_a2(sid)
            assert loc.slide_idx >= 0

    def test_all_layout_ids_resolvable(self, tm):
        """Every layout ID in the catalog lock is resolvable."""
        for lid in tm.layout_ids:
            loc = tm.resolve_layout(lid)
            assert loc.display_name


# ── Layout object resolution ─────────────────────────────────────────────


class TestLayoutObjects:
    def test_get_layout_object(self, tm):
        layout = tm.get_layout_object("methodology_overview_4")
        assert layout is not None
        assert layout.name  # Has a display name

    def test_get_layout_preserves_placeholders(self, tm):
        layout = tm.get_layout_object("team_two_members")
        assert len(layout.placeholders) > 0


# ── Slide cloning ────────────────────────────────────────────────────────


class TestSlideCloning:
    def test_clone_a1(self, tm):
        """Clone an A1 slide preserving layout linkage."""
        original_count = len(tm.presentation.slides)
        slide = tm.clone_a1("main_cover")
        assert len(tm.presentation.slides) == original_count + 1
        assert slide.slide_layout is not None

    def test_clone_a2(self, tm):
        """Clone an A2 shell."""
        original_count = len(tm.presentation.slides)
        slide = tm.clone_a2("proposal_cover")
        assert len(tm.presentation.slides) == original_count + 1

    def test_clone_divider(self, tm):
        """Clone a section divider."""
        original_count = len(tm.presentation.slides)
        slide = tm.clone_divider("01")
        assert len(tm.presentation.slides) == original_count + 1

    def test_clone_preserves_layout_name(self, tm):
        """Cloned slide has the same layout as the source."""
        loc = tm.resolve_a1("contact")
        source = tm.presentation.slides[loc.slide_idx]
        source_layout_name = source.slide_layout.name
        cloned = tm.clone_a1("contact")
        assert cloned.slide_layout.name == source_layout_name

    def test_add_slide_from_layout(self, tm):
        """Add a new slide from a semantic layout ID."""
        original_count = len(tm.presentation.slides)
        slide = tm.add_slide_from_layout("content_heading_desc")
        assert len(tm.presentation.slides) == original_count + 1
        assert slide.slide_layout is not None

    def test_clone_out_of_range_raises(self, tm):
        """Clone with invalid index raises IndexError."""
        with pytest.raises(IndexError):
            tm.clone_slide(99999)


# ── Pool access ──────────────────────────────────────────────────────────


class TestPoolAccess:
    def test_case_study_pool_disjoint_from_service_dividers(self, tm):
        """Case study and service divider pools are disjoint."""
        case_indices = set()
        for items in tm.get_case_study_pool().values():
            for item in items:
                case_indices.add(item.slide_idx)

        svc_indices = {s.slide_idx for s in tm.get_service_divider_pool()}
        assert not case_indices & svc_indices, "Pools must be disjoint"

    def test_team_pool_has_family(self, tm):
        """Every team pool entry has a team_family field."""
        for entry in tm.get_team_pool():
            assert "team_family" in entry

    def test_case_study_pool_not_empty(self, tm):
        pool = tm.get_case_study_pool()
        total = sum(len(v) for v in pool.values())
        assert total > 0

    def test_service_divider_pool_not_empty(self, tm):
        assert len(tm.get_service_divider_pool()) > 0


# ── Properties ───────────────────────────────────────────────────────────


class TestProperties:
    def test_language(self, tm):
        assert tm.language == "en"

    def test_presentation_loaded(self, tm):
        assert tm.presentation is not None
        assert len(tm.presentation.slides) > 0

    def test_a1_asset_ids(self, tm):
        ids = tm.a1_asset_ids
        assert len(ids) > 0
        assert "ksa_context" in ids
        assert "main_cover" in ids
        assert "contact" in ids

    def test_a2_shell_ids(self, tm):
        ids = tm.a2_shell_ids
        assert len(ids) > 0
        assert "proposal_cover" in ids

    def test_divider_numbers(self, tm):
        nums = tm.divider_numbers
        assert "01" in nums
        assert "06" in nums

    def test_layout_ids(self, tm):
        ids = tm.layout_ids
        assert len(ids) > 0
        assert "methodology_overview_4" in ids
        assert "team_two_members" in ids
        assert "case_study_cases" in ids


# ── Save ─────────────────────────────────────────────────────────────────


class TestSave:
    def test_save_output(self, tm, tmp_path):
        """Save produces a valid PPTX file."""
        out = tm.save(tmp_path / "output.pptx")
        assert out.exists()
        assert out.stat().st_size > 0

        # Verify it can be opened
        prs = PptxPresentation(str(out))
        assert len(prs.slides) > 0
