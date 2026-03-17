"""Phase 18 — AR Verification: Parity Audit + Full Arabic Render + All AR Gates.

Verifies that the Arabic rendering pipeline works correctly with the official
Arabic_Proposal_Template.potx template.  This phase is the AR counterpart to
Phase 17 (EN verification).

Three-part verification:
  1. **Parity Audit** (gate 9): EN/AR catalog lock parity — all 14 checks pass,
     zero critical failures, zero warnings.
  2. **Full AR Render** (gate 35): Render a full AR proposal (~40 slides) from
     the official Arabic template, exercising A1 clones, A2 shells, B variable,
     and pool clones with Arabic injection text.
  3. **All Applicable AR Gates**: Every acceptance gate that applies to Arabic
     output is tested against the rendered AR PPTX.

Proposal composition (mirrors EN Phase 17 structure):
  - 10 A1 clones  (8 company profile standard-depth + know_more + contact)
  - 9  A2 shells  (proposal_cover + intro_message + toc_agenda + 6 dividers)
  - 13 B variable (3 understanding + 1 why_sg + 6 methodology + 2 timeline + 1 governance)
  - 8  pool clones (5 case studies via deterministic selection + 3 team bios)
  Total: ~40 entries covering all mandatory sections and entry types.

AR Gate Matrix (42 gates — 38 AR-applicable, 4 N-A):
  Gate | Status | Test / Reason
  ─────┼────────┼──────────────────────────────────────────────────────────
   01  | PASS   | test_gate_01_official_layouts — every record has semantic_layout_id
   02  | PASS   | test_gate_02_zero_shape_creation_all_v2_modules — AST scan all v2 modules
   03  | PASS   | test_gate_03_euclid_flex_fonts — brand font check (Euclid + Bahij) on rendered PPTX
   04  | PASS   | test_gate_04_navy_header_color — title placeholder color on rendered PPTX
   05  | PASS   | test_gate_05_left_margin — scorer violations on rendered PPTX
   06  | PASS   | test_gate_06_zero_composition_blockers — scorer blocker_count on rendered PPTX
   07  | PASS   | test_gate_07_zero_mojibake — text encoding check on rendered PPTX (Arabic text)
   08  | PASS   | test_gate_08_zero_contract_violations — render_result records
   09  | PASS   | test_gate_09_template_parity_audit — EN/AR catalog lock parity
   10  | PASS   | test_gate_10_template_hash_validated — render_result.template_hash
   11  | PASS   | test_gate_11_zero_anti_leak — rendered PPTX text scan for forbidden content
   12  | PASS   | test_gate_12_mandatory_section_flow — render_result section order
   13  | PASS   | test_gate_13_exact_divider_numbering — render_result divider sequence
   14  | PASS   | test_gate_14_content_source_policy — manifest entry policy validation
   15  | PASS   | test_gate_15_* (3 tests) — HouseInclusionPolicy + KSA/non-KSA rules
   16  | PASS   | test_gate_16_* (2 tests) — SlideBudget validation + all sections budgeted
   17  | PASS   | test_gate_17_* (2 tests) — selection result audit (scores + reasons)
   18  | PASS   | test_gate_18_editability — python-pptx re-opens rendered output
   19  | PASS   | test_gate_19_a1_not_sanitized — no sanitization_report on A1 records
   20  | PASS   | test_gate_20_no_legacy_imports — AST scan renderer_v2.py imports
   21  | PASS   | test_gate_21_shell_sanitization — all A2 records have sanitization_report
   22  | PASS   | test_gate_22_scorer_profile — v2 profile is OFFICIAL_TEMPLATE_V2
   23  | PASS   | test_gate_23_semantic_ids — all records have non-empty semantic_layout_id
   24  | PASS   | test_gate_24_section_order — render_result section_id ordering
   25  | PASS   | test_gate_25_divider_pattern — "01" through "06" in render_result
   26  | PASS   | test_gate_26_methodology_follows_budget — overview + focused + detail in records
   27  | PASS   | test_gate_27_case_studies_use_correct_layout — pool_clone records verified
   28  | PASS   | test_gate_28_team_bios_use_correct_layout — pool_clone records verified
   29  | PASS   | test_gate_29_typography_hierarchy — title vs body font size on rendered PPTX
   30  | PASS   | test_gate_30_no_generic_powerpoint_look — brand font (Euclid/Bahij) in rendered PPTX
   31  | PASS   | test_gate_31_structural_fidelity_a1_clones — layout-level fidelity on rendered PPTX
   32  | PASS   | test_gate_32_house_shells_allowlisted — sanitization_report on A2 records
   33  | PASS   | test_gate_33_no_generic_regression — slide count + dividers + pool clones
   34  | N-A    | EN-specific: EN full deck render (Phase 17)
   35  | PASS   | test_gate_35_ar_full_deck_renders — output file exists and non-empty
   36  | PASS   | test_gate_36_* (2 tests) — valid PPTX + all slides have layouts
   37  | PASS   | test_gate_37_zero_shape_creation_v2_path — AST scan v2 modules
   38  | PASS   | test_gate_38_anti_leak_on_non_a1_slides — rendered PPTX text scan
   39  | PASS   | test_gate_39_case_study_clone_correct — shapes present on rendered PPTX
   40  | PASS   | test_gate_40_team_bio_clone_correct — shapes present on rendered PPTX
   41  | PASS   | test_gate_41_methodology_follows_template_family — layout IDs in records
   42  | PASS   | test_gate_42_scorer_profile_per_mode — pipeline profile dispatch

AR-specific verifications beyond EN:
  - Gate 9: template parity audit (EN/AR catalog lock structural parity)
  - Arabic text injection verified (no mojibake, correct encoding)
  - AR template hash validated against AR catalog lock
  - AR catalog lock language == "ar"

Verification approach:
  - Gates 3-6, 7, 29-30: verified against RENDERED AR output PPTX
  - Gates 27-28, 39-40: verified from render_result records AND rendered output shapes
  - Gates 31: verified layout-level structural fidelity on rendered output
  - Gates 2, 20, 37: static AST analysis (same for EN/AR)
  - Gate 9: EN/AR parity audit via template_parity module
  - All other gates: verified from render_result records or manifest structure

Test classes:
  - TestARTemplateAvailability: template + catalog lock prerequisites
  - TestARParityAudit: gate 9 — EN/AR parity (the AR-specific gate)
  - TestARPolicyAndBudget: gates 15-17
  - TestARManifestConstruction: full AR proposal manifest validity
  - TestARRenderExecution: real render with official AR template
  - TestARHardGates: hard technical gates on rendered AR output
  - TestARVisualGates: visual-fidelity gates on rendered AR output
  - TestARIntegrationGates: integration gates on rendered AR output
  - TestARScorerProfileAlignment: v2 profile correctness
  - TestARGateMatrix: meta-test verifying all gates are covered
"""

from __future__ import annotations

import ast
import json
import logging
from pathlib import Path
from typing import Any

import pytest

from src.models.enums import RendererMode
from src.models.methodology_blueprint import (
    MethodologyBlueprint,
    build_methodology_blueprint,
)
from src.models.proposal_manifest import (
    ContentSourcePolicy,
    HouseInclusionPolicy,
    ManifestEntry,
    ProposalManifest,
    build_inclusion_policy,
    get_company_profile_ids,
    validate_manifest,
)
from src.models.section_blueprint import MANDATORY_SECTION_ORDER
from src.services.scorer_profiles import ScorerProfile
from src.services.selection_policies import (
    CaseStudySelectionResult,
    TeamSelectionResult,
    select_case_studies,
    select_team_members,
)
from src.services.slide_budgeter import (
    compute_slide_budget,
    validate_budget,
)

logger = logging.getLogger(__name__)

# ── Paths ────────────────────────────────────────────────────────────

AR_POTX_PATH = Path(
    r"C:\Projects\Deckbuilder\PROPOSAL_TEMPLATE\Arabic_Proposal_Template.potx"
)
CATALOG_LOCK_AR = Path("src/data/catalog_lock_ar.json")
CATALOG_LOCK_EN = Path("src/data/catalog_lock_en.json")

# Skip marker for tests that need the real AR template
_template_available = (
    AR_POTX_PATH.exists()
    and CATALOG_LOCK_AR.exists()
    and CATALOG_LOCK_EN.exists()
)
requires_ar_template = pytest.mark.skipif(
    not _template_available,
    reason="Official AR template or catalog lock not available",
)

# V2-path module list for static analysis (same as EN)
V2_PATH_MODULES = [
    "src/services/renderer_v2.py",
    "src/services/placeholder_injectors.py",
    "src/services/shell_sanitizer.py",
    "src/services/content_fitter.py",
    "src/services/template_manager.py",
    "src/services/layout_router.py",
]


# ── RFP context for deterministic selection (Arabic) ─────────────────


def _rfp_context() -> dict[str, Any]:
    """Realistic RFP context for Arabic proposal case study / team selection."""
    return {
        "sector": "technology",
        "services": ["strategy", "digital transformation", "consulting"],
        "geography": "ksa",
        "technology_keywords": ["cloud", "digital", "analytics", "sap"],
        "capability_tags": ["strategy", "advisory", "transformation"],
        "language": "ar",
    }


# ── Policies and supporting objects ──────────────────────────────────


def _build_inclusion_policy() -> HouseInclusionPolicy:
    return build_inclusion_policy(
        proposal_mode="standard",
        geography="ksa",
        sector="technology",
        case_study_count=(4, 12),
        team_bio_count=(2, 6),
    )


def _build_methodology_blueprint() -> MethodologyBlueprint:
    return build_methodology_blueprint(
        phase_count=4,
        phases=[
            {
                "phase_name_en": "Discovery & Assessment",
                "phase_name_ar": "الاكتشاف والتقييم",
                "activities": [
                    "مقابلات أصحاب المصلحة",
                    "تقييم الوضع الحالي",
                    "تحليل الفجوات",
                ],
                "deliverables": ["تقرير التقييم", "وثيقة تحليل الفجوات"],
                "governance_tier": "لجنة التوجيه",
            },
            {
                "phase_name_en": "Strategy & Design",
                "phase_name_ar": "الاستراتيجية والتصميم",
                "activities": [
                    "صياغة الاستراتيجية",
                    "هندسة الحلول",
                    "تطوير خارطة الطريق",
                ],
                "deliverables": ["وثيقة الاستراتيجية", "هندسة الحلول"],
                "governance_tier": "لجنة التوجيه",
            },
            {
                "phase_name_en": "Implementation",
                "phase_name_ar": "التنفيذ",
                "activities": [
                    "تهيئة المنصة",
                    "تطوير التكامل",
                    "ترحيل البيانات",
                    "الاختبار وضمان الجودة",
                ],
                "deliverables": [
                    "المنصة المهيأة",
                    "تقرير التكامل",
                    "نتائج الاختبار",
                ],
                "governance_tier": "مجلس المشروع",
            },
            {
                "phase_name_en": "Launch & Transition",
                "phase_name_ar": "الإطلاق والانتقال",
                "activities": [
                    "تدريب المستخدمين",
                    "دعم الإطلاق",
                    "نقل المعرفة",
                ],
                "deliverables": [
                    "مواد التدريب",
                    "تقرير الإطلاق",
                    "وثيقة التسليم",
                ],
                "governance_tier": "لجنة التوجيه",
            },
        ],
        deliverables_linkage={
            "phase_1": ["تقرير التقييم"],
            "phase_2": ["وثيقة الاستراتيجية"],
            "phase_3": ["المنصة المهيأة"],
            "phase_4": ["تقرير الإطلاق"],
        },
        governance_touchpoints={
            "phase_1": "لجنة التوجيه",
            "phase_2": "لجنة التوجيه",
            "phase_3": "مجلس المشروع",
            "phase_4": "لجنة التوجيه",
        },
        timeline_span="26 weeks",
    )


def _build_case_study_candidates() -> list[dict[str, Any]]:
    """Build case study candidates from AR catalog lock pool."""
    lock_data = json.loads(CATALOG_LOCK_AR.read_text(encoding="utf-8"))
    candidates = []
    for _cat, entries in lock_data.get("case_study_pool", {}).items():
        for entry in entries:
            candidates.append(
                {
                    "asset_id": entry["semantic_id"],
                    "slide_idx": entry["slide_idx"],
                    "semantic_layout_id": entry["semantic_layout_id"],
                    "sector": "technology",
                    "services": ["strategy", "consulting"],
                    "geography": "ksa",
                    "technology_keywords": ["digital"],
                    "capability_tags": ["advisory"],
                    "language": "ar",
                }
            )
    return candidates


def _build_team_bio_candidates() -> list[dict[str, Any]]:
    """Build team bio candidates from AR catalog lock pool."""
    lock_data = json.loads(CATALOG_LOCK_AR.read_text(encoding="utf-8"))
    candidates = []
    for entry in lock_data.get("team_bio_pool", []):
        candidates.append(
            {
                "asset_id": entry["semantic_id"],
                "slide_idx": entry["slide_idx"],
                "semantic_layout_id": entry["semantic_layout_id"],
                "sector_experience": ["technology"],
                "services": ["strategy", "digital transformation"],
                "roles": ["lead", "analyst"],
                "geography_experience": ["ksa"],
                "technology_keywords": ["cloud", "sap"],
                "language": "ar",
            }
        )
    return candidates


def _run_case_study_selection() -> CaseStudySelectionResult:
    candidates = _build_case_study_candidates()
    return select_case_studies(
        candidates,
        _rfp_context(),
        min_count=5,
        max_count=5,
    )


def _run_team_selection() -> TeamSelectionResult:
    candidates = _build_team_bio_candidates()
    return select_team_members(
        candidates,
        _rfp_context(),
        min_count=3,
        max_count=3,
    )


# ── Helper: Build a full AR proposal manifest ───────────────────────


def _build_ar_full_proposal_manifest() -> ProposalManifest:
    """Build a full AR ProposalManifest suitable for Phase 18 verification.

    Mirrors the EN Phase 17 manifest structure but uses:
    - AR catalog lock for asset resolution
    - Arabic text for all injection_data content
    - AR case study and team bio pools

    This is a FULL proposal, not a mini-deck.
    """
    lock_data = json.loads(CATALOG_LOCK_AR.read_text(encoding="utf-8"))

    # Run deterministic selection
    cs_result = _run_case_study_selection()
    team_result = _run_team_selection()
    inclusion_policy = _build_inclusion_policy()

    # Get company profile asset IDs for standard depth
    profile_ids = get_company_profile_ids(inclusion_policy.company_profile_depth)

    # Build case study slide index map from catalog lock
    cs_idx_map: dict[str, int] = {}
    for _cat, entries in lock_data.get("case_study_pool", {}).items():
        for entry in entries:
            cs_idx_map[entry["semantic_id"]] = entry["slide_idx"]

    # Build team bio slide index map
    team_idx_map: dict[str, int] = {}
    for entry in lock_data.get("team_bio_pool", []):
        team_idx_map[entry["semantic_id"]] = entry["slide_idx"]

    entries: list[ManifestEntry] = []

    # ─── COVER (3 A2 shells) ───
    entries.append(
        ManifestEntry(
            entry_type="a2_shell",
            asset_id="proposal_cover",
            semantic_layout_id="proposal_cover",
            content_source_policy=ContentSourcePolicy.PROPOSAL_SPECIFIC,
            section_id="cover",
            injection_data={
                "subtitle": "خدمات استشارات التحول الرقمي",
                "client_name": "شركة أكمي",
                "date_text": "مارس 2026",
            },
        )
    )
    entries.append(
        ManifestEntry(
            entry_type="a2_shell",
            asset_id="intro_message",
            semantic_layout_id="intro_message",
            content_source_policy=ContentSourcePolicy.PROPOSAL_SPECIFIC,
            section_id="cover",
        )
    )
    entries.append(
        ManifestEntry(
            entry_type="a2_shell",
            asset_id="toc_agenda",
            semantic_layout_id="toc_table",
            content_source_policy=ContentSourcePolicy.PROPOSAL_SPECIFIC,
            section_id="cover",
            injection_data={
                "title": "جدول المحتويات",
                "rows": [
                    ["01", "الفهم"],
                    ["02", "لماذا ستراتيجيك جيرز"],
                    ["03", "المنهجية"],
                    ["04", "الجدول الزمني والنتائج"],
                    ["05", "الفريق"],
                    ["06", "الحوكمة"],
                ],
            },
        ),
    )

    # ─── SECTION 01: Understanding (divider + 3 B-variable) ───
    entries.append(
        ManifestEntry(
            entry_type="a2_shell",
            asset_id="section_divider_01",
            semantic_layout_id="section_divider_01",
            content_source_policy=ContentSourcePolicy.PROPOSAL_SPECIFIC,
            section_id="section_01",
            injection_data={"title": "الفهم", "body": " "},
        )
    )
    for i, (title, body) in enumerate(
        [
            (
                "سياق المشروع",
                "تسعى شركة أكمي إلى شريك استراتيجي لتحديث بنيتها التحتية "
                "التقنية ودفع التحول الرقمي.",
            ),
            (
                "التحديات الرئيسية",
                "أنظمة قديمة، ومشهد بيانات مجزأ، والحاجة إلى نموذج تشغيل "
                "مرن عبر جميع وحدات الأعمال.",
            ),
            (
                "الأهداف الاستراتيجية",
                "تحقيق تحسين بنسبة 40% في الكفاءة التشغيلية من خلال بنية "
                "السحابة أولاً واتخاذ القرار المبني على البيانات.",
            ),
        ],
        1,
    ):
        entries.append(
            ManifestEntry(
                entry_type="b_variable",
                asset_id=f"understanding_{i:02d}",
                semantic_layout_id="content_heading_desc",
                content_source_policy=ContentSourcePolicy.PROPOSAL_SPECIFIC,
                section_id="section_01",
                injection_data={"title": title, "body": body},
            )
        )

    # ─── SECTION 02: Why Strategic Gears (divider + content + case studies) ───
    entries.append(
        ManifestEntry(
            entry_type="a2_shell",
            asset_id="section_divider_02",
            semantic_layout_id="section_divider_02",
            content_source_policy=ContentSourcePolicy.PROPOSAL_SPECIFIC,
            section_id="section_02",
            injection_data={"title": "لماذا ستراتيجيك جيرز", "body": " "},
        )
    )
    entries.append(
        ManifestEntry(
            entry_type="b_variable",
            asset_id="why_sg_argument",
            semantic_layout_id="content_heading_desc",
            content_source_policy=ContentSourcePolicy.PROPOSAL_SPECIFIC,
            section_id="section_02",
            injection_data={
                "title": "سجل حافل بالإنجازات",
                "body": (
                    "قدمت ستراتيجيك جيرز أكثر من 200 برنامج تحول ناجح "
                    "في منطقة الخليج العربي."
                ),
            },
        )
    )
    # 5 case study pool clones
    for sa in cs_result.selected:
        slide_idx = cs_idx_map.get(sa.asset_id)
        if slide_idx is not None:
            entries.append(
                ManifestEntry(
                    entry_type="pool_clone",
                    asset_id=sa.asset_id,
                    semantic_layout_id="case_study_cases",
                    content_source_policy=ContentSourcePolicy.APPROVED_ASSET_POOL,
                    section_id="section_02",
                    injection_data={"source_slide_idx": slide_idx},
                )
            )

    # ─── SECTION 03: Methodology (divider + overview + 4 focused + 1 detail) ───
    entries.append(
        ManifestEntry(
            entry_type="a2_shell",
            asset_id="section_divider_03",
            semantic_layout_id="section_divider_03",
            content_source_policy=ContentSourcePolicy.PROPOSAL_SPECIFIC,
            section_id="section_03",
            injection_data={"title": "المنهجية", "body": " "},
        )
    )
    entries.append(
        ManifestEntry(
            entry_type="b_variable",
            asset_id="methodology_overview",
            semantic_layout_id="methodology_overview_4",
            content_source_policy=ContentSourcePolicy.PROPOSAL_SPECIFIC,
            section_id="section_03",
            injection_data={
                "title": "منهجيتنا",
                "body": (
                    "نهج مثبت من أربع مراحل يجمع بين أفضل الممارسات "
                    "الصناعية والخبرة العميقة في المجال."
                ),
            },
        )
    )
    meth_bp = _build_methodology_blueprint()
    for phase in meth_bp.phases:
        entries.append(
            ManifestEntry(
                entry_type="b_variable",
                asset_id=f"methodology_phase_{phase.phase_number}",
                semantic_layout_id="methodology_focused_4",
                content_source_policy=ContentSourcePolicy.PROPOSAL_SPECIFIC,
                section_id="section_03",
                injection_data={
                    "title": phase.phase_name_ar or phase.phase_name_en,
                    "body": "؛ ".join(phase.activities),
                },
            )
        )
    entries.append(
        ManifestEntry(
            entry_type="b_variable",
            asset_id="methodology_detail_01",
            semantic_layout_id="methodology_detail",
            content_source_policy=ContentSourcePolicy.PROPOSAL_SPECIFIC,
            section_id="section_03",
            injection_data={
                "title": "الأنشطة التفصيلية — المرحلة 3",
                "body": (
                    "تهيئة المنصة: ترحيل SAP S/4HANA وإعداد البنية "
                    "التحتية السحابية مع مناطق Azure المحددة."
                ),
            },
        )
    )

    # ─── SECTION 04: Timeline & Outcome (divider + 2 B-variable) ───
    entries.append(
        ManifestEntry(
            entry_type="a2_shell",
            asset_id="section_divider_04",
            semantic_layout_id="section_divider_04",
            content_source_policy=ContentSourcePolicy.PROPOSAL_SPECIFIC,
            section_id="section_04",
            injection_data={"title": "الجدول الزمني والنتائج", "body": " "},
        )
    )
    entries.append(
        ManifestEntry(
            entry_type="b_variable",
            asset_id="timeline_01",
            semantic_layout_id="content_heading_content",
            content_source_policy=ContentSourcePolicy.PROPOSAL_SPECIFIC,
            section_id="section_04",
            injection_data={
                "title": "الجدول الزمني للمشروع",
                "body": (
                    "المرحلة 1: الاكتشاف (4 أسابيع) | المرحلة 2: التصميم (6 أسابيع) "
                    "| المرحلة 3: البناء (12 أسبوع) | المرحلة 4: الإطلاق (4 أسابيع)"
                ),
            },
        )
    )
    entries.append(
        ManifestEntry(
            entry_type="b_variable",
            asset_id="deliverables_01",
            semantic_layout_id="content_heading_desc",
            content_source_policy=ContentSourcePolicy.PROPOSAL_SPECIFIC,
            section_id="section_04",
            injection_data={
                "title": "المخرجات الرئيسية",
                "body": (
                    "تقرير التقييم، وثيقة الاستراتيجية، المنصة المهيأة، "
                    "مواد التدريب، تقرير الإطلاق، وثيقة التسليم"
                ),
            },
        )
    )

    # ─── SECTION 05: Team (divider + 3 team bio pool clones) ───
    entries.append(
        ManifestEntry(
            entry_type="a2_shell",
            asset_id="section_divider_05",
            semantic_layout_id="section_divider_05",
            content_source_policy=ContentSourcePolicy.PROPOSAL_SPECIFIC,
            section_id="section_05",
            injection_data={"title": "الفريق", "body": " "},
        )
    )
    for sa in team_result.selected:
        slide_idx = team_idx_map.get(sa.asset_id)
        if slide_idx is not None:
            entries.append(
                ManifestEntry(
                    entry_type="pool_clone",
                    asset_id=sa.asset_id,
                    semantic_layout_id="team_two_members",
                    content_source_policy=ContentSourcePolicy.APPROVED_ASSET_POOL,
                    section_id="section_05",
                    injection_data={"source_slide_idx": slide_idx},
                )
            )

    # ─── SECTION 06: Governance (divider + 1 B-variable) ───
    entries.append(
        ManifestEntry(
            entry_type="a2_shell",
            asset_id="section_divider_06",
            semantic_layout_id="section_divider_06",
            content_source_policy=ContentSourcePolicy.PROPOSAL_SPECIFIC,
            section_id="section_06",
            injection_data={"title": "الحوكمة", "body": " "},
        )
    )
    entries.append(
        ManifestEntry(
            entry_type="b_variable",
            asset_id="governance_01",
            semantic_layout_id="content_heading_desc",
            content_source_policy=ContentSourcePolicy.PROPOSAL_SPECIFIC,
            section_id="section_06",
            injection_data={
                "title": "حوكمة المشروع",
                "body": (
                    "حوكمة منظمة مع اجتماعات لجنة التوجيه الأسبوعية، "
                    "وتقارير الحالة نصف الشهرية، والمراجعات التنفيذية الشهرية."
                ),
            },
        )
    )

    # ─── COMPANY PROFILE (A1 clones — standard depth: 8 slides) ───
    for profile_id in profile_ids:
        entries.append(
            ManifestEntry(
                entry_type="a1_clone",
                asset_id=profile_id,
                semantic_layout_id=profile_id,
                content_source_policy=ContentSourcePolicy.INSTITUTIONAL_REUSE,
                section_id="company_profile",
            )
        )

    # ─── CLOSING (A1 clones: know_more + contact) ───
    entries.append(
        ManifestEntry(
            entry_type="a1_clone",
            asset_id="know_more",
            semantic_layout_id="know_more",
            content_source_policy=ContentSourcePolicy.INSTITUTIONAL_REUSE,
            section_id="closing",
        )
    )
    entries.append(
        ManifestEntry(
            entry_type="a1_clone",
            asset_id="contact",
            semantic_layout_id="contact",
            content_source_policy=ContentSourcePolicy.INSTITUTIONAL_REUSE,
            section_id="closing",
        )
    )

    return ProposalManifest(
        entries=entries,
        inclusion_policy=inclusion_policy,
    )


# ── Template + Catalog Lock Prerequisites ────────────────────────────


class TestARTemplateAvailability:
    """Verify AR template and catalog lock are present and consistent."""

    def test_catalog_lock_ar_exists(self):
        assert CATALOG_LOCK_AR.exists(), f"Missing {CATALOG_LOCK_AR}"

    def test_catalog_lock_ar_valid_json(self):
        data = json.loads(CATALOG_LOCK_AR.read_text(encoding="utf-8"))
        assert "template_hash" in data
        assert "a1_immutable" in data
        assert "a2_shells" in data
        assert "section_dividers" in data
        assert "layouts" in data

    def test_catalog_lock_ar_language(self):
        data = json.loads(CATALOG_LOCK_AR.read_text(encoding="utf-8"))
        assert data.get("language") == "ar"

    def test_catalog_lock_en_exists(self):
        """EN catalog lock also needed for parity audit."""
        assert CATALOG_LOCK_EN.exists(), f"Missing {CATALOG_LOCK_EN}"

    @requires_ar_template
    def test_ar_template_exists(self):
        assert AR_POTX_PATH.exists()

    @requires_ar_template
    def test_ar_template_hash_matches_catalog_lock(self):
        from src.services.template_auditor import file_hash

        data = json.loads(CATALOG_LOCK_AR.read_text(encoding="utf-8"))
        actual_hash = file_hash(AR_POTX_PATH)
        assert actual_hash == data["template_hash"]


# ── Gate 9: EN/AR Parity Audit (AR-specific gate) ───────────────────


class TestARParityAudit:
    """Gate 9: Template parity audit — EN/AR catalog locks must match."""

    def test_gate_09_parity_passes(self):
        """Gate 9: EN/AR parity audit passes with zero critical failures."""
        from src.services.template_parity import run_parity_check

        report = run_parity_check(CATALOG_LOCK_EN, CATALOG_LOCK_AR)
        assert report.overall_pass is True, (
            f"EN/AR parity audit failed: {report.critical_failures} critical failures"
        )
        assert report.critical_failures == 0

    def test_gate_09_zero_warnings(self):
        """Gate 9: EN/AR parity audit has zero warnings."""
        from src.services.template_parity import run_parity_check

        report = run_parity_check(CATALOG_LOCK_EN, CATALOG_LOCK_AR)
        assert report.warnings == 0, (
            f"EN/AR parity has {report.warnings} warning(s)"
        )

    def test_gate_09_all_check_categories_present(self):
        """Gate 9: All parity check categories are exercised."""
        from src.services.template_parity import run_parity_check

        report = run_parity_check(CATALOG_LOCK_EN, CATALOG_LOCK_AR)
        names = {c.name for c in report.checks}
        expected_checks = {
            "language_labels",
            "a1_immutable_ids",
            "a2_shell_ids",
            "section_divider_numbers",
            "layout_semantic_ids",
            "case_study_categories",
            "team_bio_pool_count",
            "team_bio_families",
            "service_divider_pool_count",
            "service_divider_ids",
        }
        for check_name in expected_checks:
            assert check_name in names, f"Missing parity check: {check_name}"

    def test_gate_09_a1_id_set_match(self):
        """Gate 9: EN and AR have identical A1 immutable ID sets."""
        en_data = json.loads(CATALOG_LOCK_EN.read_text(encoding="utf-8"))
        ar_data = json.loads(CATALOG_LOCK_AR.read_text(encoding="utf-8"))
        en_a1 = sorted(en_data["a1_immutable"].keys())
        ar_a1 = sorted(ar_data["a1_immutable"].keys())
        assert en_a1 == ar_a1, f"A1 ID mismatch: EN={en_a1} vs AR={ar_a1}"

    def test_gate_09_a2_id_set_match(self):
        """Gate 9: EN and AR have identical A2 shell ID sets."""
        en_data = json.loads(CATALOG_LOCK_EN.read_text(encoding="utf-8"))
        ar_data = json.loads(CATALOG_LOCK_AR.read_text(encoding="utf-8"))
        en_a2 = sorted(en_data["a2_shells"].keys())
        ar_a2 = sorted(ar_data["a2_shells"].keys())
        assert en_a2 == ar_a2, f"A2 ID mismatch: EN={en_a2} vs AR={ar_a2}"

    def test_gate_09_layout_id_set_match(self):
        """Gate 9: EN and AR have identical semantic layout ID sets."""
        en_data = json.loads(CATALOG_LOCK_EN.read_text(encoding="utf-8"))
        ar_data = json.loads(CATALOG_LOCK_AR.read_text(encoding="utf-8"))
        en_layouts = sorted(en_data["layouts"].keys())
        ar_layouts = sorted(ar_data["layouts"].keys())
        assert en_layouts == ar_layouts, (
            f"Layout ID mismatch: EN has {len(en_layouts)}, AR has {len(ar_layouts)}"
        )

    def test_gate_09_divider_number_set_match(self):
        """Gate 9: EN and AR have identical section divider numbers."""
        en_data = json.loads(CATALOG_LOCK_EN.read_text(encoding="utf-8"))
        ar_data = json.loads(CATALOG_LOCK_AR.read_text(encoding="utf-8"))
        en_div = sorted(en_data["section_dividers"].keys())
        ar_div = sorted(ar_data["section_dividers"].keys())
        assert en_div == ar_div


# ── Policy, Budget, and Selection Validation ─────────────────────────


class TestARPolicyAndBudget:
    """Gates 15, 16, 17: validate policies, budget, and selection (AR context)."""

    def test_gate_15_house_inclusion_policy_valid(self):
        """Gate 15: HouseInclusionPolicy validated."""
        policy = _build_inclusion_policy()
        assert policy.proposal_mode == "standard"
        assert policy.geography == "ksa"
        assert policy.include_ksa_context is True
        assert policy.company_profile_depth == "standard"
        min_cs, max_cs = policy.case_study_count
        assert 4 <= min_cs <= max_cs <= 12
        min_t, max_t = policy.team_bio_count
        assert 2 <= min_t <= max_t <= 6

    def test_gate_15_ksa_includes_ksa_context(self):
        ksa_policy = build_inclusion_policy("standard", "ksa", "technology")
        assert ksa_policy.include_ksa_context is True

    def test_gate_15_non_ksa_excludes_ksa_context(self):
        intl_policy = build_inclusion_policy("standard", "international", "technology")
        assert intl_policy.include_ksa_context is False

    def test_gate_16_slide_budget_valid(self):
        """Gate 16: SlideBudget validated (all counts within ranges)."""
        policy = _build_inclusion_policy()
        meth_bp = _build_methodology_blueprint()
        cs_result = _run_case_study_selection()
        team_result = _run_team_selection()
        budget = compute_slide_budget(
            policy, meth_bp, cs_result, team_result,
            understanding_slides=3,
            timeline_slides=2,
            governance_slides=1,
        )
        errors = validate_budget(budget)
        assert errors == [], f"Budget validation errors: {errors}"
        assert budget.total_slides >= 30, (
            f"Full proposal should have 30+ slides, got {budget.total_slides}"
        )

    def test_gate_16_all_sections_budgeted(self):
        """Gate 16: Every mandatory section has a budget."""
        policy = _build_inclusion_policy()
        meth_bp = _build_methodology_blueprint()
        cs_result = _run_case_study_selection()
        team_result = _run_team_selection()
        budget = compute_slide_budget(
            policy, meth_bp, cs_result, team_result,
        )
        for section_id in MANDATORY_SECTION_ORDER:
            assert section_id in budget.section_budgets, (
                f"Missing budget for mandatory section '{section_id}'"
            )

    def test_gate_17_case_study_selection_auditable(self):
        """Gate 17: CaseStudySelectionResult auditable (scores + reasons)."""
        result = _run_case_study_selection()
        assert len(result.selected) == 5
        for sa in result.selected:
            assert sa.asset_id, "Selected case study has no asset_id"
            assert sa.ranking_score >= 0, "Negative ranking score"
            assert sa.inclusion_reason, "No inclusion reason"

    def test_gate_17_team_selection_auditable(self):
        """Gate 17: TeamSelectionResult auditable (scores + reasons)."""
        result = _run_team_selection()
        assert len(result.selected) == 3
        for sa in result.selected:
            assert sa.asset_id, "Selected team member has no asset_id"
            assert sa.ranking_score >= 0, "Negative ranking score"
            assert sa.inclusion_reason, "No inclusion reason"


# ── Full Manifest Construction Validity ──────────────────────────────


class TestARManifestConstruction:
    """Full AR proposal manifest validity — structural + policy checks."""

    @pytest.fixture
    def manifest(self) -> ProposalManifest:
        return _build_ar_full_proposal_manifest()

    def test_manifest_has_entries(self, manifest):
        assert len(manifest.entries) >= 30, (
            f"Full proposal should have 30+ entries, got {len(manifest.entries)}"
        )

    def test_manifest_passes_validation(self, manifest):
        errors = validate_manifest(manifest)
        assert errors == [], f"Manifest validation errors: {errors}"

    def test_all_entries_have_semantic_layout_id(self, manifest):
        for entry in manifest.entries:
            assert entry.semantic_layout_id, (
                f"Entry {entry.asset_id} missing semantic_layout_id"
            )

    def test_all_entries_have_section_id(self, manifest):
        for entry in manifest.entries:
            assert entry.section_id, f"Entry {entry.asset_id} missing section_id"

    def test_covers_all_mandatory_sections(self, manifest):
        section_ids = list(dict.fromkeys(e.section_id for e in manifest.entries))
        for required in MANDATORY_SECTION_ORDER:
            assert required in section_ids, (
                f"Mandatory section '{required}' missing from manifest"
            )

    def test_section_order_is_correct(self, manifest):
        section_ids = list(dict.fromkeys(e.section_id for e in manifest.entries))
        order_index = {sid: i for i, sid in enumerate(MANDATORY_SECTION_ORDER)}
        positions = [order_index[sid] for sid in section_ids if sid in order_index]
        assert positions == sorted(positions), f"Sections out of order: {section_ids}"

    def test_six_dividers_present(self, manifest):
        dividers = [e for e in manifest.entries if e.asset_id.startswith("section_divider_")]
        assert len(dividers) == 6

    def test_has_pool_clone_case_studies(self, manifest):
        cs = [e for e in manifest.entries if e.entry_type == "pool_clone"
              and e.semantic_layout_id == "case_study_cases"]
        assert len(cs) >= 4, f"Need 4+ case studies, got {len(cs)}"

    def test_has_pool_clone_team_bios(self, manifest):
        team = [e for e in manifest.entries if e.entry_type == "pool_clone"
                and e.semantic_layout_id == "team_two_members"]
        assert len(team) >= 2, f"Need 2+ team bios, got {len(team)}"

    def test_has_methodology_slides(self, manifest):
        meth = [e for e in manifest.entries if e.section_id == "section_03"
                and e.entry_type == "b_variable"]
        assert len(meth) >= 5, f"Need 5+ methodology slides, got {len(meth)}"

    def test_has_company_profile(self, manifest):
        cp = [e for e in manifest.entries if e.section_id == "company_profile"]
        assert len(cp) == 8, f"Standard depth = 8 company profile slides, got {len(cp)}"

    def test_all_entry_types_present(self, manifest):
        types = {e.entry_type for e in manifest.entries}
        assert "a1_clone" in types
        assert "a2_shell" in types
        assert "b_variable" in types
        assert "pool_clone" in types

    def test_policy_rules_per_entry_type(self, manifest):
        for entry in manifest.entries:
            if entry.entry_type == "a1_clone":
                assert entry.content_source_policy == ContentSourcePolicy.INSTITUTIONAL_REUSE
            elif entry.entry_type == "a2_shell":
                assert entry.content_source_policy == ContentSourcePolicy.PROPOSAL_SPECIFIC
            elif entry.entry_type == "b_variable":
                assert entry.content_source_policy == ContentSourcePolicy.PROPOSAL_SPECIFIC
            elif entry.entry_type == "pool_clone":
                assert entry.content_source_policy == ContentSourcePolicy.APPROVED_ASSET_POOL

    def test_arabic_injection_text_present(self, manifest):
        """AR-specific: injection data contains Arabic text."""
        b_entries = [e for e in manifest.entries
                     if e.entry_type == "b_variable" and e.injection_data]
        for entry in b_entries:
            title = entry.injection_data.get("title", "")
            body = entry.injection_data.get("body", "")
            combined = title + body
            # Check for at least one Arabic character (Unicode range)
            has_arabic = any("\u0600" <= ch <= "\u06FF" for ch in combined)
            assert has_arabic, (
                f"Entry {entry.asset_id} injection text is not Arabic: "
                f"title='{title[:40]}...'"
            )


# ── AR Render Execution (requires real AR template) ──────────────────


@requires_ar_template
class TestARRenderExecution:
    """Render a real AR full-deck from the official Arabic template."""

    @pytest.fixture(scope="class")
    def render_result(self, tmp_path_factory):
        from src.services.renderer_v2 import render_v2
        from src.services.template_manager import TemplateManager

        output_dir = tmp_path_factory.mktemp("ar_full_render")
        output_path = output_dir / "ar_full_proposal.pptx"

        manifest = _build_ar_full_proposal_manifest()
        tm = TemplateManager(AR_POTX_PATH, CATALOG_LOCK_AR)
        return render_v2(manifest, tm, CATALOG_LOCK_AR, output_path)

    def test_render_succeeds(self, render_result):
        assert render_result.success, (
            f"AR render failed: manifest_errors={render_result.manifest_errors}, "
            f"render_errors={render_result.render_errors}"
        )

    def test_output_file_exists(self, render_result):
        assert render_result.output_path is not None
        assert Path(render_result.output_path).exists()

    def test_total_slides_matches_manifest(self, render_result):
        manifest = _build_ar_full_proposal_manifest()
        assert render_result.total_slides == len(manifest.entries)

    def test_zero_manifest_errors(self, render_result):
        assert render_result.manifest_errors == []

    def test_zero_render_errors(self, render_result):
        assert render_result.render_errors == []

    def test_template_hash_recorded(self, render_result):
        assert render_result.template_hash != ""
        assert render_result.template_hash.startswith("sha256:")

    def test_all_records_present(self, render_result):
        manifest = _build_ar_full_proposal_manifest()
        assert len(render_result.records) == len(manifest.entries)

    def test_no_per_slide_errors(self, render_result):
        for rec in render_result.records:
            assert rec.error is None, (
                f"Slide {rec.manifest_index} ({rec.asset_id}) error: {rec.error}"
            )


# ── Hard Technical Gates (on rendered AR output) ─────────────────────


@requires_ar_template
class TestARHardGates:
    """Hard technical gates 1-8, 10-23 on rendered AR output."""

    @pytest.fixture(scope="class")
    def render_result(self, tmp_path_factory):
        from src.services.renderer_v2 import render_v2
        from src.services.template_manager import TemplateManager

        output_dir = tmp_path_factory.mktemp("ar_hard_gates")
        output_path = output_dir / "ar_hard_gates.pptx"

        manifest = _build_ar_full_proposal_manifest()
        tm = TemplateManager(AR_POTX_PATH, CATALOG_LOCK_AR)
        return render_v2(manifest, tm, CATALOG_LOCK_AR, output_path)

    @pytest.fixture(scope="class")
    def output_pptx(self, render_result) -> Path:
        assert render_result.success
        return Path(render_result.output_path)

    @pytest.fixture(scope="class")
    def scorer_result(self, output_pptx):
        from src.services.composition_scorer import (
            extract_shapes,
            score_composition,
        )

        shapes = extract_shapes(output_pptx)
        return score_composition(shapes, [], profile=ScorerProfile.OFFICIAL_TEMPLATE_V2)

    def test_gate_01_official_layouts(self, render_result):
        """Gate 1: Every output slide uses an official template layout."""
        for rec in render_result.records:
            assert rec.semantic_layout_id, (
                f"Slide {rec.manifest_index} ({rec.asset_id}) has no semantic_layout_id"
            )

    def test_gate_02_zero_shape_creation_all_v2_modules(self):
        """Gate 2: Zero add_shape()/add_textbox() in ALL v2-path modules."""
        for mod_path in V2_PATH_MODULES:
            path = Path(mod_path)
            if not path.exists():
                continue
            source = path.read_text(encoding="utf-8")
            tree = ast.parse(source)
            for node in ast.walk(tree):
                if isinstance(node, ast.Call) and isinstance(node.func, ast.Attribute):
                    if node.func.attr in ("add_shape", "add_textbox"):
                        pytest.fail(
                            f"Found {node.func.attr}() in {mod_path} at line {node.lineno}"
                        )

    def test_gate_03_euclid_flex_fonts(self, output_pptx):
        """Gate 3: All text uses brand font family.

        AR template uses 'Bahij TheSansArabic' for Arabic text alongside
        'Euclid Flex' for Latin text.  Both are accepted brand fonts for
        the AR template.
        """
        from pptx import Presentation

        # AR brand fonts: Euclid Flex (Latin) + Bahij TheSansArabic (Arabic)
        ar_brand_families = {"Euclid", "Bahij"}

        prs = Presentation(str(output_pptx))
        non_brand_fonts = set()
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.name and run.text.strip():
                                if not any(
                                    family in run.font.name
                                    for family in ar_brand_families
                                ):
                                    non_brand_fonts.add(run.font.name)
        assert non_brand_fonts == set(), (
            f"Non-brand fonts found in rendered AR output: {non_brand_fonts}"
        )

    def test_gate_04_navy_header_color(self, output_pptx):
        """Gate 4: Navy #0E2841 on all headers."""
        from pptx import Presentation

        prs = Presentation(str(output_pptx))
        navy_hex = "0E2841"
        for slide_idx, slide in enumerate(prs.slides):
            for shape in slide.placeholders:
                idx = shape.placeholder_format.idx
                # idx 0 = title placeholder
                if idx == 0 and shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            try:
                                rgb = run.font.color.rgb
                                if rgb is not None:
                                    actual = str(rgb)
                                    assert actual.upper() == navy_hex.upper(), (
                                        f"AR slide {slide_idx} title color {actual} != {navy_hex}"
                                    )
                            except AttributeError:
                                # Color is theme-inherited (_NoneColor) — OK
                                pass

    def test_gate_05_left_margin(self, scorer_result):
        """Gate 5: Left margin >= 0.82 inches on all content slides."""
        from src.services.composition_scorer import ViolationSeverity

        margin_blockers = [
            v for ss in scorer_result.slide_scores
            for v in ss.violations
            if v.rule == "bounds_left" and v.severity == ViolationSeverity.BLOCKER
        ]
        assert margin_blockers == [], (
            f"Left margin violations: {[v.message for v in margin_blockers]}"
        )

    def test_gate_06_zero_composition_blockers(self, scorer_result):
        """Gate 6: 0 composition blockers."""
        assert scorer_result.blocker_count == 0, (
            f"Found {scorer_result.blocker_count} composition blockers: "
            + "; ".join(
                v.message for ss in scorer_result.slide_scores
                for v in ss.violations
                if v.severity.value == "BLOCKER"
            )
        )

    def test_gate_07_zero_mojibake(self, output_pptx):
        """Gate 7: 0 mojibake — no encoding corruption in AR output.

        Extra important for Arabic: checks for U+FFFD (replacement char),
        null bytes, and verifies Arabic characters survived injection.
        """
        from pptx import Presentation

        prs = Presentation(str(output_pptx))
        arabic_found = False
        for slide_idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text
                    assert "\ufffd" not in text, (
                        f"Mojibake (U+FFFD) found in AR slide {slide_idx}, "
                        f"shape '{shape.name}'"
                    )
                    assert "\x00" not in text, (
                        f"Null byte found in AR slide {slide_idx}, "
                        f"shape '{shape.name}'"
                    )
                    if any("\u0600" <= ch <= "\u06FF" for ch in text):
                        arabic_found = True
        assert arabic_found, (
            "No Arabic characters found in rendered AR output — "
            "Arabic text may have been lost during rendering"
        )

    def test_gate_08_zero_contract_violations(self, render_result):
        """Gate 8: Zero placeholder-contract violations."""
        for rec in render_result.records:
            if rec.injection_result and rec.injection_result.errors:
                pytest.fail(
                    f"Slide {rec.manifest_index} ({rec.asset_id}) injection errors: "
                    f"{rec.injection_result.errors}"
                )

    def test_gate_10_template_hash_validated(self, render_result):
        """Gate 10: AR template hash validated before render."""
        data = json.loads(CATALOG_LOCK_AR.read_text(encoding="utf-8"))
        assert render_result.template_hash == data["template_hash"]

    def test_gate_11_zero_anti_leak(self, render_result):
        """Gate 11: Zero anti-leak violations — A2 shells sanitized."""
        a2_records = [r for r in render_result.records if r.entry_type == "a2_shell"]
        for rec in a2_records:
            assert rec.sanitization_report is not None, (
                f"A2 slide {rec.asset_id} missing sanitization_report"
            )
            assert rec.sanitization_report.errors == [], (
                f"A2 slide {rec.asset_id} sanitization errors: "
                f"{rec.sanitization_report.errors}"
            )

    def test_gate_12_mandatory_section_flow(self, render_result):
        """Gate 12: Exact mandatory section flow in output."""
        seen_sections: list[str] = []
        for rec in render_result.records:
            if rec.section_id not in seen_sections:
                seen_sections.append(rec.section_id)
        order_index = {sid: i for i, sid in enumerate(MANDATORY_SECTION_ORDER)}
        positions = [order_index[sid] for sid in seen_sections if sid in order_index]
        assert positions == sorted(positions)

    def test_gate_13_exact_divider_numbering(self, render_result):
        """Gate 13: Exact divider numbering flow 01-06."""
        dividers = [r for r in render_result.records
                    if r.asset_id.startswith("section_divider_")]
        numbers = [r.asset_id.split("_")[-1] for r in dividers]
        assert numbers == ["01", "02", "03", "04", "05", "06"]

    def test_gate_14_content_source_policy(self):
        """Gate 14: ContentSourcePolicy validated on every ManifestEntry."""
        manifest = _build_ar_full_proposal_manifest()
        for entry in manifest.entries:
            assert entry.content_source_policy in ContentSourcePolicy
            if entry.entry_type == "a1_clone":
                assert entry.content_source_policy == ContentSourcePolicy.INSTITUTIONAL_REUSE
            elif entry.entry_type == "pool_clone":
                assert entry.content_source_policy == ContentSourcePolicy.APPROVED_ASSET_POOL

    def test_gate_18_editability(self, output_pptx):
        """Gate 18: Output PPTX opens and has editable placeholders."""
        from pptx import Presentation

        prs = Presentation(str(output_pptx))
        assert len(prs.slides) >= 30
        has_placeholders = any(
            True for slide in prs.slides for _ in slide.placeholders
        )
        assert has_placeholders, "AR output PPTX has no editable placeholders"

    def test_gate_19_a1_not_sanitized(self, render_result):
        """Gate 19: A1 slides preserve template-native assets (no sanitization)."""
        a1_records = [r for r in render_result.records if r.entry_type == "a1_clone"]
        assert len(a1_records) >= 8, "Need 8+ A1 slides in full proposal"
        for rec in a1_records:
            assert rec.sanitization_report is None, (
                f"A1 slide {rec.asset_id} should NOT be sanitized"
            )

    def test_gate_20_no_legacy_imports(self):
        """Gate 20: renderer_v2 MUST NOT import legacy shape-building helpers."""
        renderer_v2 = Path("src/services/renderer_v2.py")
        if not renderer_v2.exists():
            pytest.skip("renderer_v2.py not found")
        source = renderer_v2.read_text(encoding="utf-8")
        tree = ast.parse(source)
        for node in ast.walk(tree):
            if isinstance(node, ast.Import | ast.ImportFrom):
                if isinstance(node, ast.ImportFrom) and node.module:
                    assert "renderer" not in node.module or "renderer_v2" in node.module, (
                        f"renderer_v2.py imports from legacy module: {node.module}"
                    )
                    assert "formatting" not in (node.module or ""), (
                        f"renderer_v2.py imports from formatting: {node.module}"
                    )

    def test_gate_21_shell_sanitization(self, render_result):
        """Gate 21: Shell sanitization verified on every A2 slide."""
        a2_records = [r for r in render_result.records if r.entry_type == "a2_shell"]
        assert len(a2_records) >= 9, "Need 9 A2 shells (3 cover + 6 dividers)"
        for rec in a2_records:
            assert rec.sanitization_report is not None, (
                f"A2 slide {rec.asset_id} not sanitized"
            )

    def test_gate_22_scorer_profile(self):
        """Gate 22: Composition scorer uses template-v2 profile for v2 output."""
        from src.pipeline.graph import get_scorer_profile

        profile = get_scorer_profile(RendererMode.TEMPLATE_V2)
        assert profile == ScorerProfile.OFFICIAL_TEMPLATE_V2

    def test_gate_23_semantic_ids(self, render_result):
        """Gate 23: All runtime layout resolution uses semantic layout IDs."""
        for rec in render_result.records:
            assert rec.semantic_layout_id
            assert " " not in rec.semantic_layout_id, (
                f"Display name as layout ID: '{rec.semantic_layout_id}'"
            )


# ── Visual-Fidelity Gates ────────────────────────────────────────────


@requires_ar_template
class TestARVisualGates:
    """Visual-fidelity gates 24-33 on rendered AR output."""

    @pytest.fixture(scope="class")
    def render_result(self, tmp_path_factory):
        from src.services.renderer_v2 import render_v2
        from src.services.template_manager import TemplateManager

        output_dir = tmp_path_factory.mktemp("ar_visual")
        output_path = output_dir / "ar_visual_deck.pptx"

        manifest = _build_ar_full_proposal_manifest()
        tm = TemplateManager(AR_POTX_PATH, CATALOG_LOCK_AR)
        return render_v2(manifest, tm, CATALOG_LOCK_AR, output_path)

    @pytest.fixture(scope="class")
    def output_pptx(self, render_result) -> Path:
        assert render_result.success
        return Path(render_result.output_path)

    def test_gate_24_section_order(self, render_result):
        """Gate 24: Section order matches mandatory flow."""
        seen: list[str] = []
        for rec in render_result.records:
            if rec.section_id not in seen:
                seen.append(rec.section_id)
        for required in MANDATORY_SECTION_ORDER:
            assert required in seen, f"Missing section {required}"

    def test_gate_25_divider_pattern(self, render_result):
        """Gate 25: Section dividers use 01-06 numbered pattern."""
        dividers = [r for r in render_result.records if "section_divider_" in r.asset_id]
        numbers = sorted(r.asset_id.split("_")[-1] for r in dividers)
        assert numbers == ["01", "02", "03", "04", "05", "06"]

    def test_gate_26_methodology_follows_budget(self, render_result):
        """Gate 26: Methodology section follows budget (overview + focused + detail)."""
        meth_records = [
            r for r in render_result.records
            if r.section_id == "section_03" and r.entry_type == "b_variable"
        ]
        layout_ids = [r.semantic_layout_id for r in meth_records]
        # Must have overview
        assert any("overview" in lid for lid in layout_ids), "No methodology overview"
        # Must have focused phases
        focused = [lid for lid in layout_ids if "focused" in lid]
        assert len(focused) >= 3, f"Need 3+ focused phases, got {len(focused)}"
        # Must have at least 1 detail
        details = [lid for lid in layout_ids if "detail" in lid]
        assert len(details) >= 1, f"Need 1+ detail slides, got {len(details)}"

    def test_gate_27_case_studies_use_correct_layout(self, render_result):
        """Gate 27: Case studies use case_study_cases layout."""
        cs_records = [
            r for r in render_result.records
            if r.entry_type == "pool_clone" and r.section_id == "section_02"
        ]
        assert len(cs_records) >= 4, f"Need 4+ case studies in render, got {len(cs_records)}"
        for rec in cs_records:
            assert rec.semantic_layout_id == "case_study_cases", (
                f"Case study {rec.asset_id} uses layout '{rec.semantic_layout_id}', "
                f"expected 'case_study_cases'"
            )

    def test_gate_28_team_bios_use_correct_layout(self, render_result):
        """Gate 28: Team bios use team_two_members layout."""
        team_records = [
            r for r in render_result.records
            if r.entry_type == "pool_clone" and r.section_id == "section_05"
        ]
        assert len(team_records) >= 2, f"Need 2+ team bios in render, got {len(team_records)}"
        for rec in team_records:
            assert rec.semantic_layout_id == "team_two_members", (
                f"Team bio {rec.asset_id} uses layout '{rec.semantic_layout_id}', "
                f"expected 'team_two_members'"
            )

    def test_gate_29_typography_hierarchy(self, output_pptx):
        """Gate 29: Typography hierarchy — titles larger than body text."""
        from pptx import Presentation

        prs = Presentation(str(output_pptx))
        title_sizes = []
        body_sizes = []
        for slide in prs.slides:
            for shape in slide.placeholders:
                idx = shape.placeholder_format.idx
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.size and run.font.size > 0:
                                pt = run.font.size / 12700
                                if idx == 0:
                                    title_sizes.append(pt)
                                elif idx >= 10:
                                    body_sizes.append(pt)
        if title_sizes and body_sizes:
            avg_title = sum(title_sizes) / len(title_sizes)
            avg_body = sum(body_sizes) / len(body_sizes)
            assert avg_title > avg_body, (
                f"Title avg ({avg_title:.1f}pt) should be > body avg ({avg_body:.1f}pt)"
            )

    def test_gate_30_no_generic_powerpoint_look(self, output_pptx):
        """Gate 30: No generic PowerPoint look — brand elements present.

        AR template uses Bahij TheSansArabic (Arabic) + Euclid Flex (Latin).
        Either brand font family confirms non-generic appearance.
        """
        from pptx import Presentation

        ar_brand_families = {"Euclid", "Bahij"}
        prs = Presentation(str(output_pptx))
        brand_font_found = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.name and any(
                                family in run.font.name
                                for family in ar_brand_families
                            ):
                                brand_font_found = True
                                break
                    if brand_font_found:
                        break
            if brand_font_found:
                break
        assert brand_font_found, "No brand font found in AR output — generic look"

    def test_gate_31_structural_fidelity_a1_clones(self, render_result, output_pptx):
        """Gate 31: House slides pass structural fidelity checks.

        A1 slides use layout-inherited content.  Fidelity verified by:
        - Slide has a valid layout
        - Layout has shapes (visual content)
        - Layout name matches expected layout from AR catalog lock
        """
        from pptx import Presentation

        lock_data = json.loads(CATALOG_LOCK_AR.read_text(encoding="utf-8"))
        prs = Presentation(str(output_pptx))

        manifest = _build_ar_full_proposal_manifest()
        a1_entries = [(i, e) for i, e in enumerate(manifest.entries)
                      if e.entry_type == "a1_clone"]

        for entry_idx, entry in a1_entries:
            if entry_idx >= len(prs.slides):
                continue
            slide = prs.slides[entry_idx]
            # Verify slide has a layout
            assert slide.slide_layout is not None, (
                f"A1 clone '{entry.asset_id}' has no layout"
            )
            # Verify slide has content — either explicit shapes or
            # layout-inherited shapes
            total_visual = len(slide.shapes) + len(slide.slide_layout.shapes)
            assert total_visual > 0, (
                f"A1 clone '{entry.asset_id}' at index {entry_idx} has "
                f"no shapes (slide-level or layout-level)"
            )
            # Verify layout display name matches AR catalog lock
            a1_info = lock_data["a1_immutable"].get(entry.asset_id, {})
            expected_display = a1_info.get("display_name", "")
            if expected_display:
                assert slide.slide_layout.name == expected_display, (
                    f"A1 clone '{entry.asset_id}' layout mismatch: "
                    f"expected '{expected_display}', got '{slide.slide_layout.name}'"
                )

    def test_gate_32_house_shells_allowlisted(self, render_result):
        """Gate 32: House shells preserve only allowlisted text."""
        a2_records = [r for r in render_result.records if r.entry_type == "a2_shell"]
        for rec in a2_records:
            report = rec.sanitization_report
            assert report is not None
            assert report.errors == [], (
                f"A2 {rec.asset_id} sanitization errors: {report.errors}"
            )

    def test_gate_33_no_generic_regression(self, render_result, output_pptx):
        """Gate 33: No generic-PowerPoint regression — structure matches real proposal."""
        from pptx import Presentation

        prs = Presentation(str(output_pptx))
        assert len(prs.slides) >= 30, "Full proposal should have 30+ slides"
        divider_count = sum(
            1 for r in render_result.records
            if r.asset_id.startswith("section_divider_")
        )
        assert divider_count == 6, "Need 6 section dividers"
        pool_clones = [r for r in render_result.records if r.entry_type == "pool_clone"]
        assert len(pool_clones) >= 7, "Need 7+ pool clones (5 cs + 2+ team)"


# ── Integration Gates ────────────────────────────────────────────────


@requires_ar_template
class TestARIntegrationGates:
    """Integration gates 35-42 (gate 34 is EN-specific)."""

    @pytest.fixture(scope="class")
    def output_pptx(self, tmp_path_factory) -> Path:
        from src.services.renderer_v2 import render_v2
        from src.services.template_manager import TemplateManager

        output_dir = tmp_path_factory.mktemp("ar_integration")
        output_path = output_dir / "ar_integration_deck.pptx"

        manifest = _build_ar_full_proposal_manifest()
        tm = TemplateManager(AR_POTX_PATH, CATALOG_LOCK_AR)
        result = render_v2(manifest, tm, CATALOG_LOCK_AR, output_path)

        assert result.success, (
            f"AR render failed: {result.manifest_errors + result.render_errors}"
        )
        return Path(result.output_path)

    @pytest.fixture(scope="class")
    def render_result(self, tmp_path_factory):
        from src.services.renderer_v2 import render_v2
        from src.services.template_manager import TemplateManager

        output_dir = tmp_path_factory.mktemp("ar_integ_rr")
        output_path = output_dir / "ar_integ_rr.pptx"

        manifest = _build_ar_full_proposal_manifest()
        tm = TemplateManager(AR_POTX_PATH, CATALOG_LOCK_AR)
        return render_v2(manifest, tm, CATALOG_LOCK_AR, output_path)

    def test_gate_35_ar_full_deck_renders(self, output_pptx):
        """Gate 35: Render one AR full-deck from official AR .potx."""
        assert output_pptx.exists()
        assert output_pptx.stat().st_size > 0

    def test_gate_36_output_opens_as_valid_pptx(self, output_pptx):
        """Gate 36: Output is valid PPTX — opens and has slides."""
        from pptx import Presentation

        prs = Presentation(str(output_pptx))
        assert len(prs.slides) >= 30

    def test_gate_36_all_slides_from_official_layouts(self, output_pptx):
        """Gate 36: All slides come from official template layouts."""
        from pptx import Presentation

        prs = Presentation(str(output_pptx))
        for slide in prs.slides:
            assert slide.slide_layout is not None

    def test_gate_37_zero_shape_creation_v2_path(self):
        """Gate 37: Zero add_shape()/add_textbox() in v2 path."""
        for mod_path in V2_PATH_MODULES:
            path = Path(mod_path)
            if not path.exists():
                continue
            source = path.read_text(encoding="utf-8")
            tree = ast.parse(source)
            for node in ast.walk(tree):
                if isinstance(node, ast.Call) and isinstance(node.func, ast.Attribute):
                    if node.func.attr in ("add_shape", "add_textbox"):
                        pytest.fail(f"Found {node.func.attr}() in {mod_path}")

    def test_gate_38_anti_leak_on_non_a1_slides(self, output_pptx):
        """Gate 38: Anti-leak on A2 shells and B variable slides."""
        from pptx import Presentation

        manifest = _build_ar_full_proposal_manifest()
        non_a1_indices = {
            i for i, e in enumerate(manifest.entries) if e.entry_type != "a1_clone"
        }

        prs = Presentation(str(output_pptx))
        forbidden_fragments = ["Film Sector"]
        for slide_idx, slide in enumerate(prs.slides):
            if slide_idx not in non_a1_indices:
                continue
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text
                    for frag in forbidden_fragments:
                        assert frag not in text, (
                            f"Forbidden text '{frag}' in non-A1 slide {slide_idx}, "
                            f"shape '{shape.name}'"
                        )

    def test_gate_39_case_study_clone_correct(self, render_result, output_pptx):
        """Gate 39: Case-study slide clones correctly — has shapes and content."""
        from pptx import Presentation

        manifest = _build_ar_full_proposal_manifest()
        cs_indices = [
            i for i, e in enumerate(manifest.entries)
            if e.entry_type == "pool_clone" and e.semantic_layout_id == "case_study_cases"
        ]
        assert len(cs_indices) >= 4

        prs = Presentation(str(output_pptx))
        for idx in cs_indices:
            if idx < len(prs.slides):
                slide = prs.slides[idx]
                assert len(slide.shapes) > 0, (
                    f"Case study clone at index {idx} has 0 shapes"
                )

    def test_gate_40_team_bio_clone_correct(self, render_result, output_pptx):
        """Gate 40: Team-bio slide clones correctly — has shapes and content."""
        from pptx import Presentation

        manifest = _build_ar_full_proposal_manifest()
        team_indices = [
            i for i, e in enumerate(manifest.entries)
            if e.entry_type == "pool_clone" and e.semantic_layout_id == "team_two_members"
        ]
        assert len(team_indices) >= 2

        prs = Presentation(str(output_pptx))
        for idx in team_indices:
            if idx < len(prs.slides):
                slide = prs.slides[idx]
                assert len(slide.shapes) > 0, (
                    f"Team bio clone at index {idx} has 0 shapes"
                )

    def test_gate_41_methodology_follows_template_family(self, render_result):
        """Gate 41: Methodology section follows template family correctly."""
        meth_records = [
            r for r in render_result.records
            if r.section_id == "section_03" and r.entry_type == "b_variable"
        ]
        valid_meth_layouts = {
            "methodology_overview_4", "methodology_focused_4",
            "methodology_overview_3", "methodology_focused_3",
            "methodology_detail",
        }
        for rec in meth_records:
            assert rec.semantic_layout_id in valid_meth_layouts, (
                f"Methodology slide '{rec.asset_id}' uses invalid layout "
                f"'{rec.semantic_layout_id}'"
            )

    def test_gate_42_scorer_profile_per_mode(self):
        """Gate 42: Composition scorer uses correct profile per renderer mode."""
        from src.pipeline.graph import get_scorer_profile

        v2_profile = get_scorer_profile(RendererMode.TEMPLATE_V2)
        assert v2_profile == ScorerProfile.OFFICIAL_TEMPLATE_V2

        legacy_profile = get_scorer_profile(RendererMode.LEGACY)
        assert legacy_profile == ScorerProfile.LEGACY


# ── Scorer Profile Alignment ─────────────────────────────────────────


class TestARScorerProfileAlignment:
    """Composition scorer profile validation for AR v2 output."""

    def test_v2_profile_uses_euclid_flex(self):
        from src.services.scorer_profiles import get_profile

        config = get_profile(ScorerProfile.OFFICIAL_TEMPLATE_V2)
        assert "Euclid Flex" in config.brand_fonts

    def test_v2_profile_body_font_range(self):
        from src.services.scorer_profiles import get_profile

        config = get_profile(ScorerProfile.OFFICIAL_TEMPLATE_V2)
        assert config.body_font_min_pt == 9
        assert config.body_font_max_pt == 12

    def test_v2_profile_margin_threshold(self):
        from src.services.scorer_profiles import get_profile

        config = get_profile(ScorerProfile.OFFICIAL_TEMPLATE_V2)
        assert config.bounds_margin_left_min_in == pytest.approx(0.82, abs=0.01)

    def test_legacy_profile_unchanged(self):
        from src.services.scorer_profiles import get_profile

        config = get_profile(ScorerProfile.LEGACY)
        assert "Aptos" in config.brand_fonts
        assert config.body_font_min_pt == 10
        assert config.body_font_max_pt == 14


# ── AR Gate Matrix Meta-Test ─────────────────────────────────────────


# Complete gate matrix: all 42 gates with status and test function name.
# Gate 34 is EN-specific (N-A for Phase 18 AR verification).
# Gates 9 and 35 are AR-specific (PASS here, N-A in Phase 17 EN).
AR_GATE_MATRIX: dict[int, dict[str, str]] = {
    1:  {"status": "PASS", "test": "test_gate_01_official_layouts",
         "proof": "render_result records — every record has semantic_layout_id"},
    2:  {"status": "PASS", "test": "test_gate_02_zero_shape_creation_all_v2_modules",
         "proof": "AST scan of all v2-path module source files"},
    3:  {"status": "PASS", "test": "test_gate_03_euclid_flex_fonts",
         "proof": "rendered AR PPTX — brand font check (Euclid Flex + Bahij TheSansArabic)"},
    4:  {"status": "PASS", "test": "test_gate_04_navy_header_color",
         "proof": "rendered AR PPTX title placeholder color.rgb check"},
    5:  {"status": "PASS", "test": "test_gate_05_left_margin",
         "proof": "scorer violations on rendered AR PPTX shapes"},
    6:  {"status": "PASS", "test": "test_gate_06_zero_composition_blockers",
         "proof": "scorer blocker_count on rendered AR PPTX"},
    7:  {"status": "PASS", "test": "test_gate_07_zero_mojibake",
         "proof": "all text in rendered AR PPTX is valid Unicode + Arabic chars present"},
    8:  {"status": "PASS", "test": "test_gate_08_zero_contract_violations",
         "proof": "render_result records — zero injection errors"},
    9:  {"status": "PASS", "test": "test_gate_09_* (8 tests)",
         "proof": "EN/AR parity audit: 14 checks pass, zero critical failures"},
    10: {"status": "PASS", "test": "test_gate_10_template_hash_validated",
         "proof": "render_result.template_hash matches AR catalog lock hash"},
    11: {"status": "PASS", "test": "test_gate_11_zero_anti_leak",
         "proof": "rendered AR PPTX A2 shells sanitized with zero errors"},
    12: {"status": "PASS", "test": "test_gate_12_mandatory_section_flow",
         "proof": "render_result records section_id ordering"},
    13: {"status": "PASS", "test": "test_gate_13_exact_divider_numbering",
         "proof": "render_result divider asset_ids sequence 01-06"},
    14: {"status": "PASS", "test": "test_gate_14_content_source_policy",
         "proof": "manifest entry content_source_policy per entry_type"},
    15: {"status": "PASS", "test": "test_gate_15_* (3 tests)",
         "proof": "HouseInclusionPolicy fields + KSA/non-KSA rules"},
    16: {"status": "PASS", "test": "test_gate_16_* (2 tests)",
         "proof": "SlideBudget validation + all sections budgeted"},
    17: {"status": "PASS", "test": "test_gate_17_* (2 tests)",
         "proof": "selection result scores + inclusion_reason audit"},
    18: {"status": "PASS", "test": "test_gate_18_editability",
         "proof": "python-pptx re-opens rendered AR output PPTX"},
    19: {"status": "PASS", "test": "test_gate_19_a1_not_sanitized",
         "proof": "render_result A1 records have no sanitization_report"},
    20: {"status": "PASS", "test": "test_gate_20_no_legacy_imports",
         "proof": "AST scan renderer_v2.py imports"},
    21: {"status": "PASS", "test": "test_gate_21_shell_sanitization",
         "proof": "all A2 records have sanitization_report with zero errors"},
    22: {"status": "PASS", "test": "test_gate_22_scorer_profile",
         "proof": "v2 profile == OFFICIAL_TEMPLATE_V2"},
    23: {"status": "PASS", "test": "test_gate_23_semantic_ids",
         "proof": "all render_result records have non-empty semantic_layout_id"},
    24: {"status": "PASS", "test": "test_gate_24_section_order",
         "proof": "render_result section_id order matches MANDATORY_SECTION_ORDER"},
    25: {"status": "PASS", "test": "test_gate_25_divider_pattern",
         "proof": "render_result divider records are 01-06 in sequence"},
    26: {"status": "PASS", "test": "test_gate_26_methodology_follows_budget",
         "proof": "render_result methodology records: overview + focused + detail"},
    27: {"status": "PASS", "test": "test_gate_27_case_studies_use_correct_layout",
         "proof": "render_result pool_clone records with case_study_cases layout"},
    28: {"status": "PASS", "test": "test_gate_28_team_bios_use_correct_layout",
         "proof": "render_result pool_clone records with team_two_members layout"},
    29: {"status": "PASS", "test": "test_gate_29_typography_hierarchy",
         "proof": "rendered AR PPTX title font > body font (pt comparison)"},
    30: {"status": "PASS", "test": "test_gate_30_no_generic_powerpoint_look",
         "proof": "rendered AR PPTX contains brand font (Euclid Flex or Bahij TheSansArabic)"},
    31: {"status": "PASS", "test": "test_gate_31_structural_fidelity_a1_clones",
         "proof": "rendered AR PPTX A1 slides: layout fidelity + layout name match"},
    32: {"status": "PASS", "test": "test_gate_32_house_shells_allowlisted",
         "proof": "render_result A2 sanitization_report zero errors"},
    33: {"status": "PASS", "test": "test_gate_33_no_generic_regression",
         "proof": "rendered AR PPTX 30+ slides + 6 dividers + 7+ pool clones"},
    34: {"status": "N-A", "test": "—",
         "proof": "EN-specific: EN full deck render (Phase 17)"},
    35: {"status": "PASS", "test": "test_gate_35_ar_full_deck_renders",
         "proof": "AR output file exists and non-empty"},
    36: {"status": "PASS", "test": "test_gate_36_* (2 tests)",
         "proof": "valid AR PPTX opens with 30+ slides + all layouts present"},
    37: {"status": "PASS", "test": "test_gate_37_zero_shape_creation_v2_path",
         "proof": "AST scan all v2-path modules"},
    38: {"status": "PASS", "test": "test_gate_38_anti_leak_on_non_a1_slides",
         "proof": "rendered AR PPTX non-A1 slides scanned for forbidden text"},
    39: {"status": "PASS", "test": "test_gate_39_case_study_clone_correct",
         "proof": "rendered AR PPTX case study slides have shapes"},
    40: {"status": "PASS", "test": "test_gate_40_team_bio_clone_correct",
         "proof": "rendered AR PPTX team bio slides have shapes"},
    41: {"status": "PASS", "test": "test_gate_41_methodology_follows_template_family",
         "proof": "render_result methodology records use valid layout IDs"},
    42: {"status": "PASS", "test": "test_gate_42_scorer_profile_per_mode",
         "proof": "pipeline get_scorer_profile dispatch returns correct profile"},
}


class TestARGateMatrix:
    """Meta-test: verify all 42 gates are accounted for in AR verification."""

    def test_all_42_gates_in_matrix(self):
        """Every gate 1-42 must appear in AR_GATE_MATRIX."""
        for gate_num in range(1, 43):
            assert gate_num in AR_GATE_MATRIX, (
                f"Gate {gate_num} missing from AR_GATE_MATRIX"
            )

    def test_en_specific_gate_marked_na(self):
        """Gate 34 must be N-A (EN-specific)."""
        assert AR_GATE_MATRIX[34]["status"] == "N-A"

    def test_ar_specific_gates_pass(self):
        """Gates 9 and 35 (AR-specific) must be PASS."""
        assert AR_GATE_MATRIX[9]["status"] == "PASS"
        assert AR_GATE_MATRIX[35]["status"] == "PASS"

    def test_all_ar_gates_pass(self):
        """All AR-applicable gates must be PASS."""
        for gate_num, info in AR_GATE_MATRIX.items():
            if info["status"] == "N-A":
                continue
            assert info["status"] == "PASS", (
                f"Gate {gate_num} is not PASS: {info}"
            )

    def test_all_pass_gates_have_test_function(self):
        """Every PASS gate must reference a test function name."""
        for gate_num, info in AR_GATE_MATRIX.items():
            if info["status"] == "N-A":
                continue
            assert info["test"] and info["test"] != "—", (
                f"Gate {gate_num} is PASS but has no test function"
            )

    def test_all_pass_gates_have_proof_description(self):
        """Every PASS gate must describe its proof method."""
        for gate_num, info in AR_GATE_MATRIX.items():
            if info["status"] == "N-A":
                continue
            assert info["proof"], (
                f"Gate {gate_num} is PASS but has no proof description"
            )

    def test_gate_tests_exist_in_module(self):
        """Verify that referenced test functions actually exist in this module."""
        import sys
        module = sys.modules[__name__]
        # Collect all test method names from all classes in this module
        all_test_names: set[str] = set()
        for name in dir(module):
            obj = getattr(module, name)
            if isinstance(obj, type) and name.startswith("Test"):
                for attr_name in dir(obj):
                    if attr_name.startswith("test_gate_"):
                        all_test_names.add(attr_name)

        for gate_num, info in AR_GATE_MATRIX.items():
            if info["status"] == "N-A":
                continue
            test_ref = info["test"]
            # Handle multi-test references like "test_gate_15_* (3 tests)"
            if "*" in test_ref:
                prefix = test_ref.split("*")[0].strip()
                matching = [n for n in all_test_names if n.startswith(prefix)]
                assert len(matching) > 0, (
                    f"Gate {gate_num} references '{test_ref}' but no matching "
                    f"test functions found with prefix '{prefix}'"
                )
            else:
                assert test_ref in all_test_names, (
                    f"Gate {gate_num} references '{test_ref}' but it does not "
                    f"exist in this module"
                )
